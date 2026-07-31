"""EasyPPTX渲染器"""
import os, random
from typing import Dict, Optional, List
from easypptx import Presentation as EasyPres
from ..templates.style_tokens import get_theme
try:
    from ..extract.savings_normalizer import calc_irr as _ma_irr_no_wd, calc_irr_withdraw as _ma_irr_wd
except ImportError:
    _ma_irr_no_wd = None
    _ma_irr_wd = None

def _c(h): return h.lstrip('#') if h else '000000'
def _p(v): return f"{v}%"
def _d(sl, text, x=5, y=5, w=90, h=10, size=12, color=None, bold=False):
    return sl.add_text(str(text), x=_p(x), y=_p(y), w=_p(w), h=_p(h),
                       font_size=size, bold=bold, color=_c(color) if color else '332825')

def _slide(pres, bg=None):
    sl = pres.add_slide()
    if bg:
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        sh = sl.pptx_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(12192000), Emu(6858000))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(int(_c(bg)[0:2],16), int(_c(bg)[2:4],16), int(_c(bg)[4:6],16))
        sh.line.fill.background()
        spTree = sl.pptx_slide.shapes._spTree
        spTree.remove(sh._element)
        spTree.insert(2, sh._element)
    return sl

def render_pptx_easy(data, output_path, theme='caramel', cover_image=None, logo_path=None,
                     company_images=None, scene_images=None,
                     ci_data=None, iul_data=None, ci_company=None, iul_company=None):
    C = get_theme(theme)
    summary = data.get('summary', {})
    nw = {int(k) if str(k).isdigit() else k: v for k, v in data.get('no_withdraw', {}).items()}
    wd = {int(k) if str(k).isdigit() else k: v for k, v in data.get('withdraw', {}).items()}
    pt = data.get('paid_total', 0)
    meta = data.get('meta', {})
    hw = bool(wd) and any(v.get('Annual_WD', 0) > 0 for v in wd.values())
    ia = meta.get('insured_age', 1) or 1
    py = int(meta.get('payment_years') or 0)
    cur = meta.get('currency') or meta.get('product_currency') or 'USD'
    start_wd_yr = 0
    annual_wd = 0
    if wd:
        for yk in sorted(int(k) for k in wd.keys()):
            aw = wd[yk].get('Annual_WD', 0) or 0
            if yk > 0 and aw > 0:
                start_wd_yr = yk
                annual_wd = aw
                break

    for d in [nw, wd]:
        for k, r in list(d.items()):
            t = r.get('Total', 0) or 0
            if r.get('Mult') is None: r['Mult'] = t/pt if pt else 0
            if r.get('IRR') is None and int(k) > 0:
                if d is wd and _ma_irr_wd:
                    r['IRR'] = _ma_irr_wd(int(k), (r.get('Cum_WD',0) or 0)+t, pt, py, cur, start_wd_yr, annual_wd)
                elif _ma_irr_no_wd:
                    r['IRR'] = _ma_irr_no_wd(int(k), t, pt, py, cur)
            if r.get('Simple') is None and int(k) > 0: r['Simple'] = (t-pt)/pt/float(k) if pt else None

    pres = EasyPres()
    
    # Cover
    sl = _slide(pres, C['bg_page'])
    sl.title = f"{meta.get('company_short_en','')} | {meta.get('company_short','')}"
    _d(sl, meta.get('product_name', '') or summary.get('product_name', '计划'), x=5, y=20, w=55, h=10, size=26, bold=True, color=C['dark_text'])
    _d(sl, f"{summary.get('payment_years',5)}年缴 · {summary.get('currency','USD')} {int(summary.get('annual_premium',0)):,}/年", x=5, y=33, w=55, h=5, size=11, color=C['mid_text'])
    _d(sl, f"受保人：{summary.get('insured_name','VIP')}（{summary.get('insured_age',1)}岁）", x=5, y=40, w=55, h=5, size=12, color=C['body_text'])
    if cover_image and os.path.exists(cover_image):
        sl.pptx_slide.shapes.add_picture(cover_image, 0, 0, int(9.5*914400), int(7.5*914400))

    # Company
    sl = _slide(pres, C['bg_page'])
    sl.title = f"关于我们 | {meta.get('company_name_zh','保险公司')}"
    bp = meta.get('brand_profile', {})
    for i, v in enumerate([bp.get('founded_year','—'), bp.get('rating_value','—'), bp.get('series_label','—')]):
        _d(sl, str(v), x=5+i*32, y=22, w=28, h=8, size=22, bold=True, color=C['dark_text'])
    for i, t in enumerate((bp.get('business_lines') or [])[:4]):
        _d(sl, t, x=5, y=50+i*5, w=45, h=4, size=9, color=C['body_text'])
    for i, t in enumerate((bp.get('brand_background') or [])[:4]):
        _d(sl, t, x=55, y=50+i*5, w=45, h=4, size=9, color=C['body_text'])

    # Features
    sl = _slide(pres, C['bg_page'])
    sl.title = "产品亮点"
    pb = next((y for y in sorted(int(k) for k in nw.keys()) if nw[y].get('Total',0) > pt), None)
    items = [('01','5年短期缴付','只需5年完成缴费\n后续无需再缴','5','年缴清')]
    if pb: items.append(('02',f'第{pb}年快速回本',f'Y{pb}现价 USD {nw[pb]["Total"]:,}','Y'+str(pb),'回本'))
    for i,(no,title,body,kpi,unit) in enumerate(items[:4]):
        cx, cy = 5+(i%2)*50, 22+(i//2)*28
        _d(sl, no, x=cx, y=cy, w=8, h=6, size=24, bold=True, color=C['primary'])
        _d(sl, title, x=cx+10, y=cy, w=38, h=5, size=14, bold=True, color=C['dark_text'])
        _d(sl, body, x=cx, y=cy+8, w=48, h=8, size=9, color=C['body_text'])
        _d(sl, kpi, x=cx, y=cy+18, w=48, h=6, size=20, bold=True, color=C['primary'])

    # Growth
    sl = _slide(pres, C['bg_page'])
    sl.title = "现金价值增长（不提领）"
    try:
        import matplotlib; matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        plt.rcParams['font.family'] = 'STHeiti Medium'
        yrs = sorted(int(k) for k in nw.keys() if nw[k].get('Total',0) > 0)
        vals = [nw[y]['Total']/1000 for y in yrs]
        guar = [nw[y]['Guar_CV']/1000 for y in yrs]
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5), facecolor='#'+_c(C['bg_page']))
        ax1.bar(yrs[:30], guar[:30], label='保证', color='#'+_c(C['primary']), alpha=0.7)
        ax1.bar(yrs[:30], [v-g for v,g in zip(vals[:30],guar[:30])], bottom=guar[:30], label='非保证', color='#'+_c(C['accent']), alpha=0.7)
        ax1.legend(); ax1.set_title('保证/非保证构成')
        ax2.plot(yrs, vals, color='#'+_c(C['accent_dark']), linewidth=2)
        ax2.set_title('退保总额趋势')
        plt.tight_layout()
        p = f'/tmp/c_{random.randint(0,99999)}.png'; plt.savefig(p, dpi=150); plt.close()
        sl.add_image(p, x=_p(3), y=_p(20), width=_p(94))
    except Exception as e:
        for y in [5,10,15,20,30]:
            if y in nw: _d(sl, f'Y{y}: USD {nw[y]["Total"]:,}', x=5+(y//5)*18, y=25+10*(y%10), w=16, h=4, size=9, color=C['body_text'])

    # No-withdraw table
    sl = _slide(pres, C['bg_page'])
    sl.title = "不提领情形"
    sy = sorted(y for y in [1,5,10,15,20,25,30,35,40,50,60,70] if y in nw)[:12]
    for ri, y in enumerate(sy):
        r = nw[y]; irr = f"{r['IRR']*100:.2f}%" if r.get('IRR') else '-'
        for ci, txt in enumerate([str(y), str(r.get('Age',y)), f"{r.get('Paid',pt):,}", f"{r['Guar_CV']:,}", f"{r['Total']-r['Guar_CV']:,}", f"{r['Total']:,}", irr]):
            _d(sl, txt, x=3+ci*13, y=20+ri*4.5, w=12, h=4, size=7, bold=ri==0, color=C['dark_text'] if ri==0 else C['body_text'])

    # Withdraw table
    if hw:
        sl = _slide(pres, C['bg_page'])
        ws = next((y for y in sorted(wd.keys()) if wd[y].get('Annual_WD',0) > 0), 7)
        sl.title = f"提领方案 · Y{ws}起年提 USD {wd[ws]['Annual_WD']:,}"
        sy = sorted(set([5,ws,ws+1,10,15,20,30,40,50,60] + [max(wd.keys())]))[:12]
        for ri, y in enumerate(sy):
            r = wd[y]
            for ci, txt in enumerate([str(y), str(r.get('Age',y)), f"{r.get('Annual_WD',0):,}", f"{r.get('Cum_WD',0):,}", f"{r.get('Total',0):,}"]):
                _d(sl, txt, x=3+ci*19, y=20+ri*4.5, w=18, h=4, size=7, bold=ri==0, color=C['dark_text'] if ri==0 else C['body_text'])

    # Compare
    sl = _slide(pres, C['bg_page'])
    if hw:
        sl.title = "不提领 vs 提领"
        cm = sorted(set(nw.keys()) & set(wd.keys()))[:10]
        for i, y in enumerate(cm):
            wv = wd[y].get('Total_Received', wd[y]['Cum_WD']+wd[y]['Total'])
            _d(sl, f"Y{y}", x=5+i*9, y=20, w=8, h=4, size=8, bold=True, color=C['dark_text'])
            _d(sl, f"不提领\nUSD {nw[y]['Total']:,}", x=5+i*9, y=25, w=8, h=8, size=6, color=C['primary'])
            _d(sl, f"提领\nUSD {wv:,}", x=5+i*9, y=34, w=8, h=8, size=6, color=C['accent'])
    else:
        sl.title = "关键年限数据"
        for y in [5,10,15,20,30]:
            if y in nw: _d(sl, f"Y{y}: USD {nw[y]['Total']:,}  (IRR: {nw[y].get('IRR',0)*100:.1f}%)" if nw[y].get('IRR') else '', x=5+(y//5)*18, y=25, w=16, h=5, size=10, color=C['body_text'])

    # Scenario
    if hw:
        sl = _slide(pres, C['bg_page'])
        is_edu = ia < 18
        sl.title = "教育金现金流" if is_edu else "养老金现金流"
        ws = next((y for y in sorted(wd.keys()) if wd[y].get('Annual_WD',0) > 0), 7)
        _d(sl, f"从{ia+ws-1}岁起每年 USD {wd[ws]['Annual_WD']:,}", x=5, y=18, w=90, h=5, size=11, color=C['mid_text'])
        bands = [(6,12,'小学'),(13,15,'初中'),(16,18,'高中'),(19,22,'大学'),(23,99,'研究生')] if is_edu else [(60,65,'退休初期'),(65,75,'活跃养老'),(75,85,'稳健养老'),(85,95,'享老晚年'),(95,120,'传承')]
        for i,(lo,hi,name) in enumerate(bands):
            cx = 5+i*19
            _d(sl, name, x=cx, y=25, w=17, h=5, size=12, bold=True, color=C['dark_text'])
            ys = [y for y in sorted(wd.keys()) if lo <= ia+y-1 <= hi]
            if ys: _d(sl, f"累计: USD {wd[ys[-1]]['Cum_WD']:,}\n剩余: USD {wd[ys[-1]]['Total']:,}", x=cx, y=32, w=17, h=10, size=8, color=C['body_text'])

    # CI section
    if ci_data:
        sl = _slide(pres, C['dark_text'])
        sl.title = ""; _d(sl, "家庭保障规划", x=5, y=25, w=90, h=10, size=28, bold=True, color=C['bg_light'])
        _d(sl, "筑起收入高墙，守护家人未来", x=5, y=38, w=90, h=6, size=16, color=C['accent'])
        sl = _slide(pres, C['bg_page'])
        sl.title = "危疾保障概要"
        pol = ci_data.get('policy',{})
        _d(sl, f"年缴: USD {int(pol.get('annual_premium',0)):,}  |  保额: USD {int(pol.get('sum_insured',0)):,}", x=5, y=22, w=90, h=5, size=11, color=C['body_text'])
        for i, item in enumerate((ci_data.get('coverage_items') or [])[:8]):
            _d(sl, f"· {item.get('label','')}  USD {int(item.get('amount',0)):,}", x=5, y=32+i*4, w=90, h=4, size=9, color=C['body_text'])

    # IUL section
    if iul_data:
        sl = _slide(pres, C['dark_text'])
        sl.title = ""; _d(sl, "高杠杆寿险规划", x=5, y=25, w=90, h=10, size=28, bold=True, color=C['bg_light'])
        _d(sl, "用现金流撬动千万保障", x=5, y=38, w=90, h=6, size=16, color=C['accent'])
        sl = _slide(pres, C['bg_page'])
        sl.title = "万用寿险概要"
        ip = float(iul_data.get('policy',{}).get('annual_premium',0))
        ic = float(iul_data.get('policy',{}).get('sum_insured',0))
        _d(sl, f"年缴: USD {int(ip):,}  |  身故赔偿: USD {int(ic):,}  |  杠杆: {ic/ip:.1f}x" if ip else '', x=5, y=22, w=90, h=5, size=11, color=C['body_text'])
        if hw and wd:
            sl2 = _slide(pres, C['bg_page'])
            sl2.title = "现金流 + 高杠杆"
            aw = wd.get(next((y for y in sorted(wd.keys()) if wd[y].get('Annual_WD',0) > 0), None), {}).get('Annual_WD', 0)
            _d(sl2, f"储蓄每年提领 USD {int(aw):,} → IUL保费 USD {int(ip):,}", x=5, y=22, w=90, h=6, size=11, color=C['body_text'])
            _d(sl2, f"IUL身故赔偿 USD {int(ic):,} ({int(ic/ip) if ip else 0}x杠杆)", x=5, y=32, w=90, h=6, size=16, bold=True, color=C['primary'])

    # Common pages
    sl = _slide(pres, C['bg_page']); sl.title = "重要事项声明"
    sl = _slide(pres, C['bg_page']); sl.title = "数据来源与回溯"
    _d(sl, "所有数据直接源自保险公司官方计划书", x=5, y=22, w=90, h=5, size=10, color=C['body_text'])

    sl = _slide(pres, C['bg_page'])
    sl.title = "总结与建议"
    _d(sl, f"5年缴清 USD {pt:,}", x=5, y=22, w=28, h=5, size=10, color=C['body_text'])
    if pb: _d(sl, f"第{pb}年回本", x=37, y=22, w=28, h=5, size=10, color=C['body_text'])
    if hw: _d(sl, f"Y{next(iter(wd.keys()))}起年提", x=69, y=22, w=28, h=5, size=10, color=C['body_text'])

    # Ending
    sl = _slide(pres, C['dark_text'])
    if ci_data or iul_data:
        _d(sl, "一份规划，多重守护", x=5, y=22, w=90, h=8, size=24, bold=True, color=C['bg_light'])
        _d(sl, "现金流 + 保障 + 杠杆 · 完整的家庭财务规划", x=5, y=33, w=90, h=6, size=14, color=C['accent'])
    else:
        _d(sl, "愿这份规划", x=5, y=25, w=90, h=10, size=28, bold=True, color=C['bg_light'])
        _d(sl, "陪伴您和家人稳健成长", x=5, y=38, w=90, h=8, size=20, bold=True, color=C['accent'])
    _d(sl, f"—— 致 {summary.get('insured_name','VIP')} 与家人 ——", x=5, y=50, w=90, h=5, size=12, color=C['bg_light'])

    pres.save(output_path)
    return output_path
