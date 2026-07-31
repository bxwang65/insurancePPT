"""
PPTX 渲染器 (python-pptx 原生, 完全可编辑)
支持: 5 套主题色 / 图片插入 / 折线图叠加 / 销售叙事驱动
"""
import os, math, re
import random
from typing import Dict, List, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree

from ..templates.style_tokens import get_theme, FONT_HEI, FONT_LATIN


# ── 工具函数 ──────────────────────────────────────────
def _c(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _extract_bottom_color(img_path: str, sample_height: int = 20) -> Optional[RGBColor]:
    """从图片底部采样提取主色调, 压暗后作为背景色"""
    try:
        from PIL import Image
        import colorsys
        with Image.open(img_path) as img:
            w, h = img.size
            # 采样底部 sample_height 像素区域
            box = (0, h - sample_height, w, h)
            strip = img.crop(box)
            # 缩放到1px宽取平均色
            strip = strip.resize((1, 1))
            r, g, b = strip.getpixel((0, 0))[:3]
            # 转 HSL, 降低亮度(压暗约35%), 降低饱和度(让底色更柔和)
            r_n, g_n, b_n = r/255.0, g/255.0, b/255.0
            hue, sat, light = colorsys.rgb_to_hls(r_n, g_n, b_n)
            light = max(0.08, light * 0.50)  # 压暗到50%亮度, 最低0.08
            sat = sat * 0.7                  # 降饱和度30%
            r2, g2, b2 = colorsys.hls_to_rgb(hue, light, sat)
            return RGBColor(int(r2*255), int(g2*255), int(b2*255))
    except Exception:
        return None


def _add_picture(slide, img_path: str, x, y, w, h):
    """插入图片, 文件不存在时静默跳过"""
    if not img_path or not os.path.exists(img_path):
        return None
    try:
        return slide.shapes.add_picture(img_path, x, y, w, h)
    except Exception:
        return None


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _c(fill) if isinstance(fill, str) else fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = _c(line) if isinstance(line, str) else line
        if line_width is not None: shp.line.width = Pt(line_width)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=11, color=None, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_HEI):
    if color is None: color = '#442E24'
    if isinstance(color, str): color = _c(color)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        rPr = run._r.get_or_add_rPr()
        eaFont = rPr.find(qn('a:ea'))
        if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
        eaFont.set('typeface', FONT_HEI)
        latinFont = rPr.find(qn('a:latin'))
        if latinFont is None: latinFont = etree.SubElement(rPr, qn('a:latin'))
        latinFont.set('typeface', FONT_LATIN)
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def fill_bg(slide, color='#EFE4DF'):
    bg = add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=color, line=None)
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_header(slide, meta=None, C=None, company_override=None):
    """页面头部: 公司名 + 页码 + 角标（支持公司覆盖）"""
    page_num = meta.get('_page_num', 1) if meta else 1
    total = meta.get('_total_slides', 12) if meta else 12
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])

    # 公司信息: 有覆盖则用覆盖, 否则用meta
    cid = company_override.get('id', meta.get('company_id', '')) if company_override else (meta.get('company_id', '') if meta else '')
    cname = company_override.get('name_zh', meta.get('company_short', '')) if company_override else (meta.get('company_short', '') if meta else '')
    cname_en = company_override.get('brand_profile', {}).get('short_en', meta.get('company_short_en', '')) if company_override else (meta.get('company_short_en', '') if meta else '')

    # 角标
    assets_dir = meta.get('_assets_dir', '') if meta else ''
    if cid and assets_dir:
        co_dir = os.path.join(assets_dir, 'library/companies', cid)
        for corner in ['brand-corner-left.png', 'brand-corner-right.png']:
            cp = os.path.join(co_dir, corner)
            if os.path.exists(cp):
                x = Emu(0) if 'left' in corner else Inches(12.5)
                _add_picture(slide, cp, x, Emu(0), Inches(0.8), Inches(0.5))

    brand = f"{cname_en}  |  {cname}" if cname_en and cname else 'INSURANCE  |  保险计划'
    add_text(slide, Inches(0.5), Inches(0.15), Inches(6), Inches(0.35),
             brand, size=20, color=C['primary'], bold=True)
    add_text(slide, Inches(11.5), Inches(0.18), Inches(1.5), Inches(0.32),
             f'{page_num:02d} / {total:02d}', size=9, color=C['mid_text'],
             align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None, title_size=36, sub_size=14, C=None):
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
             title, size=title_size, color=C['dark_text'], bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
                 subtitle, size=sub_size, color=C['mid_text'])
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(0.6), Inches(0.04),
             fill=C['primary'])


# ── 图标装饰: 金色方块粒子 ──────────────────────────
def _add_sparkles(slide, count=35, seed=7, x_start=0, y_start=0, w_range=13.0, h_range=7.3):
    random.seed(seed)
    for _ in range(count):
        cx = Inches(x_start + random.random() * w_range)
        cy = Inches(y_start + random.random() * h_range)
        sz = Emu(int(Inches(0.04 + random.random() * 0.12)))
        d = add_rect(slide, cx, cy, sz, sz, fill='#D4B878')
        d.fill.transparency = random.uniform(0.5, 0.85)


# ── 客户信息表 ──────────────────────────────────────────
def _build_core_features(product_data, product_type, paid_total=None, meta=None):
    """根据产品类型和数据生成产品核心功能描述"""
    policy = (product_data or {}).get('policy', {}) or {}
    insured = (product_data or {}).get('insured', {}) or {}
    meta = meta or {}
    # 储蓄险data无顶层policy，从meta回退
    if not policy.get('annual_premium'):
        policy = {
            'annual_premium': meta.get('annual_premium', 0),
            'premium_payment_period': str(meta.get('payment_years', 5)) + '年',
            'coverage_period': meta.get('coverage_period', '终身'),
        }
    # 储蓄险: BI在no_withdraw dict里; CI/IUL: BI是list
    bi_raw = (product_data or {}).get('benefit_illustration', []) or []
    no_wd = (product_data or {}).get('no_withdraw', {}) or {}
    bi = bi_raw if isinstance(bi_raw, list) else list(bi_raw.values()) if isinstance(bi_raw, dict) else []
    if not bi and no_wd:
        bi = list(no_wd.values())

    premium = float(policy.get('annual_premium', 0) or 0)
    pay_yrs = str(policy.get('premium_payment_period', '—') or '—')
    coverage = float(policy.get('sum_insured', 0) or 0)
    cvr_period = str(policy.get('coverage_period', '终身') or '终身')
    # IUL总保费单独计算
    if product_type == 'iul' and premium > 0:
        iul_pay = 1 if pay_yrs == '趸交' else int(''.join(c for c in pay_yrs if c.isdigit()) or '5')
        total_prem = premium * iul_pay
    else:
        total_prem = (paid_total or premium * 5) if premium else 0

    if product_type == 'savings':
        vals = [r for r in bi if r.get('Total') or r.get('total_surrender_value')]
        first_pay = None
        # 找第一个现价>已缴的年份
        for r in sorted(bi, key=lambda x: int(x.get('Y', x.get('policy_year', 0)))):
            y = int(r.get('Y', r.get('policy_year', 0)))
            total = float(r.get('Total', r.get('total_surrender_value', 0)))
            if y > 0 and total >= total_prem * 0.9 and total > 0:
                first_pay = y
                break
        return "储蓄 · 现金流规划", (
            f"▸ 短期投入，长期增值：{pay_yrs}缴清，资金即享复利滚存"
            + (f"，第{first_pay}年起可灵活提取" if first_pay else "")
            + "\n"
            f"▸ 稳定现金流：可根据需要灵活提取现金价值，兼顾品质生活与财富传承\n"
            f"▸ 多元货币配置：支持多币种选择，满足子女教育、退休养老等不同阶段需求\n"
            f"▸ 财富定向传承：可指定受益人，实现财富无缝传承"
        )
    elif product_type == 'ci':
        return "重疾 · 全面保障", (
            f"▸ 高额保障：基本保额{format(coverage/10000, '.0f') if coverage else '—'}万"
            + (f"，保终身" if '终身' in cvr_period else f"，保障至{cvr_period}")
            + "\n"
            f"▸ 多重赔付：覆盖多种疾病，癌症/心脏病/中风享多次赔付\n"
            f"▸ 保额增长：保额可随年份递增，抵御医疗通胀\n"
            f"▸ 保费豁免：确诊重疾后豁免剩余保费，保障继续有效"
        )
    elif product_type == 'iul':
        leverage = coverage / total_prem if total_prem > 0 and coverage > 0 else 0
        return "IUL · 高杠杆传承", (
            f"▸ 超高杠杆：总保费USD {total_prem:,.0f}，即获USD {coverage:,.0f}身故保障"
            + (f"，杠杆{leverage:.1f}倍" if leverage > 0 else "")
            + "\n"
            f"▸ 指数挂钩增长：账户挂钩指数表现，享市场潜力同时保底\n"
            f"▸ 税务优势：身故赔偿免遗产税，高效财富传承\n"
            f"▸ 保费灵活：可配合私人银行保单融资，优化资金效率"
        )
    return "", ""


def _slide_client_info(prs, product_data, company_info, meta, C, product_type,
                       paid_total=None, product_name_override=None):
    """客户信息表：在每个产品呈现前插入"""
    # 支持两种数据格式: 原始提取数据(CI/IUL) 和 归一化数据(储蓄data字典)
    policy = (product_data or {}).get('policy', {}) or {}
    insured = (product_data or {}).get('insured', {}) or {}
    # 储蓄险的 data 结构含 meta/summary, 无顶层 policy/insured, 从meta回退
    summary = (product_data or {}).get('summary', {}) or {}
    if not policy and not insured:
        policy = {
            'annual_premium': meta.get('annual_premium', 0),
            'premium_payment_period': str(meta.get('payment_years', 5)) + '年',
            'coverage_period': meta.get('coverage_period', '终身'),
            'sum_insured': meta.get('sum_insured', None),
            'product_name': meta.get('product_name', ''),
        }
        insured = {
            'name': meta.get('insured_name', '—'),
            'age': meta.get('insured_age', '—'),
        }

    product_name = product_name_override or (product_data or {}).get('product_name', '')
    if not product_name:
        product_name = summary.get('product_name', policy.get('product_name', '—'))
    company_zh = (company_info or {}).get('name_zh', meta.get('company_short', ''))
    company_en = (company_info or {}).get('brand_profile', {}).get('short_en', meta.get('company_short_en', ''))

    premium = float(policy.get('annual_premium', 0) or 0)
    raw_pay = str(policy.get('premium_payment_period', '—') or '—')
    pay_yrs = raw_pay if raw_pay.endswith('缴') else (raw_pay if raw_pay == '趸交' else raw_pay + '缴')
    cvr_period = str(policy.get('coverage_period', '终身') or '终身')
    coverage = policy.get('sum_insured', None)
    coverage_str = f"USD {coverage:,}" if coverage and float(coverage) > 0 else '—（不适用）' if product_type == 'savings' else '—'
    # IUL总保费: 从premium和缴费期计算
    if product_type == 'iul' and premium > 0:
        iul_pay_yrs = 1 if raw_pay == '趸交' else int(''.join(c for c in raw_pay if c.isdigit()) or '5')
        total_prem = premium * iul_pay_yrs
    else:
        total_prem = (paid_total or premium * 5) if premium else 0
    total_str = f"USD {total_prem:,.0f}" if total_prem else '—'
    premium_str = f"USD {premium:,.0f}" if premium else '—'

    insured_name = insured.get('name', '—') or '—'
    insured_age = insured.get('age', '—') or '—'
    insured_age_str = f"{insured_age}岁" if str(insured_age) != '—' else '—'

    # 核心功能
    feat_title, feat_body = _build_core_features(product_data, product_type, paid_total, meta)

    # 客户关键信息（投保人/受益人未知则留空）
    type_labels = {'savings': '储蓄寿险', 'ci': '危疾保障', 'iul': '指数型万用寿险'}
    type_label = type_labels.get(product_type, '—')

    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    # 标题（用"客户信息表"替代重复的公司名，header已有公司信息）
    add_text(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.45),
             "客户信息表", size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.5), Inches(1.0), Inches(12), Inches(0.3),
             "方案概览 · 投保信息 · 产品核心功能", size=11, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.35), Inches(0.6), Inches(0.04),
             fill=C['primary'])

    # 信息卡片: 2列 x 5行
    table_top = Inches(1.65)
    row_h = Inches(0.48)
    label_w = Inches(1.6)
    value_w = Inches(4.0)
    left = Inches(0.5)
    mid_x = Inches(6.8)
    gap_row = Inches(0.06)

    rows_data = [
        ("产品名称", product_name, "产品类型", type_label),
        ("投保人", '—（暂未提供）', "投保人年龄", '—（暂未提供）'),
        ("受保人", insured_name, "受保人年龄", insured_age_str),
        ("年缴保费", premium_str, "缴费年期", pay_yrs),
        ("总缴保费", total_str, "保障年期", cvr_period),
        ("基本保额", coverage_str, "受益人", '—（暂未提供）'),
    ]

    for i, (lbl1, val1, lbl2, val2) in enumerate(rows_data):
        y = table_top + i * (row_h + gap_row)
        # Left pair
        add_rect(s, left, y, label_w, row_h, fill=C['bg_card'])
        add_text(s, left, y, label_w, row_h, lbl1, size=9, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, left + label_w, y, value_w, row_h, fill=C['bg_light'])
        add_text(s, left + label_w, y, value_w, row_h, val1, size=11, color=C['dark_text'], bold=True, align=PP_ALIGN.CENTER)
        # Right pair
        add_rect(s, mid_x, y, label_w, row_h, fill=C['bg_card'])
        add_text(s, mid_x, y, label_w, row_h, lbl2, size=9, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, mid_x + label_w, y, value_w, row_h, fill=C['bg_light'])
        add_text(s, mid_x + label_w, y, value_w, row_h, val2, size=11, color=C['dark_text'], bold=True, align=PP_ALIGN.CENTER)

    # 产品核心功能
    note_y = table_top + len(rows_data) * (row_h + gap_row) + Inches(0.15)
    add_rect(s, left, note_y, Inches(12.3), Inches(1.5), fill=C['bg_card'])
    add_text(s, left + Inches(0.2), note_y + Inches(0.08), Inches(11.9), Inches(0.28),
             f"产品核心功能 · {feat_title}", size=11, color=C['accent'], bold=True)
    add_text(s, left + Inches(0.2), note_y + Inches(0.38), Inches(11.9), Inches(1.05),
             feat_body, size=9, color=C['mid_text'])


def render_pptx(data: Dict, output_path: str,
                theme: str = 'caramel',
                cover_image: Optional[str] = None,
                logo_path: Optional[str] = None,
                company_images: Optional[List[str]] = None,
                scene_images: Optional[List[str]] = None,
                ci_data: Optional[Dict] = None,
                iul_data: Optional[Dict] = None,
                ci_company: Optional[Dict] = None,
                iul_company: Optional[Dict] = None,
                all_extractions: Optional[Dict] = None,
                ) -> str:
    """
    主入口: 从normalized数据生成PPTX

    参数:
        data — build_normalized_data() 的输出 (储蓄险)
        output_path — 保存路径
        theme — caramel / broker / business / chinese / ink / minimal
        cover_image — 封面背景图 (JPG/PNG)
        logo_path — 公司 Logo PNG
        company_images — 公司介绍页照片列表
        scene_images — 教育金/养老金场景图列表
        ci_data — 重疾险数据 (可选, 用于组合方案)
        iul_data — IUL数据 (可选, 用于组合方案)
    """
    C = get_theme(theme)
    summary = data['summary']
    no_wd_raw = data['no_withdraw']
    wd_raw = data['withdraw']
    # 兼容 str/int key
    no_wd = {}
    for k, v in no_wd_raw.items():
        try: no_wd[int(k)] = v
        except: no_wd[k] = v
    wd = {}
    for k, v in wd_raw.items():
        try: wd[int(k)] = v
        except: wd[k] = v
    paid_total = data['paid_total']

    # 确保计算字段存在 (Mult, IRR, Simple)
    for _dict in [no_wd, wd]:
        for k, r in list(_dict.items()):
            if r.get('Mult') is None:
                r['Mult'] = r['Total'] / paid_total if paid_total else 0
            total_received = r.get('Total_Received', r.get('Cum_WD', 0) + r.get('Total', 0))
            if r.get('IRR') is None and total_received and paid_total and total_received > paid_total and k > 0:
                try: r['IRR'] = (total_received / paid_total) ** (1 / float(k)) - 1
                except: r['IRR'] = None
            if r.get('Total_Received') is None and 'Cum_WD' in r:
                r['Total_Received'] = r.get('Cum_WD', 0) + r.get('Total', 0)
            if r.get('IRR') is None and r.get('Total') and paid_total and r['Total'] > paid_total and k > 0:
                try: r['IRR'] = (r['Total'] / paid_total) ** (1 / float(k)) - 1
                except: r['IRR'] = None
            if r.get('Simple') is None and r.get('Total') and paid_total and k > 0:
                try: r['Simple'] = (r['Total'] - paid_total) / paid_total / float(k)
                except: r['Simple'] = None

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    meta = data.get('meta', {})
    # 注入主题色
    meta['_theme'] = theme
    meta['_C'] = C

    # 解析销售叙事 Markdown (如果有)
    narrative_md = meta.get('narrative_md', '')
    if narrative_md:
        try:
            from .narrative_parser import narrative_to_meta
            narrative_meta = narrative_to_meta(narrative_md)
            meta.update(narrative_meta)
            # 封面副标题优先用叙事
            if meta.get('narrative_subtitle') and not meta.get('chat_narrative'):
                meta['chat_narrative'] = meta['narrative_subtitle']
        except Exception:
            pass  # 叙事解析失败不影响主流程

    # ── 叙事驱动幻灯片构建 ──
    # 解析叙事 Markdown 确定结构和文案
    has_wd = bool(wd) and any(v.get('Annual_WD', 0) > 0 for v in wd.values())
    has_ci = bool(ci_data)
    has_iul = bool(iul_data)
    slides_def = []

    # 解析叙事章节
    narrative_slides = []
    narrative_md = meta.get('narrative_md', '')
    if narrative_md:
        try:
            from .narrative_parser import parse_narrative_markdown
            narrative_slides = parse_narrative_markdown(narrative_md)
        except Exception:
            narrative_slides = []

    # 检查叙事中是否包含各类型章节
    has_narrative_chapter = lambda t: any(s.get('type') == t for s in narrative_slides)
    nar_has_savings = has_narrative_chapter('savings') or any(s.get('section', '').startswith('chapter') and s.get('type') == 'savings' for s in narrative_slides)
    nar_has_ci = has_narrative_chapter('ci')
    nar_has_iul = has_narrative_chapter('iul')
    nar_has_combo = has_narrative_chapter('combo')

    # 有叙事时按叙事结构走, 否则用默认硬编码结构
    if narrative_slides:
        for ns in narrative_slides:
            sec = ns.get('section', '')
            stype = ns.get('type', '')

            if sec == 'cover':
                slides_def.append(('cover', lambda: _slide_cover(prs, summary, paid_total, meta, C, cover_image)))

            elif sec == 'company':
                slides_def.append(('company', lambda: _slide_company(prs, data, meta, C, logo_path, company_images)))

            elif stype == 'savings' or (sec.startswith('chapter') and stype == 'savings'):
                # 储蓄章节: 客户信息表 + 数据幻灯片
                slides_def.append(('savings_client_info', lambda: _slide_client_info(
                    prs, data, None, meta, C, 'savings', paid_total)))
                slides_def.append(('features', lambda: _slide_features(prs, no_wd, paid_total, meta, C)))
                slides_def.append(('growth', lambda: _slide_growth_chart(prs, no_wd, meta, C)))
                slides_def.append(('no_withdraw_table', lambda: _slide_no_withdraw_table(prs, no_wd, paid_total, meta, C)))
                if has_wd:
                    slides_def.append(('withdraw_table', lambda: _slide_withdraw_table(prs, wd, paid_total, meta, C)))
                slides_def.append(('compare_chart', lambda: _slide_compare_chart(prs, no_wd, wd, meta, C)))
                if has_wd:
                    slides_def.append(('education', lambda: _slide_education(prs, wd, meta, C, scene_images)))

            elif stype == 'ci':
                # CI 章节: 客户信息表 + CI内容
                if has_ci:
                    slides_def.append(('ci_divider', lambda: _slide_ci_divider(prs, ci_data, meta, C)))
                    slides_def.append(('ci_client_info', lambda: _slide_client_info(
                        prs, ci_data, ci_company, meta, C, 'ci')))
                    if ci_company:
                        slides_def.append(('ci_company', lambda: _slide_sub_company(prs, ci_company, '危疾保障', meta, C)))
                    slides_def.append(('ci_overview', lambda: _slide_ci_overview(prs, ci_data, meta, C)))

            elif stype == 'iul':
                # IUL 章节
                if has_iul:
                    slides_def.append(('iul_divider', lambda: _slide_iul_divider(prs, iul_data, meta, C)))
                    if iul_company:
                        iul_sub_meta = {**meta, 'product_name': (iul_data.get('summary') or {}).get('product_name', '新加坡IUL')}
                        slides_def.append(('iul_company', lambda: _slide_sub_company(prs, iul_company, '万用寿险', iul_sub_meta, C)))
                    slides_def.append(('iul_overview', lambda: _slide_iul_overview(prs, iul_data, meta, C)))
                    if has_wd:
                        slides_def.append(('iul_funding', lambda: _slide_iul_funding(prs, wd, iul_data, meta, C)))

            elif sec == 'combo':
                slides_def.append(('combo', lambda: _slide_combo_narrative(prs, meta, C, ci_data, iul_data)))

            elif sec == 'ending':
                if has_ci or has_iul:
                    slides_def.append(('ending_combined', lambda: _slide_ending_combined(prs, summary, ci_data, iul_data, meta, C)))
                else:
                    slides_def.append(('ending', lambda: _slide_ending(prs, summary, meta, C)))

        # 无论叙事如何, 都加声明和来源页
        slides_def.append(('disclaimer', lambda: _slide_disclaimer(prs, meta, C)))
        slides_def.append(('sources', lambda: _slide_sources(prs, data, meta, C)))
        slides_def.append(('summary', lambda: _slide_summary(prs, no_wd, wd, paid_total, meta, C)))

    else:
        # 无叙事时用默认硬编码结构
        has_savings = meta.get('has_savings', bool(no_wd))
        slides_def = [
            ('cover', lambda: _slide_cover(prs, summary, paid_total, meta, C, cover_image)),
            # 2026-07-16 V3.3.9+: 不再加全局 deck-level 'company' 封面,
            # 每个产品自己带 _slide_company (在 _savings_product_section 里), 否则第一个产品会双重封面
        ]
        # V3.3.6: 多储蓄产品支持 — 每个产品一份完整 section, per-product meta override
        def _savings_meta(sav_c):
            bp = (sav_c or {}).get('brand_profile') or {}
            return {
                **meta,
                # 2026-07-16 V3.3.9+: 覆盖 brand_profile 否则 deck-level profile 串到所有产品
                'brand_profile': bp or meta.get('brand_profile'),
                'company_id': (sav_c or {}).get('id') or meta.get('company_id'),
                'company_name_zh': (sav_c or {}).get('name_zh') or bp.get('name_zh') or meta.get('company_name_zh'),
                'company_short': (sav_c or {}).get('name_zh') or bp.get('short') or meta.get('company_short'),
                'company_short_en': bp.get('short_en') or meta.get('company_short_en'),
                'product_name': (sav_c or {}).get('product_name') or meta.get('product_name'),
            }

        def _savings_product_section(sav_d, sav_c):
            """Build slides_def for one savings product (normalize BI/WD, render 6-7 slides).

            2026-07-16 V3.3.9+: 每个产品前加 deck-level '关于我们' 封面页 (用产品自己的 logo_path/company_images),
            不再依赖全局 logo_path/company_images, 保证多产品场景每家公司封面都用各自 logo
            """
            sm = _savings_meta(sav_c)
            ins_age = 1
            ins = sav_d.get('insured') or {}
            try:
                ins_age = int(ins.get('age') or 1)
            except Exception:
                pass
            bi_list = sav_d.get('benefit_illustration') or []
            wd_list = sav_d.get('withdrawal_illustration') or []
            no_wd_p = _bi_to_no_wd(bi_list, prod_data=sav_d, insured_age=ins_age)
            wd_p = _wd_to_dict(wd_list, prod_data=sav_d, insured_age=ins_age) if wd_list else {}
            pt_p = _paid_total_from_bi(bi_list)
            pt_wd = _paid_total_from_bi(wd_list) if wd_list else 0
            paid_p = pt_wd if pt_wd > 0 else pt_p
            has_wd_p = bool(wd_p) and any(v.get('Annual_WD', 0) > 0 for v in wd_p.values())
            section = [
                # 2026-07-16 V3.3.9+: per-product deck-level cover (关于我们), 之前只第一个产品有
                ('savings_deck_cover', lambda c=sav_c, m=sm: _slide_company(
                    prs, sav_d, m, C,
                    logo_path=c.get('logo_path') or logo_path,
                    company_imgs=c.get('company_images') or company_images)),
                ('savings_company', lambda c=sav_c, m=sm: _slide_sub_company(prs, c, '储蓄保障', m, C)),
                ('savings_client_info', lambda d=sav_d, c=sav_c, m=sm, pt=paid_p: _slide_client_info(
                    prs, d, c, m, C, 'savings', pt)),
                ('features', lambda nw=no_wd_p, pt=paid_p, m=sm: _slide_features(prs, nw, pt, m, C)),
                ('growth', lambda nw=no_wd_p, m=sm: _slide_growth_chart(prs, nw, m, C)),
                ('no_withdraw_table', lambda nw=no_wd_p, pt=paid_p, m=sm: _slide_no_withdraw_table(prs, nw, pt, m, C)),
            ]
            if has_wd_p:
                section.append(('withdraw_table', lambda w=wd_p, pt=paid_p, m=sm: _slide_withdraw_table(prs, w, pt, m, C)))
            section.append(('compare_chart', lambda nw=no_wd_p, w=wd_p, m=sm: _slide_compare_chart(prs, nw, w, m, C)))
            return section

        if all_extractions and all_extractions.get('savings'):
            # 多储蓄: 遍历每个产品, 每个一份 section
            for sav_ext in all_extractions['savings']:
                sav_d = sav_ext.get('data') or {}
                bp = sav_ext.get('brand_profile') or {}
                # 从 PDF 名推断公司短名 (若 brand_profile 缺失)
                pdf_name = sav_ext.get('pdfName') or ''
                if not bp.get('name_zh'):
                    # 文件名通常为 "公司__产品.pdf" 或 "公司_产品.pdf"
                    inferred = pdf_name.split('__')[0].split('_')[0].strip()
                    bp = {**bp, 'name_zh': inferred or '储蓄险', 'short': inferred or 'SAV', 'short_en': inferred or 'SAV'}
                    sav_ext['brand_profile'] = bp
                if not sav_ext.get('company_short'):
                    sav_ext['company_short'] = bp.get('short') or bp.get('name_zh')
                pol = sav_d.get('policy') or {}
                prod_name = pol.get('product_name') or sav_d.get('product_name') or pdf_name
                sav_c = {
                    'brand_profile': bp,
                    # 储蓄险为香港实体, 用去掉地区前缀的 short (如"宏利"), 不用 name_zh("新加坡宏利", 那是IUL新加坡实体)
                    'name_zh': bp.get('short') or bp.get('name_zh'),
                    'id': sav_ext.get('company_id'),
                    'logo_path': sav_ext.get('logo_path'),
                    'cover_path': sav_ext.get('cover_path'),
                    'company_images': sav_ext.get('company_images') or [],
                    'product_name': prod_name,
                }
                slides_def += _savings_product_section(sav_d, sav_c)
            # V3.3.6: >=2 个储蓄加 3 张对比 slide (divider / overview / chart)
            if len(all_extractions['savings']) >= 2:
                exts = all_extractions['savings']
                # 确保 products 列表每项有 product_name + company_short 给对比函数用
                for e in exts:
                    if not e.get('product_name'):
                        pol = (e.get('data') or {}).get('policy') or {}
                        e['product_name'] = pol.get('product_name') or e.get('pdfName') or '储蓄险'
                    if not e.get('company_short'):
                        bp = e.get('brand_profile') or {}
                        e['company_short'] = bp.get('short') or bp.get('name_zh') or 'SAV'
                # 2026-07-16 V3.3.11+: 对比章节 header 用合并 meta (A VS B)
                sav_cmp_meta = _build_compare_meta(exts, meta)
                slides_def.append(('savings_compare_divider',
                    lambda e=exts, m=sav_cmp_meta: _slide_savings_compare_divider(prs, e, m, C)))
                slides_def.append(('savings_no_wd_compare_table',
                    lambda e=exts, m=sav_cmp_meta: _slide_savings_no_wd_compare_table(prs, e, m, C)))
                slides_def.append(('savings_wd_compare_table',
                    lambda e=exts, m=sav_cmp_meta: _slide_savings_wd_compare_table(prs, e, m, C)))
                slides_def.append(('savings_overview_multi',
                    lambda e=exts, m=sav_cmp_meta: _slide_savings_overview(prs, e, m, C)))
                slides_def.append(('savings_compare_chart',
                    lambda e=exts, m=sav_cmp_meta: _slide_savings_compare_chart(prs, e, m, C)))
        elif has_savings:
            # 单储蓄: 兼容旧逻辑
            slides_def.append(('savings_client_info', lambda: _slide_client_info(
                prs, data, None, meta, C, 'savings', paid_total)))
            slides_def += [
                ('features', lambda: _slide_features(prs, no_wd, paid_total, meta, C)),
                ('growth', lambda: _slide_growth_chart(prs, no_wd, meta, C)),
                ('no_withdraw_table', lambda: _slide_no_withdraw_table(prs, no_wd, paid_total, meta, C)),
            ]
            if has_wd:
                slides_def.append(('withdraw_table', lambda: _slide_withdraw_table(prs, wd, paid_total, meta, C)))
            slides_def.append(('compare_chart', lambda: _slide_compare_chart(prs, no_wd, wd, meta, C)))
        if has_wd and has_savings:
            slides_def.append(('education', lambda: _slide_education(prs, wd, meta, C, scene_images)))
        # V3.3.1+ 多CI产品支持: 每个 CI 一个完整 section, **per-product meta** 覆盖 company_short
        def _ci_meta(ci_c):
            # 合并 deck meta + 产品公司信息, 让 _slide_* 使用正确 company_short
            bp = (ci_c or {}).get('brand_profile') or {}
            return {
                **meta,
                'company_id': (ci_c or {}).get('id') or meta.get('company_id'),
                'company_short': (ci_c or {}).get('name_zh') or bp.get('short') or meta.get('company_short'),
                'company_short_en': bp.get('short_en') or meta.get('company_short_en'),
            }

        def _ci_product_section(ci_d, ci_c):
            ci_meta = _ci_meta(ci_c)
            section = [
                ('ci_divider', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_ci_divider(prs, d, m, C, c)),
                ('ci_client_info', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_client_info(prs, d, c, m, C, 'ci')),
            ]
            if ci_c:
                section.append(('ci_company', lambda c=ci_c, m=ci_meta: _slide_sub_company(prs, c, '危疾保障', m, C)))
            section += [
                ('ci_overview', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_ci_overview(prs, d, m, C, c)),
                ('ci_table', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_ci_data_table(prs, d, m, C, c)),
                ('ci_coverage', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_ci_coverage(prs, d, m, C, c)),
                ('ci_chart', lambda d=ci_d, c=ci_c, m=ci_meta: _slide_ci_premium_chart(prs, d, m, C, c)),
            ]
            return section

        if all_extractions and all_extractions.get('ci'):
            # 多CI: 遍历每个产品, 每个一份 section (含自己的 meta)
            for ci_ext in all_extractions['ci']:
                ci_d = ci_ext.get('data') or {}
                bp = ci_ext.get('brand_profile') or {}
                ci_c = {
                    'brand_profile': bp,
                    'name_zh': bp.get('name_zh') or ci_ext.get('name_zh'),
                    'id': ci_ext.get('company_id'),
                    'logo_path': ci_ext.get('logo_path'),
                    'cover_path': ci_ext.get('cover_path'),
                    'company_images': ci_ext.get('company_images') or [],
                }
                slides_def += _ci_product_section(ci_d, ci_c)
            # V3.3.5: ≥2 个 CI 加对比 slide (插在单产品之后)
            if len(all_extractions['ci']) >= 2:
                exts = all_extractions['ci']
                # 2026-07-16 V3.3.11+: 对比章节 header 用合并 meta (A VS B), 而不是单公司
                cmp_meta = _build_compare_meta(exts, meta)
                slides_def.append(('ci_compare_divider',
                    lambda e=exts, m=cmp_meta: _slide_comparison_divider(prs, ['ci'], m, C)))
                slides_def.append(('ci_overview_multi',
                    lambda e=exts, m=cmp_meta: _slide_ci_overview_multi(prs, e, m, C)))
                slides_def.append(('ci_combined_chart',
                    lambda e=exts, m=cmp_meta: _slide_ci_combined_chart(prs, e, m, C)))
                # V3.3.11+: 删掉旧 _slide_ci_coverage_compare_multi (LLM 抽 coverage_items 经常漏, 表里全是"—"),
                # 只保留 V3.3.10+ 的 _slide_ci_coverage_highlights (手工整理营销内容, 真正可用)
                slides_def.append(('ci_coverage_highlights',
                    lambda e=exts, m=cmp_meta: _slide_ci_coverage_highlights(prs, e, m, C)))
                # 最佳产品单卡推荐 (按 杠杆率 最高)
                def _pick_best_ci(exts=exts):
                    best, best_lev = None, -1
                    for e in exts:
                        d = e.get('data') or {}
                        pol = d.get('policy') or {}
                        cov = float((d.get('summary') or {}).get('sum_insured') or pol.get('sum_insured') or pol.get('basic_sum_insured') or 0)
                        prem = float((d.get('summary') or {}).get('annual_premium') or pol.get('annual_premium') or 0)
                        lev = cov / prem if prem > 0 else 0
                        if lev > best_lev:
                            best, best_lev = e, lev
                    return best
                slides_def.append(('best_ci_card',
                    lambda e=exts, m=cmp_meta: _slide_best_pick_card(prs, _pick_best_ci(e), 'ci', m, C)))
        elif has_ci:
            # 单CI: 兼容旧逻辑
            slides_def += _ci_product_section(ci_data, ci_company)
        if has_ci and has_savings and not has_iul:
            # 收入替代仅对第一CI生效 (用 ci_data / ci_company 单值即可)
            slides_def.append(('income_protection', lambda: _slide_ci_income_protection(prs, ci_data, wd, meta, C, ci_company)))
        # V3.3.1+ 多IUL产品支持: 每个 IUL 一个完整 section, per-product meta
        def _iul_meta(iul_c):
            bp = (iul_c or {}).get('brand_profile') or {}
            return {
                **meta,
                'company_id': (iul_c or {}).get('id') or meta.get('company_id'),
                'company_short': (iul_c or {}).get('name_zh') or bp.get('short') or meta.get('company_short'),
                'company_short_en': bp.get('short_en') or meta.get('company_short_en'),
            }

        def _iul_product_section(iul_d, iul_c):
            iul_meta = _iul_meta(iul_c)
            section = [
                ('iul_divider', lambda d=iul_d, c=iul_c, m=iul_meta: _slide_iul_divider(prs, d, m, C, c)),
                ('iul_client_info', lambda d=iul_d, c=iul_c, m=iul_meta: _slide_client_info(
                    prs, d, c, m, C, 'iul', product_name_override='新加坡IUL')),
            ]
            if iul_c:
                iul_sub_meta = {**iul_meta, 'product_name': (iul_d.get('summary') or {}).get('product_name', '新加坡IUL')}
                section.append(('iul_company', lambda c=iul_c, sm=iul_sub_meta: _slide_sub_company(prs, c, '万用寿险', sm, C)))
            section += [
                ('iul_overview', lambda d=iul_d, c=iul_c, m=iul_meta: _slide_iul_overview(prs, d, m, C, c)),
                ('iul_leverage_chart', lambda d=iul_d, c=iul_c, m=iul_meta: _slide_iul_leverage_chart(prs, d, m, C, c)),
                ('iul_table', lambda d=iul_d, c=iul_c, m=iul_meta: _slide_iul_data_table(prs, d, m, C, c)),
            ]
            return section

        if all_extractions and all_extractions.get('iul'):
            # 多IUL
            for iul_ext in all_extractions['iul']:
                iul_d = iul_ext.get('data') or {}
                bp = iul_ext.get('brand_profile') or {}
                iul_c = {
                    'brand_profile': bp,
                    'name_zh': bp.get('name_zh') or iul_ext.get('name_zh'),
                    'id': iul_ext.get('company_id'),
                    'logo_path': iul_ext.get('logo_path'),
                    'cover_path': iul_ext.get('cover_path'),
                    'company_images': iul_ext.get('company_images') or [],
                }
                slides_def += _iul_product_section(iul_d, iul_c)
            # V3.3.5: ≥2 个 IUL 加 4 张对比 slide (插在单产品之后)
            if len(all_extractions['iul']) >= 2:
                exts = all_extractions['iul']
                # 2026-07-16 V3.3.11+: 对比章节 header 用合并 meta (A VS B)
                iul_cmp_meta = _build_compare_meta(exts, meta)
                slides_def.append(('iul_compare_divider',
                    lambda e=exts, m=iul_cmp_meta: _slide_comparison_divider(prs, ['iul'], m, C)))
                slides_def.append(('iul_overview_multi',
                    lambda e=exts, m=iul_cmp_meta: _slide_iul_overview_multi(prs, e, m, C)))
                slides_def.append(('iul_combined_chart',
                    lambda e=exts, m=iul_cmp_meta: _slide_iul_combined_chart(prs, e, m, C)))
                slides_def.append(('iul_feature_compare',
                    lambda e=exts, m=iul_cmp_meta: _slide_iul_feature_compare_multi(prs, e, m, C)))
                def _pick_best_iul(exts=exts):
                    best, best_lev = None, -1
                    for e in exts:
                        d = e.get('data') or {}
                        pol = d.get('policy') or {}
                        cov = float(pol.get('sum_insured') or (d.get('summary') or {}).get('sum_insured') or 0)
                        prem = float(pol.get('annual_premium') or (d.get('summary') or {}).get('annual_premium') or 0)
                        pay_raw = pol.get('premium_payment_period') or (d.get('summary') or {}).get('payment_years') or 10
                        if '趸交' in str(pay_raw):
                            py = 1
                        else:
                            digits = ''.join(cc for cc in str(pay_raw) if cc.isdigit())
                            py = int(digits) if digits else 10
                        total_prem = prem * py
                        lev = cov / total_prem if total_prem > 0 else 0
                        if lev > best_lev:
                            best, best_lev = e, lev
                    return best
                slides_def.append(('best_iul_card',
                    lambda e=exts: _slide_best_pick_card(prs, _pick_best_iul(e), 'iul', meta, C)))
        elif has_iul:
            slides_def += _iul_product_section(iul_data, iul_company)
        # 方案协同页放在IUL方案摘要后面
        product_count = sum([has_savings, has_ci, has_iul])
        if product_count >= 2:
            slides_def.append(('combo', lambda: _slide_combo_narrative(prs, meta, C, ci_data, iul_data)))
        if has_iul and has_savings and not has_ci:
            slides_def.append(('savings_iul_funding', lambda: _slide_savings_iul_premium_funding(prs, wd, iul_data, meta, C, iul_company)))
        slides_def.append(('disclaimer', lambda: _slide_disclaimer(prs, meta, C)))
        is_ci_only = has_ci and not has_savings and not has_iul
        if not has_iul and not is_ci_only:
            slides_def.append(('sources', lambda: _slide_sources(prs, data, meta, C)))
        if meta.get("ai_narrative", "") or not is_ci_only:
            slides_def.append(('summary', lambda: _slide_summary(prs, no_wd, wd, paid_total, meta, C)))
        if has_ci or has_iul:
            slides_def.append(('ending_combined', lambda: _slide_ending_combined(prs, summary, ci_data, iul_data, meta, C)))
        else:
            slides_def.append(('ending', lambda: _slide_ending(prs, summary, meta, C)))

    total_slides = len(slides_def)
    meta['_total_slides'] = total_slides

    for idx, (_, slide_fn) in enumerate(slides_def, 1):
        meta['_page_num'] = idx
        slide_fn()

    prs.save(output_path)
    return output_path


# ══════════════════════════════════════════════════════
#  各页 slide 函数
# ══════════════════════════════════════════════════════

def _add_overlay(slide, color_rgb, alpha=50):
    """添加半透明覆盖层"""
    from pptx.oxml.ns import qn
    from lxml import etree
    overlay = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color_rgb
    # 设置透明度 (alpha = 0-100, 0=全透, 100=全实)
    spPr = overlay._element.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha_tag = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha_tag.set('val', str(int((100 - alpha) * 1000)))
    overlay.line.fill.background()
    return overlay

def _slide_cover(prs, summary, paid_total, meta, C, cover_img=None):
    """封面: 储蓄=顶部16:9图+底部文字, IUL=顶部banner+底部深色"""
    prod_type = (meta.get('product_type', '') or '').lower()
    is_iul = prod_type == 'iul'
    is_savings = prod_type == 'savings'
    s = prs.slides.add_slide(prs.slide_layouts[6])
    white = RGBColor(0xFF, 0xFF, 0xFF)

    if is_iul:
        # IUL: 顶部banner图 + 底部深色背景
        fill_bg(s, C['primary_dark'])
        banner_h = Inches(4.5)
        if cover_img and os.path.exists(cover_img):
            _add_picture(s, cover_img, Emu(0), Emu(0), SLIDE_W, banner_h)
    elif is_savings and cover_img and os.path.exists(cover_img):
        # 储蓄险: 16:9图片占顶部~60%, 底部自动取色放文字
        bottom_color = _extract_bottom_color(cover_img)
        fill_bg(s, bottom_color or C['primary_dark'])
        img_h = Inches(4.5)  # 顶部图片高度
        _add_picture(s, cover_img, Emu(0), Emu(0), SLIDE_W, img_h)
    else:
        # CI/兜底: 全屏图 + 遮罩
        if cover_img and os.path.exists(cover_img):
            _add_picture(s, cover_img, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
        else:
            fill_bg(s, C['primary_dark'])
        _add_overlay(s, RGBColor(0x0D, 0x1B, 0x2A), alpha=50)

    name = summary.get('insured_name', 'VIP 先生') or 'VIP 先生'
    prod = (meta.get('product_name') or summary.get('product_name') or '储蓄寿险计划').replace('\n', '')
    pay_years = summary.get('payment_years', 5)
    annual = summary.get('annual_premium') or 100000
    currency = summary.get('currency', 'USD')
    insured_age = summary.get('insured_age', 1)

    if is_iul:
        base_y = Inches(4.8)
    elif is_savings:
        # 储蓄险: 文字在底部(图片下方)
        base_y = Inches(4.6)
    else:
        base_y = Inches(2.1)

    add_text(s, Inches(0.5), base_y, Inches(6.5), Inches(0.35),
             prod, size=22, color=white, bold=True)

    gender = meta.get('insured_gender', '')
    prefix = 'MR.' if gender in ('Male', '男') else 'MS.'
    add_text(s, Inches(0.5), base_y + Inches(0.5), Inches(6.5), Inches(0.4),
             f'{prefix}{name}', size=22, color=white, bold=True)

    add_text(s, Inches(0.5), base_y + Inches(1.0), Inches(6.5), Inches(0.35),
             'PRIVATE CLIENT', size=22, color=white)

    add_text(s, Inches(0.5), base_y + Inches(1.5), Inches(6.5), Inches(0.35),
             f'{pay_years}年缴 · {currency} {annual:,.0f}', size=22, color=white, bold=True)

    add_text(s, Inches(0.5), base_y + Inches(2.0), Inches(6.5), Inches(0.4),
             f'受保人：{name}（{insured_age}岁）',
             size=22, color=white, bold=True)

    # 右侧使用纯装饰图（文字由封面图自带）


def _slide_company(prs, data, meta, C, logo_path=None, company_imgs=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    bp = meta.get('brand_profile', {})
    company_zh = meta.get('company_name_zh', '保险公司')
    add_title(s, '关于我们', f'About Us  ·  {company_zh}',
              title_size=28, sub_size=11, C=C)

    # 叙事公司介绍 (如有)
    nar_intro = meta.get('narrative_company_intro', '')
    if nar_intro:
        add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.4),
                 nar_intro, size=10, color=C['mid_text'])

    # Logo (正方形) + 公司图片区
    if logo_path and os.path.exists(logo_path):
        # 按原始宽高比显示，高度固定1.2in
        try:
            from PIL import Image
            with Image.open(logo_path) as img:
                w, h = img.size
                logo_w = Inches(1.2 * w / h)
                _add_picture(s, logo_path, Inches(0.5), Inches(2.0), logo_w, Inches(1.2))
        except Exception:
            _add_picture(s, logo_path, Inches(0.5), Inches(2.0), Inches(1.2), Inches(1.2))

    imgs = company_imgs or []
    if imgs:
        # 三张公司图片靠右排列，避免盖住左侧LOGO
        img_count = min(len(imgs), 3)
        img_start = SLIDE_W - Inches(2.2) - (img_count - 1) * Inches(2.1) - Inches(0.5)
        img_x = img_start
        for ip in imgs[:img_count]:
            _add_picture(s, ip, img_x, Inches(2.0), Inches(2.2), Inches(1.4))
            img_x += Inches(2.1)

    # 三张信息卡
    cards = [
        (bp.get('founded_year', '—'), bp.get('founded_label', '成立年份'), bp.get('founded_sub', '')),
        (bp.get('rating_value', '—'), bp.get('rating_label', '财务实力评级'), bp.get('rating_sub', '')),
        (meta.get('product_name', bp.get('series_label', '—')), bp.get('series_sub', ''), bp.get('series_products', '')),
    ]
    cy = Inches(3.6)
    cw = Inches(4.0); ch = Inches(1.5); gap = Inches(0.25); sx = Inches(0.5)
    for i, (val, lbl, sub) in enumerate(cards):
        x = sx + i * (cw + gap)
        add_rect(s, x, cy, cw, ch, fill=C['bg_card'])
        add_rect(s, x, cy + ch - Inches(0.08), cw, Inches(0.08), fill=C['primary'])
        add_text(s, x + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4),
                 Inches(0.7), str(val), size=26, color=C['dark_text'], bold=True)
        add_text(s, x + Inches(0.2), cy + Inches(1.1), cw - Inches(0.4),
                 Inches(0.3), str(lbl), size=12, color=C['body_text'], bold=True)
        add_text(s, x + Inches(0.2), cy + Inches(1.25), cw - Inches(0.4),
                 Inches(0.2), str(sub), size=8, color=C['mid_text'])

    add_rect(s, Inches(0.5), Inches(5.4), Inches(0.3), Inches(0.04), fill=C['primary'])
    add_text(s, Inches(0.85), Inches(5.3), Inches(5), Inches(0.4),
             '业务范围', size=14, color=C['dark_text'], bold=True)
    for i, t in enumerate((bp.get('business_lines', []) or [])[:4]):
        add_text(s, Inches(0.85), Inches(5.7 + i*0.3), Inches(5.5),
                 Inches(0.3), t, size=11, color=C['body_text'])

    add_rect(s, Inches(7.0), Inches(5.4), Inches(0.3), Inches(0.04), fill=C['primary'])
    add_text(s, Inches(7.35), Inches(5.3), Inches(5), Inches(0.4),
             '品牌背景', size=14, color=C['dark_text'], bold=True)
    for i, t in enumerate((bp.get('brand_background', []) or [])[:4]):
        add_text(s, Inches(7.35), Inches(5.7 + i*0.3), Inches(5.5),
                 Inches(0.3), t, size=11, color=C['body_text'])

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             f'数据来源：{bp.get("data_source", "保险公司官方资料")}',
             size=8, color=C['mid_text'])


def _slide_features(prs, no_wd, paid_total, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    # 叙事标题 (如有)
    nar_title = meta.get('narrative_title', '')
    feat_title = nar_title if nar_title and len(nar_title) < 30 else '产品亮点'
    pay_yrs = meta.get('payment_years', 5) or 5
    add_title(s, feat_title, f'{pay_yrs} 年缴 · 短期供款完成 · 终身财富复利',
              title_size=28, sub_size=11, C=C)

    paid = paid_total
    pb, y20, y30, _ = _find_milestone(no_wd, paid)

    items = [
        (f'{pay_yrs} 年短期缴付', f'只需 {pay_yrs} 年完成缴费，\n后续无需再缴；缓解供款压力', str(pay_yrs), '年缴清'),
    ]
    if pb:
        items.append((f'第 {pb} 年快速回本',
                      f'保单年度终结时退保现价\n超过已缴保费总额，本金安全\nY{pb} 现价 USD {no_wd[pb]["Total"]:,}',
                      f'Y{pb}', '回本'))
    if y20:
        items.append((f'{y20} 年财富复利',
                      f'保证 + 非保证现金价值\n长期复利，稳健增长\nY{y20} 约 USD {no_wd[y20]["Total"]:,}',
                      f'{no_wd[y20]["Mult"]:.2f}x', f'{y20} 年倍数'))
    if y30:
        items.append((f'{y30} 年财富复利',
                      f'持续滚存跨越复利临界点\n增值潜力可观\nY{y30} 约 USD {no_wd[y30]["Total"]:,}',
                      f'{no_wd[y30]["Mult"]:.2f}x', f'{y30} 年倍数'))

    cw = Inches(6.1); ch = Inches(2.2)
    gx = Inches(0.25); gy = Inches(0.2)
    sx = Inches(0.5); sy = Inches(2.4)
    for i, (title, body, kpi, unit) in enumerate(items[:4]):
        col = i % 2; row = i // 2
        x = sx + col * (cw + gx); y = sy + row * (ch + gy)
        no_text = f'{i+1:02d}'
        add_rect(s, x, y, cw, ch, fill=C['bg_card'])
        add_rect(s, x, y, Inches(0.1), ch, fill=C['primary'])
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(1.2), Inches(0.7),
                 no_text, size=28, color=C['primary'], bold=True)
        add_text(s, x + Inches(1.4), y + Inches(0.3), Inches(4.4), Inches(0.5),
                 title, size=15, color=C['dark_text'], bold=True)
        add_text(s, x + Inches(0.3), y + Inches(0.95), cw - Inches(0.5), Inches(1.1),
                 body, size=10, color=C['body_text'])
        add_text(s, x + Inches(0.3), y + Inches(1.65), cw - Inches(0.5), Inches(0.4),
                 kpi, size=22, color=C['primary'], bold=True, align=PP_ALIGN.RIGHT)
        add_text(s, x + Inches(0.3), y + Inches(1.95), cw - Inches(0.5), Inches(0.2),
                 unit, size=9, color=C['mid_text'], align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.3),
             '注：倍数基于不退保情形下「退保发还金额总额」÷「已缴保费总额」计算',
             size=8, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.3), Inches(12.5), Inches(0.2),
             '数据来源：保险公司官方计划书 第3部分', size=8, color=C['mid_text'])


def _find_milestone(no_wd, threshold_paid):
    """动态找回本年/20年/30年"""
    if not no_wd: return (None, None, None, None)
    def _get(y):
        k = str(y) if y in no_wd else y
        return no_wd.get(k) or no_wd.get(str(y))
    ys = sorted(int(k) for k in no_wd.keys())
    payback = next((y for y in ys if (_get(y) or {}).get('Total', 0) > threshold_paid), None)
    y20 = next((y for y in ys if y >= 20), None)
    y30 = next((y for y in ys if y >= 30), None)
    y70 = next((y for y in ys if y >= 70), None)
    return (payback, y20, y30, y70)


def _get_data_for_chart(no_wd, insured_age=1, max_year=30):
    """提取图表数据: 只显示前30年"""
    all_yrs = sorted(int(y) for y in no_wd.keys() if no_wd[y].get('Total', 0) > 0)
    all_yrs = [y for y in all_yrs if y <= max_year]
    yrs = [y for y in all_yrs if y <= 30]
    if not yrs: yrs = all_yrs[:15]
    guar = [no_wd[y]['Guar_CV'] for y in yrs]
    non_g = [no_wd[y]['Total'] - no_wd[y]['Guar_CV'] for y in yrs]
    insured_age_n = insured_age
    cat_labels = []
    for y in yrs:
        r = no_wd.get(y) or no_wd.get(str(y))
        age = r.get('Age', insured_age_n + y - 1) if r else insured_age_n + y - 1
        cat_labels.append(f'Y{y}\n({int(age)}岁)')
    return yrs, guar, non_g, cat_labels


def _slide_growth_chart(prs, no_wd, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    # 无数据时跳过图表，避免 python-pptx 崩溃
    if not no_wd or not any(v.get('Total', 0) > 0 for v in no_wd.values()):
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '暂无利益演示数据', size=16, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return
    add_title(s, '现金价值长期复利增长（不提领）',
              '保证 + 非保证现金价值 · 演示至保单年度 80 年（受保人 80 岁）',
              title_size=36, sub_size=14, C=C)

    insured_age = meta.get('insured_age', 1) or 1
    yrs, guar, non_g, cat_labels = _get_data_for_chart(no_wd, insured_age)

    # ── 左: 柱状图 (堆叠, 宽5.8) ──
    cd = CategoryChartData()
    cd.categories = cat_labels
    cd.add_series('保证现价 (USD)', guar)
    cd.add_series('非保证红利/分红 (USD)', non_g)

    graphic = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.3), Inches(1.85), Inches(6.0), Inches(4.8), cd
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    for plot in chart.plots:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(C['primary']) if i == 0 else _c(C['accent'])
    va = chart.value_axis
    va.tick_labels.font.size = Pt(7)
    va.tick_labels.font.color.rgb = _c(C['mid_text'])
    va.tick_labels.number_format = '#,##0,,'
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(7)
    ca.tick_labels.font.name = FONT_HEI
    ca.tick_labels.font.color.rgb = _c(C['mid_text'])
    chart.legend.font.size = Pt(8)
    chart.legend.font.name = FONT_HEI
    chart.legend.font.color.rgb = _c(C['body_text'])

    # ── 右: 折线图 (退保总额趋势, 宽5.8) ──
    total_vals = [no_wd[y]['Total'] if y is not None else 0 for y in yrs]
    cd_line = CategoryChartData()
    cd_line.categories = cat_labels
    cd_line.add_series('退保发还总额 (USD)', total_vals)

    graphic2 = s.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(6.6), Inches(1.85), Inches(6.0), Inches(4.8), cd_line
    )
    lc = graphic2.chart
    lc.has_title = False
    lc.has_legend = True
    lc.legend.position = XL_LEGEND_POSITION.BOTTOM
    lc.legend.include_in_layout = False
    lc.legend.font.size = Pt(8)
    lc.legend.font.name = FONT_HEI
    lc.legend.font.color.rgb = _c(C['body_text'])
    for plot in lc.plots:
        for ser in plot.series:
            ser.format.line.color.rgb = _c(C['accent_dark'])
            ser.format.line.width = Pt(2.5)
            try:
                ser.marker.style = XL_MARKER_STYLE.CIRCLE
                ser.marker.size = 6
            except Exception:
                pass
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(C['accent_dark'])
    lv = lc.value_axis
    lv.tick_labels.font.size = Pt(7)
    lv.tick_labels.font.color.rgb = _c(C['mid_text'])
    lv.tick_labels.number_format = '#,##0,,'
    lc_ax = lc.category_axis
    lc_ax.tick_labels.font.size = Pt(7)
    lc_ax.tick_labels.font.name = FONT_HEI
    lc_ax.tick_labels.font.color.rgb = _c(C['mid_text'])

    # ── 下方 callout 行 ──
    paid = next((no_wd[y]['Paid'] for y in yrs if y and no_wd[y].get('Paid', 0) > 0), 0)
    pb, y20, y30, y70 = _find_milestone(no_wd, paid)
    callouts = []
    if pb:
        callouts.append(('回本年', f'Y{pb}', f'USD {no_wd[pb]["Total"]:,}', C['primary']))
    if y20:
        callouts.append((f'{y20}年', f'x{no_wd[y20]["Mult"]:.2f}', f'USD {no_wd[y20]["Total"]:,}', C['primary_dark']))
    if y30:
        callouts.append((f'{y30}年', f'x{no_wd[y30]["Mult"]:.2f}', f'USD {no_wd[y30]["Total"]:,}', C['primary_dark']))
    if y70:
        callouts.append((f'{y70}年', f'x{no_wd[y70]["Mult"]:.0f}', f'USD {no_wd[y70]["Total"]:,}', C['accent_dark']))

    cx = Inches(0.3)
    cw = Inches(3.0)
    cgap = Inches(0.2)
    for i, (yr, val, sub, col) in enumerate(callouts[:4]):
        x = cx + i * (cw + cgap)
        add_rect(s, x, Inches(6.75), cw, Inches(0.55), fill=C['bg_card'])
        add_rect(s, x, Inches(6.75), Inches(0.06), Inches(0.55), fill=col)
        add_text(s, x + Inches(0.1), Inches(6.78), Inches(0.8), Inches(0.25),
                 yr, size=9, color=C['dark_text'], bold=True)
        add_text(s, x + Inches(0.85), Inches(6.78), Inches(1.0), Inches(0.25),
                 val, size=12, color=col, bold=True)
        add_text(s, x + Inches(0.1), Inches(6.98), cw - Inches(0.15), Inches(0.2),
                 sub, size=7, color=C['mid_text'])

    add_text(s, Inches(0.3), Inches(7.35), Inches(12.5), Inches(0.2),
             f'数据来源：{meta.get("company_short", "保险公司")}官方计划书 第3部分',
             size=7, color=C['mid_text'])


def _slide_no_withdraw_table(prs, no_wd, paid_total, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    # 动态生成: 40年前每5年, 40年后每10年
    max_y = max(int(k) for k in no_wd.keys()) if no_wd else 70
    show_yrs = list(range(5, min(41, max_y + 1), 5))
    show_yrs += list(range(50, max_y + 1, 10))
    if 1 in no_wd: show_yrs = [1] + show_yrs
    show_yrs = [y for y in show_yrs if y in no_wd]
    show_yrs = show_yrs[:15]
    subtitle_text = f'至保单年度 {max_y} · 40年前每5年/后每10年 · 含 IRR'
    add_title(s, '不提领情形 · 退保发还金额明细',
              subtitle_text,
              title_size=36, sub_size=14, C=C)
    headers = ['保单年度', '年龄', '已缴保费\n(USD)', '保证现价\n(USD)',
               '非保证金额\n(USD)', '退保发还总额\n(USD)', '复利 IRR', '单利IRR']
    rows = []
    for y in show_yrs:
        r = no_wd[y]
        irr = f"{r['IRR']*100:.2f}%" if r.get('IRR') else '-'
        simple = f"{r['Simple']*100:.2f}%" if r.get('Simple') else '-'
        rows.append((y, r.get('Age', y), r.get('Paid', paid_total), r['Guar_CV'],
                     r['Total'] - r['Guar_CV'], r['Total'], irr, simple))

    tx = Inches(0.5); ty = Inches(1.95)
    tw = Inches(12.333); th = Inches(4.0)
    n_rows = len(rows) + 1
    ts = s.shapes.add_table(n_rows, len(headers), tx, ty, tw, th)
    table = ts.table
    cw_target = [0.9, 0.7, 1.3, 1.5, 1.5, 1.7, 1.1, 1.1]
    total_w = sum(cw_target)
    for i, w in enumerate(cw_target):
        table.columns[i].width = Inches(w * 12.333 / total_w)
    table.rows[0].height = Inches(0.55)
    for i in range(1, n_rows):
        table.rows[i].height = Inches(0.36)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
        tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(9); run.font.name = FONT_HEI
        run.font.color.rgb = _c(C['bg_light']); run.font.bold = True
        rPr = run._r.get_or_add_rPr()
        eaFont = rPr.find(qn('a:ea'))
        if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
        eaFont.set('typeface', FONT_HEI)

    for r, row in enumerate(rows, start=1):
        bg = C['bg_card'] if r % 2 == 1 else C['bg_light']
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top = Emu(12000); tf.margin_bottom = Emu(12000)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            if c < 2:
                txt = str(val); color = C['dark_text']; bold = True
            elif c == 2:
                txt = f'{val:,}'; color = C['body_text']; bold = False
            elif c in [3, 4, 5]:
                txt = f'{val:,}'
                color = C['primary'] if c == 3 else (C['accent_dark'] if c == 4 else C['dark_text'])
                bold = (c == 5)
            else:
                txt = str(val)
                color = C['primary'] if c == 6 else C['mid_text']
                bold = (c == 6)
            run = p.add_run(); run.text = txt
            run.font.size = Pt(9); run.font.name = FONT_HEI
            run.font.color.rgb = _c(color); run.font.bold = bold
            rPr = run._r.get_or_add_rPr()
            eaFont = rPr.find(qn('a:ea'))
            if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
            eaFont.set('typeface', FONT_HEI)

    add_rect(s, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.85), fill=C['bg_card'])
    add_rect(s, Inches(0.5), Inches(6.2), Inches(0.08), Inches(0.85), fill=C['primary'])
    add_text(s, Inches(0.7), Inches(6.25), Inches(5), Inches(0.3),
             '数据口径说明', size=11, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.25),
             f'· 已缴保费：{meta.get("payment_years", 5)}年×计划 = USD {paid_total:,.0f}    · 复利IRR = (总额/保费)^(1/年数) - 1    · 单利IRR = (总额-保费)/保费/年数',
             size=9, color=C['body_text'])
    add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.25),
             '· 非保证金额含复归红利+终期分红    · 演示至保单80年；计划保障至128岁',
             size=9, color=C['body_text'])
    add_text(s, Inches(0.5), Inches(7.18), Inches(12.5), Inches(0.25),
             f'数据来源：{meta.get("company_short", "保险公司")}官方计划书 第3部分',
             size=8, color=C['mid_text'])


def _slide_withdraw_table(prs, wd, paid_total, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    wd_start = None
    for y in sorted(int(k) for k in wd.keys()):
        if wd[y].get('Annual_WD', 0) > 0:
            wd_start = y
            break
    wd_start = wd_start or 7
    annual_wd = wd.get(wd_start, {}).get('Annual_WD', 0) or 0
    insured_age = meta.get('insured_age', 0) or 0
    wd_title = '提领方案 · 退保发还金额及累计提取'
    wd_sub = f'保单年度 {wd_start}（{insured_age + wd_start - 1} 岁）起 · 每年提取 USD {annual_wd:,} · 含 IRR'
    add_title(s, wd_title, wd_sub, title_size=36, sub_size=14, C=C)

    # 动态: 40年前每5年, 后每10年; 确保提领起始年和最后一年在列
    max_y = max(int(k) for k in wd.keys()) if wd else 80
    preferred = set(range(5, min(41, max_y + 1), 5))
    preferred |= set(range(50, max_y + 1, 10))
    preferred.add(wd_start)
    if wd_start > 1: preferred.add(wd_start - 1)
    preferred.add(max_y)
    show_yrs = sorted(y for y in preferred if y in wd)
    if len(show_yrs) < 6:
        for y in sorted(int(k) for k in wd.keys()):
            if y not in show_yrs:
                show_yrs.append(y)
            if len(show_yrs) >= 12: break
    show_yrs = sorted(show_yrs)
    headers = ['保单\n年度', '年龄', '已缴保费', '年提取', '累计提取',
               '退保发还总额\n(B+C+D)', '累计已领+\n退保现价', '复利\nIRR', '单利\nIRR']
    rows = []
    for y in show_yrs:
        r = wd[y]
        irr = f"{r['IRR']*100:.2f}%" if r.get('IRR') else '-'
        simple = f"{r['Simple']*100:.2f}%" if r.get('Simple') else '-'
        rows.append((y, r['Age'], r.get('Paid', paid_total), r['Annual_WD'], r['Cum_WD'],
                     r['Total'], r.get('Total_Received', r['Cum_WD']+r['Total']),
                     irr, simple))

    tx = Inches(0.4); ty = Inches(1.9)
    tw = Inches(12.5); th = Inches(4.6)
    n_rows = len(rows) + 1
    ts = s.shapes.add_table(n_rows, len(headers), tx, ty, tw, th)
    table = ts.table
    cw_target = [0.8, 0.65, 1.15, 1.0, 1.2, 1.4, 1.55, 1.0, 1.0]
    total_w = sum(cw_target)
    for i, w in enumerate(cw_target):
        table.columns[i].width = Inches(w * 12.5 / total_w)
    table.rows[0].height = Inches(0.5)
    for i in range(1, n_rows):
        table.rows[i].height = Inches(0.32)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Emu(18000); tf.margin_right = Emu(18000)
        tf.margin_top = Emu(12000); tf.margin_bottom = Emu(12000)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(8.5); run.font.name = FONT_HEI
        run.font.color.rgb = _c(C['bg_light']); run.font.bold = True
        rPr = run._r.get_or_add_rPr()
        eaFont = rPr.find(qn('a:ea'))
        if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
        eaFont.set('typeface', FONT_HEI)

    for r, row in enumerate(rows, start=1):
        bg = C['bg_card'] if r % 2 == 1 else C['bg_light']
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Emu(12000); tf.margin_right = Emu(12000)
            tf.margin_top = Emu(6000); tf.margin_bottom = Emu(6000)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            if c < 2:
                txt = str(val); color = C['dark_text']; bold = True
            elif c == 2:
                txt = f'{val:,}'; color = C['body_text']; bold = False
            elif c == 3:
                txt = f'{val:,}' if val else '0'; color = C['primary_dark']; bold = False
            elif c == 4:
                txt = f'{val:,}'; color = C['primary_dark']; bold = True
            elif c == 5:
                txt = f'{val:,}'; color = C['dark_text']; bold = True
            elif c == 6:
                txt = f'{val:,}'; color = C['accent_dark']; bold = True
            else:
                txt = str(val)
                color = C['primary'] if c == 7 else C['mid_text']
                bold = (c == 7)
            run = p.add_run(); run.text = txt
            run.font.size = Pt(8); run.font.name = FONT_HEI
            run.font.color.rgb = _c(color); run.font.bold = bold
            rPr = run._r.get_or_add_rPr()
            eaFont = rPr.find(qn('a:ea'))
            if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
            eaFont.set('typeface', FONT_HEI)

    add_rect(s, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.55), fill=C['bg_card'])
    add_rect(s, Inches(0.4), Inches(6.65), Inches(0.08), Inches(0.55), fill=C['primary'])
    add_text(s, Inches(0.6), Inches(6.7), Inches(3), Inches(0.3),
             '数据口径', size=10, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.25),
             f'· Y{wd_start} 起每年提取 USD {annual_wd:,}    · 累计已领+退保现价 = 年提取累计 + 期末退保发还金额    · 复利IRR = (累计+退保现价)/保费^(1/年数)-1',
             size=8.5, color=C['body_text'])
    add_text(s, Inches(0.4), Inches(7.28), Inches(12.5), Inches(0.2),
             f'数据来源：{meta.get("company_short", "保险公司")}官方计划书 第5部分（现金提取）',
             size=8, color=C['mid_text'])


def _render_milestone_table(slide, no_wd, show_yrs, C, meta):
    """无提领数据时展示关键年限数据卡片"""
    inset_x = Inches(1.0)
    card_w = Inches(3.6)
    card_h = Inches(2.8)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    cy = Inches(2.2)
    for i, y in enumerate(show_yrs[:8]):
        col = i % 3
        row = i // 3
        x = inset_x + col * (card_w + gap_x)
        y_pos = cy + row * (card_h + gap_y)
        r = no_wd[y]
        add_rect(slide, x, y_pos, card_w, card_h, fill=C['bg_card'])
        add_rect(slide, x, y_pos, card_w, Inches(0.06), fill=C['primary'])
        age_val = r.get('Age', meta.get('insured_age', 1) + y - 1) if r else meta.get('insured_age', 1) + y - 1
        add_text(slide, x + Inches(0.2), y_pos + Inches(0.15), card_w - 0.4, Inches(0.3),
                 f'保单年度 {y}    |    {int(age_val)} 岁', size=12, color=C['dark_text'], bold=True)
        mult = r.get('Mult', 0)
        irr = r.get('IRR', 0)
        add_text(slide, x + Inches(0.2), y_pos + Inches(0.6), card_w - 0.4, Inches(0.3),
                 f'退保发还总额', size=9, color=C['mid_text'])
        add_text(slide, x + Inches(0.2), y_pos + Inches(0.85), card_w - 0.4, Inches(0.45),
                 f'USD {r["Total"]:,}', size=22, color=C['primary'], bold=True)
        add_text(slide, x + Inches(0.2), y_pos + Inches(1.45), card_w - 0.4, Inches(0.25),
                 f'倍数: {mult:.2f}x    |    IRR: {irr*100:.2f}%' if irr else f'倍数: {mult:.2f}x',
                 size=9, color=C['body_text'])
        add_text(slide, x + Inches(0.2), y_pos + Inches(1.8), card_w - 0.4, Inches(0.4),
                 f'保证: USD {r["Guar_CV"]:,}    非保证: USD {r["Total"] - r["Guar_CV"]:,}',
                 size=8, color=C['mid_text'])
        bar_ratio = r['Guar_CV'] / max(r['Total'], 1)
        bar_w = (card_w - 0.4) * 0.85
        add_rect(slide, x + Inches(0.2), y_pos + Inches(2.3), Inches(bar_w * 0.85), Inches(0.12), fill=C['accent_light'])
        add_rect(slide, x + Inches(0.2), y_pos + Inches(2.3), Inches(bar_w * min(bar_ratio, 1) * 0.85), Inches(0.12), fill=C['primary'])


def _slide_compare_chart(prs, no_wd, wd, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    # 无提领数据时展示不提领增长总览
    has_wd = bool(wd) and any(v.get('Annual_WD', 0) > 0 for v in wd.values())
    if not has_wd:
        add_title(s, '退保发还金额增长趋势',
                  '保证 + 非保证现金价值 · 关键年限数据',
                  title_size=36, sub_size=14, C=C)
        show_yrs = [y for y in [5, 10, 15, 20, 25, 30, 40, 50, 60, 70] if y in no_wd]
        _render_milestone_table(s, no_wd, show_yrs, C, meta)
        return

    add_title(s, '不提领 vs 提领 · 总收益对比',
              '累计已领 + 期末退保现价 · 关键年限对比',
              title_size=36, sub_size=14, C=C)

    common_yrs = sorted(set(int(y) for y in no_wd.keys() if no_wd[y].get('Total', 0) > 0) &
                        set(int(y) for y in wd.keys() if wd[y].get('Total', 0) > 0 or wd[y].get('Annual_WD', 0) > 0))
    common_yrs = [y for y in common_yrs if y <= 80]
    front = [y for y in common_yrs if y <= 30]
    tail = [y for y in common_yrs if y > 30]
    tail_filtered = []
    for y in tail:
        if not tail_filtered or y - tail_filtered[-1] >= 10:
            tail_filtered.append(y)
    if len(front) + len(tail_filtered) > 18:
        front = [y for y in front if y <= 5 or y in (10, 15, 20, 25, 30)]
        tail_filtered = tail_filtered[:3]
    has_gap = bool(front) and bool(tail_filtered) and (tail_filtered[0] - front[-1] > 1)
    if has_gap:
        yrs = front + [None] + tail_filtered
    else:
        yrs = front + tail_filtered
    # 关键: 用实际年龄, 不是保单年度
    insured_age_cc = meta.get('insured_age', 1) or 1
    def _cc_age(y):
        if y is None: return None
        r = no_wd.get(y) or wd.get(y) or {}
        return r.get('Age', insured_age_cc + y - 1) if r else insured_age_cc + y - 1
    cat_labels = []
    for y in yrs:
        if y is None:
            cat_labels.append('…')
        else:
            age = _cc_age(y)
            cat_labels.append(f'Y{y}\n({int(age)}岁)')
    real_yrs = [y for y in yrs if y is not None]
    no_tot = [no_wd[y]['Total'] for y in real_yrs]
    w_tot = [wd[y].get('Total_Received', wd[y]['Cum_WD']+wd[y]['Total']) for y in real_yrs]

    cd = CategoryChartData()
    cd.categories = cat_labels
    cd.add_series('不提领 · 退保现价 (USD)', no_tot)
    cd.add_series('提领 · 累计+退保现价 (USD)', w_tot)

    graphic = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.85), Inches(7.5), Inches(5.0), cd
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    # 线性刻度: 最大值 = 最高柱 × 2
    va = chart.value_axis
    all_vals = [v for v in no_tot + w_tot if v > 0]
    if all_vals:
        va.maximum_scale = max(all_vals) * 2.0
        va.minimum_scale = 0
    va.tick_labels.font.size = Pt(7)
    va.tick_labels.number_format = '$#,##0,, "M"'
    va.tick_labels.font.color.rgb = _c(C['mid_text'])
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(8)
    ca.tick_labels.font.name = FONT_HEI
    ca.tick_labels.font.color.rgb = _c(C['mid_text'])
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = FONT_HEI
    chart.legend.font.color.rgb = _c(C['body_text'])
    for plot in chart.plots:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(C['primary']) if i == 0 else _c(C['accent'])

    # 右侧解读文字
    cxr = Inches(8.1); rw = Inches(5.0)
    add_rect(s, cxr, Inches(1.9), Inches(0.3), Inches(0.04), fill=C['primary'])
    add_text(s, cxr + Inches(0.5), Inches(1.85), rw - 0.5, Inches(0.4),
             '不提领 vs 提领', size=14, color=C['dark_text'], bold=True)

    # 说不提领的好处
    add_rect(s, cxr, Inches(2.5), rw, Inches(1.8), fill=C['bg_card'])
    add_rect(s, cxr, Inches(2.5), Inches(0.08), Inches(1.8), fill=C['primary'])
    add_text(s, cxr + Inches(0.3), Inches(2.6), rw - 0.4, Inches(0.3),
             '📈 不提领策略', size=12, color=C['dark_text'], bold=True)
    add_text(s, cxr + Inches(0.3), Inches(2.9), rw - 0.4, Inches(1.2),
             '账户价值持续滚存，享受复利增长\n适合：长期财富传承需求\n无需现金流补充的场景',
             size=8.5, color=C['body_text'])

    # 说提领的好处
    add_rect(s, cxr, Inches(4.5), rw, Inches(1.8), fill=C['bg_card'])
    add_rect(s, cxr, Inches(4.5), Inches(0.08), Inches(1.8), fill=C['accent'])
    add_text(s, cxr + Inches(0.3), Inches(4.6), rw - 0.4, Inches(0.3),
             '💳 提领策略', size=12, color=C['dark_text'], bold=True)
    add_text(s, cxr + Inches(0.3), Inches(4.9), rw - 0.4, Inches(1.2),
             '定期提取现金补充现金流\n适合：教育金/养老金规划\n兼顾当下现金流与长期增长',
             size=8.5, color=C['body_text'])

    add_text(s, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.25),
             '注：纵轴使用对数刻度；不提领按退保发还总额计算，提领按累计已领+退保现价计算',
             size=7, color=C['mid_text'])


def _slide_education(prs, wd, meta, C, scene_imgs=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    wd_start = None
    for y in sorted(int(k) for k in wd.keys()):
        if wd[y].get('Annual_WD', 0) > 0:
            wd_start = y
            break
    wd_start = wd_start or 7
    annual_wd = wd.get(wd_start, {}).get('Annual_WD', 0) or 0
    insured_age = meta.get('insured_age', 0) or 0
    def age_of(y): return insured_age + y - 1

    # ── 根据年龄判定场景 ──
    is_education = insured_age < 18
    is_retirement = insured_age >= 55

    if is_education:
        title_text = '教育金现金流 · 与受保人年龄节点结合'
        sub_text = f'保单年度 {wd_start} 起每年提取 USD {annual_wd:,} · 从小学到大学的稳健现金流'
        narrative_prefix = f'保单年度 {wd_start}（受保人 {age_of(wd_start)} 岁）开始提领'
        narrative_suffix = '精准覆盖小学到大学教育支出 → 本金持续滚存'
        # 教育金阶段: 基于孩子年龄 6→23+
        bands = [(6, 12, '小学', '兴趣班/课外辅导'),
                 (13, 15, '初中', '升学辅导/素质拓展'),
                 (16, 18, '高中', '国际课程/留学预备'),
                 (19, 22, '大学', '学费/生活费'),
                 (23, 99, '研究生/创业', '灵活支取')]
        stage_colors = [C['primary_light'], C['primary'], C['primary_dark'], C['accent_dark'], C['dark_text']]
        scene_category = 'education'
    elif is_retirement:
        title_text = '养老金现金流 · 与退休年龄节点结合'
        sub_text = f'保单年度 {wd_start}（{age_of(wd_start)} 岁）起每年提取 USD {annual_wd:,} · 退休生活从容有品质'
        narrative_prefix = f'从 {age_of(wd_start)} 岁开始每年领取 US${annual_wd:,}'
        narrative_suffix = '退休收入稳定补充 → 剩余价值持续滚存 → 财富代际传承'
        # 养老金阶段: 基于退休年龄 60→100+
        bands = [(60, 65, '退休初期', '环球旅行/兴趣爱好'),
                 (65, 75, '活跃养老', '品质生活/健康管理'),
                 (75, 85, '稳健养老', '医疗护理/居家服务'),
                 (85, 95, '享老晚年', '专业照护/尊严养老'),
                 (95, 120, '财富传承', '遗产规划/代际传承')]
        stage_colors = [C['accent'], C['primary_light'], C['primary'], C['primary_dark'], C['dark_text']]
        scene_category = 'retirement'
    else:
        title_text = '财富增值现金流 · 与人生阶段结合'
        sub_text = f'保单年度 {wd_start}（{age_of(wd_start)} 岁）起每年提取 USD {annual_wd:,}'
        narrative_prefix = f'从 {age_of(wd_start)} 岁开始'
        narrative_suffix = '现金流灵活安排 → 财富持续增长'
        bands = [(insured_age, insured_age+5, '规划起步', '现金流启动'),
                 (insured_age+6, insured_age+15, '财富积累', '持续增长期'),
                 (insured_age+16, insured_age+25, '财富增值', '复利效应期'),
                 (insured_age+26, insured_age+35, '稳健管理', '风险控制期'),
                 (insured_age+36, 120, '财富传承', '代际传承')]
        stage_colors = [C['primary_light'], C['primary'], C['primary_dark'], C['accent_dark'], C['dark_text']]
        scene_category = 'savings'

    add_title(s, title_text, sub_text, title_size=36, sub_size=14, C=C)

    def get_yr_data(y):
        if y in wd:
            return {'cum': f'USD {wd[y]["Cum_WD"]:,}', 'remain': f'USD {wd[y]["Total"]:,}'}
        return {'cum': '-', 'remain': '-'}
    def age_band(lo, hi):
        return [y for y in sorted(int(k) for k in wd.keys()) if lo <= age_of(y) <= hi]
    def yr_label(ys):
        if not ys: return '-'
        if len(ys) == 1: return f'Y{ys[0]}'
        return f'Y{ys[0]}-Y{ys[-1]}'

    # 构建5阶段
    stage_data = []
    for lo, hi, name, detail in bands:
        ys = age_band(lo, hi)
        last_y = ys[-1] if ys else None
        d = get_yr_data(last_y) if last_y else {'cum': '-', 'remain': '-'}
        stage_data.append({
            'age': f'{lo}-{hi}岁' if hi < 120 else f'{lo}+岁',
            'name': name,
            'years': yr_label(ys),
            'cum': d['cum'],
            'detail': f'每年 {annual_wd//1000}K\n{detail}',
            'remain': d['remain'],
        })

    wd_kilo = f'{annual_wd//1000}K'

    # 时间线
    line_y = Inches(2.05)
    add_rect(s, Inches(0.5), line_y, Inches(12.333), Inches(0.03), fill=C['primary'])

    cw = Inches(2.32); ch = Inches(4.0); gx = Inches(0.15)
    sx = Inches(0.5); sy = Inches(2.4)

    for i, st in enumerate(stage_data):
        x = sx + i * (cw + gx)
        col = stage_colors[i] if i < len(stage_colors) else C['dark_text']
        add_rect(s, x + cw/2 - Inches(0.08), line_y - Inches(0.06),
                 Inches(0.16), Inches(0.16), fill=C['accent'])
        add_text(s, x, line_y - Inches(0.45), cw, Inches(0.3),
                 st['age'], size=11, color=C['dark_text'], bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x, sy, cw, ch, fill=C['bg_card'])
        add_rect(s, x, sy + ch - Inches(0.55), cw, Inches(0.55), fill=col)
        name_color = C['bg_light'] if i < 4 else C['accent_light']
        add_text(s, x, sy + ch - Inches(0.5), cw, Inches(0.45),
                 st['name'], size=15, color=name_color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, sy + ch - Inches(0.85), cw, Inches(0.3),
                 st['years'], size=9, color=C['mid_text'], align=PP_ALIGN.CENTER)
        cum_y = sy + Inches(2.5)
        add_text(s, x, cum_y, cw, Inches(0.25),
                 '累计已领', size=9, color=C['mid_text'], align=PP_ALIGN.CENTER)
        add_text(s, x, cum_y - Inches(0.6), cw, Inches(0.5),
                 st['cum'], size=15, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_rect(s, x + Inches(0.2), sy + Inches(1.5), cw - Inches(0.4), Inches(0.02), fill=C['gray_line'])
        add_text(s, x + Inches(0.15), sy + Inches(0.7), cw - Inches(0.3), Inches(0.7),
                 st['detail'], size=9, color=C['body_text'], align=PP_ALIGN.CENTER)
        add_text(s, x, sy + Inches(0.2), cw, Inches(0.3),
                 f'保单现价 {st["remain"]}', size=8.5, color=C['dark_text'],
                 align=PP_ALIGN.CENTER, bold=True)

    # 底部场景图
    if scene_imgs:
        img_x = Inches(0.5)
        for sip in scene_imgs[:3]:
            _add_picture(s, sip, img_x, Inches(6.6), Inches(2.8), Inches(0.7))
            img_x += Inches(3.0)

    # 优先用叙事文案, 回退到默认
    nar_savings = meta.get('narrative_chapter_savings', '')
    final_narrative = nar_savings if nar_savings else f'{narrative_prefix} → {narrative_suffix}'
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.35),
             final_narrative[:200], size=10, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.25),
             f'数据来源：{meta.get("company_short", "保险公司")}官方计划书 第5部分',
             size=8, color=C['mid_text'])


def _slide_disclaimer(prs, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    add_title(s, '重要事项声明', 'Important Notes · 客户须知的风险与限制',
              title_size=36, sub_size=14, C=C)
    notes = [
        ('01', '非保证利益',
         '本计划所演示之非保证金额（包括复归红利及终期分红）乃基于现时假设投资回报计算，并非保证。',
         '实际金额或会因投资市场波动而调整，在某些情况下非保证金额可能为零。'),
        ('02', '汇率风险',
         '若以保单货币（美元）以外的其他货币支付保费或收取利益，将按本公司不时厘定的汇率兑换。',
         '外币汇率波动可能影响实际支付金额。'),
        ('03', '提领限制',
         '现金提取须符合本公司最低投保单位要求；若提取导致投保单位低于最低要求，则不可提取。',
         '所演示之提取金额基于非保证红利与分红，未必可维持。'),
        ('04', '通胀风险',
         '未来生活成本可能因通胀而上升；本计划之金额为名义金额，未必能完全追上通胀。',
         '建议结合其他投资工具作综合规划。'),
    ]
    cw = Inches(6.1); ch = Inches(2.1)
    gx = Inches(0.25); gy = Inches(0.2)
    sx = Inches(0.5); sy = Inches(2.0)
    for i, (no, t, b1, b2) in enumerate(notes):
        col = i % 2; row = i // 2
        x = sx + col * (cw + gx); y = sy + row * (ch + gy)
        add_rect(s, x, y, cw, ch, fill=C['bg_card'])
        add_rect(s, x, y, cw, Inches(0.55), fill=C['dark_text'])
        add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(0.8), Inches(0.45),
                 no, size=20, color=C['accent'], bold=True)
        add_text(s, x + Inches(1.0), y + Inches(0.05), cw - Inches(1.2), Inches(0.45),
                 t, size=14, color=C['bg_light'], bold=True)
        add_text(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.7),
                 b1, size=10, color=C['body_text'])
        add_text(s, x + Inches(0.2), y + Inches(1.4), cw - Inches(0.4), Inches(0.7),
                 b2, size=10, color=C['body_text'])

    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.7), fill=C['bg_card'])
    add_rect(s, Inches(0.5), Inches(6.5), Inches(0.08), Inches(0.7), fill=C['primary'])
    add_text(s, Inches(0.7), Inches(6.55), Inches(4), Inches(0.3),
             '完整条款', size=11, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             '本文件不包含完整的计划条款，详情参阅保单文件及主要产品推销刊物。阁下应向持牌保险中介人查询。',
             size=9, color=C['body_text'])


def _slide_sources(prs, data, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    add_title(s, '数据来源与回溯',
              'Data Sources · 全部数据可逐项回溯至官方计划书页码',
              title_size=36, sub_size=14, C=C)
    rows = [
        ('保障摘要 · 保额/保费/年期', 'PDF p1', '受保人资料 + 投保资料'),
        ('不提领 · 退保发还金额（Y1-Y40）', 'PDF p2', '基本计划-说明摘要（第3部分）'),
        ('不提领 · 退保发还金额（Y41-Y80）', 'PDF p3', '基本计划-说明摘要（第3部分）'),
        ('不提领 · 身故赔偿额说明', 'PDF p17-19', '身故赔偿额-不同投资回报下的说明'),
        ('提领 · 退保发还金额（Y1-Y40）', 'PDF p42-43', '基本计划-说明摘要（续）-现金提取'),
        ('提领 · 退保发还金额（Y41-Y80）', 'PDF p44-45', '基本计划-说明摘要（续）-现金提取'),
        ('提领 · 身故赔偿额说明', 'PDF p47-50', '身故赔偿额下的说明（提取）'),
        ('注释与重要事项', 'PDF p20-22', '本计划之注释条款'),
        ('产品详细说明', 'PDF p23-41', '财富增值选项、跃进选项等'),
    ]
    headers = ['数据项', 'PDF 页码', '来源说明']
    tx = Inches(0.5); ty = Inches(1.95)
    tw = Inches(12.333); th = Inches(4.5)
    ts = s.shapes.add_table(len(rows) + 1, 3, tx, ty, tw, th)
    table = ts.table
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(6.0)
    table.rows[0].height = Inches(0.5)
    for i in range(1, len(rows) + 1):
        table.rows[i].height = Inches(0.4)

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
        tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = h
        run.font.size = Pt(10); run.font.name = FONT_HEI
        run.font.color.rgb = _c(C['bg_light']); run.font.bold = True
        rPr = run._r.get_or_add_rPr()
        eaFont = rPr.find(qn('a:ea'))
        if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
        eaFont.set('typeface', FONT_HEI)

    for r, row in enumerate(rows, start=1):
        bg = C['bg_card'] if r % 2 == 1 else C['bg_light']
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
            tf.margin_top = Emu(12000); tf.margin_bottom = Emu(12000)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            if c == 0: color = C['dark_text']; bold = True
            elif c == 1: color = C['primary_dark']; bold = True
            else: color = C['body_text']; bold = False
            run = p.add_run(); run.text = val
            run.font.size = Pt(9); run.font.name = FONT_HEI
            run.font.color.rgb = _c(color); run.font.bold = bold
            rPr = run._r.get_or_add_rPr()
            eaFont = rPr.find(qn('a:ea'))
            if eaFont is None: eaFont = etree.SubElement(rPr, qn('a:ea'))
            eaFont.set('typeface', FONT_HEI)

    add_rect(s, Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.6), fill=C['bg_card'])
    add_rect(s, Inches(0.5), Inches(6.6), Inches(0.08), Inches(0.6), fill=C['primary'])
    add_text(s, Inches(0.7), Inches(6.65), Inches(4), Inches(0.3),
             '核验方法', size=11, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3),
             f'所有数据均直接源自{meta.get("company_short", "保险公司")}官方计划书原文（pdfplumber按列精确提取）；如需进一步核验，可逐项对照 PDF 指定页码。',
             size=9, color=C['body_text'])


def _slide_summary(prs, no_wd, wd, paid_total, meta, C):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    ai_narrative = (meta.get('ai_narrative', '') or '').strip()

    if ai_narrative:
        # AI 总结模式
        add_title(s, 'AI 总结与建议', '基于对话内容生成 · Summary & Recommendation',
                  title_size=36, sub_size=14, C=C)
        add_rect(s, Inches(0.5), Inches(2.0), Inches(12), Inches(4.5), fill=C['bg_card'])
        add_rect(s, Inches(0.5), Inches(2.0), Inches(12), Inches(0.06), fill=C['accent'])
        add_text(s, Inches(0.8), Inches(2.3), Inches(11.4), Inches(4.0),
                 ai_narrative[:300], size=14, color=C['dark_text'])
        add_text(s, Inches(0.8), Inches(6.3), Inches(11.4), Inches(0.3),
                 '— AI 基于对话记录生成，仅供参考', size=9, color=C['mid_text'])
    else:
        # 标准总结（原逻辑）
        add_title(s, '总结与建议', 'Summary & Recommendation', title_size=36, sub_size=14, C=C)
        has_wd = bool(wd) and any(v.get('Annual_WD', 0) > 0 for v in wd.values())
        pb, y20, y30, _ = _find_milestone(no_wd, paid_total)
        y_pb = no_wd.get(pb, {}) if pb else {}
        y30 = no_wd.get(y30, {}) if y30 else {}
        w20 = wd.get(y20, {}) if y20 else {}
        wd_start = None
        for y in sorted(int(k) for k in wd.keys()):
            if wd[y].get('Annual_WD', 0) > 0: wd_start = y; break
        wd_start = wd_start or 7
        annual_wd = wd.get(wd_start, {}).get('Annual_WD', 0) or 0
        pb_display = pb if pb else '—'
        pb_metric = f'Y{pb}' if pb else '—'
        pb_body = f'保单第 {pb_display} 年回本（退保现价 > 已缴保费）。\n' + (f'Y{pb} 现价 USD {y_pb.get("Total", 0):,}；' if pb else '') + f'Y30 约 USD {y30.get("Total", 0):,}。'
        cards = [
            {'icon': 'A', 'title': '短期供款 · 终身受惠', 'body': f'{meta.get("payment_years",5)} 年缴清 USD {paid_total:,.0f}，后续无需再缴费。\n保单持续滚存至 128 岁，长期复利效应明显。', 'metric': str(meta.get('payment_years',5)), 'unit': '年缴清'},
            {'icon': 'B', 'title': '稳健回本 · 风险可控', 'body': pb_body, 'metric': pb_metric, 'unit': '回本'},
            {'icon': 'C', 'title': '灵活提领 · 教育无忧', 'body': (f'Y{wd_start} 起每年提取 USD {annual_wd:,}，覆盖小学到大学。\nY20 累计提取 USD {w20.get("Cum_WD", 0):,}，本金仍持续滚存。' if has_wd else '保单持续滚存，可灵活提取现金价值应对人生各阶段需求。'), 'metric': (f'{annual_wd//1000}K' if has_wd else '灵活'), 'unit': ('年提取' if has_wd else '支取')},
        ]
        cw = Inches(4.0); ch = Inches(3.7); gx = Inches(0.2); sx = Inches(0.5); sy = Inches(2.0)
        for i, c in enumerate(cards):
            x = sx + i * (cw + gx)
            add_rect(s, x, sy, cw, ch, fill=C['bg_card'])
            add_rect(s, x, sy, cw, Inches(0.7), fill=C['dark_text'])
            add_rect(s, x + Inches(0.2), sy + Inches(0.18), Inches(0.35), Inches(0.35), fill=C['primary'])
            add_text(s, x + Inches(0.2), sy + Inches(0.18), Inches(0.35), Inches(0.35), c['icon'], size=14, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, x + Inches(0.7), sy + Inches(0.18), cw - Inches(0.8), Inches(0.4), c['title'], size=13, color=C['bg_light'], bold=True)
            add_text(s, x + Inches(0.25), sy + Inches(0.95), cw - Inches(0.5), Inches(1.4), c['body'], size=10, color=C['body_text'])
            add_text(s, x, sy + Inches(2.5), cw, Inches(0.7), c['metric'], size=32, color=C['primary'], bold=True, align=PP_ALIGN.CENTER)
            add_text(s, x, sy + Inches(3.2), cw, Inches(0.3), c['unit'], size=9, color=C['mid_text'], align=PP_ALIGN.CENTER)
        add_rect(s, Inches(0.5), Inches(5.9), Inches(12.333), Inches(1.4), fill=C['bg_card'])
        add_rect(s, Inches(0.5), Inches(5.9), Inches(0.08), Inches(1.4), fill=C['primary'])
        add_text(s, Inches(0.7), Inches(5.95), Inches(4), Inches(0.3), '综合建议', size=12, color=C['dark_text'], bold=True)
        add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.3), '本计划适合具备 5 年供款能力、追求稳健长期财富传承、且有子女教育金规划需求的高净值家庭。', size=10.5, color=C['body_text'])
        add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.3), '· 短期供款压力适中  ·  长期复利跨越代际  ·  提领方案精准对接教育节点  ·  兼顾传承与流动性', size=10, color=C['body_text'])
        add_text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3), '实际收益取决于投资市场表现，建议结合整体资产配置作综合考虑。', size=8.5, color=C['mid_text'])


def _slide_ending(prs, summary, meta, C):
    """美化版尾页: 简洁居中大字号"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.4),
             f'{meta.get("company_short_en", "")}  |  {meta.get("company_short", "")}',
             size=12, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    name = summary.get('insured_name', 'VIP 先生') or 'VIP 先生'
    age = summary.get('insured_age', 1)

    add_text(s, Inches(0.5), Inches(2.2), Inches(12.333), Inches(0.8),
             '愿这份规划', size=32, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.1), Inches(12.333), Inches(0.8),
             '陪伴孩子稳健成长', size=32, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.4),
             '让爱与财富，代代相传', size=14, color=C['accent_light'], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.17), Inches(4.5), Inches(1.0), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.75), Inches(12.333), Inches(0.5),
             f'—— 致 {name} 与家人  ——', size=13, color=C['bg_light'], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.25), Inches(12.333), Inches(0.4),
             f'规划起点 · 受保人 {age} 岁', size=10, color=C['accent_light'], align=PP_ALIGN.CENTER)
    add_rect(s, Emu(0), Inches(7.2), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(7.3), Inches(8), Inches(0.2),
             f'本演示基于{meta.get("company_short", "保险公司")}官方计划书数据生成 · 实际保单条款以保单文件为准',
             size=8, color=C['accent_light'])
    add_text(s, Inches(11.5), Inches(7.3), Inches(1.5), Inches(0.2),
             f'{meta.get("_page_num", 12):02d} / {meta.get("_total_slides", 12):02d}', size=8, color=C['accent_light'], align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════
#  组合方案: 重疾(CI) / IUL 附加页面
# ══════════════════════════════════════════════════════════

def _slide_ci_divider(prs, ci_data, meta, C, ci_company=None):
    """重疾篇章封面（简洁版）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    co = ci_company.get('name_zh', meta.get('company_short', '')) if ci_company else meta.get('company_short', '')
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             '守护家庭 · 风险保障篇', size=11, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1.0),
             '家庭保障规划', size=36, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.6),
             '筑起收入高墙，守护家人未来', size=18, color=C['accent'], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             '{} · 危疾保障计划'.format(co), size=12, color=C['accent_light'], align=PP_ALIGN.CENTER)


def _slide_ci_overview(prs, ci_data, meta, C, company=None):
    """重疾概要: 交多少 · 保多少 · 保什么"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_title(s, '危疾保障概要', '交多少 · 保多少 · 保什么', title_size=36, sub_size=14, C=C)

    ci_summary = ci_data.get('summary', {})
    ci_policy = ci_data.get('policy', {})
    ci_items = ci_data.get('coverage_items', [])
    ins = ci_data.get('insured', {})

    premium = ci_summary.get('annual_premium', ci_policy.get('annual_premium', 0))
    coverage = ci_summary.get('sum_insured', ci_policy.get('sum_insured', 0))
    pay_years = ci_summary.get('payment_years', ci_policy.get('premium_payment_period', 10))

    cards_data = [
        ('年缴保费', 'USD {:,}'.format(int(premium)), '{}年缴清'.format(pay_years), C['primary']),
        ('保障总额', 'USD {:,}'.format(int(coverage)), '首十年升级保障', C['primary_dark']),
        ('保障年期', '终身', '至受保人终身', C['accent_dark']),
        ('受保人', '{}'.format(ins.get("name", "—")), '{}岁 · {}'.format(ins.get("age", "—"), ins.get("gender", "—")), C['dark_text']),
    ]

    card_w = Inches(3.0); card_h = Inches(1.8); gap = Inches(0.25); sx = Inches(0.4); sy = Inches(2.2)
    for i, (label, val, sub, col) in enumerate(cards_data):
        x = sx + i * (card_w + gap)
        add_rect(s, x, sy, card_w, card_h, fill=C['bg_card'])
        add_rect(s, x, sy, card_w, Inches(0.06), fill=col)
        add_text(s, x + Inches(0.2), sy + Inches(0.2), card_w - 0.4, Inches(0.25),
                 label, size=9, color=C['mid_text'])
        add_text(s, x + Inches(0.2), sy + Inches(0.5), card_w - 0.4, Inches(0.45),
                 val, size=20, color=col, bold=True)
        add_text(s, x + Inches(0.2), sy + Inches(1.1), card_w - 0.4, Inches(0.3),
                 sub, size=9, color=C['body_text'])

    # CI利益演示图表 (左侧保障项, 右侧折线图)
    add_rect(s, Inches(0.4), Inches(4.3), Inches(0.3), Inches(0.04), fill=C['primary'])
    add_text(s, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4),
             '保障范围 ({} 项)'.format(len(ci_items)), size=14, color=C['dark_text'], bold=True)

    item_lines = []
    for item in (ci_items or [])[:8]:
        name = item.get('label', item.get('name', '—'))
        amt = item.get('amount', item.get('sum_insured', 0))
        desc = item.get('description', '')
        line = '· ' + name
        if desc: line += ' — ' + str(desc)[:15]
        if amt: line += '  USD {:,}'.format(int(amt))
        item_lines.append(line)
    item_text = '\n'.join(item_lines) if item_lines else '· 严重疾病保障\n· 早期危疾保障\n· 保费豁免'
    add_text(s, Inches(0.8), Inches(4.7), Inches(6), Inches(2.5),
             item_text, size=9, color=C['body_text'])

    # CI 折线图已移至组合叙事项

    add_text(s, Inches(0.4), Inches(7.2), Inches(12), Inches(0.25),
             '数据来源：保险公司官方计划书', size=8, color=C['mid_text'])


def _slide_ci_data_table(prs, ci_data, meta, C, company=None):
    """CI 保单摘要表: 年龄/年度/保费/退保价值/身故赔偿"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_rect(s, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.35), fill=C['accent'])
    add_text(s, Inches(0.75), Inches(0.55), Inches(8), Inches(0.35),
             '保单摘要', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.75), Inches(0.95), Inches(10), Inches(0.25),
             '年龄 · 保单年度 · 保费 · 退保价值 · 身故赔偿', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.3), Inches(11.5), Inches(0.015), fill=C['accent_light'])

    ci_summary = ci_data.get('summary', {})
    ci_policy = ci_data.get('policy', {})
    ins = ci_data.get('insured', {})
    prem = float(ci_summary.get('annual_premium', ci_policy.get('annual_premium', 0)))
    coverage = float(ci_summary.get('sum_insured', ci_policy.get('sum_insured', 0)))
    pay_yrs = int(ci_summary.get('payment_years', str(ci_policy.get('premium_payment_period', '10')).replace('年','')))
    insured_age = int(ins.get('age', 1))
    bi = ci_data.get('benefit_illustration', [])
    upgrade_yrs = int(ci_data.get('upgrade_benefit_years', 10))
    upgrade_amt = float(ci_data.get('upgrade_benefit_amount', 0))

    # 筛选行: Y1-Y10(每年), Y10之后每10年
    selected = []
    for r in bi:
        y = int(r.get('policy_year', 0))
        age = insured_age + y - 1
        if y <= 0 or age > 120: continue
        if y <= 10 or y % 10 == 0:
            selected.append(r)

    headers = ['年龄', '保单年度', '年缴保费\n(USD)', '累计已缴\n(USD)', '基本保额\n(USD)', '总身故赔偿\n(USD)']
    upgrade_extra = upgrade_amt  # 首N年额外赔偿金额
    rows_data = []
    for r in selected:
        y = int(r['policy_year'])
        age = insured_age + y - 1
        prem_this = prem if y <= pay_yrs else 0
        total_paid = int(r.get('total_premium_paid', 0))
        db = int(r.get('death_benefit', coverage))
        base = int(coverage)
        # 总身故赔偿 = death_benefit from raw data
        total_db = db if db > base else base
        rows_data.append([age, y, prem_this, total_paid, base, total_db])

    n_rows = len(rows_data) + 1
    n_cols = len(headers)
    tbl_left = Inches(0.5); tbl_top = Inches(1.5); tbl_w = Inches(12.0)
    row_h = Inches(0.38)
    table_shape = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, row_h * n_rows)
    table = table_shape.table
    col_widths = [Inches(1.2), Inches(1.3), Inches(1.8), Inches(2.0), Inches(2.8), Inches(2.9)]
    for i, w in enumerate(col_widths): table.columns[i].width = w

    # 表头
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        p = cell.text_frame.paragraphs[0]
        p.text = h.replace('\n', ' ')
        p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['primary'])
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for ri, row in enumerate(rows_data):
        is_upgrade = row[1] <= upgrade_yrs and upgrade_amt > 0
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ci == 2:
                cell.text = '{:,.0f}'.format(val) if val > 0 else '—'
            elif isinstance(val, int) and val >= 1000:
                cell.text = '{:,.0f}'.format(val)
            else:
                cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = _c(C['accent_dark']) if ci == 5 and is_upgrade else _c(C['dark_text'])
                if ci == 5 and is_upgrade: p.font.bold = True
            cell.fill.solid()
            if ri % 2 == 0: cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
            else: cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 底部说明
    fy = tbl_top + row_h * n_rows + Inches(0.08)
    notes = []
    total_prem = int(prem * pay_yrs)
    notes.append('总缴保费: USD {:,.0f}'.format(total_prem))
    notes.append('基本保额: USD {:,.0f}'.format(int(coverage)))
    if upgrade_amt > 0 and upgrade_yrs > 0:
        notes.append('首{}年总保障: USD {:,.0f}'.format(upgrade_yrs, int(coverage + upgrade_amt)))
    add_rect(s, Inches(0.5), fy, Inches(11.5), Inches(0.3), fill=C['bg_card'])
    add_text(s, Inches(0.6), fy, Inches(11.3), Inches(0.3),
             '  '.join(notes), size=8, color=C['body_text'])
    add_text(s, Inches(0.5), fy + Inches(0.35), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])


def _slide_ci_coverage(prs, ci_data, meta, C, company=None):
    """CI 保障项目详解（基于PDF保险责任）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_rect(s, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.35), fill=C['accent'])
    add_text(s, Inches(0.75), Inches(0.55), Inches(8), Inches(0.35),
             '保障项目', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.75), Inches(0.95), Inches(10), Inches(0.25),
             '严重疾病 · 多重赔付 · 保障还原 · 保费豁免', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.3), Inches(11.5), Inches(0.015), fill=C['accent_light'])

    coverage = float(ci_data.get('policy', {}).get('sum_insured', 100000))
    upgrade_amt = float(ci_data.get('upgrade_benefit_amount', 0))
    upgrade_yrs = int(ci_data.get('upgrade_benefit_years', 20))
    if upgrade_yrs <= 0: upgrade_yrs = 20

    # 从PDF提取的「守護家倍198」实际保障项目
    all_items = [
        {'label': '严重疾病保障', 'amount': coverage, 'description': '严重程度3之危疾，赔付100%保额\n覆盖100种严重疾病'},
        {'label': '首{}年额外赔偿'.format(upgrade_yrs), 'amount': upgrade_amt if upgrade_amt > 0 else coverage * 0.6,
         'description': '首{}个保单年度内额外赔付60%保额\n身故/严重程度3之危疾均适用'.format(upgrade_yrs)},
        {'label': '早期危疾保障', 'amount': coverage * 0.5, 'description': '严重程度1-2之危疾，赔付50%保额\n涵盖早期危疾及儿童疾病'},
        {'label': '保障还原利益', 'amount': coverage, 'description': '首次严重疾病赔付后，后续可还原至原始保额\n保障至70岁'},
        {'label': '严重都市疾病额外保障', 'amount': coverage, 'description': '3次癌症+3次中风/严重心脏病\n每次额外赔付100%保额，完成后自动增值'},
        {'label': '严重都市疾病无限次增值', 'amount': coverage,
         'description': '每赔付6次后自动增值6次保障\n不限增值次数，保障至85岁'},
        {'label': '保费豁免', 'amount': 0, 'description': '确诊严重程度3之危疾后\n豁免基本计划后续保费'},
    ]

    # 补充AI提取的其他保障项（去重）
    ci_items = ci_data.get('coverage_items', [])
    existing_labels = set(i['label'] for i in all_items)
    for item in ci_items:
        label = item.get('label', '')
        if label not in existing_labels:
            all_items.append(item)
            existing_labels.add(label)

    card_w = Inches(3.6); card_h = Inches(1.7); gap = Inches(0.2)
    sx = Inches(0.4); sy = Inches(1.55)
    for i, item in enumerate(all_items[:6]):
        col_idx = i % 3
        row_idx = i // 3
        x = sx + col_idx * (card_w + gap)
        y = sy + row_idx * (card_h + gap)
        label = item.get('label', '')
        amt = float(item.get('amount', 0))
        desc = item.get('description', '').replace('\n', ' | ')
        add_rect(s, x, y, card_w, card_h, fill=C['bg_card'])
        add_rect(s, x, y, card_w, Inches(0.06), fill=C['primary'])
        add_text(s, x + Inches(0.2), y + Inches(0.1), card_w - 0.4, Inches(0.25),
                 label, size=11, color=C['dark_text'], bold=True)
        if amt > 0:
            add_text(s, x + Inches(0.2), y + Inches(0.4), card_w - 0.4, Inches(0.25),
                     'USD {:,.0f}'.format(int(amt)), size=13, color=C['primary_dark'], bold=True)
        add_text(s, x + Inches(0.2), y + Inches(0.75), card_w - 0.4, Inches(0.8),
                 desc, size=7.5, color=C['body_text'])

    add_text(s, Inches(0.4), Inches(7.2), Inches(12), Inches(0.25),
             '数据来源：保险公司官方计划书 · 实际保障以保单条款为准', size=8, color=C['mid_text'])


def _ia_irr_cap(currency: str = 'USD') -> float:
    """HK IA 演示利率封顶: HKD 6.0%, 非港元 6.5%"""
    c = (currency or 'USD').upper().strip()
    return 0.06 if c in ('HKD', '港币', '港元', 'HK$') else 0.065


def _ma_irr_bisect(cf):
    """cf: list of (t, amount) tuples; bisect 二分法 NPV=0 求解"""
    def npv(r):
        # 防止 r≈-1 且 t 极大时 (1+r)**t 下溢为 0 → 除零
        # 此时该项贡献按 a 的符号返回 ±inf, 简化处理
        s = 0.0
        denom_base = 1.0 + r
        for t, a in cf:
            if denom_base == 0:
                return float('inf') if a > 0 else float('-inf')
            try:
                d = denom_base ** t
            except (OverflowError, ValueError):
                return float('inf') if a > 0 else float('-inf')
            if d == 0:
                return float('inf') if a > 0 else float('-inf')
            s += a / d
        return s
    lo, hi = -0.99, 1.0
    f_lo, f_hi = npv(lo), npv(hi)
    if f_lo * f_hi > 0:
        return None
    for _ in range(200):
        mid = (lo + hi) / 2
        f_mid = npv(mid)
        if abs(f_mid) < 1e-6 or (hi - lo) < 1e-10:
            return mid
        if f_lo * f_mid < 0:
            hi, f_hi = mid, f_mid
        else:
            lo, f_lo = mid, f_mid
    return (lo + hi) / 2


def compute_irr_ma(annual_prem: float, pay_yrs: int, sv: float, yr: int,
                   currency: str = 'USD'):
    """M-A NPV IRR (不退保/不提取):
    现金流 = [-P at t=0..payYrs-1, +SV at t=yr]; 封顶 HK IA 6.0%/6.5%
    """
    if yr <= 0 or annual_prem <= 0 or sv <= 0 or pay_yrs < 1:
        return None
    cf = [(t, -annual_prem) for t in range(pay_yrs)]
    cf.append((yr, sv))
    r = _ma_irr_bisect(cf)
    if r is None:
        return None
    return min(r, _ia_irr_cap(currency))


def compute_irr_ma_withdraw(annual_prem: float, pay_yrs: int, sv: float, yr: int,
                            wd_start: int = 0, annual_wd: float = 0,
                            currency: str = 'USD'):
    """M-A NPV IRR (含提领):
    现金流 = [-P at t=0..payYrs-1, +aw at t=wd_start..yr-1, +(aw+SV) at t=yr]; 封顶 6.0%/6.5%
    """
    if yr <= 0 or annual_prem <= 0 or sv <= 0 or pay_yrs < 1:
        return None
    cf = [(t, -annual_prem) for t in range(pay_yrs)]
    if wd_start > 0 and annual_wd > 0 and wd_start < yr:
        for t in range(wd_start, yr):
            cf.append((t, annual_wd))
        cf.append((yr, annual_wd + sv))
    else:
        cf.append((yr, sv))
    r = _ma_irr_bisect(cf)
    if r is None:
        return None
    return min(r, _ia_irr_cap(currency))




def _slide_comparison_divider(prs, compare_types, meta, C):
    """多产品对比章节封面"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])
    type_label = {
        'savings': '储蓄险',
        'ci': '重疾险',
        'iul': 'IUL 万用寿险',
    }
    type_names = ' · '.join(type_label.get(t, t) for t in compare_types)
    add_text(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             '横向对比篇 · 数据说话', size=11, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1.0),
             '产品对比', size=40, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.6),
             type_names, size=16, color=C['accent_light'], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             '退保现价 · 保费 · 保额 · 杠杆率', size=12, color=C['accent_light'], align=PP_ALIGN.CENTER)




def _paid_total_from_bi(bi_list):
    """从 benefit_illustration 最后一行取总保费 (含提取场景优先用 paid_total)"""
    if not bi_list:
        return 0
    if isinstance(bi_list, dict):
        try:
            last = bi_list[max(bi_list.keys(), key=lambda k: int(k) if str(k).isdigit() else 0)]
            return float(last.get('paid_total') or last.get('Total_Premium_Paid') or last.get('total_premium_paid') or 0)
        except: return 0
    # list
    last = bi_list[-1]
    return float(last.get('total_premium_paid') or 0)

def _bi_to_no_wd(bi, prod_data=None, insured_age=1):
    """benefit_illustration (list) → dict[int] 形式
    同时归一化字段: Total / Mult / Guar_CV / Paid / Total_Premium_Paid / Cum_WD / Age
    (适配 _slide_features / _slide_growth_chart / _slide_no_withdraw_table / _slide_compare_chart)
    """
    if isinstance(bi, dict):
        return bi
    ins_age = insured_age
    if prod_data:
        ins = prod_data.get('insured') or {}
        try: ins_age = int(ins.get('age') or insured_age)
        except: pass
    out = {}
    # 预计算实际总保费峰值 (累计 total_premium_paid 最大值), 用于校正 M-A IRR 年缴额
    # 部分产品 annual_premium 抽错 (如宏利抽成 100K 实际 400K), 用真实累计保费反推更可靠
    _paid_final = 0.0
    for _r in bi or []:
        try:
            _pp = float(_r.get('total_premium_paid') or _r.get('Total_Premium_Paid') or 0)
        except Exception:
            _pp = 0.0
        if _pp > _paid_final:
            _paid_final = _pp
    for r in bi or []:
        try: yr = int(r.get('policy_year') or 0)
        except: continue
        normalized = {**r}
        total = float(r.get('total_surrender_value') or r.get('Total') or 0)
        prem = float(r.get('total_premium_paid') or r.get('Total_Premium_Paid') or 0)
        guar_cv = float(r.get('guaranteed_cash_value') or r.get('guaranteed_surrender_value') or r.get('Guar_CV') or 0)
        # Bug fix: 原读 paid_premium 字段但 BI 数据实际是 total_premium_paid (累计), 导致 Paid=0 → IRR/Simple 全为 None
        paid = float(r.get('paid_premium') or r.get('total_premium_paid') or r.get('Paid') or r.get('Total_Premium_Paid') or 0)
        cum_wd = float(r.get('cum_withdrawal') or r.get('Cum_WD') or 0)
        mult = total / prem if prem > 0 else 0
        age = r.get('Age') or r.get('age') or (ins_age + yr - 1 if yr else None)
        # 兜底: Guar_CV 若未提供, 假设为 Total 的 60% (港险常见保证部分占比)
        if guar_cv == 0 and total > 0:
            guar_cv = total * 0.6
        if 'Total' not in normalized or not normalized.get('Total'):
            normalized['Total'] = total
        if 'Mult' not in normalized or not normalized.get('Mult'):
            normalized['Mult'] = mult
        if 'Guar_CV' not in normalized:
            normalized['Guar_CV'] = guar_cv
        if 'Paid' not in normalized:
            normalized['Paid'] = paid
        if 'Total_Premium_Paid' not in normalized:
            normalized['Total_Premium_Paid'] = prem
        if 'Cum_WD' not in normalized:
            normalized['Cum_WD'] = cum_wd
        if 'Age' not in normalized and age:
            normalized['Age'] = int(age)
        # 兜底计算 IRR/Simple (M-A NPV IRR + 单利), 防止 TS 端未传值时表格全显示 '-'
        # M-A 优先, 从 prod_data.summary / policy 取 annual_prem / pay_yrs / currency
        if yr > 0 and total > 0 and paid > 0:
            if not normalized.get('IRR'):
                summ = (prod_data or {}).get('summary') or {}
                pol = (prod_data or {}).get('policy') or {}
                ap = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
                pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '5'
                py = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '5') or 5)
                cur = summ.get('currency') or pol.get('currency') or 'USD'
                # 用真实累计保费反推年缴额, 修正抽错的 annual_premium
                ap_eff = (_paid_final / py) if (_paid_final > 0 and py > 0) else ap
                ma = None
                if ap_eff > 0 and py > 0:
                    ma = compute_irr_ma(ap_eff, py, total, yr, cur)
                # 回退: M-A 失败时仍用 CAGR 兜底 (避免 None)
                normalized['IRR'] = ma if ma is not None else (total / paid) ** (1.0 / yr) - 1
            if not normalized.get('Simple'):
                normalized['Simple'] = (total - paid) / paid / yr
        out[yr] = normalized
    return out

def _wd_to_dict(wd_data, prod_data=None, insured_age=1):
    """withdrawal_illustration (list) → dict[int] 形式
    字段归一化: Annual_WD / Total / Cum_WD / Total_Received / Paid / Age
    Age 从 prod_data.insured.age + yr - 1 兜底计算
    """
    if isinstance(wd_data, dict):
        return wd_data
    # 优先从 prod_data.insured 拿年龄
    ins_age = insured_age
    if prod_data:
        ins = prod_data.get('insured') or {}
        try: ins_age = int(ins.get('age') or insured_age)
        except: pass
    out = {}
    # 预计算实际总保费峰值 (同 _bi_to_no_wd, 校正抽错的 annual_premium)
    _paid_final = 0.0
    for _r in wd_data or []:
        try:
            _pp = float(_r.get('total_premium_paid') or _r.get('Total_Premium_Paid') or 0)
        except Exception:
            _pp = 0.0
        if _pp > _paid_final:
            _paid_final = _pp
    for r in wd_data or []:
        try: yr = int(r.get('policy_year') or 0)
        except: continue
        normalized = {**r}
        annual_wd = float(r.get('annual_withdrawal') or r.get('Annual_WD') or 0)
        total = float(r.get('total_surrender_value') or r.get('surrender_value_after') or r.get('Total') or 0)
        cum_wd = float(r.get('cum_withdrawal') or r.get('total_withdrawn') or r.get('Cum_WD') or 0)
        # Bug fix: WI 数据用 total_premium_paid 表示累计保费, 原读 paid_premium 全为 0
        paid = float(r.get('paid_premium') or r.get('total_premium_paid') or r.get('Paid') or r.get('Total_Premium_Paid') or 0)
        age = r.get('Age') or r.get('age') or (ins_age + yr - 1 if yr else None)
        if 'Annual_WD' not in normalized:
            normalized['Annual_WD'] = annual_wd
        if 'Total' not in normalized:
            normalized['Total'] = total
        if 'Cum_WD' not in normalized:
            normalized['Cum_WD'] = cum_wd
        if 'Total_Received' not in normalized:
            normalized['Total_Received'] = cum_wd + total
        if 'Paid' not in normalized:
            normalized['Paid'] = paid
        if 'Age' not in normalized and age:
            normalized['Age'] = int(age)
        # 兜底 IRR/Simple (提领场景: M-A 用年现金流+终值SV+aw, Total_Received=Cum_WD+Total_剩余)
        if yr > 0 and paid > 0:
            tot_recv = normalized.get('Total_Received') or (cum_wd + total)
            if tot_recv > 0:
                if not normalized.get('IRR'):
                    summ = (prod_data or {}).get('summary') or {}
                    pol = (prod_data or {}).get('policy') or {}
                    ap = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
                    pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '5'
                    py = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '5') or 5)
                    cur = summ.get('currency') or pol.get('currency') or 'USD'
                    wd_start = int(summ.get('wd_start_year') or pol.get('wd_start_year') or 0)
                    # 用真实累计保费反推年缴额, 修正抽错的 annual_premium
                    ap_eff = (_paid_final / py) if (_paid_final > 0 and py > 0) else ap
                    ma = None
                    if ap_eff > 0 and py > 0 and total > 0:
                        ma = compute_irr_ma_withdraw(ap_eff, py, total, yr, wd_start, annual_wd, cur)
                    # 回退: CAGR (CV=剩余退保, 不含累计提取)
                    fallback = (total / paid) ** (1.0 / yr) - 1 if total > 0 else 0
                    normalized['IRR'] = ma if ma is not None else fallback
                if not normalized.get('Simple'):
                    normalized['Simple'] = (tot_recv - paid) / paid / yr
        out[yr] = normalized
    return out


def _slide_savings_no_wd_compare_table(prs, products, meta, C):
    """多储蓄产品 · 不提取后的现金价值横向对比表
    列: 保单年度 × 各产品 CV(USD) + 最优标记 + Y20 IRR
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    n = len(products)
    add_title(s, f'{n} 家储蓄险 · 不提取后的现金价值对比',
              '保单年度 × 各产品退保现价 (USD) · 红字标每行最高 · 含 Y20 IRR',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    rows_per_product = []
    for p in products:
        d = p.get('data') or {}
        bi = d.get('benefit_illustration') or []
        ins_age = 1
        try:
            ins_age = int((d.get('insured') or {}).get('age') or 1)
        except Exception:
            pass
        no_wd = _bi_to_no_wd(bi, prod_data=d, insured_age=ins_age)
        company = (p.get('company_short') or '').upper() or '—'
        rows_per_product.append({
            'company': company,
            'no_wd': no_wd,
        })

    all_yrs = set()
    for rp in rows_per_product:
        all_yrs.update(rp['no_wd'].keys())
    if not all_yrs:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无可用保单年度数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    max_y = max(all_yrs)
    show_yrs = [y for y in [5, 10, 15, 20, 25, 30, 35, 40] if y in all_yrs]
    show_yrs += [y for y in [50, 60, 70, 80, 90, 100] if y in all_yrs and y <= max_y]
    show_yrs = show_yrs[:12]

    headers = ['保单年度'] + [f'{rp["company"]}' for rp in rows_per_product] + ['最优']
    n_cols = len(headers)
    n_rows = 1 + len(show_yrs) + 1  # +1 for IRR row
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.4), Inches(1.95),
                                    Inches(12.9), Inches(0.45) * n_rows)
    tbl = tbl_shape.table
    col_widths = [1.0] + [2.0] * len(rows_per_product) + [1.2]
    total_cw = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * 12.9 / total_cw)
    tbl.rows[0].height = Inches(0.5)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.4)

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行
    for ri, y in enumerate(show_yrs, start=1):
        vals = []
        for rp in rows_per_product:
            r = rp['no_wd'].get(y, {})
            vals.append(float(r.get('Total') or 0))
        max_v = max([v for v in vals if v > 0], default=0)
        best_idx = (vals.index(max_v) + 1) if max_v > 0 else None
        row_bg = C['bg_light'] if ri % 2 == 0 else C['bg_card']

        # 保单年度
        cell = tbl.cell(ri, 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = f'Y{y}'
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI

        # CV 值
        for ci, v in enumerate(vals, start=1):
            cell = tbl.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = f'${v:,.0f}' if v > 0 else '—'
            run.font.size = Pt(10)
            run.font.color.rgb = _c(C['accent_dark'] if (ci == best_idx) else C['body_text'])
            run.font.bold = (ci == best_idx)
            run.font.name = FONT_HEI

        # 最优列
        cell = tbl.cell(ri, n_cols - 1)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = rows_per_product[best_idx - 1]['company'] if best_idx else '—'
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['primary_dark']); run.font.name = FONT_HEI

    # Y20 IRR 汇总行
    y20 = 20
    irr_row_idx = 1 + len(show_yrs)
    irr_vals = []
    for rp in rows_per_product:
        r = rp['no_wd'].get(y20, {})
        irr = r.get('IRR')
        irr_vals.append(irr if irr is not None else 0)
    max_irr = max([iv for iv in irr_vals if iv > 0], default=0)
    best_irr_idx = (irr_vals.index(max_irr) + 1) if max_irr > 0 else None

    cell = tbl.cell(irr_row_idx, 0)
    cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f'Y{y20} IRR'
    run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI

    for ci, iv in enumerate(irr_vals, start=1):
        cell = tbl.cell(irr_row_idx, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        run = p.add_run(); run.text = f'{iv*100:.2f}%' if iv > 0 else '—'
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = _c(C['primary'] if (ci == best_irr_idx) else C['body_text'])
        run.font.name = FONT_HEI

    cell = tbl.cell(irr_row_idx, n_cols - 1)
    cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = rows_per_product[best_irr_idx - 1]['company'] if best_irr_idx else '—'
    run.font.size = Pt(11); run.font.bold = True
    run.font.color.rgb = _c(C['primary']); run.font.name = FONT_HEI

    # 底部说明
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.3),
             '注: 退保现价 = 保证+非保证 (Total) · Y20 IRR 用 M-A NPV 求解 · 红字=该年最高',
             size=9, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             f'数据来源: {n} 款产品官方计划书 · 同档保费 · 含保证+非保证',
             size=8, color=C['mid_text'])


def _slide_savings_wd_compare_table(prs, products, meta, C):
    """多储蓄产品 · 提取后的现金价值横向对比表
    列: 保单年度 × [各产品「累计提取 + 剩余退保现价」配对] + 最优标记 + Y20 IRR
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    n = len(products)
    add_title(s, f'{n} 家储蓄险 · 提取后的现金价值对比',
              '保单年度 × 各产品 [累计提取 → 剩余退保现价] (USD) · 红字标该年最高剩余 · 含 Y20 IRR',
              title_size=30, sub_size=11, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    rows_per_product = []
    for p in products:
        d = p.get('data') or {}
        wd = d.get('withdrawal_illustration') or []
        ins_age = 1
        try:
            ins_age = int((d.get('insured') or {}).get('age') or 1)
        except Exception:
            pass
        wd_dict = _wd_to_dict(wd, prod_data=d, insured_age=ins_age) if wd else {}
        company = (p.get('company_short') or '').upper() or '—'
        rows_per_product.append({
            'company': company,
            'wd': wd_dict,
            'has_wd': bool(wd_dict) and any(v.get('Annual_WD', 0) > 0 for v in wd_dict.values()),
        })

    if not any(rp['has_wd'] for rp in rows_per_product):
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '所有产品均未提供提取场景数据',
                 size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    all_yrs = set()
    for rp in rows_per_product:
        if rp['has_wd']:
            all_yrs.update(rp['wd'].keys())
    if not all_yrs:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无可用保单年度数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    max_y = max(all_yrs)
    show_yrs = [y for y in [5, 10, 15, 20, 25, 30, 35, 40] if y in all_yrs]
    show_yrs += [y for y in [50, 60, 70, 80, 90, 100] if y in all_yrs and y <= max_y]
    show_yrs = show_yrs[:12]

    # 表头: 保单年度 + 每个产品两列(累计提取/剩余) + 最优
    headers = ['保单年度']
    for rp in rows_per_product:
        headers.append(f'{rp["company"]}\n累计提取')
        headers.append(f'{rp["company"]}\n剩余')
    headers.append('最优')
    n_cols = len(headers)
    n_rows = 1 + len(show_yrs) + 1
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.95),
                                    Inches(13.1), Inches(0.45) * n_rows)
    tbl = tbl_shape.table
    # 列宽: 保单年度0.8, 每产品1.4*2, 最优1.0
    col_widths = [0.8] + [1.3, 1.4] * len(rows_per_product) + [1.0]
    total_cw = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * 13.1 / total_cw)
    tbl.rows[0].height = Inches(0.6)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.4)

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        lines = h.split('\n')
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for li, line in enumerate(lines):
            if li > 0:
                p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = line
            run.font.size = Pt(9); run.font.bold = True
            run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行
    for ri, y in enumerate(show_yrs, start=1):
        # 收集每产品的累计提取 + 剩余 SV
        cum_vals = []  # 累计提取
        remain_vals = []  # 剩余退保现价
        for rp in rows_per_product:
            if not rp['has_wd']:
                cum_vals.append(0)
                remain_vals.append(0)
                continue
            r = rp['wd'].get(y, {})
            cum_vals.append(float(r.get('Cum_WD') or 0))
            remain_vals.append(float(r.get('Total') or 0))
        max_remain = max([v for v in remain_vals if v > 0], default=0)
        best_idx = (remain_vals.index(max_remain) + 1) if max_remain > 0 else None
        # best_idx 指向产品索引, 对应列索引 = 1 + best_idx*2 - 1 (剩余列)
        row_bg = C['bg_light'] if ri % 2 == 0 else C['bg_card']

        # 保单年度
        cell = tbl.cell(ri, 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = f'Y{y}'
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI

        # 每个产品: 累计提取列 + 剩余列
        for pi, rp in enumerate(rows_per_product):
            cum_col = 1 + pi * 2  # 累计提取列
            rem_col = 2 + pi * 2  # 剩余列
            # 累计提取列
            cell = tbl.cell(ri, cum_col)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            if not rp['has_wd']:
                run.text = 'N/A'
                run.font.color.rgb = _c(C['mid_text'])
                run.font.italic = True
            else:
                v = cum_vals[pi]
                run.text = f'${v:,.0f}' if v > 0 else '—'
                run.font.color.rgb = _c(C['body_text'])
            run.font.size = Pt(9); run.font.name = FONT_HEI
            # 剩余列
            cell = tbl.cell(ri, rem_col)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            if not rp['has_wd']:
                run.text = 'N/A'
                run.font.color.rgb = _c(C['mid_text'])
                run.font.italic = True
            else:
                v = remain_vals[pi]
                run.text = f'${v:,.0f}' if v > 0 else '—'
                is_best = (best_idx is not None and pi == best_idx - 1)
                run.font.color.rgb = _c(C['accent_dark'] if is_best else C['body_text'])
                run.font.bold = is_best
            run.font.size = Pt(10); run.font.name = FONT_HEI

        # 最优列
        cell = tbl.cell(ri, n_cols - 1)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = rows_per_product[best_idx - 1]['company'] if best_idx else '—'
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = _c(C['primary_dark']); run.font.name = FONT_HEI

    # Y20 IRR 汇总行
    y20 = 20
    irr_row_idx = 1 + len(show_yrs)
    irr_vals = []
    for rp in rows_per_product:
        if not rp['has_wd']:
            irr_vals.append(0)
            continue
        r = rp['wd'].get(y20, {})
        irr = r.get('IRR')
        irr_vals.append(irr if irr is not None else 0)
    max_irr = max([iv for iv in irr_vals if iv > 0], default=0)
    best_irr_pi = (irr_vals.index(max_irr)) if max_irr > 0 else None

    cell = tbl.cell(irr_row_idx, 0)
    cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f'Y{y20} IRR'
    run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI

    for pi, rp in enumerate(rows_per_product):
        # 累计提取列 (空白)
        cum_col = 1 + pi * 2
        cell = tbl.cell(irr_row_idx, cum_col)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
        cell.text = ''
        # IRR 列
        rem_col = 2 + pi * 2
        cell = tbl.cell(irr_row_idx, rem_col)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        if not rp['has_wd'] or irr_vals[pi] == 0:
            run.text = 'N/A' if not rp['has_wd'] else '—'
            if not rp['has_wd']:
                run.font.italic = True
                run.font.color.rgb = _c(C['mid_text'])
            else:
                run.font.color.rgb = _c(C['body_text'])
        else:
            iv = irr_vals[pi]
            run.text = f'{iv*100:.2f}%'
            is_best = (best_irr_pi is not None and pi == best_irr_pi)
            run.font.color.rgb = _c(C['primary'] if is_best else C['body_text'])
            run.font.bold = True
        run.font.size = Pt(11); run.font.name = FONT_HEI

    cell = tbl.cell(irr_row_idx, n_cols - 1)
    cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
    cell.text = ''
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = rows_per_product[best_irr_pi]['company'] if best_irr_pi is not None else '—'
    run.font.size = Pt(11); run.font.bold = True
    run.font.color.rgb = _c(C['primary']); run.font.name = FONT_HEI

    # 底部说明
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.5), Inches(0.3),
             '注: 累计提取 = 至该年累计已领取 · 剩余退保现价 = 提取后剩余 SV · Y20 IRR 用 M-A NPV 求解 · 红字=该年剩余最高',
             size=9, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             f'数据来源: {n} 款产品官方计划书提取场景演示表',
             size=8, color=C['mid_text'])


def _slide_savings_compare_divider(prs, products, meta, C):
    """储蓄险对比分隔页 — 标题 + N 个产品 chip"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])

    add_text(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             '储蓄险对比篇 · 数据驱动决策', size=11, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1.0),
             '储蓄险产品对比', size=40, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    n = len(products) if products else 0
    add_text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.6),
             f'{n} 款产品横向对比', size=16, color=C['accent_light'], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             '退保现价 · 复利 IRR · 提取金额 · 提取后剩余', size=12, color=C['accent_light'], align=PP_ALIGN.CENTER)




def _slide_savings_overview(prs, products, meta, C):
    """V3.1.5+ 多产品首屏 — 对比产品一览
    顶部"储蓄险对比"主标题 + 公司·产品名横幅, 下方表格列每产品核心数据
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    # V3.1.6: 大标题「储蓄险对比」+ 副标题列出所有公司·产品名 (用 · 分隔)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12.5), Inches(0.6),
             '储蓄险对比', size=32, color=C['dark_text'], bold=True, font=FONT_HEI)

    # 收集公司·产品名
    names = []
    for p in products:
        company = (p.get('company_short') or '').upper() or '—'
        name = (p.get('product_name') or p.get('pdfName') or '—')[:14]
        names.append(f'{company} · {name}')
    sub_text = f'共 {len(products)} 款 · ' + ' / '.join(names)
    add_text(s, Inches(0.7), Inches(1.65), Inches(12.5), Inches(0.4),
             sub_text, size=11, color=C['mid_text'], font=FONT_HEI)

    if not products:
        return

    # 提取每个产品的关键信息 + 每 10 年 IRR
    # 用 _bi_to_no_wd 归一化数据, 直接读 IRR 字段 (与单产品不提领表同源)
    irr_yrs = [10, 20, 30, 40, 50, 60]
    rows_data = []
    for p in products:
        d = p.get('data') or {}
        bi = d.get('benefit_illustration') or []
        summ = d.get('summary') or {}
        pol = d.get('policy') or {}
        # 缴费年期
        pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '5'
        pay_yrs = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '5') or 5)
        # 年缴额
        annual_prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        # 总缴保费
        total_prem = 0
        if isinstance(bi, list) and bi:
            last = bi[-1]
            total_prem = float(last.get('total_premium_paid') or 0)
        if total_prem == 0:
            total_prem = annual_prem * pay_yrs
        # 用 _bi_to_no_wd 归一化 (与单产品不提领表同一份数据)
        ins_age = 1
        try:
            ins_age = int((d.get('insured') or {}).get('age') or 1)
        except Exception:
            pass
        no_wd = _bi_to_no_wd(bi, prod_data=d, insured_age=ins_age)
        # 提取每 10 年的 IRR + 倍数 (倍数 = Total/Paid, 即使 IRR 封顶 6.5% 仍可区分)
        cell_map = {}
        for y in irr_yrs:
            r = no_wd.get(y, {})
            cell_map[y] = {
                'irr': r.get('IRR') or 0,
                'mult': r.get('Mult') or 0,
            }
        # 产品名/公司
        company = (p.get('company_short') or '').upper() or '—'
        name = (p.get('product_name') or p.get('pdfName') or '—')[:18]
        rows_data.append({
            'company': company,
            'name': name,
            'pay_yrs': pay_yrs,
            'annual_prem': annual_prem,
            'total_prem': total_prem,
            'cell': cell_map,
        })

    n = len(rows_data)
    # 8 列: 保单年度 + 3 公司 × (IRR + 倍数) + 最优 (与 slide 26 布局一致)
    headers = ['保单年度']
    for rd in rows_data:
        headers.append(f'{rd["company"]}\nIRR')
        headers.append(f'{rd["company"]}\n倍数')
    headers.append('最优')

    # 6 行数据: 每 10 年 1 行 (类似 slide 26 每行 1 年份)
    n_rows = 1 + len(irr_yrs)
    n_cols = len(headers)

    table_x = Inches(0.4); table_y = Inches(2.0)
    table_w = Inches(12.9)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, Inches(0.55) * n_rows)
    tbl = tbl_shape.table
    # 列宽: 保单年度 1.5 + 每公司 IRR 1.3 + 倍数 1.2 + 最优 1.1
    col_widths = [1.5] + [1.3, 1.2] * len(rows_data) + [1.1]
    total_cw = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * 12.9 / total_cw)
    tbl.rows[0].height = Inches(0.6)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.5)

    # 表头 (2 行文字 per cell: 公司\n指标)
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        tf.word_wrap = True
        lines = h.split('\n')
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for li, line in enumerate(lines):
            if li > 0:
                p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = line
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行: 每 10 年 1 行
    for ri, y in enumerate(irr_yrs, start=1):
        row_bg = C['bg_light'] if ri % 2 == 0 else C['bg_card']
        irr_vals = [rd['cell'].get(y, {}).get('irr', 0) for rd in rows_data]
        mult_vals = [rd['cell'].get(y, {}).get('mult', 0) for rd in rows_data]
        max_irr = max([v for v in irr_vals if v > 0], default=0)
        max_mult = max([v for v in mult_vals if v > 0], default=0)
        best_irr_pi = irr_vals.index(max_irr) if max_irr > 0 else None
        best_mult_pi = mult_vals.index(max_mult) if max_mult > 0 else None
        # 综合最优 (按 IRR + 倍数权重)
        combined = [irr_vals[i] * 0.5 + mult_vals[i] * 0.5 for i in range(len(rows_data))]
        max_c = max([v for v in combined if v > 0], default=0)
        best_pi = combined.index(max_c) if max_c > 0 else None

        # 第一列: Y{y}
        cell = tbl.cell(ri, 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = f'Y{y}'
        run.font.size = Pt(12); run.font.bold = True
        run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI

        # 每公司 IRR + 倍数
        for pi, rd in enumerate(rows_data):
            irr_ci = 1 + pi * 2
            mult_ci = 2 + pi * 2
            irr_v = irr_vals[pi]
            mult_v = mult_vals[pi]
            is_best_i = (best_irr_pi == pi)
            is_best_m = (best_mult_pi == pi)

            # IRR 列
            cell = tbl.cell(ri, irr_ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f'{irr_v*100:.2f}%' if irr_v > 0 else '—'
            run.font.size = Pt(12)
            run.font.color.rgb = _c(C['accent_dark'] if is_best_i else C['primary'])
            run.font.bold = True
            run.font.name = FONT_HEI

            # 倍数列
            cell = tbl.cell(ri, mult_ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f'{mult_v:.2f}x' if mult_v > 0 else '—'
            run.font.size = Pt(12)
            run.font.color.rgb = _c(C['accent_dark'] if is_best_m else C['body_text'])
            run.font.bold = is_best_m
            run.font.name = FONT_HEI

        # 最优列 (综合 IRR + 倍数)
        cell = tbl.cell(ri, n_cols - 1)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(row_bg)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = rows_data[best_pi]['company'] if best_pi is not None else '—'
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = _c(C['primary_dark']); run.font.name = FONT_HEI

    # 底部说明
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.5), Inches(0.3),
             '注: 倍数 = 退保现价/已缴保费 (Total/Paid) · IRR = M-A NPV=0 (HK IA 演示利率封顶 6.5%) · 红字=该指标当年最高',
             size=9, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(6.35), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书 · 同档保费 (USD 100K/年 × 5年) · 受保人 1 岁',
             size=8, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.5), Inches(0.3),
             '▸ 下一页: 储蓄险 · 退保现价柱状对比 · 直观看到收益差距',
             size=9, color=C['primary_dark'], bold=True)




def _slide_savings_compare_chart(prs, products, meta, C):
    """V3.1.5 储蓄险「不提取对比」柱状图
    X 轴: 保单年度 (Y10/Y20/Y30/Y40/Y50)
    每组柱: 各产品退保现价 (USD)
    颜色: 每个产品一个颜色 (循环 broker 主题色)
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    add_title(s, '储蓄险 · 退保现价柱状对比',
              '保单年度 × 各产品 · 直观看到收益差距',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    years = [10, 20, 30, 40, 50]
    # 收集每个产品在每个年度的 CV
    series_data = []  # [{name, values:[cv@y10,...], color}]
    # broker 主题色循环 (5 色, 适配最多 5 个产品)
    color_cycle = [
        C.get('primary', '1B5E7A'),
        C.get('accent', 'C89545'),
        C.get('accent_dark', '8B5A2B'),
        C.get('success', '5BA85B'),
        C.get('warning', 'D4915B'),
    ]
    for pi, p in enumerate(products):
        d = p.get('data') or {}
        bi = d.get('benefit_illustration') or []
        name = (p.get('product_name') or p.get('pdfName') or '—')[:14]
        vals = []
        for y in years:
            r = next((rr for rr in bi if int(rr.get('policy_year') or 0) == y), None) if isinstance(bi, list) else None
            cv = float((r or {}).get('total_surrender_value') or (r or {}).get('Total') or 0)
            vals.append(cv)
        series_data.append({'name': name, 'values': vals, 'color': color_cycle[pi % len(color_cycle)]})

    # 构造 chart data
    cd = CategoryChartData()
    cd.categories = [f'Y{y}' for y in years]
    for sd in series_data:
        cd.add_series(sd['name'], sd['values'])

    # 渲染柱状图
    chart_x = Inches(0.5)
    chart_y = Inches(1.9)
    chart_w = Inches(8.5)
    chart_h = Inches(5.2)
    graphic = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        chart_x, chart_y, chart_w, chart_h, cd
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_HEI
    chart.legend.font.color.rgb = _c(C['body_text'])

    # Y 轴刻度
    va = chart.value_axis
    all_vals = [v for sd in series_data for v in sd['values'] if v > 0]
    if all_vals:
        va.maximum_scale = max(all_vals) * 1.1
        va.minimum_scale = 0
    va.tick_labels.font.size = Pt(9)
    va.tick_labels.number_format = '$#,##0,, "M"'  # 单位 M (除以 1,000,000)
    va.tick_labels.font.color.rgb = _c(C['mid_text'])
    va.tick_labels.font.name = FONT_HEI
    # X 轴
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.bold = True
    ca.tick_labels.font.name = FONT_HEI
    ca.tick_labels.font.color.rgb = _c(C['dark_text'])

    # 设置每 series 颜色
    for plot in chart.plots:
        for i, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(series_data[i]['color'])

    # 右侧解读文字
    rx = Inches(9.3); rw = Inches(4.0)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4),
             '📊 怎么看这张图', size=14, color=C['dark_text'], bold=True)
    add_text(s, rx, Inches(2.5), rw, Inches(4.5),
             '▸ X 轴 = 保单年度 (Y10/Y20/Y30/Y40/Y50)\n'
             '▸ 每组 = 该年度各产品退保现价对比\n'
             '▸ 柱越高 = 收益越好\n'
             '▸ 颜色对应左下角产品图例\n\n'
             '💡 决策要点:\n'
             '· 短期看 (Y10-Y20): 哪家优势明显?\n'
             '· 长期看 (Y30-Y50): 谁滚存能力强?\n'
             '· 注意保费规模是否一致 (假设同档)',
             size=10, color=C['body_text'])

    # 底部说明
    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书「退保发还金额」演示表 · 含保证+非保证金额 · 单位: 百万美元 (M)',
             size=8, color=C['mid_text'])




def _slide_best_pick_card(prs, prod, type_, meta, C):
    """单张最佳产品卡片 (大数字 + 产品名 + 公司 + 推荐语)
    type_ = 'savings' | 'ci' | 'iul'
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    type_title = {'savings': '最佳储蓄险', 'ci': '最佳重疾险', 'iul': '最佳 IUL'}
    add_title(s, type_title.get(type_, '最佳产品') + ' · 单卡推荐',
              '从对比章节中选出的最优产品 · 数据驱动',
              title_size=32, sub_size=13, C=C)

    d = prod.get('data') or {}
    name = prod.get('product_name') or prod.get('pdfName') or '—'

    # 计算核心指标
    metric_label = '—'
    metric_value = '—'
    sub_text = ''
    color = C['primary']

    if type_ == 'savings':
        bi = d.get('benefit_illustration') or []
        summ = d.get('summary') or {}
        pol = d.get('policy') or {}
        row20 = next((r for r in bi if int(r.get('policy_year') or 0) == 20), bi[-1] if bi else {})
        # IRR 兜底 CAGR
        irr_raw = row20.get('IRR')
        yr20 = int(row20.get('policy_year') or 20)
        cv = float(row20.get('total_surrender_value') or 0)
        tp = float(row20.get('total_premium_paid') or 0)
        if not irr_raw:
            try:
                ap = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
                pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '5'
                py = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '5') or 5)
                cur = summ.get('currency') or pol.get('currency') or 'USD'
                ma = compute_irr_ma(ap, py, cv, yr20, cur) if ap > 0 and py > 0 else None
                if ma is not None:
                    irr_raw = ma
                elif cv > 0 and tp > 0 and yr20 > 0:
                    irr_raw = (cv / tp) ** (1.0 / yr20) - 1  # CAGR 兜底
            except: irr_raw = 0
        # 单利 (simple interest) = (CV - 已缴保费) / 已缴保费 / 年限
        simple = (cv - tp) / tp / yr20 if (cv > 0 and tp > 0 and yr20 > 0) else 0
        irr = float(irr_raw or 0) * 100
        metric_label = 'Y20 复利 IRR'
        metric_value = f'{irr:.2f}%'
        sub_text = f'单利 {simple*100:.2f}% · Y20 退保现价 ${cv:,.0f}' if cv else f'单利 {simple*100:.2f}%'
        color = C['primary']
    elif type_ == 'iul':
        summ = d.get('summary') or {}
        pol = d.get('policy') or {}
        cov = float(summ.get('sum_insured') or pol.get('sum_insured') or 0)
        prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '5'
        pay_yrs = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '5') or 5)
        total_prem = prem * pay_yrs
        lev = cov / total_prem if total_prem > 0 else 0
        metric_label = '身故杠杆率'
        metric_value = f'{lev:.2f}x'
        sub_text = f'保额 ${cov:,.0f} / 总保费 ${total_prem:,.0f}'
        color = C['accent_dark']
    elif type_ == 'ci':
        summ = d.get('summary') or {}
        pol = d.get('policy') or {}
        cov = float(summ.get('sum_insured') or pol.get('sum_insured') or pol.get('basic_sum_insured') or 0)
        prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        ratio = cov / prem if prem > 0 else 0
        cov_period = pol.get('coverage_period') or '—'
        cov_items = d.get('coverage_items') or []
        major_covs = [c for c in cov_items if (c.get('amount') or 0) > 0]
        # V3.1.6+: 显示杠杆率 (主指标) + 保额 (副)
        metric_label = '杠杆率 (保/年保费)'
        metric_value = f'{ratio:.2f}x'
        sub_text = f'保额 ${cov:,.0f} · 保障 {cov_period} · {len(major_covs)} 项主要保障'
        color = C['primary_dark']

    # 大数字卡片 (居中)
    card_w = Inches(8.5); card_h = Inches(4.0)
    sx = (SLIDE_W - card_w) / 2
    sy = Inches(2.2)
    add_rect(s, sx, sy, card_w, card_h, fill=C['bg_card'])
    add_rect(s, sx, sy, card_w, Inches(0.6), fill=color)
    add_text(s, sx + Inches(0.3), sy + Inches(0.15), card_w - 0.6, Inches(0.4),
             type_title.get(type_, '最佳产品'), size=14, color=C['bg_light'], bold=True)

    add_text(s, sx + Inches(0.3), sy + Inches(0.85), card_w - 0.6, Inches(0.4),
             name, size=20, color=C['dark_text'], bold=True, align=PP_ALIGN.CENTER)

    # V3.1.8+: 储蓄险双指标大字 (Y20 复利 IRR + Y20 单利), CI/IUL 单指标大字
    if type_ == 'savings':
        # 双指标并排显示 (同样 56pt 大字)
        half_w = (card_w - Inches(0.6)) / 2
        # 左: 复利 IRR
        left_x = sx + Inches(0.3)
        add_text(s, left_x, sy + Inches(1.45), half_w, Inches(0.3),
                 '20年末 · 复利 IRR', size=12, color=C['mid_text'], align=PP_ALIGN.CENTER, bold=True)
        add_text(s, left_x, sy + Inches(1.78), half_w, Inches(0.95),
                 metric_value, size=56, color=color, bold=True, align=PP_ALIGN.CENTER)
        # 中线分割
        sep_x = sx + Inches(0.3) + half_w
        add_rect(s, sep_x, sy + Inches(1.5), Inches(0.02), Inches(2.3), fill=C['bg_light'])
        # 右: 单利
        right_x = sep_x + Inches(0.15)
        right_w = half_w - Inches(0.15)
        add_text(s, right_x, sy + Inches(1.45), right_w, Inches(0.3),
                 '20年末 · 单利', size=12, color=C['mid_text'], align=PP_ALIGN.CENTER, bold=True)
        # 从 sub_text "单利 X.XX% · ..." 中提取单利率
        import re as _re
        sm = _re.search(r'单利\s*([\d.]+)%', sub_text or '')
        simple_str = f'{sm.group(1)}%' if sm else '—'
        add_text(s, right_x, sy + Inches(1.78), right_w, Inches(0.95),
                 simple_str, size=56, color=C['accent_dark'], bold=True, align=PP_ALIGN.CENTER)
        # 底部脚注: 仅保留 CV / 保额 文字
        cv_text = _re.sub(r'单利\s*[\d.]+%\s*·\s*', '', sub_text or '').strip(' ·')
        if cv_text:
            add_text(s, sx + Inches(0.3), sy + Inches(3.0), card_w - 0.6, Inches(0.4),
                     cv_text, size=11, color=C['body_text'], align=PP_ALIGN.CENTER)
    else:
        # CI / IUL 单指标大字 (沿用原版)
        add_text(s, sx + Inches(0.3), sy + Inches(1.5), card_w - 0.6, Inches(0.4),
                 metric_label, size=12, color=C['mid_text'], align=PP_ALIGN.CENTER)
        add_text(s, sx + Inches(0.3), sy + Inches(1.9), card_w - 0.6, Inches(1.2),
                 metric_value, size=64, color=color, bold=True, align=PP_ALIGN.CENTER)

        if sub_text:
            add_text(s, sx + Inches(0.3), sy + Inches(3.2), card_w - 0.6, Inches(0.4),
                     sub_text, size=12, color=C['body_text'], align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.3),
             f'注: 按 {metric_label} 在所有同类产品中排序最高 · 数据来源: 各产品官方计划书',
             size=8, color=C['mid_text'])




def _slide_ci_combined_chart(prs, products, meta, C):
    """V3.1.7 重疾险四合一柱状对比图 (年缴保费 / 初始保额 / 30年后保额 / 杠杆率)
    一张柱状图, 每产品一组, 每组 4 根柱 (4 个指标)
    第 3 柱 = 30 年后保额: 体现 HK CI 升级保障 + 分红带来保额增长的核心优势
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    meta['product_type'] = 'ci'

    add_title(s, '重疾险 · 核心指标四合一对比',
              '年缴保费 / 初始保额 / 30年后保额 / 杠杆率 · 一图看全核心差异',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    cats = []
    annual_prems = []
    sum_insured_initial = []
    sum_insured_y30 = []
    leverages = []
    for p in products:
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        summ = d.get('summary') or {}
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '—')[:14]
        cov = float(summ.get('sum_insured') or pol.get('sum_insured') or pol.get('basic_sum_insured') or 0)
        prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        lev = cov / prem if prem > 0 else 0
        # V3.1.7: 30 年后保额 = benefit_illustration 里 policy_year 30 (or 最接近) 的 death_benefit
        # HK CI 因升级保障 + 周年红利持续注入,death_benefit 持续增长;内地 CI 多为定额不变
        bi = d.get('benefit_illustration') or []
        cov_y30 = 0
        closest_y = None
        for r in bi:
            try:
                ry = int(r.get('policy_year') or 0)
                db = float(r.get('death_benefit') or 0)
            except Exception:
                continue
            if db <= 0:
                continue
            if ry == 30:
                cov_y30 = db
                break
            # 兜底: 取最接近 30 的年份
            if closest_y is None or abs(ry - 30) < abs(closest_y - 30):
                closest_y = ry
                cov_y30 = db
        cats.append(name)
        annual_prems.append(prem)
        sum_insured_initial.append(cov)
        sum_insured_y30.append(cov_y30)
        leverages.append(lev)

    if not cats:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无可对比数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series('年缴保费 (USD)', annual_prems)
    cd.add_series('初始保额 (USD) Y1', sum_insured_initial)
    # V3.1.7: 30 年后保额 = HK CI 升级保障 + 周年红利累计效果;内地 CI 多为定额不变
    cd.add_series('30年后保额 (USD) Y30', sum_insured_y30)
    # 杠杆率 × 1000 放到主轴 (避免被 USD 100K 压扁), 数据标签显示实际倍数
    scaled_leverages = [v * 1000 for v in leverages]
    cd.add_series('杠杆率 × 1000 (= 倍数)', scaled_leverages)

    graphic = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.85), Inches(8.5), Inches(5.0), cd,
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_HEI
    chart.legend.font.color.rgb = _c(C['body_text'])

    # Y 轴 (USD - 主轴)
    va = chart.value_axis
    va.tick_labels.font.size = Pt(9)
    va.tick_labels.font.name = FONT_HEI
    va.tick_labels.font.color.rgb = _c(C['mid_text'])
    va.tick_labels.number_format = '$#,##0'
    # X 轴
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.bold = True
    ca.tick_labels.font.name = FONT_HEI
    ca.tick_labels.font.color.rgb = _c(C['dark_text'])

    # 每 series 不同颜色: 年缴保费=primary, 初始保额=accent, 30年后保额=danger, 杠杆率=warning
    series_colors = [
        C.get('primary'),
        C.get('accent'),
        C.get('red', 'B03A3A'),
        C.get('warning', 'D4915B'),
    ]
    for plot in chart.plots:
        for si, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(series_colors[si % len(series_colors)])

    # V3.1.7: 杠杆率柱放主轴 (× 1000), 数据标签显示实际倍数
    chart.has_title = False
    for plot in chart.plots:
        for si, ser in enumerate(plot.series):
            # 数据标签: 柱顶显示实际值
            dl = ser.data_labels
            dl.show_value = True
            dl.font.size = Pt(9)
            dl.font.name = FONT_HEI
            dl.font.bold = True
            dl.font.color.rgb = _c(C['dark_text'])
            # 杠杆率 series (si==3): 数据是 actual×1000, 用 , 格式去除千分位 (= ÷1000)
            if si == 3:
                dl.number_format = '0.00,"x"'
            else:
                dl.number_format = '"$"#,##0'
            # 标签位置: 外侧上方
            try:
                from pptx.enum.chart import XL_LABEL_POSITION
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass

    # 右侧解读
    rx = Inches(9.3); rw = Inches(4.0)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4),
             '📊 四指标一图', size=16, color=C['dark_text'], bold=True)
    add_text(s, rx, Inches(2.5), rw, Inches(4.5),
             '▸ X 轴 = 各产品\n'
             '▸ 每组 4 根柱: 年缴保费 / 初始保额 /\n'
             '   30年后保额 / 杠杆率\n'
             '▸ Y 轴统一用 USD 单位 (杠杆率 ≈ 倍数)\n\n'
             '💡 怎么读:\n'
             '· 年缴保费越低 → 越便宜\n'
             '· 初始保额越高 → 起跑保障越多\n'
             '· 30年后保额 > 初始保额 → HK CI 升级\n'
             '   红利持续加保 (vs 内地 CI 定额不变)\n'
             '· 杠杆率 (保/年保费) 越高 → 性价比越好\n\n'
             '⚠️ 注意:\n'
             '· 升级红利部分为非保证,演示表仅供参考\n'
             '· 不同产品保障范围不同,杠杆率不能孤立看',
             size=9, color=C['body_text'])

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书 benefit_illustration · 杠杆率柱 = 实际倍数 × 1000 (标签显示真实值) · 年缴保费/保额=USD · 30 年后保额演示为非保证红利',
             size=8, color=C['mid_text'])




def _slide_ci_overview_multi(prs, products, meta, C):
    """V3.1.6+ 重疾险多产品首屏 — 对比产品一览
    顶部「重疾险对比」主标题 + 公司·产品名横幅, 下方表格列每产品核心数据
    列 (V3.1.8+): # / 公司 / 产品名 / 缴费年期 / 年缴额 / 总缴保费 / 初始保额 / 保额增额 / 杠杆率 / 保障年期 / 保障项数
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    # 强制 prod_type = ci (影响 header 等)
    meta['product_type'] = 'ci'

    # 大标题
    add_text(s, Inches(0.7), Inches(1.0), Inches(12.5), Inches(0.6),
             '重疾险对比', size=32, color=C['dark_text'], bold=True, font=FONT_HEI)

    # 副标题: 公司·产品名
    names = []
    for p in products:
        company = (p.get('company_short') or '').upper() or '—'
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '—')[:14]
        names.append(f'{company} · {name}')
    sub_text = f'共 {len(products)} 款 · ' + ' / '.join(names)
    add_text(s, Inches(0.7), Inches(1.65), Inches(12.5), Inches(0.4),
             sub_text, size=11, color=C['mid_text'], font=FONT_HEI)

    if not products:
        return

    # 收集每行数据
    rows_data = []
    for p in products:
        d = p.get('data') or {}
        summ = d.get('summary') or {}
        pol = d.get('policy') or {}
        cov_items = d.get('coverage_items') or []
        pay_raw = summ.get('payment_years') or pol.get('premium_payment_period') or '10'
        pay_yrs = 1 if pay_raw == '趸交' else (int(re.sub(r'[^\d]', '', str(pay_raw)) or '10') or 10)
        annual_prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        # 总缴保费 = annual * pay_yrs, 兜底用 total_premium_with_levy
        total_prem = float(pol.get('total_premium_with_levy') or (annual_prem * pay_yrs) or 0)
        cov = float(summ.get('sum_insured') or pol.get('sum_insured') or pol.get('basic_sum_insured') or 0)
        # V3.1.8+: 保额增额 = Y30 death_benefit - 初始保额 (HK CI 升级保障 + 分红累计效果)
        bi = d.get('benefit_illustration') or []
        cov_y30 = 0
        closest_y = None
        for r in bi:
            try:
                ry = int(r.get('policy_year') or 0)
                db = float(r.get('death_benefit') or 0)
            except Exception:
                continue
            if db <= 0:
                continue
            if ry == 30:
                cov_y30 = db
                break
            if closest_y is None or abs(ry - 30) < abs(closest_y - 30):
                closest_y = ry
                cov_y30 = db
        sum_increment = max(0, cov_y30 - cov) if cov_y30 > 0 else 0
        lev = cov / annual_prem if annual_prem > 0 else 0
        cov_period = pol.get('coverage_period') or '—'
        major_covs = [c for c in cov_items if (c.get('amount') or 0) > 0]
        company = (p.get('company_short') or '').upper() or '—'
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '—')[:16]
        rows_data.append({
            'company': company,
            'name': name,
            'pay_yrs': pay_yrs,
            'annual_prem': annual_prem,
            'total_prem': total_prem,
            'sum_insured': cov,
            'sum_increment': sum_increment,
            'leverage': lev,
            'cov_period': cov_period,
            'cov_count': len(major_covs),
        })

    n = len(rows_data)
    headers = ['#', '公司', '产品名称', '缴费\n年期', '年缴额\n(USD)',
               '总缴保费\n(USD)', '初始保额\n(USD)', '保额增额\n(Y30-Y1)',
               '杠杆率\n(保/年)', '保障年期', '保障项数']
    n_rows = 1 + n

    table_x = Inches(0.4); table_y = Inches(2.0)
    table_w = Inches(12.9)
    row_h = Inches(0.5)
    tbl_shape = s.shapes.add_table(n_rows, len(headers), table_x, table_y, table_w, row_h * n_rows)
    tbl = tbl_shape.table
    col_widths = [0.35, 0.6, 1.7, 0.65, 0.95, 1.15, 1.15, 1.25, 0.95, 0.85, 0.75]  # 11 列, total 10.35
    total_cw = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * 12.9 / total_cw)
    tbl.rows[0].height = Inches(0.55)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.5)

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行 + 找每列最大值
    col_max_idx = {}
    # V3.1.8+: 增加保额增额列 (sum_increment) 进 max 标红
    for ci_idx, key in enumerate(['annual_prem', 'total_prem', 'sum_insured', 'sum_increment', 'leverage', 'cov_count'], start=4):
        max_v, max_ri = -1, -1
        for ri, row in enumerate(rows_data, start=1):
            v = row.get(key, 0)
            if v > max_v: max_v = v; max_ri = ri
        if max_ri > 0: col_max_idx[ci_idx] = max_ri

    for ri, row in enumerate(rows_data, start=1):
        # 行底色
        bg = C['bg_card'] if ri % 2 == 1 else C['bg_light']
        # V3.1.8+: 保额增额列显示 "$X (Y30 比初始 +Y%)" 直观展示增长
        inc = row['sum_increment']
        if inc > 0 and row['sum_insured'] > 0:
            inc_pct = (inc / row['sum_insured']) * 100
            inc_str = f'+${inc:,.0f}\n(+{inc_pct:.0f}%)'
        else:
            inc_str = '—'
        vals = [
            str(ri - 1), row['company'], row['name'],
            f"{row['pay_yrs']}年",
            f"${row['annual_prem']:,.0f}",
            f"${row['total_prem']:,.0f}",
            f"${row['sum_insured']:,.0f}",
            inc_str,
            f"{row['leverage']:.2f}x",
            str(row['cov_period']),
            str(row['cov_count']),
        ]
        for ci, val in enumerate(vals):
            cell = tbl.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci <= 3 or ci == 10 else PP_ALIGN.RIGHT
            # 保额增额列允许多行 (\n 分隔)
            if ci == 7 and '\n' in val:
                is_max = col_max_idx.get(ci) == ri
                run = p.add_run(); run.text = val.split('\n')[0]
                run.font.size = Pt(9); run.font.bold = is_max; run.font.name = FONT_HEI
                run.font.color.rgb = _c(C['red'] if is_max else C['body_text'])
                p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT
                run2 = p2.add_run(); run2.text = val.split('\n')[1]
                run2.font.size = Pt(7)
                run2.font.color.rgb = _c(C['red'] if is_max else C['mid_text'])
                run2.font.name = FONT_HEI
                continue
            run = p.add_run(); run.text = val
            run.font.size = Pt(9)
            run.font.name = FONT_HEI
            is_max = col_max_idx.get(ci) == ri
            run.font.color.rgb = _c(C['red'] if is_max else C['body_text'])
            run.font.bold = is_max

    # 脚注 (V3.1.8+: 增加保额增额口径说明)
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.3),
             '注: 红色 = 该列最高值 · 总缴保费 = 年缴额 × 缴费年期 · 保额增额 = Y30身故赔偿 - 初始保额 (HK CI 升级保障+周年红利累计效果)',
             size=8, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.3), Inches(12.5), Inches(0.3),
             f'▸ 杠杆率 = 初始保额 / 年缴保费 · 数据来源: 各产品官方计划书',
             size=8, color=C['mid_text'])


# 2026-07-16 V3.3.11+: 对比章节的合并 meta — header 显示 "A VS B VS C" 而非单公司
def _build_compare_meta(exts: List[Dict], base_meta: Dict) -> Dict:
    """从 exts (per-product payloads) 提取每个产品的 brand_profile.short / short_en,
    合并成 "A VS B" 格式, 用作对比章节的 meta (header 显示)

    - exts: [{brand_profile, company_short, ...}, ...]
    - base_meta: 渲染器主 meta (作为 fallback)
    """
    shorts_zh: List[str] = []
    shorts_en: List[str] = []
    for e in exts or []:
        bp = e.get('brand_profile') or {}
        sz = bp.get('short') or bp.get('name_zh') or e.get('company_short') or ''
        se = bp.get('short_en') or e.get('company_short_en') or ''
        # 去重保序
        if sz and sz not in shorts_zh:
            shorts_zh.append(sz)
        if se and se not in shorts_en:
            shorts_en.append(se)
    cmp_meta = dict(base_meta or {})
    # 2 个产品: 'A  VS  B' (双空格让分隔更显眼); 3+ 个: 'A  VS  B  VS  C'
    cmp_meta['company_short'] = '  VS  '.join(shorts_zh) if shorts_zh else base_meta.get('company_short', '')
    cmp_meta['company_short_en'] = '  VS  '.join(shorts_en) if shorts_en else base_meta.get('company_short_en', '')
    cmp_meta['company_id'] = ''  # 对比模式无单一 company, 清空避免触发 logo 角标
    cmp_meta['_comparison_mode'] = True
    return cmp_meta


# 2026-07-16 V3.3.10+: 重疾险保障亮点对比 (产品营销内容, 一次性手工整理)
# 内容来自产品官方计划书, 不会随 LLM 抽取结果变化 — 每个产品独立列出独特保障优势
# (什么疾病多次赔付 / 赔多少 / 特别条款), 用于替代通用 coverage_items 表
CI_COVERAGE_HIGHLIGHTS: Dict[str, Dict] = {
    # CTF 周大福人寿 — 「守X家倍198」危疾保障计划 (HK)
    'ctf_jiaba198': {
        'product_label': '「守X家倍198」危疾保障计划',
        'company_zh': '周大福人寿',
        'key_stats': [
            ('基础保额', '$100,000'),
            ('年缴保费', '$4,949'),
            ('缴费年期', '10 年'),
            ('保障年期', '至 100 岁'),
        ],
        'tagline': '癌症/中风/心脏病多次赔 · 保额自动还原',
        'features': [
            ('保障还原利益',
             '早期危疾赔付后保额自动还原高达 100% 保额，保障至 70 岁'),
            ('严重都市疾病额外保障',
             '癌症 3 次 + 中风/严重心脏病 3 次，每次 100% 保额 (约 $100K/次)，保障至 85 岁'),
            ('严重都市疾病无限次增值保障',
             '同类重疾不限增值次数，每次增值后再赔 3 次癌症 + 3 次中风/心脏病'),
            ('额外身故/生存赔偿',
             '投保首 10 年额外 60% 保额 ($60K) 一次性给付'),
            ('早期危疾保障',
             '50% 保额 (约 $50K)，覆盖原位癌等早期病变'),
            ('保费豁免双轨',
             '危疾确诊豁免余下保费 + 配偶意外身故豁免保费'),
        ],
    },
    # AIA 友邦 — 「爱伴航」保险计划 2 (HK)
    'aia_aihang2': {
        'product_label': '「爱伴航」保险计划 2',
        'company_zh': '友邦保险',
        'key_stats': [
            ('基础保额', '$100,000'),
            ('年缴保费', '$7,008'),
            ('缴费年期', '10 年'),
            ('保障年期', '终身'),
        ],
        'tagline': '58 种危疾 + ICU + 10X 多重危疾 · 全方位守护',
        'features': [
            ('疾病覆盖广',
             '58 种严重疾病 (57 严重 + 1 非严重) + 44 种早期危疾 + ICU 深切治疗'),
            ('深切治疗保障 2 级',
             '级别一 $50K (ICU 连续 72h+) · 级别二 100% 保额 (120h+ + 复杂手术)'),
            ('10X 多重危疾赔偿',
             '癌症/心脏病/中风/脑退化/柏金逊症最多 10 次 100% 保额'),
            ('持续癌症现金选项',
             '确诊癌症每月领 5% 保额，最长 100 个月 (至 85 岁)'),
            ('脑退化/柏金逊终身年金',
             '每年 6% 保额赔付至保单终止'),
            ('双倍早期危疾',
             '65 岁后确诊早期危疾额外 1 倍赔偿 (上限 $25K，限 1 次)'),
            ('配偶身故豁免保费',
             '75 岁前配偶身故豁免基本计划余下保费'),
            ('首 10 年升级保障',
             '额外 35% 保额 (身故/严重疾病/深切治疗二级)'),
        ],
    },
}


def _slide_ci_coverage_highlights(prs, products, meta, C):
    """V3.3.10+ 重疾险保障亮点对比 (产品营销内容, 手工整理, 不依赖 LLM 抽取)

    对每个产品查 CI_COVERAGE_HIGHLIGHTS, 命中则渲染 2 列布局 (公司色块 + 关键数据 + 8 项保障亮点),
    未命中则渲染提示用旧版 coverage_items 表
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    meta['product_type'] = 'ci'

    add_title(s, '重疾险 · 保障亮点对比',
              '各产品独特保障优势 · 一次看清谁更适合你的客户',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    # 查 CI_COVERAGE_HIGHLIGHTS, 用 product_label / pdfName 模糊匹配
    def _lookup_highlight(p):
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        prod_name = pol.get('product_name') or d.get('product_name') or p.get('product_name') or ''
        pdf_name = (p.get('pdfName') or '').lower()
        # 优先按产品名匹配
        for key, info in CI_COVERAGE_HIGHLIGHTS.items():
            label = info.get('product_label') or ''
            if label and (label in prod_name or prod_name in label):
                return info
        # 兜底按 pdfName + key 关键词匹配
        if 'ctf' in pdf_name or 'jiaba' in pdf_name or '家倍' in prod_name or '守' in prod_name:
            return CI_COVERAGE_HIGHLIGHTS.get('ctf_jiaba198')
        if 'aia' in pdf_name or 'aihang' in pdf_name or '愛伴航' in prod_name or '爱伴航' in prod_name or '愛伴' in prod_name:
            return CI_COVERAGE_HIGHLIGHTS.get('aia_aihang2')
        return None

    cols_info = []
    for p in products:
        info = _lookup_highlight(p)
        if info:
            cols_info.append(info)
        else:
            # 没找到: 用产品名当 placeholder, 让用户知道哪些产品未配置
            d = p.get('data') or {}
            pol = d.get('policy') or {}
            cols_info.append({
                'product_label': (pol.get('product_name') or p.get('pdfName') or '—')[:18],
                'company_zh': (p.get('brand_profile') or {}).get('name_zh') or '—',
                'key_stats': [],
                'tagline': '（该产品保障亮点待手工整理）',
                'features': [],
            })

    # 2 列布局: 每列 ~6.2" 宽, 留 0.5" 左右边距 + 0.2" 列间隔
    n = len(cols_info)
    col_w = 6.2
    col_gap = 0.2
    total_w = col_w * n + col_gap * (n - 1)
    left_margin = (13.33 - total_w) / 2

    for idx, info in enumerate(cols_info):
        col_x = left_margin + idx * (col_w + col_gap)
        # 顶部色带: 公司色 (CTF=dark green, AIA=red, 其他=primary)
        company_color = C.get('primary_dark', '#A0522D')
        if '周大福' in info.get('company_zh', '') or 'CTF' in info.get('product_label', ''):
            company_color = '#1F6B47'  # CTF green
        elif '友邦' in info.get('company_zh', '') or 'AIA' in info.get('product_label', ''):
            company_color = '#B91C1C'  # AIA red
        # 顶部色块: 公司名 + 产品名
        add_rect(s, Inches(col_x), Inches(2.0), Inches(col_w), Inches(0.7), fill=company_color)
        # 公司名 (左上小字)
        add_text(s, Inches(col_x + 0.15), Inches(2.05), Inches(col_w - 0.3), Inches(0.25),
                 info.get('company_zh', ''), size=10, color='#FFFFFF', bold=True)
        # 产品名 (中下大字)
        add_text(s, Inches(col_x + 0.15), Inches(2.30), Inches(col_w - 0.3), Inches(0.4),
                 info.get('product_label', ''), size=14, color='#FFFFFF', bold=True)
        # tagline (色带下小字)
        if info.get('tagline'):
            add_text(s, Inches(col_x), Inches(2.75), Inches(col_w), Inches(0.3),
                     info['tagline'], size=10, color=C['mid_text'], align=PP_ALIGN.CENTER)

        # 关键数据 4 列 (保额/保费/年期/保障期)
        if info.get('key_stats'):
            stat_w = col_w / len(info['key_stats'])
            for si, (lbl, val) in enumerate(info['key_stats']):
                sx = col_x + si * stat_w
                add_rect(s, Inches(sx), Inches(3.1), Inches(stat_w - 0.05), Inches(0.55),
                         fill=C['bg_card'])
                add_text(s, Inches(sx), Inches(3.12), Inches(stat_w - 0.05), Inches(0.22),
                         lbl, size=8, color=C['mid_text'], align=PP_ALIGN.CENTER)
                add_text(s, Inches(sx), Inches(3.32), Inches(stat_w - 0.05), Inches(0.3),
                         val, size=11, color=C['dark_text'], bold=True, align=PP_ALIGN.CENTER)

        # 保障亮点列表: 从 3.85" 起, 每项 ~0.42"
        feats = info.get('features') or []
        feat_y = 3.85
        feat_h = 0.42
        if not feats:
            add_text(s, Inches(col_x), Inches(feat_y), Inches(col_w), Inches(0.4),
                     '（暂无结构化保障亮点，可在该产品配置 CI_COVERAGE_HIGHLIGHTS）',
                     size=10, color=C['mid_text'])
        else:
            for fi, (ftitle, fdesc) in enumerate(feats):
                fy = feat_y + fi * feat_h
                # 编号徽章
                badge_color = company_color
                add_rect(s, Inches(col_x), Inches(fy), Inches(0.32), Inches(0.32),
                         fill=badge_color)
                add_text(s, Inches(col_x), Inches(fy + 0.02), Inches(0.32), Inches(0.3),
                         f'{fi+1}', size=11, color='#FFFFFF', bold=True, align=PP_ALIGN.CENTER)
                # 标题
                add_text(s, Inches(col_x + 0.4), Inches(fy - 0.02), Inches(col_w - 0.4), Inches(0.22),
                         ftitle, size=10, color=C['dark_text'], bold=True)
                # 描述
                add_text(s, Inches(col_x + 0.4), Inches(fy + 0.18), Inches(col_w - 0.4), Inches(0.24),
                         fdesc, size=8.5, color=C['body_text'])

    # 底部注释
    add_text(s, Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.3),
             '注: 保障亮点已根据各产品官方计划书手工整理 · 详细条款请参阅保单契约',
             size=8, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.3), Inches(12.5), Inches(0.3),
             '▸ 未配置的产品会展示占位提示, 后续可在 CI_COVERAGE_HIGHLIGHTS 字典中添加',
             size=8, color=C['mid_text'])




def _slide_ci_coverage_compare_multi(prs, products, meta, C):
    """V3.1.6+ 重疾险主要保障项目横向对比表 (多产品)
    列出每产品的保障项目 (label · amount · percentage), 红色标最高保额
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    meta['product_type'] = 'ci'

    add_title(s, '重疾险 · 主要保障项目对比',
              '各产品主要保障项目保额对比 (USD) · 红色 = 同项目最高保额',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    # 收集每产品主要保障项目 (amount > 0)
    prod_covs = []  # [{name, items:[{label, amount, pct}], major_count}]
    all_labels = []  # 全部出现过的 label (按出现顺序, 取首次)
    label_seen = set()
    for p in products:
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '—')[:14]
        items = d.get('coverage_items') or []
        major = [c for c in items if (c.get('amount') or 0) > 0]
        prod_covs.append({'name': name, 'items': major})
        for it in major:
            lbl = it.get('label') or ''
            if lbl and lbl not in label_seen:
                all_labels.append(lbl)
                label_seen.add(lbl)

    # 表头: 保障项目 | 各产品保额 (取前 5 个最常见 label, 避免表过宽)
    MAX_LABELS = 5
    show_labels = all_labels[:MAX_LABELS]
    headers = ['保障项目'] + [pc['name'] for pc in prod_covs]

    # 行数: 1 表头 + N labels
    n_rows = 1 + len(show_labels)
    n_cols = 1 + len(prod_covs)

    table_x = Inches(0.4); table_y = Inches(2.0)
    table_w = Inches(12.9)
    row_h = Inches(0.55)
    tbl_shape = s.shapes.add_table(n_rows, n_cols, table_x, table_y, table_w, row_h * n_rows)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.5)
    rem_inches = (12.9 - 3.5) / max(1, len(prod_covs))
    for ci in range(1, n_cols):
        tbl.columns[ci].width = Inches(rem_inches)
    tbl.rows[0].height = Inches(0.5)
    for ri in range(1, n_rows):
        tbl.rows[ri].height = Inches(0.55)

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = h
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行: 每个 label 对应一行
    for ri, lbl in enumerate(show_labels, start=1):
        bg = C['bg_card'] if ri % 2 == 1 else C['bg_light']
        # 该 label 在每产品的金额
        col_vals = []
        for pc in prod_covs:
            v = 0
            for it in pc['items']:
                if it.get('label') == lbl:
                    v = float(it.get('amount') or 0)
                    break
            col_vals.append(v)
        max_v = max(col_vals) if col_vals else 0
        # 保障项目列
        cell = tbl.cell(ri, 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = lbl
        run.font.size = Pt(10); run.font.color.rgb = _c(C['dark_text'])
        run.font.name = FONT_HEI; run.font.bold = True
        # 各产品金额列
        for ci, v in enumerate(col_vals, start=1):
            cell = tbl.cell(ri, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(bg)
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f'${v:,.0f}' if v > 0 else '—'
            run.font.size = Pt(11); run.font.name = FONT_HEI
            is_max = (v == max_v and v > 0 and len(prod_covs) >= 2)
            run.font.color.rgb = _c(C['red'] if is_max else C['body_text'])
            run.font.bold = is_max

    # 脚注
    add_text(s, Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.3),
             f'注: 红色 = 该保障项目保额最高的产品 · 仅展示各产品前 {MAX_LABELS} 项主要保障 · 数据来源: 各产品官方计划书',
             size=8, color=C['mid_text'])
    add_text(s, Inches(0.5), Inches(7.2), Inches(12.5), Inches(0.3),
             f'▸ 完整保障项目请见各产品详情页 · 部分产品早期危疾/保费豁免等条款性保障未在表中体现',
             size=8, color=C['mid_text'])




def _slide_iul_combined_chart(prs, products, meta, C):
    """V3.1.7 IUL 四合一柱状对比图 (年缴保费 / 身故赔偿 / 杠杆率 / 第10年现金价值)
    一张柱状图, 每产品一组, 每组 4 根柱
    第 4 柱 = 第 10 年现金价值: 体现 IUL 区别于普通寿险的关键差异 (现金价值积累)
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    meta['product_type'] = 'iul'

    add_title(s, 'IUL 万用寿险 · 核心指标四合一对比',
              '年缴保费 / 身故赔偿 / 杠杆率 / 第10年现金价值 · 一图看全核心差异',
              title_size=30, sub_size=12, C=C)

    if not products:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无产品数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    cats = []
    annual_prems = []
    sum_insureds = []
    leverages = []
    cv_y10s = []
    for p in products:
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        summ = d.get('summary') or {}
        ins = d.get('insured') or {}
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '新加坡IUL')[:16]
        age = ins.get('age') or pol.get('insured_age') or 0
        cov = float(summ.get('sum_insured') or pol.get('sum_insured') or 0)
        prem = float(summ.get('annual_premium') or pol.get('annual_premium') or 0)
        # 缴费年期: 趸交=1, 否则从字符提数字
        pay_raw = pol.get('premium_payment_period') or summ.get('payment_years') or 10
        if '趸交' in str(pay_raw) or 'Single' in str(pay_raw) or str(pay_raw).lower() == 'sgl':
            pay_yrs = 1
        else:
            digits = ''.join(c for c in str(pay_raw) if c.isdigit())
            pay_yrs = int(digits) if digits else 10
        total_prem = prem * pay_yrs
        lev = (cov / total_prem) if total_prem > 0 else 0
        # 第 10 年现金价值: 取 policy_year=10 (or 最接近) 的 non_guaranteed_cash_value
        bi = d.get('benefit_illustration') or []
        cv_y10 = 0
        closest_y = None
        for r in bi:
            try:
                ry = int(r.get('policy_year') or 0)
                cv = float(r.get('non_guaranteed_cash_value') or r.get('cash_value') or r.get('non_guaranteed_account_value') or r.get('account_value') or 0)
            except Exception:
                continue
            if cv <= 0:
                continue
            if ry == 10:
                cv_y10 = cv
                break
            if closest_y is None or abs(ry - 10) < abs(closest_y - 10):
                closest_y = ry
                cv_y10 = cv
        cats.append(name)
        annual_prems.append(prem)
        sum_insureds.append(cov)
        leverages.append(lev)
        cv_y10s.append(cv_y10)

    if not cats:
        add_text(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.5),
                 '无可对比数据', size=14, color=C['mid_text'], align=PP_ALIGN.CENTER)
        return

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series('年缴保费 (USD)', annual_prems)
    cd.add_series('身故赔偿 (USD)', sum_insureds)
    # 杠杆率 × 1000 放到主轴 (避免被 USD 压扁), 数据标签显示实际倍数
    scaled_leverages = [v * 1000 for v in leverages]
    cd.add_series('杠杆率 × 1000 (= 倍数)', scaled_leverages)
    cd.add_series('第10年现金价值 (USD)', cv_y10s)

    graphic = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.85), Inches(8.5), Inches(5.0), cd,
    )
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_HEI
    chart.legend.font.color.rgb = _c(C['body_text'])

    va = chart.value_axis
    va.tick_labels.font.size = Pt(9)
    va.tick_labels.font.name = FONT_HEI
    va.tick_labels.font.color.rgb = _c(C['mid_text'])
    va.tick_labels.number_format = '$#,##0'
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.bold = True
    ca.tick_labels.font.name = FONT_HEI
    ca.tick_labels.font.color.rgb = _c(C['dark_text'])

    series_colors = [
        C.get('primary'),
        C.get('accent'),
        C.get('warning', 'D4915B'),
        C.get('red', 'B03A3A'),
    ]
    for plot in chart.plots:
        for si, ser in enumerate(plot.series):
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = _c(series_colors[si % len(series_colors)])

    for plot in chart.plots:
        for si, ser in enumerate(plot.series):
            dl = ser.data_labels
            dl.show_value = True
            dl.font.size = Pt(9)
            dl.font.name = FONT_HEI
            dl.font.bold = True
            dl.font.color.rgb = _c(C['dark_text'])
            if si == 2:  # 杠杆率 series
                dl.number_format = '0.00,"x"'
            else:
                dl.number_format = '"$"#,##0'
            try:
                from pptx.enum.chart import XL_LABEL_POSITION
                dl.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass

    rx = Inches(9.3); rw = Inches(4.0)
    add_text(s, rx, Inches(2.0), rw, Inches(0.4),
             '📊 四指标一图', size=16, color=C['dark_text'], bold=True)
    add_text(s, rx, Inches(2.5), rw, Inches(4.5),
             '▸ X 轴 = 各产品\n'
             '▸ 每组 4 根柱: 年缴保费 / 身故赔偿 /\n'
             '   杠杆率 / 第10年现金价值\n'
             '▸ Y 轴统一用 USD 单位 (杠杆率 ≈ 倍数)\n\n'
             '💡 怎么读:\n'
             '· 年缴保费越低 → 越便宜\n'
             '· 身故赔偿越高 → 传承杠杆越大\n'
             '· 杠杆率 (身故/总保费) 越高 → 越划算\n'
             '· 第10年现金价值 → 流动性储备\n\n'
             '⚠️ 注意:\n'
             '· 现金价值部分为非保证演示 (指数表现)\n'
             '· 不同受保人年龄/性别 → 保费差异显著',
             size=11, color=C['body_text'])

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书 benefit_illustration · 杠杆率柱 = 实际倍数 × 1000 (标签显示真实值) · 现金价值演示为非保证红利',
             size=8, color=C['mid_text'])




def _slide_iul_overview_multi(prs, products, meta, C):
    """V3.1.7 IUL 多产品首屏 — 横向对比一览 (10 列)
    列: # / 公司 / 产品名 / 受保人 / 缴费年期 / 年缴额 / 总缴保费 / 身故赔偿 / 杠杆率 / 保障年期
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    meta['product_type'] = 'iul'

    add_text(s, Inches(0.7), Inches(1.0), Inches(12.5), Inches(0.6),
             'IUL 万用寿险对比', size=32, color=C['dark_text'], bold=True, font=FONT_HEI)

    names = []
    for p in products:
        company = (p.get('company_short') or '').upper() or '—'
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '新加坡IUL')[:18]
        names.append(f'{company} · {name}')
    # V3.1.7+: 副标题只显示公司+款数, 避免产品名重复冗余 (公司已含在产品名 column 里)
    companies = sorted(set((p.get('company_short') or '').upper() for p in products if p.get('company_short')))
    sub_text = f'共 {len(products)} 款 · ' + (' / '.join(companies) if companies else '') + f' · 关键参数对比'
    add_text(s, Inches(0.7), Inches(1.65), Inches(12.5), Inches(0.4),
             sub_text, size=11, color=C['mid_text'], font=FONT_HEI)

    rows_data = []
    for p in products:
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        summ = d.get('summary') or {}
        ins = d.get('insured') or {}
        pay_raw = pol.get('premium_payment_period') or summ.get('payment_years') or 10
        if '趸交' in str(pay_raw) or 'Single' in str(pay_raw) or str(pay_raw).lower() == 'sgl':
            pay_str = '趸交'
            pay_yrs = 1
        else:
            digits = ''.join(c for c in str(pay_raw) if c.isdigit())
            pay_yrs = int(digits) if digits else 10
            pay_str = f'{pay_yrs}年缴'
        annual_prem = float(pol.get('annual_premium') or summ.get('annual_premium') or 0)
        total_prem = annual_prem * pay_yrs
        cov = float(pol.get('sum_insured') or summ.get('sum_insured') or 0)
        lev = (cov / total_prem) if total_prem > 0 else 0
        cov_period = pol.get('coverage_period') or '终身'
        company = (p.get('company_short') or '').upper() or '—'
        name = (pol.get('product_name') or p.get('product_name') or p.get('pdfName') or '—')[:18]
        rows_data.append({
            'company': company, 'name': name,
            'insured': f'{ins.get("name","—")} {ins.get("age","—")}岁'.strip(),
            'pay_str': pay_str, 'pay_yrs': pay_yrs,
            'annual_prem': annual_prem, 'total_prem': total_prem,
            'sum_insured': cov, 'leverage': lev, 'cov_period': cov_period,
        })

    n = len(rows_data)
    headers = ['#', '公司', '产品名称', '受保人', '缴费\n年期',
               '年缴额\n(USD)', '总缴保费\n(USD)', '身故赔偿\n(USD)', '杠杆率\n(身/总)', '保障年期']
    n_rows = 1 + n

    table_x = Inches(0.4); table_y = Inches(2.0)
    table_w = Inches(12.9)
    row_h = Inches(0.5)
    tbl_shape = s.shapes.add_table(n_rows, len(headers), table_x, table_y, table_w, row_h * n_rows)
    tbl = tbl_shape.table
    col_widths = [0.4, 0.7, 2.0, 1.4, 0.8, 1.0, 1.4, 1.5, 1.1, 0.9]
    total_cw = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w * 12.9 / total_cw)
    tbl.rows[0].height = Inches(0.55)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(0.5)

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 找最大值用于标红
    max_prem_idx = max(range(n), key=lambda i: rows_data[i]['annual_prem']) if n else 0
    max_sum_idx = max(range(n), key=lambda i: rows_data[i]['sum_insured']) if n else 0
    max_lev_idx = max(range(n), key=lambda i: rows_data[i]['leverage']) if n else 0

    for ri, r in enumerate(rows_data):
        cells = [
            str(ri + 1), r['company'], r['name'], r['insured'], r['pay_str'],
            f"${r['annual_prem']:,.0f}", f"${r['total_prem']:,.0f}",
            f"${r['sum_insured']:,.0f}", f"{r['leverage']:.2f}x",
            str(r['cov_period']),
        ]
        for ci, v in enumerate(cells):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'] if ri % 2 == 0 else C['bg_page'])
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci != 2 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = v
            run.font.size = Pt(10)
            run.font.bold = False
            run.font.name = FONT_HEI
            # 标红最大值
            if ci == 5 and ri == max_prem_idx and rows_data[ri]['annual_prem'] > 0:
                run.font.color.rgb = _c(C['red'])
                run.font.bold = True
            elif ci == 7 and ri == max_sum_idx and rows_data[ri]['sum_insured'] > 0:
                run.font.color.rgb = _c(C['red'])
                run.font.bold = True
            elif ci == 8 and ri == max_lev_idx and rows_data[ri]['leverage'] > 0:
                run.font.color.rgb = _c(C['red'])
                run.font.bold = True
            else:
                run.font.color.rgb = _c(C['dark_text'])

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书 · 受保人=保费试算年龄 · 总缴保费=年缴额×缴费年期 · 杠杆率=身故赔偿÷总缴保费',
             size=8, color=C['mid_text'])




def _slide_iul_feature_compare_multi(prs, products, meta, C):
    """V3.1.7 IUL 多产品核心特性对比 (跨产品功能/亮点表)
    列: 产品特性 / 受保人 / 各产品对照值
    包含: 缴费年期, 现金价值积累, 身故赔偿方式, 指数账户选项
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)

    meta['product_type'] = 'iul'

    add_title(s, 'IUL 万用寿险 · 核心特性对比',
              '缴费年期 · 现金价值 · 身故赔偿 · 指数账户 · 一图看全产品差异',
              title_size=28, sub_size=11, C=C)

    # 收集数据
    rows = []
    for p in products:
        d = p.get('data') or {}
        pol = d.get('policy') or {}
        summ = d.get('summary') or {}
        ins = d.get('insured') or {}
        # 缴费年期
        pay_raw = pol.get('premium_payment_period') or summ.get('payment_years') or 10
        if '趸交' in str(pay_raw) or 'Single' in str(pay_raw) or str(pay_raw).lower() == 'sgl':
            pay_str = '趸交'
        else:
            digits = ''.join(c for c in str(pay_raw) if c.isdigit())
            pay_yrs = int(digits) if digits else 10
            pay_str = f'{pay_yrs}年缴'
        # 现金价值 Y10
        bi = d.get('benefit_illustration') or []
        cv_y10 = 0
        for r in bi:
            try:
                ry = int(r.get('policy_year') or 0)
            except Exception:
                continue
            if ry == 10:
                cv_y10 = float(r.get('non_guaranteed_cash_value') or r.get('cash_value') or r.get('non_guaranteed_account_value') or r.get('account_value') or 0)
                break
        cv_y30 = 0
        for r in bi:
            try:
                ry = int(r.get('policy_year') or 0)
            except Exception:
                continue
            if ry == 30:
                cv_y30 = float(r.get('non_guaranteed_cash_value') or r.get('cash_value') or r.get('non_guaranteed_account_value') or r.get('account_value') or 0)
                break
        cov = float(pol.get('sum_insured') or summ.get('sum_insured') or 0)
        prem = float(pol.get('annual_premium') or summ.get('annual_premium') or 0)
        # 受保人
        age = ins.get('age') or pol.get('insured_age') or '—'
        gender = ins.get('gender') or pol.get('insured_gender') or ''
        insured_label = f'{gender} {age}岁' if gender else f'{age}岁'
        rows.append({
            'name': (pol.get('product_name') or p.get('product_name') or '—')[:18],
            'insured': insured_label,
            'pay_str': pay_str,
            'cov': cov, 'prem': prem,
            'cv_y10': cv_y10, 'cv_y30': cv_y30,
            'cov_period': pol.get('coverage_period') or '终身',
            'currency': pol.get('currency') or 'USD',
        })

    # 表头 (固定 4 列: 特性名 / 产品1 / 产品2 / 产品3)
    n = min(len(rows), 4)  # 限制 4 列产品
    headers = ["产品特性"] + [(rows[i].get("name","") or "—")[:12] + "\n" + (products[i].get("company_short") or "—") for i in range(n)]

    feature_rows = [
        ('产品名称', [r['name'] for r in rows[:n]]),
        ('受保人', [r['insured'] for r in rows[:n]]),
        ('缴费年期', [r['pay_str'] for r in rows[:n]]),
        ('年缴保费', [f"{r['currency']} {r['prem']:,.0f}" for r in rows[:n]]),
        ('身故赔偿', [f"{r['currency']} {r['cov']:,.0f}" for r in rows[:n]]),
        ('保障年期', [r['cov_period'] for r in rows[:n]]),
        ('第10年现金价值', [f"{r['currency']} {r['cv_y10']:,.0f}" if r['cv_y10'] > 0 else '—' for r in rows[:n]]),
        ('第30年现金价值', [f"{r['currency']} {r['cv_y30']:,.0f}" if r['cv_y30'] > 0 else '—' for r in rows[:n]]),
        ('产品亮点 (摘要)', [
            ('趸交省心 · 一次性投入' if '趸交' in r['pay_str'] else '分期缴费 · 现金流灵活')
            + '\n' + f"{r['currency']} {r['cov']/max(1,r['prem']*(1 if '趸交' in r['pay_str'] else 10)):.1f}x 杠杆"
            for r in rows[:n]
        ]),
    ]

    n_rows_f = 1 + len(feature_rows)
    # V3.1.7+: 表格上移 + 压紧行高, 确保不超出画布 (10 行 × 0.6 = 6" 太挤)
    table_x = Inches(0.4); table_y = Inches(1.55)
    table_w = Inches(12.9)
    row_h = Inches(0.42)

    # 列宽: 特性列窄, 产品列平均
    feature_w_ratio = 0.18
    col_w = [feature_w_ratio] + [(1 - feature_w_ratio) / n] * n
    tbl_shape = s.shapes.add_table(n_rows_f, len(headers), table_x, table_y, table_w, row_h * n_rows_f)
    tbl = tbl_shape.table
    for i, w in enumerate(col_w):
        tbl.columns[i].width = Inches(w * 12.9)
    tbl.rows[0].height = Inches(0.45)  # 表头稍高
    for i in range(1, n_rows_f):
        tbl.rows[i].height = Inches(0.5)  # 数据行 0.5 (原 0.6)

    # 表头
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['dark_text'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = _c(C['bg_light']); run.font.name = FONT_HEI

    # 数据行
    for ri, (label, vals) in enumerate(feature_rows):
        row = tbl.rows[ri + 1]
        # 特性名
        cell = tbl.cell(ri + 1, 0)
        cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_page'])
        cell.text = ''
        tf = cell.text_frame
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = label
        run.font.size = Pt(10); run.font.bold = True
        run.font.color.rgb = _c(C['dark_text']); run.font.name = FONT_HEI
        # 数据列
        for ci, v in enumerate(vals):
            cell = tbl.cell(ri + 1, ci + 1)
            cell.fill.solid(); cell.fill.fore_color.rgb = _c(C['bg_card'])
            cell.text = ''
            tf = cell.text_frame
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            tf.word_wrap = True
            lines = str(v).split('\n')
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run(); run.text = line
                run.font.size = Pt(9.5)
                run.font.color.rgb = _c(C['body_text']); run.font.name = FONT_HEI

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源: 各产品官方计划书 · 现金价值演示为非保证红利 · 杠杆率仅基于演示保额,实际可能因指数表现不同',
             size=8, color=C['mid_text'])






def _slide_ci_premium_chart(prs, ci_data, meta, C, company=None):
    """CI 保费 · 基本保额 · 总身故赔偿 对比图
    V3.2.4 拆成 2 子图 (左:保费+保额 / 右:总身故赔偿) + X轴限制 Y50
    V3.2.3 关键年(Y1/Y10/Y30/Y50)显示数据标签 + Y10 升级截止标注
    V3.2.2 高对比度 3 色 (蓝/金/红) + 6 项关键数据卡
    V3.2.1+ 加入基本保额线 (col 7 身故保证 = 保单年度内的固定保额)
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_text(s, Inches(0.4), Inches(0.5), Inches(12.5), Inches(0.45),
             '已缴保费 · 基本保额 · 总身故赔偿', size=20, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.28),
             '保单年度 · 保费投入 · 保障杠杆 (含终期红利) — X 轴: Y1-Y50', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.4), Inches(1.28), Inches(0.5), Inches(0.025), fill=C['accent'])

    ci_policy = ci_data.get('policy', {})
    prem = float(ci_policy.get('annual_premium') or 0)
    coverage = float(ci_policy.get('sum_insured') or ci_policy.get('basic_sum_insured') or 100000)
    base_si = float(ci_policy.get('basic_sum_insured') or coverage)
    pay_yrs = int(str(ci_policy.get('premium_payment_period', '10')).replace('年',''))
    bi = ci_data.get('benefit_illustration', [])

    labels = []; tp_series = []; si_series = []; db_series = []
    for r in bi:
        y = int(r.get('policy_year', 0))
        if y <= 0 or y > 50: continue  # V3.2.4: X 轴限制 Y50
        if y > 100: continue
        labels.append('Y{}'.format(y) if (y == 1 or y % 5 == 0) else '')
        tp = int(r.get('total_premium_paid') or 0)
        if not tp:
            tp = min(int(y * prem), int(prem * pay_yrs)) if prem > 0 else 0
        else:
            tp = min(tp, int(prem * pay_yrs)) if prem > 0 else tp
        tp_series.append(tp)
        si = int(r.get('basic_sum_insured') or 0)
        si_series.append(si if si > 0 else int(base_si))
        db_val = int(r.get('death_benefit') or 0)
        db_series.append(db_val if db_val > 0 else int(coverage))

    n = len(labels)
    last_idx = n - 1
    y1_idx = 0
    y10_idx = None
    y30_idx = None
    y50_idx = None
    for idx, lab in enumerate(labels):
        if lab == 'Y10': y10_idx = idx
        if lab == 'Y30': y30_idx = idx
        if lab == 'Y50': y50_idx = idx
    if y10_idx is None: y10_idx = min(9, n - 1)
    if y30_idx is None: y30_idx = min(29, n - 1)
    if y50_idx is None: y50_idx = last_idx
    key_idx_set = {y1_idx, y10_idx, y30_idx, y50_idx, last_idx}

    blue = _c(C.get('blue', C['primary']))
    gold = _c(C['accent'])
    red = _c(C['red'])

    # === 子图 A (左): 已缴保费 + 基本保额 — Y 轴 $0-$200K ===
    cd_a = CategoryChartData()
    cd_a.categories = labels
    cd_a.add_series('已缴保费', tp_series)
    cd_a.add_series('基本保额', si_series)
    cf_a = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                              Inches(0.3), Inches(1.5), Inches(4.4), Inches(4.8), cd_a)
    cA = cf_a.chart
    cA.has_title = True
    cA.chart_title.text_frame.text = '保费 vs 保额 (USD)'
    for p in cA.chart_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = _c(C['dark_text'])
    cA.has_legend = True
    cA.legend.position = XL_LEGEND_POSITION.BOTTOM
    cA.legend.font.size = Pt(9); cA.legend.font.bold = True
    # 自适应 Y 轴上限 (按 max of 保额 * 1.2)
    max_a = max(max(si_series) if si_series else 0, max(tp_series) if tp_series else 0)
    if max_a > 0:
        cA.value_axis.maximum_scale = max_a * 1.2
    cA.value_axis.minimum_scale = 0
    cA.value_axis.tick_labels.font.size = Pt(8)
    cA.value_axis.tick_labels.number_format = '"$"#,##0'
    cA.category_axis.tick_labels.font.size = Pt(8)
    for i, ser in enumerate(cA.series):
        ser.format.line.width = Pt(2.5)
        ser.format.line.color.rgb = [blue, gold][i]
        ser.smooth = False
        dl = ser.data_labels
        dl.show_value = False
        dl.font.size = Pt(7)
        dl.font.bold = True
        dl.font.color.rgb = [blue, gold][i]
        dl.number_format = '"$"#,##0'
        for pt_idx in key_idx_set:
            if pt_idx < len(ser.points):
                pt = ser.points[pt_idx]
                pt.data_label.show_value = True
                pt.data_label.font.size = Pt(7)
                pt.data_label.font.bold = True
                pt.data_label.font.color.rgb = [blue, gold][i]
                pt.data_label.number_format = '"$"#,##0'

    # === 子图 B (右): 总身故赔偿 — Y 轴 $0-max ===
    cd_b = CategoryChartData()
    cd_b.categories = labels
    cd_b.add_series('总身故赔偿', db_series)
    cf_b = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                              Inches(4.75), Inches(1.5), Inches(4.4), Inches(4.8), cd_b)
    cB = cf_b.chart
    cB.has_title = True
    cB.chart_title.text_frame.text = '总身故赔偿 (USD, 含终期红利)'
    for p in cB.chart_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = _c(C['dark_text'])
    cB.has_legend = False
    max_b = max(db_series) if db_series else 0
    if max_b > 0:
        cB.value_axis.maximum_scale = max_b * 1.15
    cB.value_axis.minimum_scale = 0
    cB.value_axis.tick_labels.font.size = Pt(8)
    cB.value_axis.tick_labels.number_format = '"$"#,##0'
    cB.category_axis.tick_labels.font.size = Pt(8)
    for ser in cB.series:
        ser.format.line.width = Pt(2.5)
        ser.format.line.color.rgb = red
        ser.smooth = False
        dl = ser.data_labels
        dl.show_value = False
        dl.font.size = Pt(8)
        dl.font.bold = True
        dl.font.color.rgb = red
        dl.number_format = '"$"#,##0'
        for pt_idx in key_idx_set:
            if pt_idx < len(ser.points):
                pt = ser.points[pt_idx]
                pt.data_label.show_value = True
                pt.data_label.font.size = Pt(8)
                pt.data_label.font.bold = True
                pt.data_label.font.color.rgb = red
                pt.data_label.number_format = '"$"#,##0'

    # Y10 升级截止年的视觉标注 (左子图右上)
    if y10_idx is not None and y10_idx < n and y10_idx > 0:
        y10_si_value = si_series[y10_idx] if si_series[y10_idx] else si_series[y10_idx - 1]
        y11_si_value = si_series[y10_idx + 1] if y10_idx + 1 < n else y10_si_value
        if y10_si_value > 0 and y11_si_value > 0 and y10_si_value != y11_si_value:
            drop_pct = int((y10_si_value - y11_si_value) * 100 / y10_si_value)
            add_text(s, Inches(0.3), Inches(6.35), Inches(8.85), Inches(0.32),
                     f'※ Y{10} 升级结束: ${y10_si_value:,} → ${y11_si_value:,} (回归 -{drop_pct}%)',
                     size=10, color=C['red'], bold=True, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.3), Inches(6.7), Inches(8.85), Inches(0.28),
             '※ 前 10 年含 60% 升级保障 (基本保额 100K → 实际 160K)，Y10 后回归基本保额水平',
             size=9, color=C['mid_text'], align=PP_ALIGN.CENTER)

    # === 右侧 6 项关键数据卡 ===
    ax, aw = Inches(9.25), Inches(3.6)
    box = s.shapes.add_shape(1, ax, Inches(1.5), aw, Inches(5.5))
    box.fill.solid()
    box.fill.fore_color.rgb = _c(C['bg_card'])
    box.line.color.rgb = _c(C['mid_text']); box.line.width = Pt(1)
    add_text(s, ax + Inches(0.2), Inches(1.65), aw - Inches(0.4), Inches(0.3),
             '关键数据', size=13, color=C['dark_text'], bold=True)
    add_rect(s, ax + Inches(0.2), Inches(1.95), aw - Inches(0.4), Inches(0.02), fill=C['accent'])

    final_idx = n - 1
    final_tp = tp_series[final_idx] if tp_series else 0
    final_si = si_series[final_idx] if si_series else 0
    final_db = db_series[final_idx] if db_series else 0
    y1_si = si_series[0] if si_series else int(base_si)
    final_yr = n  # 实际年度数 (≤50)

    my = Inches(2.15)
    items = [
        ('年缴保费', 'USD {:,.0f}'.format(int(prem)), C.get('blue', C['primary'])),
        ('总保费 ({}年缴)'.format(pay_yrs), 'USD {:,.0f}'.format(int(prem * pay_yrs)), C.get('blue', C['primary'])),
        ('基本保额 (Y1 升级后)', 'USD {:,.0f}'.format(int(y1_si)), C['accent']),
        ('基本保额 (Y10 后回归)', 'USD {:,.0f}'.format(int(final_si) if final_si else int(base_si)), C['accent']),
        ('总身故赔偿 (Y{})'.format(final_yr), 'USD {:,.0f}'.format(int(final_db)), C['red']),
        ('赔偿/保费倍数', '{:.1f}x'.format(final_db/(prem*pay_yrs) if prem*pay_yrs else 0), C['red']),
    ]
    for label, val, col in items:
        tb = s.shapes.add_textbox(ax + Inches(0.2), my, aw - Inches(0.4), Inches(0.6))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(8.5); p.font.color.rgb = _c(C['mid_text'])
        p2 = tf.add_paragraph(); p2.text = val; p2.font.size = Pt(13); p2.font.bold = True; p2.font.color.rgb = _c(col)
        my += Inches(0.65)
    add_text(s, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.2),
             '数据来源：保险公司官方计划书 · col 6 退保总额 / col 7 身故保证 (基本保额) / col 9 身故总额', size=7, color=C['mid_text'])




def _slide_iul_divider(prs, iul_data, meta, C, iul_company=None):
    """IUL篇章封面（简洁版）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    co = iul_company.get('name_zh', meta.get('company_short', '')) if iul_company else meta.get('company_short', '')
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
             '财富杠杆 · 寿险保障篇', size=11, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(1.0),
             '高杠杆寿险规划', size=36, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.5), Inches(12), Inches(0.6),
             '用现金流撬动千万保障', size=18, color=C['accent'], align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
             '{} · 指数型万用寿险'.format(co), size=12, color=C['accent_light'], align=PP_ALIGN.CENTER)


def _slide_iul_overview(prs, iul_data, meta, C, company=None):
    """IUL概要: 保费 · 身故赔偿 · 杠杆"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_title(s, '万用寿险概要', '保费 · 身故赔偿 · 杠杆倍数', title_size=36, sub_size=14, C=C)

    iul_summary = iul_data.get('summary', {})
    iul_policy = iul_data.get('policy', {})
    ins = iul_data.get('insured', {})

    premium = float(iul_summary.get('annual_premium', iul_policy.get('annual_premium', 0)))
    coverage = float(iul_summary.get('sum_insured', iul_policy.get('sum_insured', 0)))
    pay_yrs = int(iul_summary.get('payment_years', meta.get('payment_years', 5)))
    total_prem = premium * pay_yrs
    leverage = (coverage / total_prem) if total_prem > 0 else 0

    cards_data = [
        ('年缴保费', 'USD {:,} × {}年'.format(int(premium), pay_yrs), '共USD {:,}'.format(int(total_prem)), C['primary']),
        ('身故赔偿', 'USD {:,}'.format(int(coverage)), '保障至终身', C['primary_dark']),
        ('杠杆倍数', '{:.1f}x'.format(leverage), '身故赔偿/总保费', C['accent_dark']),
        ('受保人', '{}'.format(ins.get("name", "—")), '{}岁'.format(ins.get("age", "—")), C['dark_text']),
    ]

    card_w = Inches(3.0); card_h = Inches(1.8); gap = Inches(0.25); sx = Inches(0.4); sy = Inches(2.2)
    for i, (label, val, sub, col) in enumerate(cards_data):
        x = sx + i * (card_w + gap)
        add_rect(s, x, sy, card_w, card_h, fill=C['bg_card'])
        add_rect(s, x, sy, card_w, Inches(0.06), fill=col)
        add_text(s, x + Inches(0.2), sy + Inches(0.2), card_w - 0.4, Inches(0.25),
                 label, size=9, color=C['mid_text'])
        add_text(s, x + Inches(0.2), sy + Inches(0.5), card_w - 0.4, Inches(0.45),
                 val, size=20, color=col, bold=True)
        add_text(s, x + Inches(0.2), sy + Inches(1.1), card_w - 0.4, Inches(0.3),
                 sub, size=9, color=C['body_text'])

    add_rect(s, Inches(0.4), Inches(4.5), Inches(0.3), Inches(0.04), fill=C['primary'])
    add_text(s, Inches(0.8), Inches(4.4), Inches(10), Inches(0.4),
             '产品特点', size=14, color=C['dark_text'], bold=True)
    features = '· 指数账户挂钩市场表现，享有增长潜力\n· 保证最低派息率，下行风险可控\n· 保费灵活缴付，可根据现金流调整\n· 身故赔偿免遗产税，财富定向传承'
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    for i, line in enumerate(features.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12.5)
        run.font.color.rgb = _c(C['body_text'])
        run.font.name = FONT_HEI
        p.line_spacing = Pt(19)
    add_text(s, Inches(0.4), Inches(7.2), Inches(12), Inches(0.25),
             '数据来源：保险公司官方计划书', size=8, color=C['mid_text'])


def _slide_iul_leverage_chart(prs, iul_data, meta, C, company=None):
    """IUL全页杠杆折线图: 已缴总保费 vs 身故赔偿总保额"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    add_text(s, Inches(0.5), Inches(0.6), Inches(12), Inches(0.5),
             '已缴总保费 vs 身故赔偿总保额', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.3),
             '保单年度 · 金额 · 杠杆效应趋势', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.45), Inches(0.5), Inches(0.03), fill=C['accent'])

    # 提取 IUL 数据
    iul_summary = iul_data.get('summary', {})
    iul_policy = iul_data.get('policy', {})
    iul_insured = iul_data.get('insured', {})
    premium = float(iul_summary.get('annual_premium', iul_policy.get('annual_premium', 0)))
    coverage = float(iul_summary.get('sum_insured', iul_policy.get('sum_insured', 0)))
    pay_yrs = int(iul_summary.get('payment_years', meta.get('payment_years', 5)))
    insured_age = int(iul_insured.get('age', meta.get('insured_age', 46)))
    max_age = 120
    bi = iul_data.get('benefit_illustration', [])

    # 构建图表数据（演示到受保人120岁）
    labels = []
    tp_series = []
    db_series = []
    av_series = []  # 退保价值 (与表格列保持一致)
    for r in bi:
        y = int(r.get('policy_year', 0))
        age = insured_age + y - 1
        if y <= 0 or age > max_age: continue
        labels.append('Y{}'.format(y) if (y == 1 or y % 10 == 0) else '')
        paid = min(y * premium, premium * pay_yrs)
        tp_series.append(paid)
        db_val = r.get('non_guaranteed_death_benefit', r.get('death_benefit', coverage))
        db_series.append(int(db_val if db_val else coverage))
        # 关键: 与表格列保持一致, 优先读退保价值 (non_guaranteed_cash_value), 兜底读 account_value
        av_val = r.get('non_guaranteed_cash_value', r.get('cash_value', 0))
        if not av_val:
            av_val = r.get('non_guaranteed_account_value', r.get('account_value', 0))
        av_series.append(int(av_val if av_val else 0))

    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series('已交总保费', tp_series)
    cd.add_series('退保价值', av_series)
    cd.add_series('身故赔偿总保额', db_series)

    cf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.3), Inches(1.7), Inches(8.5), Inches(5.0), cd)
    c = cf.chart
    c.has_title = False
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    colors = [_c(C['primary']), _c(C['primary_light']), _c(C['accent_dark'])]
    for i, ser in enumerate(c.series):
        ser.format.line.width = Pt(2.5)
        ser.format.line.color.rgb = colors[i]
        ser.smooth = True

    va = c.value_axis
    va.has_title = True
    va.axis_title.text_frame.paragraphs[0].text = '金额 (USD)'
    va.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    va.major_gridlines.format.line.color.rgb = _c(C['mid_text'])
    va.format.line.color.rgb = _c(C['mid_text'])
    va.minimum_scale = 0
    va.maximum_scale = coverage * 2.0  # 纵轴上限 = 初始保额 × 2

    ca = c.category_axis
    ca.has_title = True
    ca.axis_title.text_frame.paragraphs[0].text = '保单年度'
    ca.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    ca.tick_labels.font.size = Pt(7)

    # 右侧注解面板
    ax, aw = Inches(9.0), Inches(3.5)
    box = s.shapes.add_shape(1, ax, Inches(1.7), aw, Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = _c(C['bg_card'])
    box.line.color.rgb = _c(C['mid_text'])
    box.line.width = Pt(1)

    add_text(s, ax + Inches(0.2), Inches(1.85), aw - Inches(0.4), Inches(0.3),
             '📊 关键数据', size=13, color=C['dark_text'], bold=True)
    add_rect(s, ax + Inches(0.2), Inches(2.15), aw - Inches(0.4), Inches(0.02), fill=C['accent'])

    total_prem = premium * pay_yrs
    leverage = (coverage / total_prem) if total_prem > 0 else 0
    metrics = [
        ('年缴保费', 'USD {:,.0f} × {}年'.format(int(premium), pay_yrs)),
        ('总保费', 'USD {:,.0f}'.format(int(total_prem))),
        ('身故保额', 'USD {:,.0f}'.format(int(coverage))),
        ('保费杠杆', '{:.1f}x'.format(leverage)),
        ('保额锁定', '缴费期内锁定额度'),
    ]
    my = Inches(2.4)
    for label, value in metrics:
        tb = s.shapes.add_textbox(ax + Inches(0.2), my, aw - Inches(0.4), Inches(0.5))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = label
        p.font.size = Pt(9); p.font.color.rgb = _c(C['mid_text'])
        p2 = tf.add_paragraph(); p2.text = value
        p2.font.size = Pt(14); p2.font.bold = True; p2.font.color.rgb = _c(C['dark_text'])
        my += Inches(0.55)

    add_text(s, ax + Inches(0.2), Inches(5.8), aw - Inches(0.4), Inches(0.8),
             '蓝色为累计已缴保费（{}年缴清后持平），金色为身故赔偿保额。两条线的间距体现保费杠杆效应。'.format(pay_yrs),
             size=8, color=C['body_text'])
    add_text(s, Inches(0.4), Inches(7.0), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])


def _slide_iul_data_table(prs, iul_data, meta, C, company=None):
    """IUL保单摘要表: 年龄/年度/保费/退保价值/保额（美化版）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)

    # — 标题区 + 顶部装饰线 —
    add_rect(s, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.35), fill=C['accent'])
    add_text(s, Inches(0.75), Inches(0.55), Inches(8), Inches(0.35),
             '保单摘要', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.75), Inches(0.95), Inches(10), Inches(0.25),
             '年龄 · 保单年度 · 保费 · 退保价值 · 保额', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.3), Inches(11.5), Inches(0.015), fill=C['accent_light'])

    # — 提取数据 —
    iul_summary = iul_data.get('summary', {})
    iul_policy = iul_data.get('policy', {})
    iul_insured = iul_data.get('insured', {})
    premium = float(iul_summary.get('annual_premium', iul_policy.get('annual_premium', 0)))
    coverage = float(iul_summary.get('sum_insured', iul_policy.get('sum_insured', 0)))
    pay_yrs = int(iul_summary.get('payment_years', meta.get('payment_years', 5)))
    insured_age = int(iul_insured.get('age', meta.get('insured_age', 46)))
    bi = iul_data.get('benefit_illustration', [])

    # — 筛选行: Y1-Y10(每年), Y10之后每10年 —
    selected = []
    for r in bi:
        y = int(r.get('policy_year', 0))
        if y <= 0: continue
        age = insured_age + y - 1
        if age > 120: continue
        if y <= 10 or y % 10 == 0:
            selected.append(r)

    # — 构建行列数据 —
    # 关键: 列名 = "退保价值" (surrender value), 数据取 non_guaranteed_cash_value
    # 旧版用 non_guaranteed_account_value (账户价值) 但实际数据流里该字段常被映射成 cash_value,
    # 导致 label 与数字语义不一致. 现统一: label=退保价值, data=non_guaranteed_cash_value/cash_value
    headers = ['年龄', '保单年度', '年缴保费\n(USD)', '累计已缴\n(USD)', '退保价值\n(USD)', '保额\n(USD)']
    rows_data = []
    be_y = None  # breakeven year: 退保价值首次超过累计已缴
    # 关键: Sunlife IUL 等万用寿险保费每年可变 (Year 1 趸交, Year 2-10 续期, Year 11+ 缴清)
    # 优先用 r.annual_premium (per-year), 兜底用 scalar premium 仅在 pay_yrs 内
    cum = 0  # 本地累计 (r.total_premium_paid 不可信时兜底)
    has_per_year = any(r.get('annual_premium') for r in selected)
    for r in selected:
        y = int(r['policy_year'])
        age = insured_age + y - 1
        if has_per_year and 'annual_premium' in r and r['annual_premium'] is not None:
            prem_this = int(r['annual_premium'])
        else:
            prem_this = int(premium) if y <= pay_yrs else 0
        # 累计: 优先用 r.total_premium_paid, 否则本行 prem_this 累加
        cum += prem_this
        total_paid_raw = r.get('total_premium_paid')
        if total_paid_raw is not None and float(total_paid_raw) > 0:
            total_paid = int(float(total_paid_raw))
        else:
            total_paid = cum
        # 关键: 优先读退保价值字段 (non_guaranteed_cash_value), 兜底读 account_value
        # Manulife 等 extractor 直接输出 surrender_value, 由 server.ts 映射到 cash_value
        surr_val = int(r.get('non_guaranteed_cash_value', r.get('cash_value', 0)) or 0)
        if surr_val == 0:
            surr_val = int(r.get('non_guaranteed_account_value', r.get('account_value', 0)) or 0)
        db = int(r.get('non_guaranteed_death_benefit', r.get('death_benefit', coverage)) or coverage)
        if be_y is None and y > pay_yrs and surr_val > total_paid:
            be_y = y
        rows_data.append([age, y, prem_this, total_paid, surr_val, db])

    n_rows = len(rows_data) + 1
    n_cols = len(headers)

    # — 计算表格尺寸 —
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.5)
    tbl_w = Inches(12.0)
    row_h = Inches(0.38)
    tbl_h = row_h * n_rows

    table_shape = s.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table

    # — 列宽 —
    col_widths = [Inches(1.2), Inches(1.3), Inches(1.8), Inches(2.0), Inches(2.8), Inches(2.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # — 表头样式：深色底色 + 金色装饰条 —
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        p.text = h.replace('\n', ' ')
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = _c(C['primary'])
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # 表头底部分隔线 (金色细条)
        tcPr = cell._tc.get_or_add_tcPr()
        lnB = etree.SubElement(tcPr, qn('a:lnB'))
        lnB.set('w', '12700')  # 1pt
        solidFill = etree.SubElement(lnB, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(*_c(C['accent'])[:3]))

    # — 数据行 —
    accent_rgb = _c(C['accent'])
    primary_light_rgb = _c(C['primary_light'])
    for ri, row in enumerate(rows_data):
        is_be = row[1] == be_y  # breakeven year highlight
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 格式化文本
            if ci == 2:
                cell.text = '{:,.0f}'.format(val) if val > 0 else '—'
            elif isinstance(val, int) and val >= 1000:
                cell.text = '{:,.0f}'.format(val)
            else:
                cell.text = str(val)

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.alignment = PP_ALIGN.CENTER
                if is_be:
                    p.font.bold = True
                    p.font.color.rgb = _c(C['accent_dark'])
                elif ci == 4:  # 退保价值用品牌色
                    p.font.color.rgb = accent_rgb
                    p.font.bold = True
                else:
                    p.font.color.rgb = _c(C['dark_text'])

            # 行背景: breakeven行高亮, 其余交替
            cell.fill.solid()
            if is_be:
                cell.fill.fore_color.rgb = _c(C['accent_light'])
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # — 底部说明 —
    footer_y = tbl_top + tbl_h + Inches(0.08)
    summary_parts = []
    total_prem = int(premium * pay_yrs)
    summary_parts.append('总缴保费: USD {:,.0f}'.format(total_prem))
    if be_y:
        summary_parts.append('预期回本年度: 第{}年'.format(be_y))
    last_acct = rows_data[-1][4] if rows_data else 0
    if total_prem > 0:
        summary_parts.append('末年倍数: {:.1f}x'.format(last_acct / total_prem))
    add_rect(s, Inches(0.5), footer_y, Inches(11.5), Inches(0.3), fill=C['bg_card'])
    add_text(s, Inches(0.6), footer_y, Inches(11.3), Inches(0.3),
             '  '.join(summary_parts), size=8, color=C['body_text'])

    add_text(s, Inches(0.5), footer_y + Inches(0.35), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])


def _slide_savings_iul_premium_funding(prs, wd, iul_data, meta, C, company=None):
    """储蓄+IUL组合: 储蓄提领缴纳IUL保费（全页版）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)

    add_rect(s, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.35), fill=C['accent'])
    add_text(s, Inches(0.75), Inches(0.55), Inches(8), Inches(0.35),
             '现金流 + 高杠杆 · 完美组合', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.75), Inches(0.95), Inches(10), Inches(0.25),
             '储蓄提领缴纳IUL保费 → 稳定现金流 + 高额寿险保障', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.3), Inches(11.5), Inches(0.015), fill=C['accent_light'])

    iul_premium = float(iul_data.get('policy', {}).get('annual_premium', 0))
    iul_coverage = float(iul_data.get('policy', {}).get('sum_insured', 0))
    iul_pay_yrs = int(iul_data.get('summary', {}).get('payment_years', meta.get('payment_years', 5)))
    iul_total_prem = iul_premium * iul_pay_yrs
    leverage = (iul_coverage / iul_total_prem) if iul_total_prem > 0 else 0

    wd_annual = 0; wd_start = None
    if wd:
        for y in sorted(int(k) for k in wd.keys()):
            if wd[y].get('Annual_WD', 0) > 0:
                wd_start = y; wd_annual = wd[y]['Annual_WD']; break
    insured_age = meta.get('insured_age', 1) or 1
    wd_age = insured_age + (wd_start or 7) - 1
    cover_pct = (wd_annual / iul_premium * 100) if iul_premium > 0 else 0

    lx = Inches(0.4); lw = Inches(6.0)

    # 上: 储蓄
    add_rect(s, lx, Inches(1.6), lw, Inches(2.2), fill=C['bg_card'])
    add_rect(s, lx, Inches(1.6), lw, Inches(0.06), fill=C['primary'])
    add_text(s, lx + Inches(0.2), Inches(1.75), lw - 0.4, Inches(0.3),
             '💰 储蓄险 · 稳定现金流', size=13, color=C['dark_text'], bold=True)
    if wd_annual > 0:
        add_text(s, lx + Inches(0.2), Inches(2.2), lw - 0.4, Inches(0.4),
                 '从 {} 岁起，每年提取 USD {:,.0f}'.format(int(wd_age), int(wd_annual)),
                 size=14, color=C['primary_dark'], bold=True)
        add_text(s, lx + Inches(0.2), Inches(2.7), lw - 0.4, Inches(0.8),
                 '这笔现金流可用于支付IUL保费，同时储蓄账户继续滚存增长\n'
                 '覆盖比例：约 {:.0f}%（IUL年缴USD {:,.0f}）'.format(cover_pct, int(iul_premium)),
                 size=9, color=C['body_text'])
    else:
        add_text(s, lx + Inches(0.2), Inches(2.2), lw - 0.4, Inches(1.2),
                 '储蓄计划长期复利增长，保单价值可灵活提取\n'
                 '提取的资金可用于补充IUL保费，实现以钱生钱',
                 size=10, color=C['body_text'])

    add_text(s, Inches(2.5), Inches(3.9), Inches(1.5), Inches(0.5),
             '⬇ 保费支付', size=14, color=C['accent_dark'], bold=True, align=PP_ALIGN.CENTER)

    # 下: IUL
    add_rect(s, lx, Inches(4.4), lw, Inches(2.6), fill=C['bg_card'])
    add_rect(s, lx, Inches(4.4), lw, Inches(0.06), fill=C['accent_dark'])
    add_text(s, lx + Inches(0.2), Inches(4.55), lw - 0.4, Inches(0.3),
             '📈 IUL · 高杠杆寿险', size=13, color=C['dark_text'], bold=True)
    add_text(s, lx + Inches(0.2), Inches(5.0), lw - 0.4, Inches(0.3),
             '年缴保费 USD {:,.0f} × {}年'.format(int(iul_premium), iul_pay_yrs),
             size=12, color=C['accent_dark'], bold=True)
    iul_d = ('• 身故赔偿 USD {:,.0f}（约 {:.1f} 倍总保费杠杆）\n'
             '• 指数账户增长潜力，下行风险可控\n'
             '• 保证最低派息率，保障本金安全\n'
             '• 身故赔偿免遗产税，财富定向传承').format(int(iul_coverage), leverage)
    add_text(s, lx + Inches(0.2), Inches(5.4), lw - 0.4, Inches(1.4),
             iul_d, size=9, color=C['body_text'])

    # 右侧数据
    rx = Inches(7.0); rw = Inches(5.5)
    right_data = [
        ('储蓄年提领', 'USD {:,.0f}'.format(int(wd_annual)) if wd_annual else '—', C['primary']),
        ('IUL年缴保费', 'USD {:,.0f}'.format(int(iul_premium)), C['accent_dark']),
        ('IUL总保费', 'USD {:,.0f}'.format(int(iul_total_prem)), C['mid_text']),
        ('身故赔偿', 'USD {:,.0f}'.format(int(iul_coverage)), C['primary_dark']),
        ('保费杠杆', '{:.1f}x'.format(leverage), C['accent']),
    ]
    if wd_annual > 0:
        right_data.insert(2, ('保费覆盖', '{:.0f}%'.format(cover_pct), C['accent']))
    ry = Inches(1.6)
    for rl, rv, rc in right_data[:5]:
        add_rect(s, rx, ry, rw, Inches(1.0), fill=C['bg_card'])
        add_rect(s, rx, ry, rw, Inches(0.06), fill=rc)
        add_text(s, rx + Inches(0.2), ry + Inches(0.1), rw - 0.4, Inches(0.2),
                 rl, size=9, color=C['mid_text'])
        add_text(s, rx + Inches(0.2), ry + Inches(0.4), rw - 0.4, Inches(0.35),
                 rv, size=18, color=rc, bold=True)
        ry += Inches(1.1)

    add_rect(s, Inches(0.4), Inches(6.5), Inches(12.2), Inches(0.008), fill=C['accent_light'])
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             '💡 规划理念：储蓄险创造稳定现金流，IUL提供高额寿险保障。'
             '用储蓄收益缴纳IUL保费，一份投入双重收获。',
             size=9, color=C['body_text'])
    add_text(s, Inches(0.5), Inches(7.15), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])


def _slide_ending_combined(prs, summary, ci_data, iul_data, meta, C):
    """组合方案结束页（简洁版）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['dark_text'])
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.4),
             '{}  |  {}'.format(meta.get("company_short_en", ""), meta.get("company_short", "")),
             size=12, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)

    tagline = '现金流 + 保障 + 杠杆 · 完整的家庭财务规划'
    add_text(s, Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.9),
             '一份规划，三重守护', size=30, color=C['bg_light'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.333), Inches(0.9),
             tagline, size=22, color=C['accent'], bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.17), Inches(3.9), Inches(1.0), Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5),
             '—— 致 {} 与家人  ——'.format(summary.get("insured_name", "VIP 先生")),
             size=13, color=C['bg_light'], align=PP_ALIGN.CENTER)
    add_rect(s, Emu(0), Inches(7.2), SLIDE_W, Inches(0.04), fill=C['accent'])
    add_text(s, Inches(0.5), Inches(7.3), Inches(8), Inches(0.2),
             '本演示基于保险公司官方计划书数据生成 · 实际保单条款以保单文件为准',
             size=8, color=C['accent_light'])
    add_text(s, Inches(11.5), Inches(7.3), Inches(1.5), Inches(0.2),
             '{:02d} / {:02d}'.format(meta.get("_page_num", 12), meta.get("_total_slides", 12)),
             size=8, color=C['accent_light'], align=PP_ALIGN.RIGHT)

def _slide_sub_company(prs, company, product_label, meta, C):
    """次要产品公司介绍页 (用于CI/IUL等附加产品)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)
    bp = company.get('brand_profile', company)
    add_title(s, '关于{}'.format(company.get('name_zh', '保险公司')),
              '{} · {}'.format(product_label, company.get('name_zh', '')),
              title_size=36, sub_size=14, C=C)

    cards = [
        (bp.get('founded_year', '—'), bp.get('founded_label', '成立年份'), bp.get('founded_sub', '')),
        (bp.get('rating_value', '—'), bp.get('rating_label', '财务实力评级'), bp.get('rating_sub', '')),
        (meta.get('product_name', bp.get('series_label', '—')), bp.get('series_sub', ''), bp.get('series_products', '')),
    ]
    cy = Inches(2.4); cw = Inches(4.0); ch = Inches(1.5); gap = Inches(0.25); sx = Inches(0.5)
    for i, (val, lbl, sub) in enumerate(cards):
        x = sx + i * (cw + gap)
        add_rect(s, x, cy, cw, ch, fill=C['bg_card'])
        add_rect(s, x, cy + ch - Inches(0.08), cw, Inches(0.08), fill=C['primary'])
        add_text(s, x + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), Inches(0.7),
                 str(val), size=26, color=C['dark_text'], bold=True)
        add_text(s, x + Inches(0.2), cy + Inches(1.1), cw - Inches(0.4), Inches(0.3),
                 str(lbl), size=12, color=C['body_text'], bold=True)
        add_text(s, x + Inches(0.2), cy + Inches(1.25), cw - Inches(0.4), Inches(0.2),
                 str(sub), size=8, color=C['mid_text'])

    biz = bp.get('business_lines', [])
    if biz:
        add_rect(s, Inches(0.5), Inches(4.3), Inches(0.3), Inches(0.04), fill=C['primary'])
        add_text(s, Inches(0.85), Inches(4.2), Inches(5), Inches(0.4), '业务范围', size=14, color=C['dark_text'], bold=True)
        for i, t in enumerate(biz[:4]):
            add_text(s, Inches(0.85), Inches(4.7 + i*0.3), Inches(5.5), Inches(0.3), t, size=11, color=C['body_text'])

    brand = bp.get('brand_background', [])
    if brand:
        add_rect(s, Inches(7.0), Inches(4.3), Inches(0.3), Inches(0.04), fill=C['primary'])
        add_text(s, Inches(7.35), Inches(4.2), Inches(5), Inches(0.4), '品牌背景', size=14, color=C['dark_text'], bold=True)
        for i, t in enumerate(brand[:4]):
            add_text(s, Inches(7.35), Inches(4.7 + i*0.3), Inches(5.5), Inches(0.3), t, size=11, color=C['body_text'])

    add_text(s, Inches(0.5), Inches(7.15), Inches(12.5), Inches(0.3),
             '数据来源：{}'.format(bp.get('data_source', '保险公司官方资料')), size=8, color=C['mid_text'])

def _slide_ci_income_protection(prs, ci_data, wd, meta, C, company=None):
    """储蓄+重疾组合: 现金流 + 全面保障 · 安心组合"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = _c(C['bg_page'])
    add_header(s, meta=meta, C=C, company_override=company)

    # 标题区（类似IUL组合页风格）
    add_rect(s, Inches(0.5), Inches(0.55), Inches(0.08), Inches(0.35), fill=C['accent'])
    add_text(s, Inches(0.75), Inches(0.55), Inches(8), Inches(0.35),
             '现金流 + 全面保障 · 安心组合', size=22, color=C['dark_text'], bold=True)
    add_text(s, Inches(0.75), Inches(0.95), Inches(10), Inches(0.25),
             '储蓄现金流覆盖重疾保费，收入中断有保障', size=10, color=C['mid_text'])
    add_rect(s, Inches(0.5), Inches(1.3), Inches(11.5), Inches(0.015), fill=C['accent_light'])

    # 计算CI数据
    ci_policy = ci_data.get('policy', {})
    ci_summary = ci_data.get('summary', {})
    ci_insured = ci_data.get('insured', {})
    premium = float(ci_policy.get('annual_premium', 0))
    coverage = float(ci_policy.get('sum_insured', 100000))
    pay_yrs = int(str(ci_policy.get('premium_payment_period', '10')).replace('年',''))
    insured_age = int(ci_insured.get('age', 1))
    upgrade_amt = float(ci_data.get('upgrade_benefit_amount', 0))

    # 查找储蓄提领数据
    wd_annual = 0
    wd_start = None
    if wd:
        for y in sorted(int(k) for k in wd.keys()):
            if wd[y].get('Annual_WD', 0) > 0:
                wd_start = y
                wd_annual = wd[y]['Annual_WD']
                break

    total_db = coverage + (upgrade_amt if upgrade_amt > 0 else coverage * 0.6)
    daily_cost = premium * pay_yrs / 365 if premium > 0 else 0

    # 左侧: 风险分析面板
    left_x = Inches(0.4); box_w = Inches(5.8)
    # 风险警示框
    add_rect(s, left_x, Inches(1.6), box_w, Inches(2.2), fill=C['bg_card'])
    add_rect(s, left_x, Inches(1.6), box_w, Inches(0.06), fill=C['accent_dark'])
    add_text(s, left_x + Inches(0.2), Inches(1.75), box_w - 0.4, Inches(0.3),
             '⚡ 风险场景', size=14, color=C['dark_text'], bold=True)
    risk_lines = [
        '若受保人不幸确诊严重疾病（如癌症/中风/心脏病）',
        '• 治疗期间可能无法工作 → 收入中断',
        '• 康复期通常需要1-5年 → 持续支出',
        '• 家庭日常开支、房贷、子女教育仍需支付',
        '• 储蓄计划可提供稳定现金流，缓解康复期经济压力',
    ]
    add_text(s, left_x + Inches(0.2), Inches(2.15), box_w - 0.4, Inches(1.5),
             '\n'.join(risk_lines), size=9, color=C['body_text'])

    # 解决方案框
    add_rect(s, left_x, Inches(4.0), box_w, Inches(3.0), fill=C['bg_card'])
    add_rect(s, left_x, Inches(4.0), box_w, Inches(0.06), fill=C['primary'])
    add_text(s, left_x + Inches(0.2), Inches(4.15), box_w - 0.4, Inches(0.3),
             '✅ 组合解决方案', size=14, color=C['dark_text'], bold=True)

    # CI保障信息
    ci_items_info = [
        ('危疾保障', 'USD {:,.0f}'.format(int(total_db)), '确诊即赔付，覆盖治疗及康复费用'),
    ]
    if wd_annual > 0:
        ci_items_info.append(('储蓄现金流', 'USD {:,.0f}/年'.format(int(wd_annual)),
                              '从第{}年起每年提取，补充康复期收入'.format(wd_start)))
    ci_items_info.append(('组合优势', '现金流+保障', '治疗有资金，康复有收入，家庭有保障'))

    iy = Inches(4.55)
    for label, val, desc in ci_items_info:
        add_rect(s, left_x + Inches(0.2), iy, Inches(0.08), Inches(0.55), fill=C['accent'])
        add_text(s, left_x + Inches(0.4), iy, Inches(2.0), Inches(0.25),
                 label, size=10, color=C['mid_text'])
        add_text(s, left_x + Inches(0.4), iy + Inches(0.25), box_w - 0.6, Inches(0.25),
                 val, size=13, color=C['dark_text'], bold=True)
        add_text(s, left_x + Inches(3.0), iy + Inches(0.1), box_w - 3.2, Inches(0.5),
                 desc, size=8, color=C['body_text'])
        iy += Inches(0.7)

    # 右侧: 关键数据面板
    rx = Inches(6.8); rw = Inches(5.5)
    right_data = [
        ('年缴保费(CI)', 'USD {:,.0f}'.format(int(premium)), C['primary']),
        ('危疾保障总额', 'USD {:,.0f}'.format(int(total_db)), C['accent_dark']),
    ]
    if wd_annual > 0:
        right_data.append(('储蓄年提领', 'USD {:,.0f}'.format(int(wd_annual)), C['accent']))
    right_data.append(('保障杠杆', '{:.1f}x'.format(total_db / (premium * pay_yrs) if premium > 0 else 0), C['primary']))

    ry = Inches(1.6)
    for label, val, col in right_data:
        add_rect(s, rx, ry, rw, Inches(1.0), fill=C['bg_card'])
        add_rect(s, rx, ry, rw, Inches(0.06), fill=col)
        add_text(s, rx + Inches(0.2), ry + Inches(0.1), rw - 0.4, Inches(0.2),
                 label, size=9, color=C['mid_text'])
        add_text(s, rx + Inches(0.2), ry + Inches(0.4), rw - 0.4, Inches(0.35),
                 val, size=18, color=col, bold=True)
        ry += Inches(1.15)

    # 底部叙事（含保费覆盖比例）
    cover_multiple = (wd_annual / premium) if premium > 0 and wd_annual > 0 else 0
    cover_text = ''
    if wd_annual > 0:
        if cover_multiple >= 5:
            cover_text = f'（年提USD {wd_annual:,.0f}，覆盖重疾保费{cover_multiple:.1f}倍），'
        else:
            cover_text = f'（年提USD {wd_annual:,.0f}，覆盖重疾保费约{cover_multiple*100:.0f}%），'
    add_rect(s, Inches(0.4), Inches(5.85), Inches(12.2), Inches(0.008), fill=C['accent_light'])
    add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.7),
             '💡 规划理念：<b>储蓄险</b>创造稳定现金流' + cover_text
             + '<b>危疾保障</b>提供高额医疗及收入保障。'
             '两者结合，让您既有财富增值又有风险兜底。',
             size=9, color=C['body_text'])
    add_text(s, Inches(0.5), Inches(6.8), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])


def _slide_combo_narrative(prs, meta, C, ci_data=None, iul_data=None):
    """组合叙事页: 展示产品协同关系"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C['bg_page'])
    add_header(s, meta=meta, C=C)
    combo_text = meta.get('narrative_combo', '')
    add_title(s, '方案协同 · 完整的家庭财务规划' if not combo_text else combo_text[:40],
              '攻守兼备，进退自如',
              title_size=36, sub_size=14, C=C)

    # 产品角色与详细优势
    has_savings_meta = meta.get('has_savings', True)
    has_ci_data = ci_data is not None
    has_iul_data = iul_data is not None

    # 定义每个产品的优势
    savings_advantages = [
        ('📈 复利增长', '长期 IRR 5.5-6.5%，30年翻5倍+'),
        ('💱 多元货币', '支持7-10种货币自由转换'),
        ('✂️ 保单拆分', '无限次拆分，多子女传承'),
        ('🏛️ 税务递延', '香港免税增长，无资本利得税'),
        ('👑 定向传承', '无限次变更受保人，类信托赔付'),
        ('💧 提取灵活', '提取现金价值缴纳IUL保费'),
    ]
    iul_advantages = [
        ('🛡️ 下行保护', '0%保证下限，市场跌不亏'),
        ('📊 指数挂钩', '挂钩标普500，上限率10%+'),
        ('🇸🇬 主权独立', '新加坡MAS监管，司法独立'),
        ('💰 美元资产', '美元计价，跨境资产配置'),
        ('🏛️ 免税传承', '身故赔偿免遗产税，隔离债务'),
        ('👨‍👩‍👧‍👦 多代规划', '保障至120岁+，跨代传承'),
    ]
    ci_advantages = [
        ('🛡️ 收入保障', '确诊即赔，弥补收入中断'),
        ('💰 杠杆赔付', '小投入高保额，最高100%'),
        ('🔄 多重赔付', '癌症+中风+心脏病多次赔'),
        ('🏛️ 豁免保费', '确诊后免缴后续保费'),
    ]

    # 计算每列宽度
    prod_count = sum([has_savings_meta, has_ci_data, has_iul_data])
    if prod_count == 2:
        cw = Inches(5.5); gap = Inches(0.3)
    else:
        cw = Inches(3.8); gap = Inches(0.2)
    sx = Inches(0.4); sy = Inches(1.9)
    col_idx = 0

    if has_savings_meta:
        x = sx + col_idx * (cw + gap)
        add_rect(s, x, sy, cw, Inches(0.5), fill=C['primary'])
        add_text(s, x + Inches(0.2), sy + Inches(0.1), cw - 0.4, Inches(0.35),
                 '💰 现金流引擎 · 香港储蓄险', size=12, color=C['bg_light'], bold=True)
        adv_y = sy + Inches(0.7)
        for label, desc in savings_advantages[:6]:
            add_text(s, x + Inches(0.15), adv_y, Inches(1.3), Inches(0.22),
                     label, size=8, color=C['primary'], bold=True)
            add_text(s, x + Inches(0.15), adv_y + Inches(0.22), cw - 0.3, Inches(0.35),
                     desc, size=7.5, color=C['body_text'])
            adv_y += Inches(0.55)
        col_idx += 1

    if has_ci_data:
        x = sx + col_idx * (cw + gap)
        add_rect(s, x, sy, cw, Inches(0.5), fill=C['primary_dark'])
        add_text(s, x + Inches(0.2), sy + Inches(0.1), cw - 0.4, Inches(0.35),
                 '🛡️ 风险防御 · 重疾险', size=12, color=C['bg_light'], bold=True)
        adv_y = sy + Inches(0.7)
        for label, desc in ci_advantages:
            add_text(s, x + Inches(0.15), adv_y, Inches(1.3), Inches(0.22),
                     label, size=8, color=C['primary_dark'], bold=True)
            add_text(s, x + Inches(0.15), adv_y + Inches(0.22), cw - 0.3, Inches(0.35),
                     desc, size=7.5, color=C['body_text'])
            adv_y += Inches(0.6)

    if has_iul_data:
        x = sx + col_idx * (cw + gap)
        add_rect(s, x, sy, cw, Inches(0.5), fill=C['accent_dark'])
        add_text(s, x + Inches(0.2), sy + Inches(0.1), cw - 0.4, Inches(0.35),
                 '📈 财富杠杆 · 新加坡IUL', size=12, color=C['bg_light'], bold=True)
        adv_y = sy + Inches(0.7)
        for label, desc in iul_advantages[:6]:
            add_text(s, x + Inches(0.15), adv_y, Inches(1.3), Inches(0.22),
                     label, size=8, color=C['accent_dark'], bold=True)
            add_text(s, x + Inches(0.15), adv_y + Inches(0.22), cw - 0.3, Inches(0.35),
                     desc, size=7.5, color=C['body_text'])
            adv_y += Inches(0.55)

    # 箭头连接
    if prod_count >= 2:
        arrow_y = sy + Inches(2.0)
        for i in range(prod_count - 1):
            ax = sx + (i + 1) * (cw + gap) - Inches(0.4)
            add_text(s, ax, arrow_y, Inches(0.4), Inches(0.4),
                     '⟶', size=24, color=C['primary'], bold=True, align=PP_ALIGN.CENTER)

    # IUL/CI 右侧数据图表
    iul_bi = (iul_data or {}).get('benefit_illustration', []) if iul_data else []
    iul_has_data = len(iul_bi) >= 5
    ci_has_data = ci_data and len(ci_data.get('benefit_illustration', [])) >= 5

    if iul_has_data:
        try:
            iul_summary = iul_data.get('summary', {}) or {}
            iul_policy = iul_data.get('policy', {}) or {}
            iul_prem = float(iul_summary.get('annual_premium', iul_policy.get('annual_premium', 0)))
            iul_cov = float(iul_summary.get('sum_insured', iul_policy.get('sum_insured', 0)))
            iul_total = iul_prem * iul_pay_yrs
            iul_lev = (iul_cov / iul_total) if iul_total > 0 else 0
            iul_pay_yrs = int(iul_summary.get('payment_years', 5))

            # IUL 关键指标卡片 (右上)
            chart_x = Inches(8.0)
            chart_w = Inches(5.0)
            if not ci_has_data:
                # 无CI: IUL 指标卡片
                metrics_data = [
                    ('年缴保费', 'USD {:,.0f}'.format(iul_prem), C['primary']),
                    ('保额', 'USD {:,.0f}'.format(iul_cov), C['primary_dark']),
                    ('缴费年期', '{}年 · 共USD {:,.0f}'.format(iul_pay_yrs, iul_prem * iul_pay_yrs), C['mid_text']),
                    ('杠杆倍数', '{:.1f}x'.format(iul_lev), C['accent_dark']),
                ]
                card_y = Inches(2.0); card_h = Inches(1.3)
                for mi, (ml, mv, mc) in enumerate(metrics_data):
                    mx = chart_x + (mi % 2) * (chart_w / 2 + Inches(0.05))
                    my = card_y + (mi // 2) * (card_h + Inches(0.1))
                    mw = chart_w / 2 - Inches(0.05)
                    add_rect(s, mx, my, mw, card_h, fill=C['bg_card'])
                    add_rect(s, mx, my, mw, Inches(0.06), fill=mc)
                    add_text(s, mx + Inches(0.15), my + Inches(0.15), mw - 0.3, Inches(0.2),
                             ml, size=9, color=C['mid_text'])
                    add_text(s, mx + Inches(0.15), my + Inches(0.45), mw - 0.3, Inches(0.35),
                             mv, size=16, color=mc, bold=True)
                chart_y = Inches(4.9); chart_h = Inches(2.2)
            else:
                chart_y = Inches(6.5); chart_h = Inches(1.5)

            # IUL 折线图: 已缴总保费 vs 身故赔偿(杠杆)
            cd = CategoryChartData()
            labels = []; total_paid = []; ng_db = []
            for row in iul_bi[:81]:
                py = int(row.get('policy_year', 0))
                if py <= 0: continue
                labels.append('Y{}'.format(py) if py % 10 == 0 or py == 1 else '')
                paid = min(iul_prem * py, iul_prem * iul_pay_yrs)
                total_paid.append(paid)
                ng_db.append(int(row.get('non_guaranteed_death_benefit', row.get('death_benefit', iul_cov))))
            if labels and total_paid:
                cd.categories = labels
                cd.add_series('已缴总保费', total_paid)
                cd.add_series('身故赔偿(杠杆)', ng_db)
                cf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, chart_x, chart_y, chart_w, chart_h, cd)
                c = cf.chart; c.has_title = False; c.has_legend = True
                c.legend.position = XL_LEGEND_POSITION.BOTTOM
                for p in c.plots:
                    for i, ser in enumerate(p.series):
                        ser.format.line.width = Pt(2)
                        ser.format.line.color.rgb = _c(C['accent_dark'] if i else C['primary'])
        except Exception:
            pass

    # CI/IUL 指标卡片已合并到上方产品优势列中，此处不再重复显示

    if combo_text:
        add_text(s, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.0),
                 combo_text[:200], size=9, color=C['body_text'])

    add_text(s, Inches(0.4), Inches(7.2), Inches(12), Inches(0.2),
             '数据来源：保险公司官方计划书', size=7, color=C['mid_text'])
