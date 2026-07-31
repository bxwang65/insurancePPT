"""
混合渲染器: 模板外观 + 数据内容

1. 用券商风模板做封面/公司/排版框架
2. 对数据页(表格/图表), 在模板基础上叠加python-pptx内容
3. 保留模板的精致视觉
"""
import os, copy
from typing import Dict, List, Optional, Any
from pptx import Presentation
try:
    from ..extract.savings_normalizer import calc_irr as _ma_irr_no_wd, calc_irr_withdraw as _ma_irr_wd
except ImportError:
    _ma_irr_no_wd = None
    _ma_irr_wd = None
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree

from ..templates.style_tokens import get_theme, FONT_HEI, FONT_LATIN

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def _c(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _clone_slide(prs, slide):
    """克隆一个幻灯片"""
    slide_layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        try:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.append(el)
        except:
            pass
    return new_slide

def render_hybrid(data: Dict, output_path: str, template_path: str,
                   theme: str = 'broker',
                   logo_path: Optional[str] = None,
                   company_images: Optional[List[str]] = None,
                   scene_images: Optional[List[str]] = None,
                   ci_data: Optional[Dict] = None,
                   iul_data: Optional[Dict] = None) -> str:
    """混合渲染入口"""
    C = get_theme(theme)
    summary = data.get('summary', {})
    nw = {int(k) if str(k).isdigit() else k: v for k,v in data.get('no_withdraw',{}).items()}
    wd = {int(k) if str(k).isdigit() else k: v for k,v in data.get('withdraw',{}).items()}
    pt = data.get('paid_total', 0)
    meta = data.get('meta', {})
    hw = bool(wd) and any(v.get('Annual_WD',0) > 0 for v in wd.values())
    ia = meta.get('insured_age', 1) or 1

    # 补齐计算字段 (IRR 已由 savings_normalizer 用 M-A 计算; 这里兜底也用 M-A)
    py = int(meta.get('payment_years') or 0)
    cur = meta.get('currency') or meta.get('product_currency') or 'USD'
    # 提领起始年 (供 M-A 提领公式)
    start_wd_yr = 0
    annual_wd = 0
    if wd:
        for yk in sorted(int(k) for k in wd.keys()):
            aw = wd[yk].get('Annual_WD', 0) or 0
            if yk > 0 and aw > 0:
                start_wd_yr = yk
                annual_wd = aw
                break
    for d in [nw, wd]:
        for k, r in list(d.items()):
            t = r.get('Total', 0) or 0
            if r.get('Mult') is None: r['Mult'] = t/pt if pt else 0
            if r.get('IRR') is None and int(k) > 0:
                if d is wd and _ma_irr_wd:
                    r['IRR'] = _ma_irr_wd(int(k), (r.get('Cum_WD',0) or 0)+t, pt, py, cur, start_wd_yr, annual_wd)
                elif _ma_irr_no_wd:
                    r['IRR'] = _ma_irr_no_wd(int(k), t, pt, py, cur)
            if r.get('Simple') is None and int(k) > 0: r['Simple'] = (t-pt)/pt/float(k) if pt else None

    # 打开模板
    tpl = Presentation(template_path)

    # === 数据替换映射 ===
    replace_map = {
        "2026/06/01": meta.get("date", "2026/06"),
        "2026/06": meta.get("date", "2026/06"),
        "家族传承财富保障方案": summary.get('product_name', '家庭财富保障方案'),
        "汇报人：": f"汇报人：{meta.get('company_short','')} · 财富管理团队",
    }

    # 公司数据替换
    company_data = [
        (meta.get('brand_profile',{}).get('founded_year','—'), '成立年份'),
        (meta.get('brand_profile',{}).get('rating_value','—'), '财务实力评级'),
        ("5000位", "高净值客户"),
    ]
    comp_map = {}
    for val, label in company_data:
        comp_map[str(val)] = str(val)
        comp_map[label] = label

    # 执行文本替换 (所有页)
    for slide in tpl.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replace_map.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    # === 选择要保留的模板页面 ===
    # 保留: 1封面, 2公司, 3服务体系, 5产品定位, 6增长, 7提领, 25结束
    keep_indices = [0, 1, 2, 4, 5, 6, 24]  # 0-indexed
    all_slides = list(range(len(tpl.slides)))
    to_delete = [i for i in all_slides if i not in keep_indices]
    for idx in sorted(to_delete, reverse=True):
        rId = tpl.slides._sldIdLst[idx].rId
        tpl.part.drop_rel(rId)
        tpl.slides._sldIdLst.remove(tpl.slides._sldIdLst[idx])

    # 重映射索引
    kept_slides = list(tpl.slides)

    # === 用 python-pptx 创建数据幻灯片并追加 ===
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def add_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def add_text(slide, text, x=0.5, y=0.3, w=12, h=0.6, size=10, color=None, bold=False, align=PP_ALIGN.LEFT):
        if color is None: color = '#332825'
        if isinstance(color, str): color = _c(color)
        txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txbox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT_HEI
        rPr = run._r.get_or_add_rPr()
        for tag, font in [('a:ea', FONT_HEI), ('a:latin', FONT_LATIN)]:
            el = rPr.find(qn(tag))
            if el is None: el = etree.SubElement(rPr, qn(tag))
            el.set('typeface', font)
        return txbox

    def add_table(slide, rows, headers, x, y, w, h):
        """添加数据表"""
        rows_data = [headers] + rows
        table_shape = slide.shapes.add_table(len(rows_data), len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
        table = table_shape.table
        for ri, row in enumerate(rows_data):
            for ci, cell_val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = str(cell_val)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(8)
                        run.font.name = FONT_HEI
                        if ri == 0:
                            run.font.bold = True
                            run.font.color.rgb = _c('#FFFFFF')
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _c('#0D1B2A')
                        else:
                            run.font.color.rgb = _c('#1A2A3A')
                            bg = '#F0F2F5' if ri % 2 == 0 else '#FFFFFF'
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _c(bg)
        return table_shape

    # 复制模板页到输出
    for idx in keep_indices:
        if idx < len(tpl.slides):
            _clone_slide(prs, tpl.slides[idx])

    # === 添加数据页 ===
    # 增长图表页
    s = add_slide()
    yrs = sorted(int(k) for k in nw.keys() if nw[k].get('Total',0) > 0)
    add_text(s, "现金价值增长趋势", x=0.5, y=0.3, w=12, h=0.6, size=22, bold=True, color=C['dark_text'])
    add_text(s, "保证+非保证 · 不提领情形", x=0.5, y=0.9, w=12, h=0.4, size=11, color=C['mid_text'])

    if yrs:
        front = [y for y in yrs if y <= 30]
        tail = [y for y in yrs if y > 30][:4]
        show_yrs = front + tail
        guar = [nw[y]['Guar_CV']/1000 for y in show_yrs]
        non_g = [(nw[y]['Total']-nw[y]['Guar_CV'])/1000 for y in show_yrs]
        total_v = [nw[y]['Total']/1000 for y in show_yrs]
        labels = [f'Y{y}' for y in show_yrs]

        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series('保证现价', guar)
        cd.add_series('非保证', non_g)
        chart_frame = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.5), Inches(6), Inches(5), cd)
        chart = chart_frame.chart
        chart.has_title = False; chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        for plot in chart.plots:
            for i, ser in enumerate(plot.series):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = _c(C['primary']) if i == 0 else _c(C['accent'])

        cd2 = CategoryChartData()
        cd2.categories = labels; cd2.add_series('退保总额', total_v)
        chart_frame2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(6.7), Inches(1.5), Inches(6), Inches(5), cd2)
        lc = chart_frame2.chart; lc.has_title = False; lc.has_legend = False
        for plot in lc.plots:
            for ser in plot.series:
                ser.format.line.color.rgb = _c(C['accent_dark']); ser.format.line.width = Pt(2.5)

    # 不提领表
    s = add_slide()
    add_text(s, "不提领情形 · 退保发还金额明细", x=0.5, y=0.3, w=12, h=0.6, size=20, bold=True, color=C["dark_text"])
    show = [y for y in [1,5,10,15,20,25,30,35,40,50,60] if y in nw][:12]
    if show:
        rows = []
        for y in show:
            r = nw[y]
            irr = f"{r['IRR']*100:.2f}%" if r.get('IRR') else '-'
            rows.append([str(y), str(r.get('Age',y)), f"{r.get('Paid',pt):,}", f"{r['Guar_CV']:,}", f"{r['Total']:,}", irr])
        add_table(s, rows, ['年度','年龄','已缴保费','保证现价','退保总额','IRR'], x=0.5, y=1.5, w=12, h=0.4*len(rows))

    # 提领表
    if hw:
        s = add_slide()
        ws = next((y for y in sorted(wd.keys()) if wd[y].get('Annual_WD',0) > 0), 7)
        add_text(s, f"提领方案 · Y{ws}起年提 USD {wd[ws]['Annual_WD']:,}", x=0.5, y=0.3, w=12, h=0.6, size=20, bold=True, color=C['dark_text'])
        show2 = sorted(set([5,ws,ws+1,10,15,20,30,40,50,60]))[:12]
        if show2:
            rows2 = []
            for y in show2:
                if y in wd:
                    r = wd[y]
                    rows2.append([str(y), str(r.get('Age',y)), f"{r.get('Annual_WD',0):,}", f"{r.get('Cum_WD',0):,}", f"{r.get('Total',0):,}"])
            add_table(s, rows2, ['年度','年龄','年提取','累计','退保后价值'], x=0.5, y=1.5, w=12, h=0.35*len(rows2))

    # 保存
    prs.save(output_path)
    return output_path
