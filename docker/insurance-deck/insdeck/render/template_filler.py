"""
模板填充器: 用专业PPTX模板 + 保险数据 → 输出精美计划书

原理:
1. 打开专业设计的PPTX模板
2. 找到每页的文本占位
3. 替换为实际数据
4. 保留原始设计(字体、颜色、布局、图片)
"""
import os, copy, re, json
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

LETTER_ADDR = {}

def fill_template(template_path: str, output_path: str, data: Dict[str, Any],
                   slides_to_keep: Optional[List[int]] = None) -> str:
    """
    用数据填充模板

    Args:
        template_path: 模板 PPTX 路径
        output_path: 输出 PPTX 路径
        data: 数据 dict, 格式见下方
        slides_to_keep: 要保留的模板页码列表 (1-based), None=全部

    data 格式:
    {
        "date": "2026/06",
        "title": "封面标题",
        "reporter": "汇报人",
        "company_name": "公司名",
        "company_stats": [("15年","行业深耕"), ("5000位","高净值客户"), ...],
        "product_name": "产品名",
        "features": [{"title":"亮点1","body":"描述"}, ...],
        "insured_name": "客户名",
        "insured_age": "年龄",
        "premium": "$100,000",
        "pay_years": "5年",
        "table_data": [["Y1","100","200"], ...],
        ...
    }
    """
    prs = Presentation(template_path)

    # 构建替换映射: {原文本 → 新文本}
    text_map = _build_text_map(data)

    # 处理每一页
    slides_to_process = slides_to_keep if slides_to_keep else list(range(1, len(prs.slides) + 1))
    slides_to_process = [i for i in slides_to_process if 1 <= i <= len(prs.slides)]

    for slide_idx in slides_to_process:
        slide = prs.slides[slide_idx - 1]
        _process_slide(slide, text_map, data, slide_idx)

    # 删除不需要的页 (从后往前删, 避免索引变化)
    all_slides = set(range(1, len(prs.slides) + 1))
    to_delete = all_slides - set(slides_to_process)
    for idx in sorted(to_delete, reverse=True):
        rId = prs.slides._sldIdLst[idx - 1].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx - 1])

    prs.save(output_path)
    return output_path


def _build_text_map(data: Dict) -> Dict[str, str]:
    """构建简单的文本替换映射"""
    m = {}

    # 日期
    m["2026/06/01"] = data.get("date", "2026/06")
    m["2026/06"] = data.get("date", "2026/06")

    # 标题
    title = data.get("title", "")
    if title:
        m["家族传承财富保障方案"] = title
        m["家族财富传承方案"] = title

    # 汇报人
    reporter = data.get("reporter", "")
    if reporter:
        m["汇报人："] = f"汇报人：{reporter}"

    # 产品名
    pname = data.get("product_name", "")
    if pname:
        m["兼具保障与传承功能的终身寿险"] = pname
        m["终身寿险"] = pname

    # 受保人
    iname = data.get("insured_name", "")
    if iname:
        m["VIP 先生"] = iname

    return m


def _process_slide(slide, text_map: Dict, data: Dict, slide_num: int):
    """处理单个幻灯片: 替换文本"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                orig = run.text.strip()
                if orig in text_map:
                    run.text = run.text.replace(orig, text_map[orig])


def quick_fill_savings(template_path: str, output_path: str, savings_data: Dict):
    """
    一键填充储蓄险方案

    savings_data:
    {
        "title": "封面标题",
        "date": "2026/06",
        "reporter": "团队名称",
        "company_name": "友邦保险",
        "product_name": "环宇盈活",
        "insured_name": "VIP 先生",
        "insured_age": "1",
        "premium": "USD 105,000",
        "pay_years": "5年",
        "basic_info": {...}
    }
    """
    data = {
        "date": savings_data.get("date", "2026/06"),
        "title": savings_data.get("title", "家庭财富保障方案"),
        "reporter": savings_data.get("reporter", "财富管理团队"),
        "product_name": savings_data.get("product_name", ""),
        "insured_name": savings_data.get("insured_name", "VIP 先生"),
    }

    # 券商风模板映射: 保留的页码 → 内容
    # P1=封面, P2=公司, P3=服务体系/亮点, P5=产品定位, P6=增长
    # P7=提领, P10=增长示例, P13=收益表, P14=提领对比
    # P25=结束
    keep = [1, 2, 3, 5, 6, 7, 10, 13, 14, 25]

    return fill_template(template_path, output_path, data, keep)
