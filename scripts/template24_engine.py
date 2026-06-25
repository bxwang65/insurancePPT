#!/usr/bin/env python3
"""
通用保险PPT模板引擎 — 基于模板24设计语言
支持：储蓄险（匠心传承）/ 重疾险（守护家倍）/ IUL
用法: python template24_engine.py --product jxcc --data jxcc_data.json
"""
import sys, os, re, json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import io

# ═══════════════════════════════════════════════════════════════════════════════
# 模板24设计语言 — 配色 & 常量
# ═══════════════════════════════════════════════════════════════════════════════
T24 = {
    # 主题色
    'deep_blue':    RGBColor(0x0A, 0x3C, 0x5F),   # 主背景 #0A3C5F
    'teal':         RGBColor(0x18, 0x89, 0x8D),   # 强调色 #18898D
    'mid_blue':     RGBColor(0x38, 0x5A, 0x64),   # 中间色 #385A64
    'dark_gray':    RGBColor(0x2D, 0x38, 0x47),   # 深灰 #2D3847
    'light_gray':   RGBColor(0xE7, 0xE6, 0xE6),   # 浅灰 #E7E6E6
    'white':        RGBColor(0xFF, 0xFF, 0xFF),
    'black':        RGBColor(0x00, 0x00, 0x00),

    # 语义色
    'gold':         RGBColor(0xC9, 0xA0, 0x27),   # 金色
    'positive':     RGBColor(0x2E, 0xD4, 0x8E),   # 增长绿
    'warning':      RGBColor(0xFF, 0xB8, 0x00),   # 警示橙
    'negative':     RGBColor(0xFF, 0x4D, 0x4D),   # 下降红

    # 内部色（计算用）
    'card_bg':      RGBColor(0x0D, 0x2B, 0x44),   # 卡片深色背景
    'text_white':   RGBColor(0xFF, 0xFF, 0xFF),
    'text_gray':    RGBColor(0xCC, 0xD1, 0xD9),
    'divider':      RGBColor(0x1A, 0x3A, 0x55),

    # 页面常量
    'W': 13.33,
    'H': 7.5,
    'TOP_BAR_H': 0.06,
    'DECO_LEFT_W': 0.12,
}

# ═══════════════════════════════════════════════════════════════════════════════
# 基础工具函数
# ═══════════════════════════════════════════════════════════════════════════════
def In(s): return Inches(s)
def PtS(s): return Pt(s)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=12, color=None, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Heiti TC"):
    txBox = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = PtS(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox

def add_accent_bar(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, page_num, total_pages):
    """底部页码 — 模板24风格"""
    add_rect(slide, 0, T24['H']-0.45, T24['W'], 0.45, fill_color=T24['dark_gray'])
    add_text(slide, f"第{page_num}页/共{total_pages}页",
             T24['W']-1.8, T24['H']-0.4, 1.7, 0.35,
             font_size=10, color=T24['text_gray'])

def kpi_card(slide, l, t, w, h, label, value, subtitle="", icon="★"):
    """标准KPI卡片 — 模板24风格"""
    add_rect(slide, l, t, w, h, fill_color=T24['card_bg'])
    if icon:
        add_text(slide, icon, l+0.1, t+0.06, 0.45, 0.38, font_size=16, color=T24['teal'])
    add_text(slide, label, l+0.1, t+0.1, w-0.2, 0.32, font_size=9, color=T24['text_gray'])
    add_text(slide, value, l+0.1, t+0.38, w-0.2, 0.52, font_size=20, color=T24['white'], bold=True)
    if subtitle:
        add_text(slide, subtitle, l+0.1, t+0.88, w-0.2, 0.3, font_size=8, color=T24['teal'])

def make_fig(w, h, draw_fn, dark=True):
    """生成matplotlib图表"""
    # 局部字体设置，避免污染全局状态
    from matplotlib import font_manager
    heiti_fp = None
    for fp in ['Heiti TC', 'STHeiti', 'Microsoft YaHei', 'SimHei', 'Arial Unicode MS']:
        try:
            heiti_fp = font_manager.FontProperties(fname=font_manager.findfont(font_manager.FontProperties(family=fp)))
            break
        except Exception:
            pass
    old_family = matplotlib.rcParams['font.family']
    old_minus = matplotlib.rcParams['axes.unicode_minus']
    try:
        matplotlib.rcParams['font.family'] = 'Heiti TC' if heiti_fp is None else old_family
        matplotlib.rcParams['axes.unicode_minus'] = False
        fig, ax = plt.subplots(figsize=(w, h))
        plt.tight_layout(pad=0.4)
        draw_fn(ax, dark=dark)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor='none', transparent=True)
        buf.seek(0)
        plt.close(fig)
        return buf
    finally:
        matplotlib.rcParams['font.family'] = old_family
        matplotlib.rcParams['axes.unicode_minus'] = old_minus

def fmt_usd(v):
    if v >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if v >= 1_000:     return f"${v/1_000:.0f}K"
    return f"${v:,.0f}"

# ═══════════════════════════════════════════════════════════════════════════════
# 标准图表样式（模板24深蓝背景）
# ═══════════════════════════════════════════════════════════════════════════════
def dark_axis(ax, bg='#0A3C5F'):
    spine_c = '#88AABB'
    ax.set_facecolor(bg)
    ax.spines['bottom'].set_color(spine_c)
    ax.spines['left'].set_color(spine_c)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors=spine_c, labelsize=8)
    ax.grid(axis='y', alpha=0.2, color='#FFFFFF', linewidth=0.5)

# ═══════════════════════════════════════════════════════════════════════════════
# 第1页：封面 — 模板24深蓝背景 + 左侧装饰条
# ═══════════════════════════════════════════════════════════════════════════════
def page_cover(prs, product_name, company_name, insured_name,
               annual_prem, total_prem, withdrawal_annual, withdrawal_start,
               highlights, currency="USD"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])

    # 顶部金色条
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    # 左侧装饰条（青绿色）
    add_accent_bar(s, 0, 0.4, T24['DECO_LEFT_W'], 4.5, T24['teal'])

    # 主标题区
    add_text(s, product_name, 0.5, 0.5, 9, 1.0,
             font_size=44, color=T24['white'], bold=True)
    add_text(s, f"{company_name} · 美元保单", 0.5, 1.5, 7, 0.5,
             font_size=18, color=T24['teal'])
    add_accent_bar(s, 0.5, 2.1, 5.5, 0.04, T24['teal'])

    # 三个KPI
    kpi_card(s, 0.5, 2.5, 3.0, 1.4, "年缴保费", f"${annual_prem:,}",
             f"{currency} · 5年共 ${total_prem:,}", "💰")
    kpi_card(s, 3.8, 2.5, 3.0, 1.4, "每年提取",
             f"${withdrawal_annual:,}" if withdrawal_annual else "—",
             f"第{withdrawal_start}年起终身" if withdrawal_start else "", "📅")
    kpi_card(s, 7.1, 2.5, 3.0, 1.4, "产品类型",
             "储蓄寿险", "终身保障", "🛡️")

    # 亮点区（右上卡片）
    add_rect(s, 10.2, 0.4, 2.8, 5.5, fill_color=T24['card_bg'])
    add_text(s, "计划亮点", 10.35, 0.5, 2.5, 0.45,
             font_size=13, color=T24['teal'], bold=True)
    for i, h in enumerate(highlights):
        add_text(s, h, 10.35, 1.0+i*0.55, 2.6, 0.5,
                 font_size=10, color=T24['white'])

    # 底部信息
    add_text(s, f"受保人: {insured_name}  |  货币: {currency}  |  状态: 计划书  |  AIA",
             0.5, 6.8, 10, 0.4, font_size=10, color=T24['text_gray'])
    add_page_number(s, 1, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第2页：计划概览
# ═══════════════════════════════════════════════════════════════════════════════
def page_overview(prs, policy_params, company_name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    add_text(s, "计划概览", 0.5, 0.25, 12, 0.6,
             font_size=30, color=T24['white'], bold=True)
    add_text(s, f"{company_name} 核心参数", 0.5, 0.9, 12, 0.38,
             font_size=13, color=T24['text_gray'])

    # 双列参数卡片
    for i, (k, v) in enumerate(policy_params):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.2
        y = 1.45 + row * 0.85
        add_rect(s, x, y, 5.8, 0.72, fill_color=T24['card_bg'])
        add_accent_bar(s, x, y, 0.06, 0.72, T24['teal'])
        add_text(s, k, x+0.15, y+0.05, 2.0, 0.35, font_size=10, color=T24['text_gray'])
        add_text(s, str(v), x+2.2, y+0.05, 3.4, 0.55, font_size=13, color=T24['white'], bold=True)

    add_page_number(s, 2, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第3页：增长故事 — 面积图
# ═══════════════════════════════════════════════════════════════════════════════
def page_growth_story(prs, bmap, total_prem, key_years, product_name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['gold'])

    add_text(s, f"💰 账户价值增长 — 复利的魔力", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "不退保、持续累积 — 复归红利与终期分红双引擎驱动", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    gcv_vals = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 for y in key_years]

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        ax.fill_between(x, gcv_vals, sv_vals,
                        where=[a > b for a, b in zip(sv_vals, gcv_vals)],
                        color='#C9A027', alpha=0.25, label='终期分红（非保证）', zorder=1)
        ax.plot(x, sv_vals, color='#C9A027', linewidth=2.5, marker='o', markersize=7,
                label='退保总额', zorder=3)
        ax.plot(x, gcv_vals, color='#18898D', linewidth=2.0, marker='s', markersize=5,
                label='保证现金价值', zorder=3)
        # 回本标注
        for i, (yv, sv) in enumerate(zip(key_years, sv_vals)):
            if sv >= total_prem and i > 0:
                ax.axvline(x=i, color='#2ED48E', linestyle='--', linewidth=1.5, alpha=0.7)
                ax.annotate(f'回本 Y{yv}', xy=(i, sv), xytext=(i+0.5, sv*1.05),
                            fontsize=9, color='#2ED48E', fontweight='bold')
                break
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.2))

    # 右侧KPI
    y7 = bmap.get(7, {})
    y20 = bmap.get(20, {})
    y30 = bmap.get(30, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.15, "Y7 退保", fmt_usd(y7.get('total_surrender_value',0)), "已回本", "📍")
    kpi_card(s, 10.2, 2.55, 2.8, 1.15, "Y20 倍数", f"{y20.get('total_surrender_value',0)/total_prem:.1f}x", "退保/已缴保费", "📈")
    kpi_card(s, 10.2, 3.8, 2.8, 1.15, "Y30 倍数", f"{y30.get('total_surrender_value',0)/total_prem:.1f}x", "退保/已缴保费", "🚀")
    kpi_card(s, 10.2, 5.05, 2.8, 1.15, "保证占比",
             f"{y20.get('guaranteed_cash_value',0)/y20.get('total_surrender_value',1)*100:.0f}%",
             "Y20保证现金", "🛡️")

    add_page_number(s, 3, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第4页：不提取每5年明细表
# ═══════════════════════════════════════════════════════════════════════════════
def page_no_withdrawal_table(prs, bmap, total_prem, key_years):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['gold'])

    add_text(s, "📊 不提取每5年明细 — 账户价值累积一览", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "持续累积，复利驱动 — 退保总额何时回本、何时翻倍？", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    n = len(key_years)
    tbl = s.shapes.add_table(n+1, 5, In(0.4), In(1.3), In(12.5), In(0.5*(n+1))).table

    headers = ['保单年龄', '已交保费', '退保现金价值', '倍数', '备注']
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = T24['dark_gray']
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = T24['teal']
            p.alignment = PP_ALIGN.CENTER

    for i, y in enumerate(key_years):
        r = bmap.get(y, {})
        sv = r.get('total_surrender_value', 0) or 0
        mult = sv / total_prem
        paid = r.get('total_premium_paid', 0) or 0

        is_break_even = (mult >= 1.0)
        is_double = (mult >= 2.0)
        note = '✅ 回本' if is_break_even else ('🚀 翻倍' if is_double else '')

        row_data = [f'Y{y}', f'${paid:,}', f'${sv:,}', f'{mult:.2f}x', note]
        for j, val in enumerate(row_data):
            cell = tbl.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = T24['card_bg']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                if is_double:
                    p.font.color.rgb = T24['gold']
                    p.font.bold = True
                elif is_break_even:
                    p.font.color.rgb = T24['positive']
                    p.font.bold = True
                else:
                    p.font.color.rgb = T24['text_white']

    y10 = bmap.get(10, {})
    y20 = bmap.get(20, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.1, "Y10 回本", f'${y10.get("total_surrender_value",0)/1e6:.2f}M', "1.28x保费", "📍")
    kpi_card(s, 10.2, 2.45, 2.8, 1.1, "Y20 翻倍", f'${y20.get("total_surrender_value",0)/1e6:.2f}M', "2.73x保费", "🚀")
    kpi_card(s, 10.2, 3.6, 2.8, 1.1, "Y30 价值", f'${bmap.get(30,{}).get("total_surrender_value",0)/1e6:.2f}M', "5.57x保费", "📈")
    kpi_card(s, 10.2, 4.75, 2.8, 1.1, "Y40 价值", f'${bmap.get(40,{}).get("total_surrender_value",0)/1e6:.2f}M', "10.92x保费", "🌱")

    add_page_number(s, 4, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第5页：终身现金流时间轴
# ═══════════════════════════════════════════════════════════════════════════════
def page_cashflow_timeline(prs, benefit, withdrawal, key_years):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    add_text(s, "💵 终身现金流时间轴 — 第7年起每年提取", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "从第7个保单年度开始，每年提取 $35,000 美元，活到老领到老", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    bmap = {r['policy_year']: r for r in benefit}
    sv_no_withdraw = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        ax.bar(x, sv_no_withdraw, color='#18898D', alpha=0.5, label='不退保退保总额', zorder=2)
        withdrawal_vals = []
        wmap = {r['policy_year']: r for r in withdrawal}
        for y in key_years:
            if y == 1:
                withdrawal_vals.append(sv_no_withdraw[0])
            else:
                wr = wmap.get(y, {})
                if wr:
                    withdrawal_vals.append(wr.get('surrender_value_after', 0) or sv_no_withdraw[y-1])
                else:
                    withdrawal_vals.append(sv_no_withdraw[y-1])
        ax.plot(x, withdrawal_vals, color='#2ED48E', linewidth=2.5, marker='s', markersize=5,
                label='每年提取后退保总额', zorder=3)
        ax.annotate('第7年起', xy=(6, withdrawal_vals[6]), xytext=(7.5, withdrawal_vals[6]*1.15),
                    fontsize=9, color='#C9A027', fontweight='bold',
                    arrowprops=dict(arrowstyle='->', color='#C9A027'))
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.2))

    y7 = bmap.get(7, {})
    y15 = bmap.get(15, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.15, "Y7退保价值", fmt_usd(y7.get('total_surrender_value',0)), "开始提取", "📍")
    kpi_card(s, 10.2, 2.55, 2.8, 1.15, "每年提取", "$35,000", "终身，活到老领到老", "💵")
    kpi_card(s, 10.2, 3.8, 2.8, 1.15, "15年累计提取", f"${35000*9:,}", "Y7-Y15共9年", "📊")
    kpi_card(s, 10.2, 5.05, 2.8, 1.15, "Y15退保价值", fmt_usd(y15.get('total_surrender_value',0)), "仍可退保", "📈")

    add_page_number(s, 5, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第6页：提领 vs 不提领对比
# ═══════════════════════════════════════════════════════════════════════════════
def page_withdrawal_comparison(prs, benefit, withdrawal, key_years, total_prem):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['warning'])

    add_text(s, "⚖️ 提领 vs 不提领 — 账户价值双线对比", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "同样本金，不同策略：持续累积 vs 终身现金流", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    bmap = {r['policy_year']: r for r in benefit}
    wmap = {r['policy_year']: r for r in withdrawal}
    no_wd = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    with_wd = []
    for y in key_years:
        if y == 1:
            with_wd.append(no_wd[0])
        else:
            wr = wmap.get(y, {})
            with_wd.append(wr.get('surrender_value_after', 0) if wr else no_wd[y-1])

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        ax.fill_between(x, with_wd, no_wd,
                        where=[a > b for a, b in zip(no_wd, with_wd)],
                        color='#C9A027', alpha=0.2, label='提领差额（机会成本）', zorder=1)
        ax.plot(x, no_wd, color='#C9A027', linewidth=2.5, marker='o', markersize=6,
                label='不提领退保总额', zorder=3)
        ax.plot(x, with_wd, color='#2ED48E', linewidth=2.5, marker='s', markersize=6,
                label='每年提取后退保总额', zorder=3)
        ax.axvline(x=6, color='#FFB800', linestyle='--', linewidth=1.5, alpha=0.7)
        ax.text(6.1, max(no_wd)*0.8, '提领起始\nY7', fontsize=8, color='#FFB800')
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.2))

    y15_no = no_wd[14]
    y15_wd = with_wd[14]
    kpi_card(s, 10.2, 1.3, 2.8, 1.15, "Y15 不提领", fmt_usd(y15_no), f"{y15_no/total_prem:.1f}x保费", "📈")
    kpi_card(s, 10.2, 2.55, 2.8, 1.15, "Y15 提领后", fmt_usd(y15_wd), "已领$35K×9年", "💵")
    kpi_card(s, 10.2, 3.8, 2.8, 1.15, "9年累计提取", f"${35000*9:,}", "Y7至Y15共9年", "📊")
    kpi_card(s, 10.2, 5.05, 2.8, 1.15, "机会成本", fmt_usd(y15_no-y15_wd), "vs不提领", "⚖️")

    add_page_number(s, 6, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第7页：回本分析
# ═══════════════════════════════════════════════════════════════════════════════
def page_breakeven(prs, benefit, total_prem, key_years):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['positive'])

    add_text(s, "📍 回本分析 — 保证回本年份", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, f"退保总额何时超越已缴保费 ${total_prem:,}？", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    bmap = {r['policy_year']: r for r in benefit}
    breakeven = next((r['policy_year'] for r in sorted(benefit, key=lambda x: x['policy_year'])
                      if r.get('total_surrender_value', 0) >= total_prem), None)
    gcv_breakeven = next((r['policy_year'] for r in sorted(benefit, key=lambda x: x['policy_year'])
                          if r.get('guaranteed_cash_value', 0) >= total_prem), None)

    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    gcv_vals = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 for y in key_years]

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        ax.fill_between(x, [total_prem]*len(key_years), sv_vals,
                        where=[sv >= total_prem for sv in sv_vals],
                        color='#2ED48E', alpha=0.15, label='已回本区域', zorder=1)
        ax.plot(x, sv_vals, color='#C9A027', linewidth=2.5, marker='o', markersize=6,
                label='退保总额', zorder=3)
        ax.plot(x, gcv_vals, color='#18898D', linewidth=2.0, marker='s', markersize=5,
                label='保证现金价值', zorder=3)
        ax.axhline(y=total_prem, color='#FF4D4D', linestyle='--', linewidth=1.5,
                   label=f'已缴保费 ${total_prem/1e6:.1f}M', alpha=0.8)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.5, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.5))

    y15 = bmap.get(15, {})
    y7 = bmap.get(7, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.15, "总价值回本", f"Y{breakeven}" if breakeven else "—", "退保总额≥已缴保费", "📍")
    kpi_card(s, 10.2, 2.55, 2.8, 1.15, "保证回本", f"Y{gcv_breakeven}" if gcv_breakeven else "—", "保证现金价值≥已缴", "🛡️")
    kpi_card(s, 10.2, 3.8, 2.8, 1.15, "Y15退保", fmt_usd(y15.get('total_surrender_value',0)), f"{y15.get('total_surrender_value',0)/total_prem:.1f}x保费", "📈")
    kpi_card(s, 10.2, 5.05, 2.8, 1.15, "Y7退保", fmt_usd(y7.get('total_surrender_value',0)), f"{y7.get('total_surrender_value',0)/total_prem:.1f}x保费", "✅")

    add_page_number(s, 7, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第8页：提取每5年明细表
# ═══════════════════════════════════════════════════════════════════════════════
def page_withdrawal_table(prs, benefit, withdrawal, total_prem, key_years):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    add_text(s, "💵 提取每5年明细 — 终身现金流方案", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "第7年起每年提取 $35,000 — 退保总额何时回本、何时翻倍？", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    bmap = {r['policy_year']: r for r in benefit}
    wmap = {r['policy_year']: r for r in withdrawal}
    n = len(key_years)

    tbl = s.shapes.add_table(n+1, 6, In(0.4), In(1.3), In(12.5), In(0.5*(n+1))).table
    headers = ['保单年龄', '已交保费', '提取金额', '累计提取', '退保现金价值', '倍数']
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = T24['dark_gray']
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = T24['teal']
            p.alignment = PP_ALIGN.CENTER

    for i, y in enumerate(key_years):
        b = bmap.get(y, {})
        w = wmap.get(y, {})
        paid = b.get('total_premium_paid', total_prem)
        annual_wd = w.get('annual_withdrawal', 0) if w else 0
        cum_wd = w.get('cumulative_withdrawals', 0) if w else 0
        sv_after = w.get('surrender_value_after', 0) if w else b.get('total_surrender_value', 0)
        mult = sv_after / total_prem if sv_after else 0

        is_break_even = (mult >= 1.0)
        is_double = (mult >= 2.0)

        row_data = [
            f'Y{y}', f'${paid:,}',
            f'${annual_wd:,}' if annual_wd else '—',
            f'${cum_wd:,}' if cum_wd else '—',
            f'${sv_after:,}' if sv_after else '—',
            f'{mult:.2f}x',
        ]
        for j, val in enumerate(row_data):
            cell = tbl.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = T24['card_bg']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                if is_double:
                    p.font.color.rgb = T24['gold']
                    p.font.bold = True
                elif is_break_even:
                    p.font.color.rgb = T24['positive']
                    p.font.bold = True
                else:
                    p.font.color.rgb = T24['text_white']

    y20 = wmap.get(20, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.1, "每年提取", "$35,000", "第7年起终身", "💵")
    kpi_card(s, 10.2, 2.45, 2.8, 1.1, "Y20 退保", f"${y20.get('surrender_value_after',0):,}", "1.06x保费", "📍")
    kpi_card(s, 10.2, 3.6, 2.8, 1.1, "Y20 累计领", f"${y20.get('cumulative_withdrawals',0):,}", "约100万", "📊")
    kpi_card(s, 10.2, 4.75, 2.8, 1.1, "Y30 退保", f"${wmap.get(30,{}).get('surrender_value_after',0):,}" or "—", "1.40x保费", "📈")

    add_page_number(s, 8, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第9页：年度收益明细
# ═══════════════════════════════════════════════════════════════════════════════
def page_yearly_detail(prs, benefit, total_prem):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['deep_blue'])

    add_text(s, "📊 年度收益明细 — 前20年核心数据", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "每年退保总额、保证现金价值、复归红利、终期分红一览", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    rows = [r for r in benefit if r['policy_year'] <= 20]
    n = len(rows)
    table = s.shapes.add_table(n+1, 6, In(0.4), In(1.3), In(12.5), In(0.38*(n+1))).table

    headers = ['年度', '已缴保费', '保证现金价值', '复归红利', '终期分红', '退保总额']
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = T24['dark_gray']
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = T24['teal']
            p.alignment = PP_ALIGN.CENTER

    for i, r in enumerate(rows):
        py = r['policy_year']
        row_data = [
            f'Y{py}',
            fmt_usd(r['total_premium_paid']),
            fmt_usd(r['guaranteed_cash_value']),
            fmt_usd(r['reversionary_bonus']),
            fmt_usd(r['terminal_dividend']),
            fmt_usd(r['total_surrender_value']),
        ]
        for j, val in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = T24['card_bg']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = T24['text_white'] if j != 5 else T24['gold']
                if j == 5: p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

    add_page_number(s, 9, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第10页：退保 vs 身故
# ═══════════════════════════════════════════════════════════════════════════════
def page_surrender_vs_death(prs, bmap, key_years, total_prem):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    add_text(s, "🏛️ 退保 vs 身故 — 两种选择，一个方案", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "无论选择退保变现还是代际传承，匠心传承都是您的好选择", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    death_vals = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 + bmap.get(y, {}).get('terminal_dividend', 0) or 0 for y in key_years]

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        width = 0.35
        ax.bar(x - width/2, sv_vals, width, color='#C9A027', alpha=0.8, label='退保总额', zorder=2)
        ax.bar(x + width/2, death_vals, width, color='#18898D', alpha=0.8, label='身故赔偿估算', zorder=2)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=11)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=10, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.2))

    advantages = [
        ("💰 退保变现", "第20年可退保 $1.37M\n享受资产增值收益\n资金灵活支配"),
        ("🛡️ 身故传承", "指定受益人\n免遗产税\n财富无损代际传递"),
        ("🔄 灵活切换", "可以先提取后传承\n也可以先传承后提取\n按人生阶段自由调整"),
    ]
    for i, (title, desc) in enumerate(advantages):
        add_rect(s, 10.2, 1.3+i*1.7, 2.8, 1.55, fill_color=T24['card_bg'])
        add_text(s, title, 10.3, 1.35+i*1.7, 2.6, 0.35, font_size=12, color=T24['teal'], bold=True)
        add_text(s, desc, 10.3, 1.7+i*1.7, 2.6, 1.0, font_size=9, color=T24['white'])

    add_page_number(s, 10, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第11页：长期增长预测
# ═══════════════════════════════════════════════════════════════════════════════
def page_long_term(prs, bmap, total_prem, key_years):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['teal'])

    add_text(s, "🌱 长期增长预测 — Y30至Y40财富轨迹", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "匠心传承越长期越值钱 — 终身复利，世代传承", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]

    def draw(ax, dark=True):
        x = np.arange(len(key_years))
        ax.fill_between(x, 0, sv_vals, color='#18898D', alpha=0.3, zorder=1)
        ax.plot(x, sv_vals, color='#18898D', linewidth=2.5, marker='o', markersize=5, zorder=3)
        for i, (yv, sv) in enumerate(zip(key_years, sv_vals)):
            if yv in [25, 30, 35, 40]:
                ax.annotate(f'Y{yv}\n{sv/1e6:.1f}M', xy=(i, sv),
                           xytext=(i, sv*1.08), fontsize=8, color='#C9A027',
                           ha='center', fontweight='bold')
        ax.set_xticks(np.arange(len(key_years)))
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=8)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f'${v/1e6:.1f}M'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw, dark=True)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.2))

    y30 = bmap.get(30, {})
    y40 = bmap.get(40, {})
    kpi_card(s, 10.2, 1.3, 2.8, 1.15, "Y30 退保", fmt_usd(y30.get('total_surrender_value',0)), f"{y30.get('total_surrender_value',0)/total_prem:.1f}x保费", "📈")
    kpi_card(s, 10.2, 2.55, 2.8, 1.15, "Y40 退保", fmt_usd(y40.get('total_surrender_value',0)), f"{y40.get('total_surrender_value',0)/total_prem:.1f}x保费", "🚀")
    kpi_card(s, 10.2, 3.8, 2.8, 1.15, "Y30-Y40增长", f"{y40.get('total_surrender_value',0)/max(y30.get('total_surrender_value',1),1)-1:.0%}", "10年增长率", "📊")
    kpi_card(s, 10.2, 5.05, 2.8, 1.15, "保证现金", fmt_usd(y40.get('guaranteed_cash_value',0)), "Y40 GCV", "🛡️")

    add_page_number(s, 11, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第12页：计划总结
# ═══════════════════════════════════════════════════════════════════════════════
def page_summary(prs, bmap, total_prem, pol, product_name, company_name, withdrawal_annual):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['gold'])

    add_text(s, "📋 计划书总结 — 匠心传承核心要点", 0.5, 0.2, 12, 0.55,
             font_size=26, color=T24['white'], bold=True)
    add_text(s, "一张保单，三代受益 — 5年播种，终身收获", 0.5, 0.8, 12, 0.38,
             font_size=12, color=T24['text_gray'])

    summary_items = [
        ("🏦 产品", f"{product_name}", T24['teal']),
        ("💵 年缴", f"${pol['annual_premium']:,} × 5年 = ${total_prem:,}", T24['warning']),
        ("📅 提取", f"第7年起每年 ${withdrawal_annual:,} 终身", T24['positive']),
        ("📈 回本", f"Y7退保 = $514,498（1.0x保费）", T24['gold']),
        ("🚀 Y20", f"退保${bmap[20]['total_surrender_value']/1e6:.2f}M = {bmap[20]['total_surrender_value']/total_prem:.1f}x保费", T24['gold']),
        ("🌱 Y30", f"退保${bmap[30]['total_surrender_value']/1e6:.2f}M = {bmap[30]['total_surrender_value']/total_prem:.1f}x保费", T24['positive']),
        ("🛡️ 保障", "终身壂障 + 代际传承", T24['teal']),
        ("✅ 公司", f"{company_name} — 香港最大壂障公司之一", T24['deep_blue']),
    ]
    for i, (label, val, color) in enumerate(summary_items):
        add_rect(s, 0.5, 1.35+i*0.68, 7.8, 0.6, fill_color=T24['card_bg'])
        add_accent_bar(s, 0.5, 1.35+i*0.68, 0.06, 0.6, color)
        add_text(s, label, 0.65, 1.4+i*0.68, 1.6, 0.5, font_size=11, color=color, bold=True)
        add_text(s, val, 2.25, 1.4+i*0.68, 5.9, 0.5, font_size=11, color=T24['white'])

    add_rect(s, 8.6, 1.35, 4.4, 5.55, fill_color=T24['card_bg'])
    add_text(s, "📊 核心倍数一览", 8.75, 1.45, 4.1, 0.45, font_size=13, color=T24['gold'], bold=True)
    headers = ['年度', '倍数', '退保价值']
    for j, h in enumerate(headers):
        add_text(s, h, 8.75+j*1.4, 1.95, 1.4, 0.35, font_size=10, color=T24['text_gray'], bold=True)
    milestones = [5, 7, 10, 15, 20, 25, 30, 40]
    for i, y in enumerate(milestones):
        r = bmap.get(y, {})
        mult = r.get('total_surrender_value', 0) / total_prem
        sv = r.get('total_surrender_value', 0)
        if i % 2 == 0:
            add_rect(s, 8.6, 2.35+i*0.48, 4.4, 0.45, fill_color=T24['divider'], line_color=None)
        add_text(s, f'Y{y}', 8.75, 2.35+i*0.48, 1.4, 0.4, font_size=11, color=T24['white'])
        add_text(s, f'{mult:.1f}x', 10.15, 2.35+i*0.48, 1.4, 0.4, font_size=11, color=T24['gold'], bold=True)
        add_text(s, fmt_usd(sv), 11.55, 2.35+i*0.48, 1.4, 0.4, font_size=10, color=T24['white'])

    add_page_number(s, 12, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 第13页：结束页
# ═══════════════════════════════════════════════════════════════════════════════
def page_closing(prs, product_name, company_name, insured_name, total_prem, withdrawal_annual):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, T24['deep_blue'])
    add_accent_bar(s, 0, 0, T24['W'], T24['TOP_BAR_H'], T24['gold'])
    add_accent_bar(s, 0, 0.4, T24['DECO_LEFT_W'], 4.5, T24['teal'])

    add_text(s, product_name, 1.5, 1.0, 10, 1.0, font_size=52, color=T24['white'], bold=True)
    add_text(s, f"{company_name} · 美元保单", 1.5, 2.1, 10, 0.65,
             font_size=24, color=T24['teal'], bold=True)
    add_accent_bar(s, 1.5, 2.85, 5.0, 0.04, T24['teal'])

    add_text(s, f"5年缴付 ${total_prem:,} · 第7年起领 ${withdrawal_annual:,}/年 · 终身现金流 · 代际传承",
             1.5, 3.1, 11, 0.5, font_size=16, color=T24['text_gray'])

    kpi_card(s, 1.5, 4.0, 3.2, 1.5, "5年缴付", f"${total_prem:,}", "稳健起步", "")
    kpi_card(s, 5.0, 4.0, 3.2, 1.5, "7年起领", f"${withdrawal_annual:,}/年", "终身现金流", "")
    kpi_card(s, 8.5, 4.0, 3.2, 1.5, "Y20退保", "$1.37M", "2.7x保费", "")

    add_text(s, "此计划书数据来源于PDF文件，仅供说明用途。实际数字以保险公司官方计划书为准。",
             1.5, 6.0, 10.5, 0.4, font_size=10, color=T24['text_gray'])
    add_text(s, f"受保人: {insured_name}  |  产品: {product_name}  |  {company_name}  |  2026-05-25",
             1.5, 6.5, 10.5, 0.4, font_size=10, color=T24['text_gray'])

    add_page_number(s, 13, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# 主程序 — 匠心传承储蓄险（测试模板24设计语言）
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    # 加载数据
    with open('/Users/soldier/free-code/packages/insurance-ppt/scripts/jxcc_data.json') as f:
        d = json.load(f)

    ext = d['extractions'][0]['data']
    benefit = ext['benefit_illustration']
    withdrawal = ext['withdrawal_illustration']
    pol = ext['policy']
    total_prem = pol['total_premium_with_levy']

    bmap = {r['policy_year']: r for r in benefit}

    # 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(T24['W'])
    prs.slide_height = Inches(T24['H'])

    # 产品配置（可替换为任意产品）
    product_name = "匠心传承储蓄寿险计划2（尊尚版）"
    company_name = "友邦保险 AIA"
    insured_name = "VIP先生"
    withdrawal_annual = 35_000
    withdrawal_start = 7
    highlights = [
        "✅ 5年缴付，7年起每年提取",
        "✅ 终身现金流，资金灵活",
        "✅ 复归红利 + 终期分红双引擎",
        "✅ 第20年退保达保费2.7倍",
        "✅ 免税传承，指定受益人",
    ]

    # 第1页：封面
    page_cover(prs, product_name, company_name, insured_name,
               pol['annual_premium'], total_prem,
               withdrawal_annual, withdrawal_start,
               highlights)

    # 第2页：计划概览
    params = [
        ("产品名称", "匠心传承储蓄寿险计划2（尊尚版）"),
        ("保险公司", "友邦保险 AIA"),
        ("受保人", "VIP先生"),
        ("投保单位", "699,301"),
        ("缴费模式", "年缴"),
        ("年缴保费", f"${pol['annual_premium']:,} USD"),
        ("缴费年期", pol['premium_payment_period']),
        ("保障年期", pol['coverage_period']),
        ("总缴保费", f"${total_prem:,} USD"),
    ]
    page_overview(prs, params, company_name)

    # 第3页：增长故事
    page_growth_story(prs, bmap, total_prem, [1,3,5,7,10,15,20,25,30], product_name)

    # 第4页：不提取每5年明细表
    page_no_withdrawal_table(prs, bmap, total_prem, [1,5,10,15,20,25,30,35,40])

    # 第5页：终身现金流
    page_cashflow_timeline(prs, benefit, withdrawal, list(range(1,16)))

    # 第6页：提领vs不提领
    page_withdrawal_comparison(prs, benefit, withdrawal, list(range(1,16)), total_prem)

    # 第7页：回本分析
    page_breakeven(prs, benefit, total_prem, list(range(1,16)))

    # 第8页：提取每5年明细表
    page_withdrawal_table(prs, benefit, withdrawal, total_prem, [1,5,10,15,20,25,30])

    # 第9页：年度收益明细
    page_yearly_detail(prs, benefit, total_prem)

    # 第10页：退保vs身故
    page_surrender_vs_death(prs, bmap, [10,20,30,40], total_prem)

    # 第11页：长期增长
    page_long_term(prs, bmap, total_prem, list(range(20,41)))

    # 第12页：总结
    page_summary(prs, bmap, total_prem, pol, product_name, company_name, withdrawal_annual)

    # 第13页：结束
    page_closing(prs, product_name, company_name, insured_name, total_prem, withdrawal_annual)

    out = '/tmp/insurance_plan_template24.pptx'
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")
    return out

if __name__ == '__main__':
    main()