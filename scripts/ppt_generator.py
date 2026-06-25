#!/usr/bin/env python3
"""
Insurance Proposal PPT Generator — Gamma/Google Stitch 风格 (WPS 兼容版)
======================================================================
设计原则:
  • 大字体标题层级 (标题 44-52pt, 副标题 24-28pt, 正文 14-16pt)
  • 全屏深色背景 + 渐变色块
  • 每页聚焦一个信息点，留白充足
  • KPI 卡片突出核心数字（大字号 + 强调色）
  • 折线图/区域图替代柱状图（长期增长叙事）
  • 叙事文案框（Sales Narrative Box）
  • Apple 风格圆角 (corner_radius: 16-20)
  • WPS 兼容：使用 msos + 避免高级 API

Usage: python3 ppt_generator.py --data '{"extractions":[...]}' --style modern --output /path/out.pptx
"""

import json
import re
import sys
import os
import argparse
import warnings
import io
import math
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
warnings.filterwarnings("ignore", message=".*Glyph.*missing from font.*")
import matplotlib.ticker as mticker
import matplotlib.font_manager as fm
import matplotlib.patches as mpatches
import numpy as np
import unicodedata

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from ppt_styles import get_style
from wiki_knowledge import get_wiki

def normalize_slide_text(text):
    """移除emoji与不可见控制字符，降低机器感并提升跨平台渲染稳定性。"""
    if text is None:
        return ""
    s = str(text)
    out = []
    for ch in s:
        cat = unicodedata.category(ch)
        # So(Other Symbol) 多数是 emoji / 图形符号，过滤掉
        if cat == "So":
            continue
        if cat.startswith("C") and ch not in ("\n", "\t"):
            continue
        out.append(ch)
    return "".join(out)

# ─── 中文字体配置 ──────────────────────────────────────────
_CHINA_FONT = None
for fname in [
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    "C:/Windows/Fonts/msyh.ttc",   # 微软雅黑 (WPS)
    "C:/Windows/Fonts/simhei.ttf", # 黑体
    "/usr/share/fonts/truetype/wqy-microhei/wqy-microhei.ttc",
]:
    if os.path.exists(fname):
        _CHINA_FONT = fname
        break
if _CHINA_FONT:
    font_prop = fm.FontProperties(fname=_CHINA_FONT)
    plt.rcParams["font.family"] = font_prop.get_name()
    fm.fontManager.addfont(_CHINA_FONT)
    plt.rcParams["font.sans-serif"] = [font_prop.get_name()]
else:
    plt.rcParams["font.sans-serif"] = ["Arial"]
plt.rcParams["axes.unicode_minus"] = False

# ─── 颜色工具 ──────────────────────────────────────────

def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def rgb_to_hex(r, g, b):
    return f"#{r:02X}{g:02X}{b:02X}"

def darken(hex_color, factor=0.8):
    r, g, b = hex_to_rgb(hex_color)
    return rgb_to_hex(int(r * factor), int(g * factor), int(b * factor))

def alphaBlend(c1_hex, c2_hex, alpha=0.5):
    """alpha 混合两个颜色"""
    r1, g1, b1 = hex_to_rgb(c1_hex)
    r2, g2, b2 = hex_to_rgb(c2_hex)
    r = int(r1 * (1 - alpha) + r2 * alpha)
    g = int(g1 * (1 - alpha) + g2 * alpha)
    b = int(b1 * (1 - alpha) + b2 * alpha)
    return rgb_to_hex(r, g, b)

# ─── 基础绘图工具 ──────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None,
             radius=None, gradient_start=None, gradient_end=None):
    """绘制矩形，支持渐变和圆角"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    elif gradient_start and gradient_end:
        shape.fill.patterned()
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = RGBColor(*hex_to_rgb(line_color))
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = min(0.25, radius / 200)
    return shape


def add_text(slide, text, left, top, width, height, font_size=14, color="#FFFFFF",
             bold=False, alignment=PP_ALIGN.LEFT, font_name=None, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    """文本框 — 兼容 WPS"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = normalize_slide_text(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    tf.vertical_anchor = anchor
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, lines, left, top, width, height, font_size=12,
                 color="#FFFFFF", line_spacing=1.4, bold_first=False):
    """多行文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = normalize_slide_text(line)
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(font_size * 0.3)
    return txBox


def add_narrative_box(slide, left, top, width, height, style, text,
                     brand_color=None, show_accent=True):
    """
    叙事文案框 — 多段落版
    - 支持 \\n 分段
    - 首段：品牌色+加粗（产品定位/公司背书）
    - 正文：常规颜色
    - 可选左侧 accent 条
    - 修复：圆角背景在 accent bar 下方（不覆盖），避免字体重叠
    """
    c = style["colors"]
    lines = normalize_slide_text(text).split("\n")
    if not lines:
        return

    accent_color = brand_color or c.get("accent_gold") or c.get("accent") or "#C8963E"

    # ① 先画圆角背景（rectRadius 生效）
    add_rect(slide, left, top, width, height,
             fill_color=c["bg_dark"],
             radius=style["layout"]["corner_radius"])

    # ② 左侧 accent bar（在背景上方，覆盖左侧区域）
    if show_accent:
        add_rect(slide, left, top, 0.08, height,
                 fill_color=accent_color, radius=0)

    # ③ 多段落渲染（宽度略收窄，避免碰到accent bar）
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.3), Inches(top + 0.12),
        Inches(width - 0.45), Inches(height - 0.22))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = style["fonts"]["body"]
        p.font.size = Pt(12)
        p.space_after = Pt(5)
        if i == 0:
            # 首段：品牌色+加粗
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_to_rgb(accent_color))
        else:
            p.font.bold = False
            p.font.color.rgb = RGBColor(*hex_to_rgb(c.get("text_body", "#E8E8E8")))


def add_kpi_card(slide, x, y, w, h, label, value, style, icon=None):
    """Gamma 风格 KPI 卡片 — 大数字 + 小标签"""
    c = style["colors"]
    accent_color = c.get("accent") or c.get("accent_gold", "#C8963E")
    card = add_rect(slide, x, y, w, h,
                    fill_color=c["bg_card"],
                    radius=style["layout"]["card_radius"])
    # 顶部细线强调
    add_rect(slide, x, y, w, 0.05, fill_color=accent_color, radius=0)
    # 大数字
    add_text(slide, value, x + 0.15, y + 0.15, w - 0.3, h * 0.55,
             font_size=28, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    # 小标签
    label_text = f"{icon} {label}" if icon else label
    add_text(slide, label_text, x + 0.1, y + h * 0.65, w - 0.2, h * 0.3,
             font_size=10, color=c["text_gray"], alignment=PP_ALIGN.CENTER)


def add_section_header(slide, text, left, top, width, style, icon=None):
    """大号分节标题"""
    c = style["colors"]
    label_text = f"{icon} {text}" if icon else text
    add_text(slide, label_text, left, top, width, 0.6,
             font_size=26, color=c["bg_dark"], bold=True)


def make_gradient_fig(figsize, draw_fn):
    """创建 matplotlib 图表（返回 BytesIO）"""
    fig, ax = plt.subplots(figsize=figsize, facecolor="#FFFFFF")
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#FAFAFA")
    draw_fn(ax)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def fmt(n):
    if n >= 1_000_000:
        return f"${n/1_000_000:.2f}M"
    if n >= 1_000:
        return f"${n/1_000:.1f}K"
    return f"${n:.0f}"


# ─── Gamma 风格配色 ──────────────────────────────────────────

GAMMA_COLORS = {
    "bg_dark":      "#0D1B2A",
    "bg_card":      "#14273E",
    "bg_light":     "#F8F9FA",
    "accent_gold":  "#C8963E",
    "accent_teal":  "#00D4AA",
    "accent_blue":  "#4FC3F7",
    "text_white":   "#FFFFFF",
    "text_body":    "#E8E8E8",
    "text_gray":    "#8899A6",
    "text_dark":    "#1A1A2E",
    "success":      "#34A853",
    "warning":      "#FBBC04",
    "error":        "#EA4335",
}


def gamma_style():
    """Gamma/Google Stitch 风格配置"""
    return {
        "colors": GAMMA_COLORS,
        "fonts": {
            "title": "Arial",
            "body": "Microsoft YaHei",
            "size_title": 48,
            "size_heading": 28,
            "size_subheading": 18,
            "size_body": 14,
            "size_small": 11,
            "size_kpi": 32,
        },
        "layout": {
            "corner_radius": 18,
            "card_radius": 16,
            "header_bar": True,
        }
    }


# ─── 辅助函数 ──────────────────────────────────────────

def findBreakeven(ext):
    """查找回本年份"""
    d = ext.get("data", {})
    if ext.get("plan_type") == "savings":
        years = d.get("benefit_illustration", [])
        for r in years:
            if r.get("total_surrender_value", 0) >= r.get("total_premium_paid", 0):
                return r.get("policy_year")
    elif ext.get("plan_type") in ("ci", "iul"):
        years = d.get("benefit_illustration", []) or d.get("yearly_data", [])
        for r in years:
            v = r.get("surrender_value_total", 0) or r.get("cash_value_guaranteed", 0)
            if v and v > 0:
                return r.get("policy_year") or r.get("year")
    return None


def findMultiple(ext, year):
    """查找指定年份倍数"""
    d = ext.get("data", {})
    if ext.get("plan_type") == "savings":
        years = d.get("benefit_illustration", [])
        if not years:
            return None
        try:
            row = next(r for r in years if r.get("policy_year") == year)
        except StopIteration:
            row = None
        if not row:
            return None
        total_prem = max((r.get("total_premium_paid", 0) for r in years), default=1)
        return row.get("total_surrender_value", 0) / total_prem
    return None


# ─── 幻灯片构建 ──────────────────────────────────────────

def build_cover(prs, data, style):
    """Gamma 风格封面 — 全屏深色 + 大字标题 + KPI 卡片"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = style["colors"]
    add_bg(s, c["bg_dark"])

    # 右上装饰圆形
    circle1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(*hex_to_rgb(c["accent_gold"]))
    circle1.line.fill.background()

    circle2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(10.5), Inches(-0.5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(*hex_to_rgb(c["accent_teal"]))
    circle2.line.fill.background()

    # 客户名称（大字）
    customer = data.get("customer_name", "尊貴客戶")
    add_text(s, customer, 0.6, 0.8, 12, 1.0,
             font_size=style["fonts"]["size_title"],
             color=c["text_white"], bold=True)

    # 副标题
    title = data.get("title", "家庭保障與財富方案")
    add_text(s, title, 0.6, 2.0, 10, 0.6,
             font_size=style["fonts"]["size_subheading"],
             color=c["accent_gold"])

    # 金色装饰线
    add_rect(s, 0.6, 2.7, 2.5, 0.06, fill_color=c["accent_gold"])

    # 产品类型标签
    types = data.get("plan_types", [])
    badges = "   ".join(
        t.replace("savings", "💰 储蓄").replace("ci", "🛡️ 重疾").replace("iul", "📈 IUL")
        for t in types
    )
    add_text(s, badges, 0.6, 2.9, 10, 0.3,
             font_size=12, color=c["text_gray"])

    # KPI 卡片（核心数字）
    exts = data.get("extractions", [])
    metrics = []
    for ext in exts[:3]:
        d = ext.get("data", {})
        pol = d.get("policy", {})
        if ext.get("plan_type") == "savings":
            prem = pol.get("annual_premium", 0)
            years = d.get("benefit_illustration", [])
            tp = max((r.get("total_premium_paid", 0) for r in years), default=0)
            metrics.append(("年缴保费", fmt(prem), "💰"))
            metrics.append(("总投入", fmt(tp), "📊"))
        elif ext.get("plan_type") == "ci":
            metrics.append(("危疾保额", fmt(pol.get("sum_insured", 0)), "🛡️"))
            daily = (pol.get("annual_premium", 0) or 0) / 365
            metrics.append(("每天成本", f"${daily:.1f}", "⏱️"))
        elif ext.get("plan_type") == "iul":
            si = pol.get("sum_insured", 0)
            ip = pol.get("initial_premium", 0)
            lev = f"{(si/ip):.1f}x" if ip > 0 else "—"
            metrics.append(("身故保障", fmt(si), "📈"))
            metrics.append(("杠杆比例", lev, "🎯"))

    card_y = 3.6
    card_h = 1.4
    for i, (lab, val, icon) in enumerate(metrics[:4]):
        bx = 0.6 + i * 3.1
        add_kpi_card(s, bx, card_y, 2.8, card_h, lab, val, style, icon)

    # 公司标签（wiki知识增强）
    wk = get_wiki()
    company_hint = re.sub(r'[\W]+', '', str(exts[0].get("data", {}).get("product_name", "aia")))[:8]
    wiki_overview = wk.get_company_overview(company_hint) or wk.get_company_overview("aia")
    brand_label = f"📋 {wiki_overview['name']}" if wiki_overview else ""
    if brand_label:
        add_text(s, brand_label, 0.6, 2.85, 10, 0.3,
                 font_size=11, color=c["accent_gold"])
    
    # 日期
    add_text(s, f"生成日期: {data.get('date', '')}", 0.6, 5.3, 5, 0.3,
             font_size=10, color=c["text_gray"])


def build_overview(prs, data, style):
    """三层架构概览 — Gamma 风格卡片"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = style["colors"]
    add_bg(s, c["bg_light"])

    # 顶部深色条
    add_rect(s, 0, 0, 13.33, 0.12, fill_color=c["bg_dark"])
    add_rect(s, 0, 0.12, 13.33, 1.0, fill_color=c["bg_dark"])

    add_text(s, "🏛️ 方案总览", 0.5, 0.25, 8, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)

    add_text(s, "从健康保障到财富累积，再到资产传承 — 三层架构层层递进",
             0.5, 1.3, 12, 0.4,
             font_size=style["fonts"]["size_subheading"],
             color=c["text_dark"])

    exts = data.get("extractions", [])

    layers = [
        ("ci",      "🛡️", "风险防护层", "危疾保障，抵御健康风险",          "#4FC3F7"),
        ("savings", "💰", "财富累积层", "储蓄增值，长期复利增长",          c["accent_gold"]),
        ("iul",     "📈", "传承规划层", "高杠杆传承，指数增长",            c["accent_teal"]),
    ]

    card_positions = [
        (0.5, 2.0, 3.8, 4.0),
        (4.7, 2.0, 3.8, 4.0),
        (8.9, 2.0, 3.8, 4.0),
    ]

    for idx, (pt, icon, name, desc, accent) in enumerate(layers):
        ext = next((e for e in exts if e.get("plan_type") == pt), None)
        x, y, w, h = card_positions[idx]

        add_rect(s, x, y, w, h, fill_color=c["bg_dark"],
                 radius=style["layout"]["card_radius"])
        add_rect(s, x, y, w, 0.08, fill_color=accent)
        add_text(s, icon, x + 0.2, y + 0.2, 1.2, 0.8, font_size=36, color=accent)
        add_text(s, name, x + 0.2, y + 1.0, w - 0.4, 0.5,
                 font_size=18, color=c["text_white"], bold=True)
        add_text(s, desc, x + 0.2, y + 1.5, w - 0.4, 0.8,
                 font_size=12, color=c["text_gray"], line_spacing=1.4)

        if ext:
            pn = ext.get("data", {}).get("product_name", "—")
            add_text(s, f"➜ {pn}", x + 0.2, y + 2.4, w - 0.4, 0.4,
                     font_size=13, color=accent)
            pol = ext.get("data", {}).get("policy", {})
            if pt == "savings":
                prem = pol.get("annual_premium", 0)
                add_text(s, f"年缴 {fmt(prem)}", x + 0.2, y + 2.9, w - 0.4, 0.4,
                         font_size=14, color=c["text_white"], bold=True)
            elif pt == "ci":
                si = pol.get("sum_insured", 0)
                add_text(s, f"保额 {fmt(si)}", x + 0.2, y + 2.9, w - 0.4, 0.4,
                         font_size=14, color=c["text_white"], bold=True)
            elif pt == "iul":
                si = pol.get("sum_insured", 0)
                add_text(s, f"身故保障 {fmt(si)}", x + 0.2, y + 2.9, w - 0.4, 0.4,
                         font_size=14, color=c["text_white"], bold=True)


def build_savings_growth_story(prs, data, style):
    """储蓄险增长故事页 — 折线图 + 叙事框"""
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "savings"), None)
    if not ext:
        return
    d = ext.get("data", {})
    years_data = d.get("benefit_illustration", [])
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "💰 账户价值增长分析", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "时间是最好的朋友，复利是最大的杠杆",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    pn = d.get("product_name", "")
    add_text(s, pn, 10.5, 0.3, 2.5, 0.4,
             font_size=12, color=c["accent_gold"],
             alignment=PP_ALIGN.RIGHT)

    # 折线图数据（只用有数据的年份，缺失年份跳过）
    ym = {r.get("policy_year"): r for r in years_data}
    present_years = sorted(ym.keys())
    key_years = [y for y in [1, 3, 5, 7, 10, 15, 20, 25, 30] if ym.get(y)]

    surrender_vals  = [ym.get(y, {}).get("total_surrender_value", 0) for y in key_years]
    guaranteed_vals = [ym.get(y, {}).get("guaranteed_cash_value", 0) for y in key_years]
    non_guaranteed  = [
        (ym.get(y, {}).get("reversionary_bonus", 0) + ym.get(y, {}).get("terminal_dividend", 0))
        for y in key_years
    ]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.plot(x, guaranteed_vals, color=c["accent_blue"], linewidth=2.5,
                marker='o', markersize=6, label="保证现金价值", zorder=3)
        ax.fill_between(x, guaranteed_vals, non_guaranteed,
                       where=[a > b for a, b in zip(non_guaranteed, guaranteed_vals)],
                       color=c["accent_gold"], alpha=0.3, label="非保证终期分红", zorder=2)
        ax.plot(x, surrender_vals, color=c["accent_gold"], linewidth=2.5,
                marker='s', markersize=6, label="退保总额", zorder=3)

        total_prem = years_data[0].get("total_premium_paid", 0) if years_data else 0
        for i, (yv, sv) in enumerate(zip(key_years, surrender_vals)):
            if sv >= total_prem and i > 0:
                ax.axvline(x=i, color=c["accent_teal"], linestyle="--",
                           linewidth=1.5, alpha=0.7)
                ax.annotate(f"回本 (Y{yv})", xy=(i, sv),
                            xytext=(i + 0.5, sv * 1.05),
                            fontsize=8, color=c["accent_teal"], fontweight='bold')
                break

        ax.set_xticks(x)
        ax.set_xticklabels([f"Y{y}" for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f"${v/1_000_000:.1f}M" if v >= 1_000_000 else f"${v/1_000:.0f}K"))
        ax.legend(fontsize=9, loc="upper left", framealpha=0.8)
        # 浅色/深色风格自适应
        chart_bg = "#F8F9FA" if c["bg_dark"] == "#FFFFFF" else "#0D1B2A"
        spine_c  = "#666666" if chart_bg == "#F8F9FA" else "#8899A6"
        grid_c   = "#CCCCCC" if chart_bg == "#F8F9FA" else "#FFFFFF"
        ax.grid(axis='y', alpha=0.25, color=grid_c, linewidth=0.5)
        ax.set_facecolor(chart_bg)
        ax.spines['bottom'].set_color(spine_c)
        ax.spines['left'].set_color(spine_c)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(colors=spine_c)
        ax.set_ylabel("金额 (USD)", fontsize=9, color=spine_c)

    buf = make_gradient_fig((9, 4.2), draw)
    s.shapes.add_picture(buf, Inches(0.4), Inches(1.35), Inches(9.5), Inches(4.2))

    # 右侧关键数字
    total_prem = max((r.get("total_premium_paid", 0) for r in years_data), default=0)
    pol = d.get("policy", {})
    annual_prem = pol.get("annual_premium", 0)
    for i, y in enumerate([5, 10, 20, 30]):
        row = ym.get(y, {})
        if not row:
            continue
        sv = row.get("total_surrender_value", 0)
        mult = f"{(sv/total_prem):.1f}x" if total_prem > 0 else "—"
        add_kpi_card(s, 10.0, 1.4 + i * 1.0, 2.8, 0.85,
                     f"Y{y} 倍数", mult, style)

    # 底部叙事框（wiki知识增强）
    wk = get_wiki()
    pn_raw = d.get("product_name", "")
    # Try to extract company from product name
    company_hint = re.sub(r'[\W]+', '', pn_raw)[:8]
    wiki_overview = wk.get_company_overview(company_hint) or wk.get_company_overview("aia")
    brand = wiki_overview["brand_line"] if wiki_overview else "友邦保险"
    wiki_narrative = wk.generate_narrative(
        wiki_overview["slug"] if wiki_overview else "aia",
        "savings",
        {"annual_premium": annual_prem, "sum_insured": 0}
    ) if wiki_overview else ""
    y20_row = ym.get(20, {})
    y20_mult = f"{(y20_row.get('total_surrender_value', 0) / max(total_prem, 1)):.1f}x" if y20_row else "2-3x"
    base_narrative = (
        f"「{pn_raw}」— 为您实现教育金、退休规划和长期财富增值。\n"
        f"账户价值随时间持续增长。保证现金价值提供本金安全垫，"
        f"非保证终期分红带来额外增长潜力。第20年退保总额约为已缴保费的 {y20_mult}。\n"
        f"复利的力量：时间越长，账户增长越快 — 第20年的增长斜率远超前10年。"
    )
    add_narrative_box(s, 0.4, 5.7, 12.5, 1.35, style, base_narrative,
                     brand_color=c["accent_gold"])


def build_savings_withdrawal_comparison(prs, data, style):
    """
    储蓄险提领 vs 不提领对比页
    单图双线：左Y轴不提领(蓝金)，右Y轴提领后(青绿)，标注提领起始+累计金额
    所有年份从数据中自动取，绘制到Y30上限。
    """
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "savings"), None)
    if not ext:
        return
    d = ext.get("data", {})
    benefit = d.get("benefit_illustration", [])
    withdrawal = d.get("withdrawal_illustration", [])
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "💰 提领 vs 不提领 — 账户价值对比", 0.5, 0.2, 12, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "同样本金，不同策略，体验差异立现",
             0.5, 0.85, 12, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    pn = d.get("product_name", "")
    add_text(s, pn, 10.5, 0.3, 2.5, 0.4,
             font_size=12, color=c["accent_gold"],
             alignment=PP_ALIGN.RIGHT)

    # 建立年份并集，数据上限30年；无数据年份不绘制（避免$0污染）
    all_years = sorted(set(
        [r.get("policy_year") for r in benefit] +
        [r.get("policy_year") for r in withdrawal]
    ))
    key_years = [y for y in all_years if 1 <= y <= 30]

    bm = {r.get("policy_year"): r for r in benefit}
    wm = {r.get("policy_year"): r for r in withdrawal}

    # 跳过没有任何数据的年份（Y7在withdrawal里缺失→提领线不画Y7点，用虚线连接邻近点）
    no_withdrawal_vals = [(bm.get(y, {}).get("total_surrender_value") or 0) for y in key_years]
    withdrawal_vals = [(wm.get(y, {}).get("surrender_value_after") or 0) for y in key_years]

    # 判断提领从第几年开始：第一年有 withdrawal 数据说明从Y1就开始提
    # 否则找 withdrawal 里有数据的最小年份作为起始年
    if withdrawal and wm:
        withdrawal_start_year = min(wm.keys())
    else:
        withdrawal_start_year = None

    # 累计提领金额（从有数据的年份累加 withdrawal_amount）
    cumulative = 0
    cum_list = []
    for y in key_years:
        row = wm.get(y, {})
        amt = row.get("withdrawal_amount") or 0
        cum = row.get("cumulative_withdrawals") or 0
        cumulative = max(cumulative, cum)  # 用字段本身的值更准确
        cum_list.append(cumulative)

    def draw_combined(ax):
        x = np.arange(len(key_years))
        labels = [f"Y{y}" for y in key_years]

        # 浅色/深色风格自适应
        chart_bg = "#F8F9FA" if c["bg_dark"] == "#FFFFFF" else "#0D1B2A"
        spine_c  = "#666666" if chart_bg == "#F8F9FA" else "#8899A6"
        grid_c   = "#CCCCCC" if chart_bg == "#F8F9FA" else "#FFFFFF"

        # 区域填充：提领后 vs 不提领 之间的差异（金色面积）
        # 安全比较：None → 0
        ax.fill_between(x, withdrawal_vals, no_withdrawal_vals,
                        where=[(a or 0) > (b or 0) for a, b in zip(no_withdrawal_vals, withdrawal_vals)],
                        color=c["accent_gold"], alpha=0.2,
                        label="提领差额（累计损失）", zorder=1)

        # 提领后折线（青绿色）
        ax.plot(x, withdrawal_vals, color=c["accent_teal"], linewidth=2.5,
                marker='s', markersize=5, label="提领后退保总额", zorder=3)

        # 不提领折线（金色）
        ax.plot(x, no_withdrawal_vals, color=c["accent_gold"], linewidth=2.5,
                marker='o', markersize=5, label="不提领退保总额", zorder=3)

        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=8, rotation=45)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f"${v/1_000_000:.1f}M" if v >= 1_000_000 else f"${v/1_000:.0f}K"))
        ax.legend(fontsize=8, loc="upper left", framealpha=0.8)
        ax.grid(axis='y', alpha=0.25, color=grid_c, linewidth=0.5)
        ax.set_facecolor(chart_bg)
        ax.spines['bottom'].set_color(spine_c)
        ax.spines['left'].set_color(spine_c)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(colors=spine_c)
        ax.set_ylabel("金额 (USD)", fontsize=9, color=spine_c)

        # 标注提领起始年
        title_color = "#37474F" if chart_bg == "#F8F9FA" else c["text_gray"]
        if withdrawal_start_year and withdrawal_start_year in key_years:
            idx = key_years.index(withdrawal_start_year)
            ax.axvline(x=idx, color=c["accent_teal"], linestyle="--", linewidth=1.5, alpha=0.7)
            ax.text(idx, ax.get_ylim()[1] * 0.85,
                    f"提领起始\n  Y{withdrawal_start_year}",
                    fontsize=8, color=c["accent_teal"],
                    ha='center', va='top',
                    bbox=dict(boxstyle='round,pad=0.3', facecolor=chart_bg, alpha=0.85))

        ax.set_title("双线对比 — 提领策略 vs 持续累积", fontsize=11,
                     color=title_color, pad=8)

    buf = make_gradient_fig((8.5, 4.5), draw_combined)
    s.shapes.add_picture(buf, Inches(0.5), Inches(1.35), Inches(12.5), Inches(4.0))

    # 底部 KPI：Y15 / Y20 对比
    y15_surr = bm.get(15, {}).get("total_surrender_value") or 0
    y15_with = wm.get(15, {}).get("surrender_value_after") or 0
    y15_cum = wm.get(15, {}).get("cumulative_withdrawals") or 0
    y20_surr = bm.get(20, {}).get("total_surrender_value") or 0
    y20_with = wm.get(20, {}).get("surrender_value_after") or 0
    y20_cum = wm.get(20, {}).get("cumulative_withdrawals") or 0

    add_kpi_card(s, 0.5, 5.55, 3.8, 1.0,
                 f"Y15 不提领  {fmt(y15_surr)}", "", style, "💰")
    add_kpi_card(s, 4.5, 5.55, 3.8, 1.0,
                 f"Y15 提领后  {fmt(y15_with)}  (已提 {fmt(y15_cum)})", "", style, "✅")
    add_kpi_card(s, 8.5, 5.55, 3.8, 1.0,
                 f"Y15 差异  {fmt(y15_surr - y15_with)}", "", style, "📊")


def build_ci_story(prs, data, style):
    """重疾险故事页 — 每天成本 + 保障矩阵"""
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "ci"), None)
    if not ext:
        return
    d = ext.get("data", {})
    items = d.get("coverage_items", [])
    pol = d.get("policy", {})
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "🛡️ 危疾保障方案", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "每天一杯咖啡的价格，换全年无盲点的危疾保障",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    annual_prem = pol.get("annual_premium", 0)
    sum_insured = pol.get("sum_insured", 0)
    daily_cost = annual_prem / 365
    coverage_period = pol.get("coverage_period", "终身")

    add_text(s, f"${daily_cost:.1f}", 0.5, 1.5, 3.5, 1.2,
             font_size=60, color=c["accent_gold"], bold=True,
             alignment=PP_ALIGN.LEFT)
    add_text(s, "/天", 3.8, 1.9, 1.5, 0.6,
             font_size=24, color=c["text_gray"])
    add_text(s, f"= 全年 {fmt(sum_insured)} 全面危疾保障",
             0.5, 2.7, 6, 0.4,
             font_size=style["fonts"]["size_subheading"],
             color=c["text_white"])
    add_text(s, f"保障期限: {coverage_period}  |  缴费期: {pol.get('premium_payment_period', '')}",
             0.5, 3.2, 6, 0.3,
             font_size=12, color=c["text_gray"])

    # 右侧保障项目
    half = math.ceil(len(items) / 2)
    add_text(s, "保障范围", 7.5, 1.3, 5, 0.4,
             font_size=16, color=c["accent_teal"], bold=True)

    for col_idx, start in enumerate([0, half]):
        for i, item in enumerate(items[start:start + half]):
            bx = 7.5 + col_idx * 2.7
            by = 1.85 + i * 0.65
            if by > 5.5:
                continue
            add_rect(s, bx, by, 2.5, 0.55,
                     fill_color=c["bg_card"],
                     radius=style["layout"]["corner_radius"])
            add_text(s, item.get("name", ""), bx + 0.1, by + 0.05, 1.8, 0.25,
                     font_size=10, color=c["text_body"])
            amt = item.get("amount")
            if amt:
                add_text(s, fmt(amt), bx + 0.1, by + 0.28, 2.2, 0.25,
                         font_size=10, color=c["accent_gold"],
                         bold=True, alignment=PP_ALIGN.RIGHT)

    # 多次赔付
    mc = d.get("multi_claim", [])
    if mc:
        add_text(s, "🔄 多次赔付设计", 0.5, 3.7, 5, 0.35,
                 font_size=14, color=c["accent_teal"], bold=True)
        for i, m in enumerate(mc[:3]):
            add_text(s, f"• {m['condition']} — 可赔付 {m['claim_count']} 次",
                     0.5, 4.1 + i * 0.35, 5, 0.3,
                     font_size=12, color=c["text_body"])

    # 底部叙事框（wiki知识增强）
    wk = get_wiki()
    pn_raw_ci = d.get("product_name", "此危疾保障计划")
    wiki_overview = wk.get_company_overview("aia")
    wiki_narrative = wk.generate_narrative(
        wiki_overview["slug"] if wiki_overview else "aia",
        "ci",
        {"annual_premium": annual_prem, "sum_insured": sum_insured}
    ) if wiki_overview else ""
    base_narrative = (
        f"「{pn_raw_ci}」— 保障全面、赔付及时，是家庭风险管理的基石。\n"
        f"健康是最大的财富，而一份全面的危疾保障，是对家人最负责任的投资。"
        f"这份计划书覆盖 {len(items)} 个危疾类别，癌症/中风/心脏病均可多次赔付。"
    )
    add_narrative_box(s, 0.4, 5.7, 12.5, 1.3, style, base_narrative,
                     brand_color=c["accent_teal"])


def build_ci_family_security(prs, data, style):
    """
    【新增】家庭多方位保障页
    展示 CI 如何覆盖家庭成员的不同风险维度
    """
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "ci"), None)
    if not ext:
        return
    d = ext.get("data", {})
    pol = d.get("policy", {})
    items = d.get("coverage_items", [])
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "🛡️ 家庭多方位保障规划", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "危疾不是一个人的事 — 全家都需要保障网",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    sum_insured = pol.get("sum_insured", 0)
    annual_prem = pol.get("annual_premium", 0)
    daily_cost = annual_prem / 365
    coverage_period = pol.get("coverage_period", "终身")

    # 左列：保障维度雷达卡
    dimensions = [
        ("🫀 心脏保障",  sum_insured * 0.3,  c["accent_gold"]),
        ("🧠 中风保障",  sum_insured * 0.25, c["accent_teal"]),
        ("🦀 癌症保障",  sum_insured * 0.35, c["accent_blue"]),
        ("👶 子女加成",  sum_insured * 0.1,  "#FF7043"),
    ]
    y_pos = 1.4
    for label, amt, color in dimensions:
        add_rect(s, 0.5, y_pos, 6.0, 0.7,
                 fill_color=c["bg_card"], radius=style["layout"]["corner_radius"])
        add_rect(s, 0.5, y_pos, 0.08, 0.7, fill_color=color)
        add_text(s, label, 0.7, y_pos + 0.05, 2.5, 0.3,
                 font_size=12, color=c["text_white"], bold=True)
        add_text(s, f"{fmt(amt)}", 3.2, y_pos + 0.05, 3.0, 0.3,
                 font_size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
        y_pos += 0.85

    # 右列：保障项目详情
    add_text(s, "保障范围一览", 7.0, 1.3, 5.5, 0.4,
             font_size=14, color=c["accent_gold"], bold=True)

    row_y = 1.85
    for item in items[:8]:
        name = item.get("name", "")
        amt = item.get("amount")
        add_rect(s, 7.0, row_y, 5.8, 0.55,
                 fill_color=c["bg_card"], radius=8)
        add_text(s, name, 7.15, row_y + 0.08, 3.5, 0.25,
                 font_size=11, color=c["text_body"])
        if amt:
            add_text(s, fmt(amt), 10.5, row_y + 0.08, 2.2, 0.25,
                     font_size=11, color=c["accent_teal"],
                     bold=True, alignment=PP_ALIGN.RIGHT)
        row_y += 0.65

    # 底部 KPI
    add_kpi_card(s, 0.5, 5.3, 3.8, 1.1, "总保额", fmt(sum_insured), style, "🛡️")
    add_kpi_card(s, 4.5, 5.3, 3.8, 1.1, "每天成本", f"${daily_cost:.1f}", style, "⏱️")
    add_kpi_card(s, 8.5, 5.3, 3.8, 1.1, "保障期限", coverage_period, style, "📅")


def build_iul_story(prs, data, style):
    """IUL 故事页 — 保证 vs 非保证双账户"""
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "iul"), None)
    if not ext:
        return
    d = ext.get("data", {})
    yearly = d.get("benefit_illustration", [])
    pol = d.get("policy", {})
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "📈 指数型万用寿险", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "保底 0% 不亏损，分享美国 S&P 500 经济增长红利",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    # 双账户架构
    add_rect(s, 0.5, 1.4, 5.8, 1.1,
             fill_color=c["bg_card"],
             radius=style["layout"]["corner_radius"])
    add_text(s, "🏦 固定账户（保证）", 0.7, 1.45, 2.5, 0.3,
             font_size=12, color=c["accent_blue"], bold=True)
    add_text(s, f"保证利率 {pol.get('fixed_account_rate', '—')}% | 保底 0%",
             0.7, 1.75, 5.5, 0.3, font_size=11, color=c["text_body"])
    add_text(s, "📊 指数账户（增长）", 0.7, 2.05, 2.5, 0.3,
             font_size=12, color=c["accent_teal"], bold=True)
    add_text(s, f"假设利率 {pol.get('index_account_rate', '—')}% | 对标 S&P 500",
             0.7, 2.35, 5.5, 0.3, font_size=11, color=c["text_body"])

    # 折线图
    ym = {r.get("policy_year"): r for r in yearly}
    key_years = [5, 10, 15, 20, 25, 30]
    gvals  = [ym.get(y, {}).get("guaranteed_cash_value", 0) for y in key_years]
    ngvals = [ym.get(y, {}).get("non_guaranteed_cash_value", 0) for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.plot(x, gvals, color=c["accent_blue"], linewidth=2.5,
                marker='o', markersize=7, label="保证现金价值", zorder=3)
        ax.fill_between(x, gvals, ngvals,
                       where=[a > b for a, b in zip(ngvals, gvals)],
                       color=c["accent_teal"], alpha=0.25,
                       label="指数账户增长", zorder=2)
        ax.plot(x, ngvals, color=c["accent_teal"], linewidth=2.5,
                marker='s', markersize=7, label="非保证账户价值", zorder=3)
        ax.set_xticks(x)
        ax.set_xticklabels([f"Y{y}" for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f"${v/1_000_000:.1f}M" if v >= 1_000_000 else f"${v/1_000:.0f}K"))
        ax.legend(fontsize=9, loc="upper left", framealpha=0.8)
        # 浅色/深色风格自适应
        chart_bg = "#F8F9FA" if c["bg_dark"] == "#FFFFFF" else "#0D1B2A"
        spine_c  = "#666666" if chart_bg == "#F8F9FA" else "#8899A6"
        grid_c   = "#CCCCCC" if chart_bg == "#F8F9FA" else "#FFFFFF"
        ax.grid(axis='y', alpha=0.25, color=grid_c, linewidth=0.5)
        ax.set_facecolor(chart_bg)
        ax.spines['bottom'].set_color(spine_c)
        ax.spines['left'].set_color(spine_c)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(colors=spine_c)
        ax.set_ylabel("金额 (USD)", fontsize=9, color=spine_c)

    buf = make_gradient_fig((9, 3.8), draw)
    s.shapes.add_picture(buf, Inches(0.4), Inches(2.7), Inches(9.5), Inches(3.8))

    si = pol.get("sum_insured", 0)
    ip = pol.get("initial_premium", 0)
    lev = f"{(si/ip):.1f}x" if ip > 0 else "—"
    add_kpi_card(s, 10.0, 1.5, 2.8, 1.2, "身故保障", fmt(si), style, "🛡️")
    add_kpi_card(s, 10.0, 2.85, 2.8, 1.2, "杠杆比例", lev, style, "🎯")

    # 底部叙事框（top=6.5 修复越界）
    add_narrative_box(s, 0.4, 6.5, 12.5, 0.8, style,
        f"攻守兼备：保证账户 0% 保底确保本金安全，指数账户让您参与美国经济增长。"
        f"身故保障 {fmt(si)}，是初始供款的 {lev} 倍，实现保障与投资的双重目标。")


def build_iul_succession(prs, data, style):
    """
    【新增】定向传承 IUL 页
    身故保障杠杆 + 指定受益人 + 税务优化叙事
    """
    exts = data.get("extractions", [])
    ext = next((e for e in exts if e.get("plan_type") == "iul"), None)
    if not ext:
        return
    d = ext.get("data", {})
    pol = d.get("policy", {})
    yearly = d.get("benefit_illustration", [])
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "🏛️ 定向传承规划", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "用高杠杆寿险，将财富无损传递给下一代",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    si = pol.get("sum_insured", 0)
    ip = pol.get("initial_premium", 0)
    lev = f"{(si/ip):.1f}x" if ip > 0 else "—"

    # 左：杠杆传承原理
    add_rect(s, 0.5, 1.4, 5.8, 3.2,
             fill_color=c["bg_card"], radius=style["layout"]["corner_radius"])
    add_text(s, "🎯 传承杠杆原理", 0.7, 1.5, 5.5, 0.4,
             font_size=16, color=c["accent_gold"], bold=True)

    # 投入 vs 传承
    add_text(s, "初始投入", 0.8, 2.0, 2.0, 0.4,
             font_size=12, color=c["text_gray"])
    add_text(s, fmt(ip), 0.8, 2.4, 2.5, 0.5,
             font_size=28, color=c["text_white"], bold=True)

    add_text(s, "→", 3.0, 2.35, 0.5, 0.5,
             font_size=28, color=c["accent_teal"], alignment=PP_ALIGN.CENTER)

    add_text(s, "定向传承", 3.6, 2.0, 2.5, 0.4,
             font_size=12, color=c["text_gray"])
    add_text(s, fmt(si), 3.6, 2.4, 2.5, 0.5,
             font_size=28, color=c["accent_gold"], bold=True)

    add_text(s, f"杠杆比例 {lev}x", 0.8, 3.1, 5.0, 0.4,
             font_size=14, color=c["accent_teal"], bold=True)

    # 传承优势列表
    advantages = [
        ("✅ 指定受益人", "财富直接传递，无遗产争议"),
        ("✅ 税务优化", "寿险理赔金一般免税"),
        ("✅ 私密性强", "无需公开遗嘱或 probated"),
        ("✅ 杠杆放大", "少量保费撬动大额传承"),
    ]
    ay = 3.6
    for title, desc in advantages:
        add_text(s, title, 0.8, ay, 2.5, 0.3,
                 font_size=11, color=c["text_white"], bold=True)
        add_text(s, desc, 3.3, ay, 2.8, 0.3,
                 font_size=11, color=c["text_gray"])
        ay += 0.38

    # 右：身故保障增长图（只用有数据的年份）
    ym = {r.get("policy_year"): r for r in yearly}
    key_years = [y for y in [5, 10, 15, 20, 25, 30] if ym.get(y)]
    dbvals = [ym.get(y, {}).get("guaranteed_death_benefit", 0) for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.plot(x, dbvals, color=c["accent_gold"], linewidth=2.5,
                marker='o', markersize=7, label="身故保障", zorder=3)
        ax.set_xticks(x)
        ax.set_xticklabels([f"Y{y}" for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f"${v/1_000_000:.1f}M" if v >= 1_000_000 else f"${v/1_000:.0f}K"))
        ax.legend(fontsize=9, loc="upper left", framealpha=0.8)
        # 浅色/深色风格自适应
        chart_bg = "#F8F9FA" if c["bg_dark"] == "#FFFFFF" else "#0D1B2A"
        spine_c  = "#666666" if chart_bg == "#F8F9FA" else "#8899A6"
        grid_c   = "#CCCCCC" if chart_bg == "#F8F9FA" else "#FFFFFF"
        ax.grid(axis='y', alpha=0.25, color=grid_c, linewidth=0.5)
        ax.set_facecolor(chart_bg)
        ax.spines['bottom'].set_color(spine_c)
        ax.spines['left'].set_color(spine_c)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(colors=spine_c)
        ax.set_ylabel("金额 (USD)", fontsize=9, color=spine_c)
        ax.set_title("身故保障增长趋势", fontsize=11, color=c["text_gray"], pad=6)

    buf = make_gradient_fig((6.2, 3.5), draw)
    s.shapes.add_picture(buf, Inches(6.7), Inches(1.3), Inches(6.0), Inches(3.5))

    # 底部叙事（wiki知识增强）
    wk = get_wiki()
    pn_raw_iul = d.get("product_name", "此指数型万用寿险计划")
    wiki_overview = wk.get_company_overview("yflife")
    wiki_narrative = wk.generate_narrative(
        wiki_overview["slug"] if wiki_overview else "yflife",
        "iul",
        {"initial_premium": ip, "sum_insured": si}
    ) if wiki_overview else ""
    base_narrative = (
        f"「{pn_raw_iul}」— 双账户结构兼顾保证与增长潜力，对标 S&P 500 指数，适合长期传承规划。\n"
        f"高杠杆寿险是财富传承的高效工具。投入 {fmt(ip)}，即可撬动 {fmt(si)} 的身故保障，"
        f"杠杆 {lev}x，让您的财富精准传递给下一代。"
    )
    add_narrative_box(s, 0.4, 5.0, 12.5, 1.3, style, base_narrative,
                     brand_color=c["accent_teal"])


def build_combined_plan_page(prs, data, style):
    """
    【新增】1+2+3 组合方案综合规划页
    三层架构协同：风险防护 + 财富累积 + 定向传承
    """
    exts = data.get("extractions", [])
    c = style["colors"]

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "🏛️ 组合方案：1+2+3 全面规划", 0.5, 0.2, 12, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "一张保单保健康，一张保单增值，一张保单传承 — 三层架构，守护家庭一生",
             0.5, 0.85, 12, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    # 三层架构图
    layers_cfg = [
        ("🛡️", "第一层：风险防护",    "危疾险", "ci",
         "覆盖家庭成员健康风险，癌症/中风/心脏病多重保障", "#4FC3F7"),
        ("💰", "第二层：财富累积",    "储蓄险", "savings",
         "长期复利增长，第6年回本，20年2.7倍收益", c["accent_gold"]),
        ("📈", "第三层：定向传承",    "万用寿险", "iul",
         "高杠杆身故保障，指数账户增长，精准传承", c["accent_teal"]),
    ]
    for idx, (icon, title, product, pt, desc, color) in enumerate(layers_cfg):
        ext = next((e for e in exts if e.get("plan_type") == pt), None)
        x = 0.5 + idx * 4.2
        y = 1.5

        # 连接箭头
        if idx > 0:
            arrow_x = x - 0.5
            add_text(s, "→", arrow_x, y + 0.8, 0.5, 0.5,
                     font_size=24, color=color, alignment=PP_ALIGN.CENTER)

        add_rect(s, x, y, 3.8, 3.8,
                 fill_color=c["bg_card"], radius=style["layout"]["card_radius"])
        add_rect(s, x, y, 3.8, 0.08, fill_color=color)

        add_text(s, icon, x + 0.2, y + 0.2, 1.0, 0.6,
                 font_size=32, color=color)
        add_text(s, title, x + 0.2, y + 0.85, 3.4, 0.4,
                 font_size=15, color=c["text_white"], bold=True)
        add_text(s, product, x + 0.2, y + 1.25, 3.4, 0.3,
                 font_size=12, color=color)

        if ext:
            pol = ext.get("data", {}).get("policy", {})
            d = ext.get("data", {})
            if pt == "ci":
                si = pol.get("sum_insured", 0)
                add_text(s, f"保额 {fmt(si)}", x + 0.2, y + 1.7, 3.4, 0.4,
                         font_size=20, color=c["text_white"], bold=True)
            elif pt == "savings":
                prem = pol.get("annual_premium", 0)
                add_text(s, f"年缴 {fmt(prem)}", x + 0.2, y + 1.7, 3.4, 0.4,
                         font_size=20, color=c["text_white"], bold=True)
            elif pt == "iul":
                si = pol.get("sum_insured", 0)
                ip = pol.get("initial_premium", 0)
                lev = f"{(si/ip):.1f}x" if ip > 0 else "—"
                add_text(s, f"身故 {fmt(si)} | {lev}x杠杆", x + 0.2, y + 1.7, 3.4, 0.4,
                         font_size=18, color=c["text_white"], bold=True)

        add_text(s, desc, x + 0.2, y + 2.3, 3.4, 1.0,
                 font_size=11, color=c["text_gray"], line_spacing=1.4)

    # 底部综合总结
    add_narrative_box(s, 0.4, 5.5, 12.5, 1.5, style,
        "三层架构层层递进：\n"
        "① 危疾险 — 守护健康，不让疾病摧毁家庭财务\n"
        "② 储蓄险 — 复利增值，让财富持续增长\n"
        "③ 万用寿险 — 高杠杆传承，将财富精准传递给下一代\n"
        "三张保单协同，构建完整的家庭保障与财富传承体系。")


def build_comparison_table(prs, data, style):
    """多产品对比表"""
    exts = data.get("extractions", [])
    if len(exts) < 2:
        return

    c = style["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    add_text(s, "📊 产品对比分析", 0.5, 0.2, 10, 0.6,
             font_size=style["fonts"]["size_heading"],
             color=c["text_white"], bold=True)
    add_text(s, "用数据说话，让选择更清晰",
             0.5, 0.85, 10, 0.35,
             font_size=style["fonts"]["size_body"],
             color=c["text_gray"])

    # 数据行
    type_map = {"savings": "储蓄险", "ci": "重疾险", "iul": "万用寿险"}

    metrics = [
        ("产品类型", lambda ext: type_map.get(ext.get("plan_type", ""), "")),
        ("年缴保费", lambda ext: fmt((ext.get("data", {}).get("policy", {}).get("annual_premium", 0)
                                    or ext.get("data", {}).get("policy", {}).get("initial_premium", 0)
                                    or 0))),
        ("回本年份", lambda ext: f"第{findBreakeven(ext)}年" if findBreakeven(ext) else "—"),
        ("20年倍数", lambda ext: f"{(findMultiple(ext, 20)):.1f}x" if findMultiple(ext, 20) else "—"),
    ]

    rows = 1 + len(metrics)
    cols = 1 + len(exts)
    table = s.shapes.add_table(
        rows, cols,
        Inches(0.4), Inches(1.3),
        Inches(12.5), Inches(0.6 * rows)
    ).table

    # 表头
    headers = ["指标"] + [
        ext.get("data", {}).get("product_name", f"产品{i+1}")
        for i, ext in enumerate(exts)
    ]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*hex_to_rgb(c["text_white"]))
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*hex_to_rgb(c["accent_gold"]))

    for i, (metric, fn) in enumerate(metrics):
        row = table.cell(i + 1, 0)
        row.text = metric
        for p in row.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(*hex_to_rgb(c["accent_teal"]))
            p.font.bold = True

        for j, ext in enumerate(exts):
            cell = table.cell(i + 1, j + 1)
            val = fn(ext)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(*hex_to_rgb(c["text_body"]))
                p.alignment = PP_ALIGN.CENTER

        if i % 2 == 0:
            for j in range(1, cols):
                table.cell(i + 1, j).fill.solid()
                table.cell(i + 1, j).fill.fore_color.rgb = RGBColor(0x14, 0x27, 0x3E)


def build_closing(prs, data, style):
    """Gamma 风格结束页"""
    c = style["colors"]
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, c["bg_dark"])

    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(-1), Inches(4), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*hex_to_rgb(c["accent_teal"]))
    circle.line.fill.background()

    add_text(s, "谢谢", 0.5, 2.0, 12, 1.2,
             font_size=56, color=c["text_white"], bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s, "感谢您的时间，期待为您提供专业服务",
             0.5, 3.3, 12, 0.5,
             font_size=style["fonts"]["size_subheading"],
             color=c["accent_gold"],
             alignment=PP_ALIGN.CENTER)
    add_rect(s, 5.5, 4.0, 2.5, 0.06, fill_color=c["accent_gold"])
    add_text(s,
             f"本文件仅供参考，不构成要约或建议。\n"
             f"非保证金额并非保证，实际可能高于或低于预期。\n"
             f"生成日期: {data.get('date', '')}",
             0.5, 4.3, 12, 1.2,
             font_size=10, color=c["text_gray"],
             alignment=PP_ALIGN.CENTER)


# ─── 主生成器 ──────────────────────────────────────────

def generate_ppt(data, style_name, output_path):
    style = gamma_style() if style_name in ("modern", "gamma") else get_style(style_name)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    types = list(set(e.get("plan_type") for e in data.get("extractions", [])))

    build_cover(prs, {**data, "plan_types": types, "date": data.get("date", "")}, style)
    build_overview(prs, data, style)

    if "savings" in types:
        build_savings_growth_story(prs, data, style)
        build_savings_withdrawal_comparison(prs, data, style)  # 新增：提领对比

    if "ci" in types:
        build_ci_story(prs, data, style)
        build_ci_family_security(prs, data, style)              # 新增：家庭保障

    if "iul" in types:
        build_iul_story(prs, data, style)
        build_iul_succession(prs, data, style)                  # 新增：定向传承

    if len(types) >= 2:
        build_comparison_table(prs, data, style)

    if len(types) >= 3:
        build_combined_plan_page(prs, data, style)              # 新增：1+2+3 组合

    build_closing(prs, data, style)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True, help="JSON string with extraction data")
    parser.add_argument("--style", default="modern", help="Style name (modern/gamma/professional/fresh/minimal/warm)")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    data = json.loads(args.data)
    result = generate_ppt(data, args.style, args.output)
    print(json.dumps({"status": "done", "path": result}))
