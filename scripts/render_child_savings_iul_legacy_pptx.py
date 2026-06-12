#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FONT_CN = "Microsoft YaHei"
BG = RGBColor(247, 244, 239)
PANEL = RGBColor(255, 255, 255)
PRIMARY = RGBColor(34, 34, 34)
MUTED = RGBColor(95, 95, 95)
LINE = RGBColor(224, 216, 202)
ACCENT = RGBColor(209, 151, 38)
ACCENT_LIGHT = RGBColor(245, 234, 204)
GREEN = RGBColor(36, 111, 86)
GREEN_LIGHT = RGBColor(187, 219, 206)
BLUE = RGBColor(61, 95, 148)
ROSE = RGBColor(180, 92, 92)


def money(value: float | int) -> str:
    return f"{round(float(value)):,}"


def multiple(value: float) -> str:
    return f"{value:.2f}x"


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, line_spacing=1.4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, left, top, width, height, lines, size=14.5, bullet=True, color=PRIMARY, line_spacing=1.45):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.text = f"• {line}" if bullet else line
        if p.runs:
            run = p.runs[0]
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=""):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.48), title, size=24, bold=True, line_spacing=1.15)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.78), Inches(10.0), Inches(0.28), subtitle, size=11.5, color=MUTED, line_spacing=1.15)


def add_card(slide, left, top, width, height, title, value, subtitle="", accent=False, value_size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT if accent else PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.22), title, size=11.8, bold=True, color=MUTED, line_spacing=1.15)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.38), width - Inches(0.28), Inches(0.42), value, size=value_size, bold=True, color=ACCENT if accent else PRIMARY, line_spacing=1.15)
    if subtitle:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.82), width - Inches(0.28), Inches(0.22), subtitle, size=10.2, color=MUTED, line_spacing=1.15)


def add_fact_card(slide, left, top, width, height, title, body, fill_color=PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.22), title, size=11.5, bold=True, color=MUTED, line_spacing=1.15)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.36), width - Inches(0.28), height - Inches(0.48), body, size=12.8, color=PRIMARY, line_spacing=1.42)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.62), Inches(7.0), Inches(11.8), Inches(0.2), text, size=10, color=MUTED, line_spacing=1.1)


def add_picture(slide, image_path: Path | None, left, top, width=None, height=None):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def chart(slide, left, top, width, height, categories, series_spec):
    data = CategoryChartData()
    data.categories = [str(x) for x in categories]
    for name, values in series_spec:
        data.add_series(name, values)
    c = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data).chart
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    c.legend.font.size = Pt(10)
    c.value_axis.has_major_gridlines = True
    c.value_axis.tick_labels.font.size = Pt(10)
    c.category_axis.tick_labels.font.size = Pt(10)
    c.value_axis.tick_labels.number_format = '$#,##0'
    palette = [ACCENT, GREEN, BLUE, ROSE]
    for idx, series in enumerate(c.series):
        series.format.line.width = Pt(2.2)
        series.format.line.color.rgb = palette[idx % len(palette)]
    return c


def timeline_marker(slide, left, title, lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.7), Inches(2.0), Inches(1.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.12), Inches(2.84), Inches(1.75), Inches(0.22), title, size=12.5, bold=True, color=ACCENT, line_spacing=1.15)
    add_paragraphs(slide, left + Inches(0.12), Inches(3.08), Inches(1.75), Inches(1.25), lines, size=10.8, bullet=False, line_spacing=1.3)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--savings", required=True)
    parser.add_argument("--iul", required=True)
    parser.add_argument("--savings-company", required=True)
    parser.add_argument("--iul-company", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    savings = json.loads(Path(args.savings).read_text(encoding="utf-8"))
    iul = json.loads(Path(args.iul).read_text(encoding="utf-8"))
    savings_company = json.loads(Path(args.savings_company).read_text(encoding="utf-8"))
    iul_company = json.loads(Path(args.iul_company).read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover = ROOT / "public/assets/library/themes/family/family-evening-01.jpg"
    father_child = ROOT / "public/assets/library/themes/family/father-child-01.jpg"
    child_growth = ROOT / "public/assets/library/themes/education/child-growth-01.jpg"
    graduation = ROOT / "public/assets/library/themes/education/graduation-01.jpg"
    city = ROOT / "public/assets/library/themes/business/cityline-01.jpg"
    savings_growth = ROOT / "public/assets/library/themes/savings/long-term-growth-01.jpg"
    savings_cash = ROOT / "public/assets/library/themes/savings/cashflow-future-01.jpg"

    sw = savings["withdrawalRows"]
    sb = savings["benefitRows"]
    ib = iul["benefitRows"]
    first_draw = next((row for row in sw if float(row["annualWithdrawal"]) > 0), sw[0])
    s18 = next((row for row in sw if int(row["policyYear"]) == 18), sw[-1])
    s21 = next((row for row in sw if int(row["policyYear"]) == 21), sw[-1])
    s25 = next((row for row in sw if int(row["policyYear"]) == 25), sw[-1])
    s20_nv = next(row for row in sb if int(row["policyYear"]) == 20)
    s30_nv = next(row for row in sb if int(row["policyYear"]) == 30)
    i20 = next(row for row in ib if int(row["policyYear"]) == 20)
    i30 = next(row for row in ib if int(row["policyYear"]) == 30)
    i100 = next((row for row in ib if int(row["age"]) == 100), ib[-1])
    total_iul_paid = next((row["totalPremiumPaid"] for row in ib if int(row["policyYear"]) == 5), ib[-1]["totalPremiumPaid"])
    first_year_leverage = iul["policy"]["sumInsured"] / iul["policy"]["initialPremium"]
    planned_leverage = iul["policy"]["sumInsured"] / total_iul_paid

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_picture(slide, cover, Inches(7.2), Inches(0.0), width=Inches(6.13), height=Inches(7.5))
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.62), Inches(6.1), Inches(5.85))
    panel.fill.solid(); panel.fill.fore_color.rgb = PANEL; panel.line.color.rgb = LINE
    add_textbox(slide, Inches(0.92), Inches(0.98), Inches(5.2), Inches(0.7), "教育金储备池与传承杠杆定制方案", size=24, bold=True, line_spacing=1.15)
    add_textbox(slide, Inches(0.94), Inches(1.72), Inches(5.3), Inches(0.38), f"{savings_company['companyName']} 储蓄险 + {iul_company['companyName']} IUL", size=14.5, color=MUTED, line_spacing=1.15)
    add_card(slide, Inches(1.0), Inches(2.65), Inches(2.3), Inches(1.18), "孩子储蓄险", f"US${money(savings['policy']['contractualTotalPremium'])}", "教育金储备池", accent=True, value_size=16.5)
    add_card(slide, Inches(3.55), Inches(2.65), Inches(2.3), Inches(1.18), "IUL 保额", f"US${money(iul['policy']['sumInsured'])}", "家族传承底盘", value_size=16.5)
    add_card(slide, Inches(1.0), Inches(3.95), Inches(2.3), Inches(1.18), "教育金起领", f"第{int(first_draw['policyYear'])}年", f"每年 US${money(first_draw['annualWithdrawal'])}", value_size=16.5)
    add_card(slide, Inches(3.55), Inches(3.95), Inches(2.3), Inches(1.18), "IUL 首年杠杆", multiple(first_year_leverage), "保额 ÷ 首年保费", value_size=16.5)
    add_footer(slide, "一张为孩子准备未来现金流，一张为家庭建立长期传承杠杆。")

    # 2 family goal
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "家庭目标", "储蓄险讲时间节奏，IUL 讲传承杠杆")
    add_picture(slide, father_child, Inches(7.35), Inches(1.35), width=Inches(5.2), height=Inches(4.2))
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(2.2), Inches(1.75), f"孩子 · {int(savings['insured']['age']) + 1} 岁", "未来教育金准备对象\n需要按年龄阶段释放现金流")
    add_fact_card(slide, Inches(3.25), Inches(1.45), Inches(2.2), Inches(1.75), f"家庭持有人 · {iul['insured']['age']} 岁", "长期家族资产安排对象\n需要更高身故杠杆和传承效率")
    add_fact_card(slide, Inches(5.7), Inches(1.45), Inches(1.35), Inches(1.75), "主线", "未来\n+传承")
    add_paragraphs(slide, Inches(0.86), Inches(3.55), Inches(6.05), Inches(2.3), [
        "孩子储蓄险负责 18/21/25 岁的教育金和成长现金流节奏。",
        "IUL 负责更高的身故杠杆与中长期现金价值弹性。",
        "两者不是重复配置，而是功能分层。"
    ], size=15)

    # 3 structure
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "双层结构总览", "先看现金流层，再看杠杆层")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(5.45), Inches(2.0), "第一层：孩子教育金储备池", f"年缴 US${money(savings['policy']['annualPremium'])} × {savings['policy']['payYears']} 年\n第 {int(first_draw['policyYear'])} 年开始每年领约 US${money(first_draw['annualWithdrawal'])}\n18/21/25 岁累计约 US${money(s18['cumulativeWithdrawal'])} / US${money(s21['cumulativeWithdrawal'])} / US${money(s25['cumulativeWithdrawal'])}", fill_color=ACCENT_LIGHT)
    add_fact_card(slide, Inches(6.95), Inches(1.45), Inches(5.45), Inches(2.0), "第二层：家族传承杠杆层", f"首年保费 US${money(iul['policy']['initialPremium'])}\n基础保额 US${money(iul['policy']['sumInsured'])}\n首年杠杆 {multiple(first_year_leverage)}，按5年总保费杠杆 {multiple(planned_leverage)}", fill_color=GREEN_LIGHT)
    add_paragraphs(slide, Inches(0.86), Inches(4.0), Inches(11.4), Inches(1.85), [
        "储蓄险解决的是孩子未来教育与成长现金流。",
        "IUL 解决的是长期身故保障底盘与代际传承效率。",
    ], size=16)

    # 4 savings role
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "为什么先配孩子储蓄险", "教育金是时间确定、节点清晰的刚性目标")
    add_picture(slide, child_growth, Inches(7.45), Inches(1.3), width=Inches(5.0), height=Inches(4.25))
    add_paragraphs(slide, Inches(0.86), Inches(1.45), Inches(6.05), Inches(3.6), [
        f"总保费约 US${money(savings['policy']['contractualTotalPremium'])}，回本较早，长期退保价值可持续抬升。",
        f"第 20 年不提取现价约 US${money(s20_nv['totalSurrenderValue'])}，第 30 年约 US${money(s30_nv['totalSurrenderValue'])}。",
        f"第 {int(first_draw['policyYear'])} 年开始每年约 US${money(first_draw['annualWithdrawal'])}，可顺着孩子成长阶段释放教育金。",
        "它负责未来的确定节奏，不负责家庭传承杠杆。"
    ], size=14.8)

    # 5 savings no-withdraw chart
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "储蓄险不提取曲线", "先看长期累积能力")
    years = [1, 5, 10, 20, 30, 40, 50]
    s_map = {int(r['policyYear']): r for r in sb}
    chart(slide, Inches(0.65), Inches(1.45), Inches(6.3), Inches(4.25), years, [
        ("总退保价值", [float(s_map[y]['totalSurrenderValue']) for y in years]),
        ("保证现金价值", [float(s_map[y]['guaranteedCashValue']) for y in years]),
        ("已交总保费", [float(min(y, savings['policy']['payYears']) * savings['policy']['annualPremium']) for y in years]),
    ])
    add_picture(slide, savings_growth, Inches(7.35), Inches(1.45), width=Inches(5.05), height=Inches(2.3))
    add_fact_card(slide, Inches(7.35), Inches(3.95), Inches(5.05), Inches(1.7), "解读", f"第20年总退保价值约 US${money(s20_nv['totalSurrenderValue'])}\n第30年总退保价值约 US${money(s30_nv['totalSurrenderValue'])}\n这是教育金不提前动用时的长期资金池。")

    # 6 savings withdraw chart
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "储蓄险提取曲线", "再看教育金如何按节点释放")
    wyears = [1, 6, 10, 18, 21, 25, 30]
    w_map = {int(r['policyYear']): r for r in sw}
    chart(slide, Inches(0.65), Inches(1.45), Inches(6.3), Inches(4.25), wyears, [
        ("累计领取金额", [float(w_map[y]['cumulativeWithdrawal']) for y in wyears]),
        ("提后现价", [float(w_map[y]['surrenderValueAfter']) for y in wyears]),
        ("累计经济总值", [float(w_map[y]['cumulativeWithdrawal']) + float(w_map[y]['surrenderValueAfter']) for y in wyears]),
    ])
    add_picture(slide, graduation, Inches(7.35), Inches(1.45), width=Inches(5.05), height=Inches(2.3))
    add_fact_card(slide, Inches(7.35), Inches(3.95), Inches(5.05), Inches(1.7), "教育金节点", f"18岁累计约 US${money(s18['cumulativeWithdrawal'])}\n21岁累计约 US${money(s21['cumulativeWithdrawal'])}\n25岁累计约 US${money(s25['cumulativeWithdrawal'])}")

    # 7 iul role
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "为什么再配 IUL", "IUL 解决的是教育金之外的长期杠杆层")
    add_picture(slide, city, Inches(7.35), Inches(1.35), width=Inches(5.2), height=Inches(4.15))
    add_paragraphs(slide, Inches(0.86), Inches(1.48), Inches(6.1), Inches(3.4), [
        f"基础保额 US${money(iul['policy']['sumInsured'])}，用相对有限的计划保费建立更高身故保障。",
        f"首年杠杆约 {multiple(first_year_leverage)}，按总计划保费杠杆约 {multiple(planned_leverage)}。",
        f"第20年当前假设现金值约 US${money(i20['nonGuaranteedCashValue'])}，第30年约 US${money(i30['nonGuaranteedCashValue'])}。",
        "它不是给孩子领教育金，而是给家庭做长期传承和现金价值缓冲。"
    ], size=14.8)

    # 8 iul chart
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "IUL 长期利益曲线", "把现金值和身故赔偿一起看")
    iyears = [1, 5, 10, 20, 30, 40, 50]
    i_map = {int(r['policyYear']): r for r in ib}
    chart(slide, Inches(0.65), Inches(1.45), Inches(6.3), Inches(4.25), iyears, [
        ("保证现金值", [float(i_map[y]['guaranteedCashValue']) for y in iyears]),
        ("当前假设现金值", [float(i_map[y]['nonGuaranteedCashValue']) for y in iyears]),
        ("当前假设身故赔偿", [float(i_map[y]['nonGuaranteedDeathBenefit']) for y in iyears]),
    ])
    add_fact_card(slide, Inches(7.35), Inches(1.6), Inches(5.0), Inches(4.0), "关键数字", f"第20年当前假设现金值：US${money(i20['nonGuaranteedCashValue'])}\n第30年当前假设现金值：US${money(i30['nonGuaranteedCashValue'])}\n100岁当前假设身故赔偿：US${money(i100['nonGuaranteedDeathBenefit'])}\n客户定位：{iul['insured']['age']}岁，{iul['insured']['gender']}，{iul['insured']['smoker']}")

    # 9 dual company
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "公司与服务分工", "储蓄险与 IUL 分别承担不同金融功能")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(5.45), Inches(3.5), f"{savings_company['companyName']} · 储蓄险", f"{savings_company['companyIntro']}\n\n定位：以孩子未来教育金、成长现金流和长期累积为核心。", fill_color=ACCENT_LIGHT)
    add_fact_card(slide, Inches(6.95), Inches(1.45), Inches(5.45), Inches(3.5), f"{iul_company['companyName']} · IUL", f"{iul_company['companyIntro']}\n\n定位：以高杠杆身故保障、长期现金价值和家族传承为核心。", fill_color=GREEN_LIGHT)
    add_footer(slide, "组合的价值不是叠加产品数量，而是把功能拆到最清楚。")

    # 10 synergy
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "协同关系", "教育金层与传承杠杆层不抢职责")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(5.45), Inches(2.05), "孩子储蓄险负责", "未来 18/21/25 岁教育金节奏\n成长阶段现金流\n未来启动资金", fill_color=ACCENT_LIGHT)
    add_fact_card(slide, Inches(6.95), Inches(1.45), Inches(5.45), Inches(2.05), "IUL 负责", "更高身故保障底盘\n家族传承效率\n晚年流动性与应急缓冲", fill_color=GREEN_LIGHT)
    add_paragraphs(slide, Inches(0.86), Inches(4.0), Inches(11.35), Inches(1.8), [
        f"孩子储蓄险从第 {int(first_draw['policyYear'])} 年开始每年约领取 US${money(first_draw['annualWithdrawal'])}，到 18/21/25 岁已累计约 US${money(s18['cumulativeWithdrawal'])} / US${money(s21['cumulativeWithdrawal'])} / US${money(s25['cumulativeWithdrawal'])}，为教育阶段提供长期稳健现金流。",
        f"IUL 自首年开始就建立约 US${money(iul['policy']['sumInsured'])} 的高基本保额，以更高杠杆为家庭筑起长期保障与财富传承底盘。",
        "两张保单一张管未来现金流，一张管家族资产护城河，功能清楚、互不冲突。"
    ], size=15.5)

    # 11 timeline
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "家庭时间轴", "一个看孩子年龄，一个看家族资产层级")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(2.3), Inches(10.8), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = LINE; line.line.color.rgb = LINE
    timeline_marker(slide, Inches(0.95), "现在", ["孩子储蓄开始累积", "IUL 杠杆结构就位"])
    timeline_marker(slide, Inches(3.20), "孩子 6 岁", [f"每年约 US${money(first_draw['annualWithdrawal'])}", "教育金开始释放"])
    timeline_marker(slide, Inches(5.45), "孩子 18/21 岁", [f"累计约 US${money(s18['cumulativeWithdrawal'])}", f"累计约 US${money(s21['cumulativeWithdrawal'])}"])
    timeline_marker(slide, Inches(7.70), "IUL 第20年", [f"现值约 US${money(i20['nonGuaranteedCashValue'])}", "传承底盘持续在位"])
    timeline_marker(slide, Inches(9.95), "IUL 第30年+", [f"现值约 US${money(i30['nonGuaranteedCashValue'])}", "身故保障继续服务家族"])
    add_footer(slide, "孩子的未来现金流与家庭长期传承底盘，分别沿两条时间轴发挥作用。")

    # 12 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "结论", "一张负责未来现金流，一张负责长期传承杠杆")
    add_paragraphs(slide, Inches(0.86), Inches(1.45), Inches(6.0), Inches(3.2), [
        "孩子储蓄险保证未来教育金节奏明确，不需要临时为升学资金找现金流。",
        "IUL 用更高杠杆把家族长期保障和传承底盘提前锁定下来。",
        "这套组合不是重复配置，而是把家庭资产目标拆成两个不同的功能层。"
    ], size=15.2)
    add_card(slide, Inches(0.9), Inches(4.95), Inches(2.55), Inches(1.15), "教育金池", f"US${money(savings['policy']['contractualTotalPremium'])}", "孩子未来现金流")
    add_card(slide, Inches(3.75), Inches(4.95), Inches(2.55), Inches(1.15), "IUL 杠杆", multiple(first_year_leverage), "首年保额杠杆", accent=True)
    add_card(slide, Inches(6.6), Inches(4.95), Inches(2.55), Inches(1.15), "IUL 首年保额", f"US${money(iul['policy']['sumInsured'])}", "财富护城河")
    add_picture(slide, savings_cash, Inches(9.45), Inches(1.55), width=Inches(3.3), height=Inches(4.8))
    add_footer(slide, "后续若再叠加重疾险，就会形成：教育金层 + 传承杠杆层 + 家庭风险防火墙。")

    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
