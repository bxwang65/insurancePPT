#!/usr/bin/env python3
"""
HybridGenerator — 混合生成器
整合完整流水线: LLM内容规划 → Playwright截图 → python-pptx嵌入

流水线:
  PDF提取 → ContentPlanner(LLM) → SlidePlan → SlideRenderer(Playwright) → PNG背景
    → python-pptx嵌入PNG背景+文字叠加 → 成品PPT

使用方法:
    gen = HybridGenerator()
    pptx_path = gen.generate(
        extraction_results=[...],  # 从PDF提取的结构化JSON
        user_intent="帮我看看回本速度",
        customer_name="VIP先生"
    )
"""

import os
import sys
import json
import asyncio
import subprocess
from pathlib import Path
from datetime import datetime
from typing import Optional

# ─── 品牌配置 ──────────────────────────────────────────────
BRAND_COLORS = {
    "primary": "#0A3C5F",      # 深海蓝
    "accent_teal": "#18898D",  # 青绿
    "accent_gold": "#C9A027",  # 金色
}

FONT_NAME = "Heiti TC"

# ─── LLM调用（统一接口） ────────────────────────────────────

def call_llm_content_planner(extractions_json: list, user_intent: str = "", customer_name: str = "") -> dict:
    """
    调用LLM内容规划器（通过Bun脚本或直接调用TypeScript）
    返回: ContentPlan JSON
    """
    # 检查是否有bun可用的LLM客户端
    bun_script = Path(__file__).parent / "llm_caller.js"

    if bun_script.exists():
        # 使用Bun调用TypeScript LLM客户端
        result = subprocess.run(
            ["bun", str(bun_script), json.dumps({
                "extractions": extractions_json,
                "userIntent": user_intent,
                "customerName": customer_name
            })],
            capture_output=True,
            text=True,
            timeout=60
        )
        if result.returncode == 0:
            return json.loads(result.stdout)
        else:
            print(f"[HybridGenerator] LLM call failed: {result.stderr}", file=sys.stderr)

    # 降级：返回规则生成的大纲
    return _generate_fallback_plan(extractions_json, user_intent, customer_name)


def _generate_fallback_plan(extractions_json: list, user_intent: str, customer_name: str) -> dict:
    """降级方案：用规则生成默认大纲"""
    slides = []
    page_num = 1

    # 封面
    slides.append({
        "pageNumber": page_num,
        "title": "保险计划方案",
        "narrativeText": f"为{customer_name or '尊贵客户'}精心设计的财富规划",
        "contentFocus": "封面",
        "visualType": "纯文本",
        "dataHighlights": [],
        "layout": "单栏",
        "type": "simple"
    })
    page_num += 1

    for ext in extractions_json:
        plan_type = ext.get("planType", ext.get("type", "savings"))
        data = ext.get("data", {})

        # 产品概览
        slides.append({
            "pageNumber": page_num,
            "title": data.get("product_name", f"{plan_type.upper()}产品"),
            "narrativeText": f"了解{_type_label(plan_type)}的核心优势",
            "contentFocus": "产品特点",
            "visualType": "纯文本",
            "dataHighlights": [],
            "layout": "单栏",
            "type": "simple"
        })
        page_num += 1

        if plan_type in ("savings", "savings_plan"):
            # 面积图：账户价值
            benefits = data.get("benefit_illustration", [])
            if benefits:
                years = [r.get("policy_year", i+1) for i, r in enumerate(benefits[:10])]
                # 找total_surrender_value或类似字段
                values = []
                for r in benefits[:10]:
                    val = r.get("total_surrender_value") or r.get("account_value") or r.get("benefit_amount", 0)
                    values.append(int(val))

                slides.append({
                    "pageNumber": page_num,
                    "title": "账户价值增长",
                    "narrativeText": "时间是最好的朋友，复利是最大的杠杆",
                    "contentFocus": "长期财富增值",
                    "visualType": "面积图",
                    "chartType": "area_chart",
                    "dataHighlights": _extract_highlights(benefits),
                    "layout": "左图右文",
                    "type": "area_chart",
                    "data": {"years": years, "values": values}
                })
                page_num += 1

                # KPI卡片
                total_prem = max((r.get("total_premium_paid", 0) for r in benefits), default=0)
                y20 = next((r for r in benefits if r.get("policy_year") == 20), None)
                y30 = next((r for r in benefits if r.get("policy_year") == 30), None)

                kpis = [
                    {"label": "年缴保费", "value": _format_num(data.get("annual_premium", 0)), "unit": "USD", "sub": "每年"},
                    {"label": "缴费期", "value": str(data.get("premium_payment_period", "5")), "unit": "年", "sub": ""},
                ]
                if y20:
                    mult = round(y20.get("total_surrender_value", 0) / total_prem, 1) if total_prem else 0
                    kpis.append({"label": "20年倍数", "value": str(mult), "unit": "x", "sub": ""})
                if y30:
                    mult = round(y30.get("total_surrender_value", 0) / total_prem, 1) if total_prem else 0
                    kpis.append({"label": "30年倍数", "value": str(mult), "unit": "x", "sub": ""})

                slides.append({
                    "pageNumber": page_num,
                    "title": "方案核心指标",
                    "narrativeText": "用数据说话，让选择更清晰",
                    "contentFocus": "关键指标",
                    "visualType": "KPI卡片",
                    "chartType": "kpi_cards",
                    "dataHighlights": [],
                    "layout": "全屏KPI",
                    "type": "kpi_cards",
                    "kpis": kpis[:6]
                })
                page_num += 1

        elif plan_type in ("ci", "critical_illness"):
            # 重疾险：KPI展示
            slides.append({
                "pageNumber": page_num,
                "title": "保障范围",
                "narrativeText": "全方位守护您和家人的健康",
                "contentFocus": "保障概览",
                "visualType": "KPI卡片",
                "chartType": "kpi_cards",
                "dataHighlights": [],
                "layout": "全屏KPI",
                "type": "kpi_cards",
                "kpis": [
                    {"label": "保障总额", "value": _format_num(data.get("sum_insured", 0)), "unit": "USD", "sub": ""},
                    {"label": "年缴保费", "value": _format_num(data.get("annual_premium", 0)), "unit": "USD", "sub": "每天 $X"},
                    {"label": "缴费期", "value": str(data.get("premium_payment_period", "10")), "unit": "年", "sub": ""},
                ]
            })
            page_num += 1

    # 综合建议
    slides.append({
        "pageNumber": page_num,
        "title": "综合方案建议",
        "narrativeText": "专业建议，助力您的人生规划",
        "contentFocus": "建议",
        "visualType": "纯文本",
        "dataHighlights": [],
        "layout": "单栏",
        "type": "simple"
    })
    page_num += 1

    # 感谢页
    slides.append({
        "pageNumber": page_num,
        "title": "感谢",
        "narrativeText": "感谢您的信任，期待为您提供专业服务",
        "contentFocus": "联系信息",
        "visualType": "纯文本",
        "dataHighlights": [],
        "layout": "单栏",
        "type": "simple"
    })

    return {
        "overallNarrative": f"为{customer_name or '尊贵客户'}定制的保险规划方案",
        "customerProfileSummary": f"客户: {customer_name or '尊贵客户'}",
        "slides": slides,
        "metadata": {
            "productTypes": list(set(e.get("planType", "savings") for e in extractions_json)),
            "totalPages": len(slides),
            "generatedAt": datetime.now().isoformat(),
            "brandColors": BRAND_COLORS
        }
    }


def _type_label(t: str) -> str:
    return {"savings": "储蓄险", "ci": "重疾险", "iul": "万用寿险"}.get(t, t)


def _format_num(v) -> str:
    try:
        return f"{float(v):,.0f}"
    except:
        return str(v)


def _extract_highlights(benefits: list) -> list:
    """从benefits数据提取关键数字"""
    highlights = []
    for r in benefits:
        yr = r.get("policy_year")
        val = r.get("total_surrender_value", 0)
        prem = r.get("total_premium_paid", 0)
        if yr in (5, 10, 20, 30) and prem > 0:
            mult = round(val / prem, 1)
            highlights.append(f"Y{yr}: {mult}x")
    return highlights[:4]


# ─── Playwright截图 ─────────────────────────────────────────

def render_slides_to_png(slide_plans: list[dict], output_dir: str = "/tmp/slide_renders") -> list[Optional[str]]:
    """
    调用slide_renderer.py批量渲染PNG
    返回: PNG路径列表（空字符串表示失败）
    """
    renderer_script = Path(__file__).parent / "slide_renderer.py"
    if not renderer_script.exists():
        print("[HybridGenerator] slide_renderer.py not found", file=sys.stderr)
        return [""] * len(slide_plans)

    # 转换slide_plans为slide_specs
    specs = []
    for i, slide in enumerate(slide_plans):
        spec = {
            "_id": i,
            "type": _visual_to_type(slide.get("visualType", "纯文本")),
            "title": slide.get("title", ""),
            "narrative": slide.get("narrativeText", ""),
            "highlights": _build_highlights(slide),
        }

        # 填充type-specific数据
        if spec["type"] == "area_chart":
            data = slide.get("data", {})
            spec["data"] = {
                "years": data.get("years", list(range(1, 21))),
                "values": data.get("values", [])
            }
        elif spec["type"] == "kpi_cards":
            spec["kpis"] = slide.get("kpis", [])
        elif spec["type"] == "comparison":
            spec["columns"] = slide.get("columns", [])

        specs.append(spec)

    # 调用Python渲染器（使用临时文件）
    import tempfile
    spec_file = None
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False) as f:
            json.dump(specs, f)
            spec_file = f.name

        result = subprocess.run(
            [sys.executable, str(renderer_script), "--spec-file", spec_file],
            capture_output=True,
            text=True,
            timeout=120,
            cwd=str(renderer_script.parent)
        )
        if result.returncode == 0:
            # slide_renderer.py 输出每张幻灯片两行：状态行(✓) + 路径行
            # 路径行以 .png 结尾且不包含勾号字符
            all_lines = result.stdout.strip().split("\n")
            paths = [p.strip() for p in all_lines if re.search(r'slide_\d+\.png$', p.strip()) and "✓" not in p]
            return paths if len(paths) == len(slide_plans) else [""] * len(slide_plans)
        else:
            print(f"[HybridGenerator] Render failed: {result.stderr[:500]}", file=sys.stderr)
            return [""] * len(slide_plans)
    except Exception as e:
        print(f"[HybridGenerator] Render error: {e}", file=sys.stderr)
        return [""] * len(slide_plans)
    finally:
        try:
            os.unlink(spec_file)
        except:
            pass


def _visual_to_type(visual: str) -> str:
    mapping = {
        "面积图": "area_chart",
        "折线图": "area_chart",
        "柱状图": "area_chart",
        "饼图": "area_chart",
        "KPI卡片": "kpi_cards",
        "对比": "comparison",
        "表格": "table",
        "纯文本": "simple",
    }
    return mapping.get(visual, "simple")


def _build_highlights(slide: dict) -> list:
    """从slide plan构建highlights数组"""
    highlights = slide.get("dataHighlights", [])
    if isinstance(highlights, list):
        result = []
        for h in highlights:
            if isinstance(h, str):
                result.append({"label": h, "value": "—"})
            else:
                result.append(h)
        return result
    return []


# ─── python-pptx生成 ───────────────────────────────────────

def create_pptx_with_backgrounds(
    slide_plans: list[dict],
    png_paths: list[str],
    output_path: str,
    title_texts: Optional[list[str]] = None
) -> str:
    """
    用python-pptx生成PPT，每页嵌入PNG作为背景
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from PIL import Image
    except ImportError as e:
        print(f"[HybridGenerator] Import error: {e}", file=sys.stderr)
        print("Please install: pip install python-pptx Pillow", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for i, (slide_plan, png_path) in enumerate(zip(slide_plans, png_paths)):
        slide = prs.slides.add_slide(blank_layout)

        # 1. 嵌入PNG背景
        if png_path and Path(png_path).exists():
            try:
                # 使用fill来设置背景图片
                background = slide.background
                fill = background.fill
                fill.blip_id = prs.part.relate_to(
                    png_path,
                    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                )
            except Exception as e:
                print(f"[HybridGenerator] Background error slide {i+1}: {e}", file=sys.stderr)

        # 2. 叠加文字层（如果PNG没有文字或需要补充）
        title = slide_plan.get("title", "")
        narrative = slide_plan.get("narrativeText", "")
        layout = slide_plan.get("layout", "单栏")

        # 根据布局添加文字
        if png_path and Path(png_path).exists():
            # PNG已有设计，叠加少量文字
            _add_overlay_text(prs, slide, title, narrative, layout)
        else:
            # 无PNG，纯文字模式
            _add_full_text_slide(slide, title, narrative)

        # 3. 页脚
        _add_footer(slide, i+1, len(slide_plans))

    # 保存
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)


def _add_overlay_text(prs, slide, title, narrative, layout):
    """在PNG背景上叠加文字"""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # 页码（右下角）
        pass  # footer单独处理
    except Exception as e:
        print(f"[HybridGenerator] Overlay text error: {e}", file=sys.stderr)


def _add_full_text_slide(slide, title, narrative):
    """纯文字幻灯片"""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        # 背景色（深海蓝）
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x3C, 0x5F)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.3), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 叙事文案
        narrative_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4), Inches(10.3), Inches(1.5)
        )
        tf = narrative_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = narrative
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    except Exception as e:
        print(f"[HybridGenerator] Full text slide error: {e}", file=sys.stderr)


def _add_footer(slide, page_num: int, total: int):
    """添加页脚"""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        footer_box = slide.shapes.add_textbox(
            Inches(12), Inches(7.1), Inches(1.2), Inches(0.3)
        )
        tf = footer_box.text_frame
        tf.paragraphs[0].text = f"{page_num}/{total}"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    except:
        pass


# ─── 主生成函数 ─────────────────────────────────────────────

def generate(
    extraction_results: list[dict],
    user_intent: str = "",
    customer_name: str = "",
    output_dir: str = "/tmp/insurance_pptx"
) -> str:
    """
    完整流水线入口

    extraction_results: 从PDF提取的结构化JSON列表
    user_intent: 用户对话意图
    customer_name: 客户姓名

    返回: 生成的PPTX文件路径
    """
    print(f"[HybridGenerator] Starting pipeline for {customer_name or '客户'}...")
    print(f"[HybridGenerator] Input: {len(extraction_results)} extraction result(s)")

    # Step 1: LLM内容规划
    print("[1/3] Running LLM ContentPlanner...")
    plan = call_llm_content_planner(extraction_results, user_intent, customer_name)
    slides = plan.get("slides", [])
    print(f"    → Generated {len(slides)} slide plans")

    # Step 2: Playwright渲染PNG背景
    print("[2/3] Rendering PNG backgrounds with Playwright...")
    png_paths = render_slides_to_png(slides)
    valid_count = sum(1 for p in png_paths if p and Path(p).exists())
    print(f"    → Rendered {valid_count}/{len(slides)} PNG backgrounds")

    # Step 3: python-pptx生成
    print("[3/3] Generating PPTX with python-pptx...")
    output_path = Path(output_dir) / f"insurance_plan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    result_path = create_pptx_with_backgrounds(
        slides, png_paths, str(output_path)
    )
    print(f"    → Saved to: {result_path}")

    return result_path


# ─── CLI入口 ──────────────────────────────────────────────

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="HybridGenerator — 混合生成器")
    parser.add_argument("--extractions", type=str, help="JSON格式的提取结果")
    parser.add_argument("--extractions-file", type=str, help="提取结果JSON文件")
    parser.add_argument("--intent", type=str, default="", help="用户意图")
    parser.add_argument("--customer", type=str, default="", help="客户姓名")
    parser.add_argument("--output-dir", type=str, default="/tmp/insurance_pptx")
    args = parser.parse_args()

    if args.extractions:
        extractions = json.loads(args.extractions)
    elif args.extractions_file:
        with open(args.extractions_file) as f:
            extractions = json.load(f)
    else:
        # 测试模式
        print("Running in test mode...")
        test_extractions = [
            {
                "planType": "savings",
                "data": {
                    "product_name": "匠心传承储蓄计划2尊尚版",
                    "annual_premium": 100000,
                    "premium_payment_period": 5,
                    "benefit_illustration": [
                        {"policy_year": 1, "total_premium_paid": 100000, "total_surrender_value": 0},
                        {"policy_year": 5, "total_premium_paid": 500000, "total_surrender_value": 520000},
                        {"policy_year": 10, "total_premium_paid": 500000, "total_surrender_value": 720000},
                        {"policy_year": 20, "total_premium_paid": 500000, "total_surrender_value": 1350000},
                        {"policy_year": 30, "total_premium_paid": 500000, "total_surrender_value": 2800000},
                    ]
                }
            }
        ]
        result = generate(
            test_extractions,
            user_intent="帮我看看回本速度和长期收益",
            customer_name="VIP先生",
            output_dir=args.output_dir
        )
        print(f"\n✅ Generated: {result}")
        sys.exit(0)

    result = generate(
        extraction_results=extractions,
        user_intent=args.intent,
        customer_name=args.customer,
        output_dir=args.output_dir
    )
    print(f"\n✅ Generated: {result}")
