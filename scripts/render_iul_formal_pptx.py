#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FONT_CN = "Microsoft YaHei"
BG = RGBColor(247, 244, 239)
PANEL = RGBColor(255, 255, 255)
PRIMARY = RGBColor(34, 34, 34)
MUTED = RGBColor(92, 92, 92)
LINE = RGBColor(221, 214, 201)
ACCENT = RGBColor(228, 162, 41)
ACCENT_LIGHT = RGBColor(248, 236, 201)
GREEN = RGBColor(39, 108, 84)
BLUE = RGBColor(58, 97, 152)


def money(value: float | int) -> str:
    return f"{round(float(value)):,}"


def multiple(value: float) -> str:
    return f"{value:.2f}x"


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, line_spacing=1.4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, left, top, width, height, lines, size=14.5, bullet=True, color=PRIMARY, line_spacing=1.45):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.text = f"• {line}" if bullet else line
        if p.runs:
            run = p.runs[0]
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=""):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.48), title, size=24, bold=True, line_spacing=1.15)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.78), Inches(10.0), Inches(0.28), subtitle, size=11.5, color=MUTED, line_spacing=1.15)


def add_card(slide, left, top, width, height, title, value, subtitle="", accent=False, value_size=19):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT if accent else PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.22), title, size=12, bold=True, color=MUTED, line_spacing=1.15)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.38), width - Inches(0.28), Inches(0.42), value, size=value_size, bold=True, color=ACCENT if accent else PRIMARY, line_spacing=1.15)
    if subtitle:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.82), width - Inches(0.28), Inches(0.22), subtitle, size=10.5, color=MUTED, line_spacing=1.15)


def add_fact_card(slide, left, top, width, height, title, body, fill_color=PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.22), title, size=11.5, bold=True, color=MUTED, line_spacing=1.15)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.36), width - Inches(0.28), height - Inches(0.48), body, size=13, color=PRIMARY, line_spacing=1.45)


def add_picture(slide, image_path: Path | None, left, top, width=None, height=None):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.62), Inches(7.0), Inches(11.8), Inches(0.2), text, size=10, color=MUTED, line_spacing=1.1)


def chart(slide, left, top, width, height, categories, series_spec):
    data = CategoryChartData()
    data.categories = [str(x) for x in categories]
    for name, values in series_spec:
        data.add_series(name, values)
    c = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, data).chart
    c.has_legend = True
    c.legend.position = XL_LEGEND_POSITION.BOTTOM
    c.legend.font.size = Pt(10)
    c.value_axis.has_major_gridlines = True
    c.value_axis.tick_labels.font.size = Pt(10)
    c.category_axis.tick_labels.font.size = Pt(10)
    c.value_axis.tick_labels.number_format = '$#,##0'
    palette = [ACCENT, GREEN, BLUE]
    for idx, series in enumerate(c.series):
        series.format.line.width = Pt(2.2)
        series.format.line.color.rgb = palette[idx % len(palette)]
    return c


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--normalized", required=True)
    parser.add_argument("--company-context", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    normalized = json.loads(Path(args.normalized).read_text(encoding="utf-8"))
    company = json.loads(Path(args.company_context).read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    company_id = company.get("companyId", "")
    logo = ROOT / f"public/assets/library/companies/{company_id}/logo.png" if company_id else None
    if logo is not None and not logo.exists():
        logo = None
    city = ROOT / "public/assets/library/themes/business/cityline-01.jpg"
    retirement = ROOT / "public/assets/library/themes/retirement/senior-life-01.jpg"

    benefit = normalized["benefitRows"]
    year20 = next(row for row in benefit if int(row["policyYear"]) == 20)
    year30 = next(row for row in benefit if int(row["policyYear"]) == 30)
    age100 = next((row for row in benefit if int(row["age"]) == 100), benefit[-1])
    total_planned = year20["totalPremiumPaid"]
    first_year_leverage = normalized["policy"]["sumInsured"] / normalized["policy"]["initialPremium"]
    planned_leverage = normalized["policy"]["sumInsured"] / total_planned
    age_based_role = (
        "对这位客户而言，它更像“家族财富传承工具”，重点解决晚年流动性和代际传承效率。"
        if normalized["insured"]["age"] >= 60
        else "对这位客户而言，它更像“家族长期杠杆资产”，重点解决早期高杠杆保障和中长期传承积累。"
    )

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_picture(slide, city, Inches(7.45), Inches(0), width=Inches(5.88), height=Inches(7.5))
    cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.62), Inches(6.25), Inches(5.85))
    cover.fill.solid(); cover.fill.fore_color.rgb = PANEL; cover.line.color.rgb = LINE
    add_picture(slide, logo, Inches(0.92), Inches(0.92), width=Inches(1.8))
    add_textbox(slide, Inches(0.95), Inches(1.55), Inches(5.35), Inches(0.7), f"{normalized['insured']['name']} 家庭新加坡 IUL 定制方案", size=24, bold=True, line_spacing=1.15)
    add_textbox(slide, Inches(0.95), Inches(2.24), Inches(5.4), Inches(0.4), f"{company['companyName']} · {normalized['productName']}", size=14.5, color=MUTED, line_spacing=1.15)
    add_card(slide, Inches(0.95), Inches(3.05), Inches(2.35), Inches(1.18), "首年保费", f"US${money(normalized['policy']['initialPremium'])}", "第1保单年度", accent=True, value_size=16.5)
    add_card(slide, Inches(3.48), Inches(3.05), Inches(2.35), Inches(1.18), "后续年缴", f"US${money(normalized['policy']['annualPremium'])}", "第2-10保单年度", value_size=16.5)
    add_card(slide, Inches(0.95), Inches(4.35), Inches(2.35), Inches(1.18), "基础保额", f"US${money(normalized['policy']['sumInsured'])}", "身故保障底盘", value_size=16.5)
    add_card(slide, Inches(3.48), Inches(4.35), Inches(2.35), Inches(1.18), "首年杠杆", multiple(first_year_leverage), "保额 ÷ 首年保费", value_size=16.5)
    add_footer(slide, "定位：高杠杆身故保障 + 现金价值增长弹性 + 家族传承工具。")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "公司与产品定位", "先讲公司实力，再讲 IUL 的家庭功能")
    add_picture(slide, logo, Inches(0.75), Inches(1.45), width=Inches(2.4))
    add_textbox(slide, Inches(0.78), Inches(2.45), Inches(5.2), Inches(0.7), company["companyIntro"], size=14, line_spacing=1.45)
    facts = company.get("companyFacts", [])[:4]
    positions = [(Inches(6.15), Inches(1.45)), (Inches(9.2), Inches(1.45)), (Inches(6.15), Inches(3.55)), (Inches(9.2), Inches(3.55))]
    for idx, fact in enumerate(facts):
        left, top = positions[idx]
        add_fact_card(slide, left, top, Inches(2.7), Inches(1.65), fact["label"], fact["value"], fill_color=ACCENT_LIGHT if idx == 0 else PANEL)
    add_footer(slide, "公司事实来自本地 Company Factbook 与内部公司资料库。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "保单核心参数", "先把客户最关心的姓名、年龄、缴费和保额讲清楚")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(3.1), Inches(1.7), "客户信息", f"姓名：{normalized['insured']['name']}\n年龄：{normalized['insured']['age']} 岁\n性别：{normalized['insured']['gender']}\n类别：{normalized['insured']['smoker']}")
    add_fact_card(slide, Inches(4.15), Inches(1.45), Inches(3.9), Inches(1.7), "缴费结构", f"首年：US${money(normalized['policy']['initialPremium'])}\n第2-10年：每年 US${money(normalized['policy']['annualPremium'])}\n累计计划保费：US${money(total_planned)}")
    add_fact_card(slide, Inches(8.35), Inches(1.45), Inches(3.9), Inches(1.7), "保障与杠杆", f"基础保额：US${money(normalized['policy']['sumInsured'])}\n首年杠杆：{multiple(first_year_leverage)}\n按10年总计划保费：{multiple(planned_leverage)}")
    add_paragraphs(slide, Inches(0.85), Inches(3.65), Inches(11.4), Inches(2.2), [
        "这份保单不是等额年缴：首年投入较高，后续9年降到 11.88 万美元。",
        "它的核心不是短期收益，而是用 300 万美元身故保障先建立家族传承底盘。",
        "客户 66 岁时投保，适合用来解决晚年流动性和代际传承效率。"
    ], size=15.2)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "指数账户与利率结构", "现金价值增长来自账户结构，不是单一路径")
    accounts = normalized["indexAccounts"]
    for idx, account in enumerate(accounts):
        add_fact_card(
            slide,
            Inches(0.8 + idx * 4.15), Inches(1.5), Inches(3.7), Inches(2.65),
            account["name"],
            f"配置比例：{account['allocation']}%\n当前假设利率：{account['assumedRate']}\n保底/下限：{account['floorRate']}\n封顶/参与：{account['capRate'] or '-'} / {account['participationRate'] or '-'}",
            fill_color=ACCENT_LIGHT if idx == 0 else PANEL,
        )
    add_paragraphs(slide, Inches(0.85), Inches(4.55), Inches(11.4), Inches(1.7), [
        "本方案 100% 配置在倍数指数账户，当前假设利率 7.00%，下限 0%。",
        "固定收益账户当前派息率 4.20%，第 11 个保单年度起另有 0.50% 保证忠诚红利派息率。",
        "这张 IUL 真正要解释的是“保底 + 指数弹性”的组合，而不是单看一个收益数字。 "
    ], size=14.5)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "长期利益演示", "用官方表格里的当前假设与保证值解释长期路径")
    years = [1, 5, 10, 20, 30, 40, 50]
    year_map = {int(row["policyYear"]): row for row in benefit}
    chart(
        slide,
        Inches(0.65), Inches(1.45), Inches(6.4), Inches(4.25),
        years,
        [
            ("保证现金值", [float(year_map[y]["guaranteedCashValue"]) for y in years]),
            ("当前假设现金值", [float(year_map[y]["nonGuaranteedCashValue"]) for y in years]),
            ("当前假设身故赔偿", [float(year_map[y]["nonGuaranteedDeathBenefit"]) for y in years]),
        ],
    )
    add_fact_card(slide, Inches(7.45), Inches(1.55), Inches(5.15), Inches(4.05), "关键节点", f"第20年当前假设现金值：US${money(year20['nonGuaranteedCashValue'])}\n第20年保证现金值：US${money(year20['guaranteedCashValue'])}\n第30年当前假设现金值：US${money(year30['nonGuaranteedCashValue'])}\n100岁当前假设身故赔偿：US${money(age100['nonGuaranteedDeathBenefit'])}")
    add_footer(slide, "关键理解：保证值与当前假设值差异很大，客户沟通必须先讲风险承受，再讲收益区间。")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
    add_title(slide, "它为家庭起到什么作用", "IUL 在家庭资产结构里不负责教育金，而负责杠杆与传承")
    add_picture(slide, retirement, Inches(7.55), Inches(1.45), width=Inches(5.0), height=Inches(3.5))
    add_paragraphs(slide, Inches(0.85), Inches(1.55), Inches(6.0), Inches(3.4), [
        "第一层：用较高的身故杠杆，给家庭留下确定性的传承底盘。",
        "第二层：随着现金价值累积，未来可作为晚年流动性和应急资金来源。",
        "第三层：如果市场长期表现良好，非保证现金值和身故赔偿都会继续抬升。",
        age_based_role
    ], size=15.2)
    add_fact_card(slide, Inches(0.85), Inches(5.25), Inches(2.8), Inches(1.15), "首年杠杆", multiple(first_year_leverage), fill_color=ACCENT_LIGHT)
    add_fact_card(slide, Inches(3.85), Inches(5.25), Inches(2.8), Inches(1.15), "10年总缴杠杆", multiple(planned_leverage))
    add_footer(slide, "后续若与储蓄险、重疾险叠加，IUL 适合作为第三层：高杠杆定向传承层。")

    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
