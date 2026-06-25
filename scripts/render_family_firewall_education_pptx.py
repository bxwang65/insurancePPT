#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FONT_CN = "Microsoft YaHei"
BG = RGBColor(248, 245, 239)
PANEL = RGBColor(255, 255, 255)
PRIMARY = RGBColor(38, 38, 38)
ACCENT = RGBColor(206, 138, 44)
ACCENT_LIGHT = RGBColor(240, 230, 208)
MUTED = RGBColor(96, 96, 96)
LINE = RGBColor(223, 214, 195)
GOOD = RGBColor(24, 108, 79)
GOOD_LIGHT = RGBColor(165, 204, 189)
ALERT = RGBColor(173, 68, 52)
BLUE = RGBColor(54, 94, 140)
BLUE_LIGHT = RGBColor(197, 214, 234)


def money(value: float | int) -> str:
    return f"{round(float(value)):,}"


def pct(value: float) -> str:
    return f"{value:.2f}%"


def paid_premium_for_year(year: int, annual_premium: float, pay_years: int) -> float:
    return float(annual_premium) * min(int(year), int(pay_years))


def simple_return(value: float, premium: float) -> float:
    if premium <= 0:
        return 0.0
    return (float(value) / float(premium) - 1.0) * 100.0


def compound_return(value: float, premium: float, year: int) -> float:
    if premium <= 0 or value <= 0 or year <= 0:
        return 0.0
    return ((float(value) / float(premium)) ** (1.0 / float(year)) - 1.0) * 100.0


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, line_spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, left, top, width, height, lines, size=15, color=PRIMARY, bullet=False, line_spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.text = f"• {line}" if bullet else line
        if p.runs:
            run = p.runs[0]
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.45), title, size=24, bold=True, line_spacing=1.2)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.78), Inches(9.2), Inches(0.28), subtitle, size=11, color=MUTED, line_spacing=1.2)


def add_card(slide, left, top, width, height, title, value, subtitle="", accent=False, title_size=12, value_size=20, subtitle_size=10.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT if accent else PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.22), title, size=title_size, bold=True, color=MUTED, line_spacing=1.2)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.38), width - Inches(0.28), Inches(0.4), value, size=value_size, bold=True, color=ACCENT if accent else PRIMARY, line_spacing=1.2)
    if subtitle:
        add_textbox(slide, left + Inches(0.14), top + Inches(0.82), width - Inches(0.28), Inches(0.22), subtitle, size=subtitle_size, color=MUTED, line_spacing=1.2)


def add_fact_card(slide, left, top, width, height, label, value, fill_color=PANEL, title_size=11.5, body_size=13.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.12), width - Inches(0.28), Inches(0.2), label, size=title_size, bold=True, color=MUTED, line_spacing=1.2)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.38), width - Inches(0.28), height - Inches(0.5), value, size=body_size, color=PRIMARY, line_spacing=1.5)


def add_panel(slide, left, top, width, height, title, lines, fill_color=PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.16), top + Inches(0.14), width - Inches(0.32), Inches(0.24), title, size=13, bold=True, color=PRIMARY, line_spacing=1.2)
    add_paragraphs(slide, left + Inches(0.16), top + Inches(0.46), width - Inches(0.32), height - Inches(0.6), lines, size=12.5, bullet=True, line_spacing=1.45)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.65), Inches(7.0), Inches(12.0), Inches(0.2), text, size=10, color=MUTED, line_spacing=1.2)


def add_timeline_marker(slide, left, title, lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.75), Inches(2.1), Inches(2.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), Inches(2.9), Inches(1.82), Inches(0.26), title, size=13, bold=True, color=ACCENT, line_spacing=1.2)
    add_paragraphs(slide, left + Inches(0.14), Inches(3.18), Inches(1.82), Inches(1.25), lines, size=11, bullet=False, line_spacing=1.35)


def set_chart_style(chart):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.number_format = '$#,##0'


def add_line_chart(slide, left, top, width, height, categories, series_spec):
    chart_data = CategoryChartData()
    chart_data.categories = [str(x) for x in categories]
    for name, values in series_spec:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data).chart
    set_chart_style(chart)
    colors = [ACCENT, GOOD, BLUE, ALERT, GOOD_LIGHT]
    for idx, series in enumerate(chart.series):
        series.format.line.width = Pt(2.2)
        series.format.line.color.rgb = colors[idx % len(colors)]
    return chart


def add_picture(slide, image_path: Path | None, left, top, width, height):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def style_table(table, widths, header_fill=ACCENT_LIGHT):
    for idx, width in enumerate(widths):
        table.columns[idx].width = width
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            text = cell.text
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.05
            run = p.add_run()
            run.text = text
            run.font.name = FONT_CN
            run.font.size = Pt(12)
            run.font.bold = r == 0
            run.font.color.rgb = PRIMARY
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else PANEL


def decade_rows_no_withdraw(benefit_rows, annual_premium, pay_years):
    targets = [1, 10, 20, 30, 40, 50, 60, 70, 80, 90]
    rows = []
    for year in targets:
        row = next((r for r in benefit_rows if int(r["policyYear"]) == year), None)
        if not row:
            continue
        paid = paid_premium_for_year(year, annual_premium, pay_years)
        value = float(row["totalSurrenderValue"])
        rows.append([
            str(int(row["age"])),
            str(year),
            money(paid),
            money(row["guaranteedCashValue"]),
            money(value),
            pct(simple_return(value, paid)),
            pct(compound_return(value, paid, year)),
        ])
    return rows


def decade_rows_withdraw(withdraw_rows, annual_premium, pay_years):
    targets = [1, 10, 20, 30, 40, 50, 60, 70, 80, 90]
    rows = []
    for year in targets:
        row = next((r for r in withdraw_rows if int(r["policyYear"]) == year), None)
        if not row:
            continue
        paid = paid_premium_for_year(year, annual_premium, pay_years)
        economic_total = float(row["cumulativeWithdrawal"]) + float(row["surrenderValueAfter"])
        rows.append([
            str(int(row["age"])),
            str(year),
            money(paid),
            money(row["annualWithdrawal"]),
            money(row["cumulativeWithdrawal"]),
            money(row["surrenderValueAfter"]),
            money(economic_total),
            pct(simple_return(economic_total, paid)),
            pct(compound_return(economic_total, paid, year)),
        ])
    return rows


def add_table_slide(slide, title, subtitle, cols, rows, note_lines):
    add_bg(slide)
    add_title(slide, title, subtitle)
    table = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(0.55), Inches(1.38), Inches(9.55), Inches(5.3)).table
    widths = [Inches(0.65), Inches(0.8), Inches(1.15), Inches(1.0), Inches(1.2), Inches(1.3), Inches(1.15), Inches(0.9), Inches(0.9)]
    if len(cols) == 7:
        widths = [Inches(0.75), Inches(0.9), Inches(1.2), Inches(1.25), Inches(1.35), Inches(1.0), Inches(1.0)]
    for i, col in enumerate(cols):
        table.cell(0, i).text = col
    for r_idx, row in enumerate(rows, 1):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = str(value)
    style_table(table, widths[:len(cols)])
    add_panel(slide, Inches(10.4), Inches(1.65), Inches(2.25), Inches(3.25), "解读", note_lines, fill_color=ACCENT_LIGHT)
    add_footer(slide, "表格口径：关键年度每10年展示，数字直接来自官方计划书标准化结果。")


def pick_existing_asset(*relative_candidates: str) -> Path | None:
    for rel in relative_candidates:
        path = ROOT / rel
        if path.exists():
            return path
    return None


def find_payback_year(benefit_rows, annual_premium, pay_years):
    for row in benefit_rows:
        year = int(row["policyYear"])
        if float(row["totalSurrenderValue"]) >= paid_premium_for_year(year, annual_premium, pay_years):
            return year
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--savings", required=True)
    parser.add_argument("--ci", required=True)
    parser.add_argument("--company-context", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    savings = json.loads(Path(args.savings).read_text(encoding="utf-8"))
    ci = json.loads(Path(args.ci).read_text(encoding="utf-8"))
    company = json.loads(Path(args.company_context).read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover_img = pick_existing_asset(
        "outputs/formal_環宇盈活儲蓄保險計劃_1780288726854_pipeline/assets/cover.jpg",
        "outputs/9f164714_api_pipeline/assets/cover.jpg",
    )
    family_img = pick_existing_asset(
        "outputs/formal_環宇盈活儲蓄保險計劃_1780288726854_pipeline/assets/closing.jpg",
        "outputs/9f164714_api_pipeline/assets/closing.jpg",
    )
    company_img = pick_existing_asset(
        "outputs/formal_環宇盈活儲蓄保險計劃_1780288726854_pipeline/assets/company.jpg",
        "outputs/9f164714_api_pipeline/assets/company.jpg",
    )

    child_age = int(savings["insured"]["age"]) + 1
    mom_age = int(ci["insured"]["age"])
    child_premium = float(savings["policy"]["annualPremium"])
    child_total = float(savings["policy"]["contractualTotalPremium"])
    child_pay_years = int(savings["policy"]["payYears"])
    ci_premium = float(ci["policy"]["annualPremium"])
    ci_total = float(ci["policy"]["totalPremium"])
    ci_pay_years = int(ci["policy"]["payYears"])
    ci_cover = float(ci["policy"]["sumInsured"]) + float(ci["policy"]["upgradeBenefitAmount"])
    company_facts = company.get("companyFacts", [])[:4]

    benefit_rows = savings["benefitRows"]
    payback_year = find_payback_year(benefit_rows, child_premium, child_pay_years)
    s20 = next((row for row in benefit_rows if int(row["policyYear"]) == 20), benefit_rows[-1])
    s30 = next((row for row in benefit_rows if int(row["policyYear"]) == 30), benefit_rows[-1])
    sw = savings["withdrawalRows"]
    first_draw = next((row for row in sw if float(row["annualWithdrawal"]) > 0), sw[0])
    w18 = next((row for row in sw if int(row["policyYear"]) == 18), sw[-1])
    w21 = next((row for row in sw if int(row["policyYear"]) == 21), sw[-1])
    w25 = next((row for row in sw if int(row["policyYear"]) == 25), sw[-1])
    ci_rows = ci["benefitRows"]
    ci10 = next((row for row in ci_rows if int(row["policyYear"]) == 10), ci_rows[-1])
    ci20 = next((row for row in ci_rows if int(row["policyYear"]) == 20), ci_rows[-1])
    ci30 = next((row for row in ci_rows if int(row["policyYear"]) == 30), ci_rows[-1])

    # Slide 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_picture(slide, cover_img, Inches(7.35), Inches(0.0), Inches(5.98), Inches(7.5))
    cover_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.6), Inches(6.2), Inches(5.2))
    cover_panel.fill.solid()
    cover_panel.fill.fore_color.rgb = PANEL
    cover_panel.line.color.rgb = LINE
    add_textbox(slide, Inches(0.9), Inches(0.92), Inches(5.7), Inches(0.95), "家庭防火墙与教育金定制方案", size=24, bold=True, line_spacing=1.15)
    add_textbox(slide, Inches(0.92), Inches(1.72), Inches(5.8), Inches(0.42), f"{company['companyName']} · 爱伴航重疾险 + 环宇盈活储蓄险", size=15.5, color=MUTED, line_spacing=1.2)
    # Cover KPI cards need extra upward offset in WPS/PPT viewers; otherwise the
    # bottom row visually reads as escaping the white panel even when raw bounds fit.
    add_card(slide, Inches(1.0), Inches(2.70), Inches(2.3), Inches(1.18), "妈妈重疾险", f"US${money(ci_cover)}", "前10年总保障", accent=True, value_size=16.5, subtitle_size=10)
    add_card(slide, Inches(3.55), Inches(2.70), Inches(2.3), Inches(1.18), "孩子储蓄险", f"US${money(child_total)}", "5年教育金储备", value_size=16.5, subtitle_size=10)
    add_card(slide, Inches(1.0), Inches(4.00), Inches(2.3), Inches(1.18), "妈妈年缴", f"US${money(ci_premium)}", "10年缴", value_size=16.5, subtitle_size=10)
    add_card(slide, Inches(3.55), Inches(4.00), Inches(2.3), Inches(1.18), "孩子年缴", f"US${money(child_premium)}", "5年缴", value_size=16.5, subtitle_size=10)
    add_footer(slide, "主线：妈妈负责家庭稳定，孩子负责未来规划。")

    # Slide 2 family portrait
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第一章 家庭画像", "先按家庭角色讲，再按产品落位")
    add_picture(slide, family_img, Inches(7.35), Inches(1.32), Inches(5.35), Inches(4.42))
    add_fact_card(slide, Inches(0.88), Inches(1.46), Inches(2.18), Inches(1.72), f"妈妈 · {mom_age} 岁", "家庭收入与照顾责任核心\n需要先锁住重大健康风险", body_size=12.3)
    add_fact_card(slide, Inches(3.24), Inches(1.46), Inches(2.18), Inches(1.72), f"孩子 · {child_age} 岁", "未来教育金准备对象\n需要长期现金流按节奏累积", body_size=12.3)
    add_fact_card(slide, Inches(5.6), Inches(1.46), Inches(1.4), Inches(1.72), "目标", "防风险\n备教育金", body_size=12.3)
    add_paragraphs(slide, Inches(0.85), Inches(3.35), Inches(5.95), Inches(2.6), [
        "如果妈妈先出风险，家庭收入和照顾能力会受影响。",
        "这时候不能先动孩子的教育金，所以重疾险必须放在第一层。",
        "孩子储蓄险负责未来 18/21/25 岁阶段的教育资金节奏。"
    ], size=15, bullet=True)
    add_footer(slide, "家庭主语先统一，后续每一页都围绕“谁承担什么责任”展开。")

    # Slide 3 company page
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第二章 公司实力", "公司页先交代品牌与财务底座，再讲产品")
    add_picture(slide, company_img, Inches(0.7), Inches(1.45), Inches(5.1), Inches(4.35))
    add_paragraphs(slide, Inches(0.85), Inches(5.95), Inches(5.0), Inches(0.8), [company["companyIntro"]], size=12.5, bullet=False, line_spacing=1.45)
    positions = [
        (Inches(6.15), Inches(1.45)),
        (Inches(9.35), Inches(1.45)),
        (Inches(6.15), Inches(3.55)),
        (Inches(9.35), Inches(3.55)),
    ]
    for idx, fact in enumerate(company_facts[:4]):
        left, top = positions[idx]
        add_fact_card(slide, left, top, Inches(2.75), Inches(1.7), fact["label"], fact["value"], fill_color=ACCENT_LIGHT if idx == 0 else PANEL)
    add_footer(slide, "公司事实采用内部 Company Factbook 权威口径。")

    # Slide 4 overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第三章 方案总览", "两张保单不是并排，而是各自承担不同职责")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(3.75), Inches(1.7), "左：妈妈重疾险", f"年缴 US${money(ci_premium)}\n总保费 US${money(ci_total)}\n首十年总保障 US${money(ci_cover)}")
    add_fact_card(slide, Inches(4.85), Inches(1.45), Inches(3.75), Inches(1.7), "右：孩子储蓄险", f"年缴 US${money(child_premium)}\n总保费 US${money(child_total)}\n第 {int(first_draw['policyYear'])} 年开始提领")
    add_fact_card(slide, Inches(8.9), Inches(1.45), Inches(3.45), Inches(1.7), "中间一句话", "妈妈负责家庭稳定\n孩子负责未来规划")
    add_paragraphs(slide, Inches(0.85), Inches(3.55), Inches(11.4), Inches(2.5), [
        "爱伴航解决的是大病、收入中断和医疗支出冲击。",
        "环宇盈活解决的是孩子未来 18/21/25 岁阶段的教育金节奏。",
        "先防风险，再做未来，是这套组合的核心结构。"
    ], size=15.5, bullet=True)

    # Slide 5 why CI first
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第四章 为什么先做妈妈重疾险", "第一层先处理家庭健康与现金流风险")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(3.8), Inches(1.65), "如果妈妈发生重疾", "收入减少\n照顾能力下降\n治疗与康复支出上升")
    add_fact_card(slide, Inches(4.9), Inches(1.45), Inches(3.8), Inches(1.65), "这时候不能做什么", "不能先动孩子教育金\n不能把未来准备拿去填当前风险")
    add_fact_card(slide, Inches(9.0), Inches(1.45), Inches(3.2), Inches(1.65), "所以爱伴航的角色", "家庭防火墙\n先扛风险支出", fill_color=ACCENT_LIGHT)
    add_paragraphs(slide, Inches(0.85), Inches(3.55), Inches(11.4), Inches(2.5), [
        f"基础保额 US${money(ci['policy']['baseSumInsured'])}，前 10 年升级保障 US${money(ci['policy']['upgradeBenefitAmount'])}，把前期保障抬到 US${money(ci_cover)}。",
        f"重疾 {ci['coverageSummary']['majorCiCount']} 项、早期危疾 {ci['coverageSummary']['earlyCiCount']} 项，加上 ICU 和豁免责任，核心是先把家庭底线守住。",
        "这张单的任务不是替代储蓄，而是保护储蓄不被提前动用。"
    ], size=15.5, bullet=True)

    # Slide 6 CI chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第五章 重疾险数据曲线", "把保障力度和现金价值拆开看，比纯文字更直观")
    ci_years = [int(r["policyYear"]) for r in ci_rows]
    add_line_chart(
        slide,
        Inches(0.7), Inches(1.45), Inches(6.5), Inches(4.3),
        ci_years,
        [
            ("重疾赔付额", [float(r["ciBenefit"]) for r in ci_rows]),
            ("退保价值", [float(r["totalSurrenderValue"]) for r in ci_rows]),
            ("已交总保费", [float(r["totalPremiumPaid"]) for r in ci_rows]),
        ],
    )
    add_panel(slide, Inches(7.55), Inches(1.55), Inches(5.1), Inches(4.1), "图表解读", [
        f"10 年缴总保费 US${money(ci_total)}，第 10 年重疾赔付约 US${money(ci10['ciBenefit'])}。",
        f"第 20 年重疾赔付约 US${money(ci20['ciBenefit'])}，第 30 年约 US${money(ci30['ciBenefit'])}。",
        f"现金价值在第 10/20/30 年约为 US${money(ci10['totalSurrenderValue'])} / US${money(ci20['totalSurrenderValue'])} / US${money(ci30['totalSurrenderValue'])}。",
        "这张保单的重点是家庭风险缓冲，而不是长期财富增值。"
    ])
    add_footer(slide, "图表口径：重疾赔付额、退保价值、已交保费均来自官方计划书标准化字段。")

    # Slide 7 why savings
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第六章 为什么给孩子做储蓄险", "第二层解决未来教育金与成长现金流")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(3.8), Inches(1.65), "环宇盈活不解决什么", "它不解决家庭健康风险\n不承担妈妈重疾后的支出压力")
    add_fact_card(slide, Inches(4.9), Inches(1.45), Inches(3.8), Inches(1.65), "它解决什么", f"18 岁累计约 US${money(w18['cumulativeWithdrawal'])}\n21 岁累计约 US${money(w21['cumulativeWithdrawal'])}\n25 岁累计约 US${money(w25['cumulativeWithdrawal'])}")
    add_fact_card(slide, Inches(9.0), Inches(1.45), Inches(3.2), Inches(1.65), "第二层角色", "教育金储备池\n为孩子未来预留节奏", fill_color=ACCENT_LIGHT)
    add_paragraphs(slide, Inches(0.85), Inches(3.55), Inches(11.4), Inches(2.4), [
        f"总保费约 US${money(child_total)}，回本约第 {payback_year} 年，第 20 年退保价值约 US${money(s20['totalSurrenderValue'])}，第 30 年约 US${money(s30['totalSurrenderValue'])}。",
        "如果家庭平稳，这张单就按成长阶段承担教育金、升学金、未来启动资金。",
        "所以它是第二层，不是第一层。"
    ], size=15.2, bullet=True)

    # Slide 8 savings no-withdraw chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第七章 储蓄险不提取曲线", "先看长期累积，再看提领取用")
    savings_years = [1, 5, 10, 20, 30, 40, 50]
    value_map = {int(r["policyYear"]): r for r in benefit_rows}
    add_line_chart(
        slide,
        Inches(0.7), Inches(1.45), Inches(6.5), Inches(4.3),
        savings_years,
        [
            ("总退保价值", [float(value_map[y]["totalSurrenderValue"]) for y in savings_years]),
            ("保证现金价值", [float(value_map[y]["guaranteedCashValue"]) for y in savings_years]),
            ("已交总保费", [paid_premium_for_year(y, child_premium, child_pay_years) for y in savings_years]),
        ],
    )
    add_panel(slide, Inches(7.55), Inches(1.55), Inches(5.1), Inches(4.1), "图表解读", [
        f"回本约第 {payback_year} 年，之后总退保价值开始明显拉开。",
        f"第 20 年总退保价值约 US${money(s20['totalSurrenderValue'])}，约为总保费的 {round(float(s20['totalSurrenderValue']) / child_total, 2)} 倍。",
        f"第 30 年总退保价值约 US${money(s30['totalSurrenderValue'])}，约为总保费的 {round(float(s30['totalSurrenderValue']) / child_total, 2)} 倍。",
        "这张曲线对应的是“不提取、让教育金继续积累”的路径。"
    ])
    add_footer(slide, "不提取口径：展示保证现金价值、总退保价值与已交保费的长期差距。")

    # Slide 9 savings withdrawal chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第八章 储蓄险提取曲线", "提领与剩余价值要一起看")
    withdraw_years = [1, 6, 10, 20, 30, 40, 50]
    withdraw_map = {int(r["policyYear"]): r for r in sw}
    add_line_chart(
        slide,
        Inches(0.7), Inches(1.45), Inches(6.5), Inches(4.3),
        withdraw_years,
        [
            ("累计领取金额", [float(withdraw_map[y]["cumulativeWithdrawal"]) for y in withdraw_years]),
            ("提取后退保价值", [float(withdraw_map[y]["surrenderValueAfter"]) for y in withdraw_years]),
            ("累计经济总值", [float(withdraw_map[y]["cumulativeWithdrawal"]) + float(withdraw_map[y]["surrenderValueAfter"]) for y in withdraw_years]),
        ],
    )
    add_panel(slide, Inches(7.55), Inches(1.55), Inches(5.1), Inches(4.1), "图表解读", [
        f"第 {int(first_draw['policyYear'])} 年开始每年约 US${money(first_draw['annualWithdrawal'])}。",
        f"到第 20 年累计领取约 US${money(w18['cumulativeWithdrawal'])}，提取后仍保留现价。",
        f"到第 30 年累计领取约 US${money(next(r['cumulativeWithdrawal'] for r in sw if int(r['policyYear']) == 30))}，经济总值继续抬升。",
        "这一页对应“教育阶段持续提取，但资金池没有被一次性抽空”的路径。"
    ])
    add_footer(slide, "提取口径：同时观察累计领取、提取后现价和累计经济总值。")

    # Slide 10 synergy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第九章 协同关系", "这页是整套方案的灵魂")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(5.45), Inches(2.05), "有重疾险", "风险来时，用重疾险扛支出\n储蓄险继续留给孩子\n教育金目标不被打断", fill_color=GOOD_LIGHT)
    add_fact_card(slide, Inches(6.95), Inches(1.45), Inches(5.45), Inches(2.05), "没有重疾险", "风险来时，家里可能先动教育金\n孩子未来目标被迫让位给当前支出", fill_color=ACCENT_LIGHT)
    add_paragraphs(slide, Inches(0.9), Inches(3.95), Inches(11.3), Inches(2.0), [
        "先用爱伴航守住家庭不失速，再用环宇盈活准备孩子未来不落空。",
        "这不是两张单相加，而是一个家庭资产与风险分层方案。"
    ], size=17, bullet=True)

    # Slide 11 policy summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第十章 两张保单分别展开", "先妈妈，后孩子，不要来回跳")
    add_fact_card(slide, Inches(0.8), Inches(1.45), Inches(5.55), Inches(2.15), "妈妈重疾险：爱伴航", f"基础保额 US${money(ci['policy']['baseSumInsured'])}\n年缴 US${money(ci_premium)}，缴费 {ci_pay_years} 年\n首 10 年总保障 US${money(ci_cover)}\n重疾 {ci['coverageSummary']['majorCiCount']} 项，早期危疾 {ci['coverageSummary']['earlyCiCount']} 项")
    add_fact_card(slide, Inches(6.95), Inches(1.45), Inches(5.55), Inches(2.15), "孩子储蓄险：环宇盈活", f"年缴 US${money(child_premium)}，缴费 {child_pay_years} 年\n第 {int(first_draw['policyYear'])} 年开始每年提领 US${money(first_draw['annualWithdrawal'])}\n第 20 年总退保价值 US${money(s20['totalSurrenderValue'])}\n教育金按阶段释放")
    add_paragraphs(slide, Inches(0.85), Inches(4.15), Inches(11.4), Inches(1.75), [
        "顺序固定：先讲谁扛风险，再讲谁承接未来目标。",
        "客户听到这里，已经知道两张单各自承担什么角色。"
    ], size=16, bullet=True)

    # Slide 12 no-withdraw table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    no_withdraw_rows = decade_rows_no_withdraw(benefit_rows, child_premium, child_pay_years)
    add_table_slide(
        slide,
        "第十一章 储蓄险不提取数据表（每10年）",
        "口径：官方不提取现金价值路径",
        ["年龄", "保单年度", "已交保费", "保证现价", "总退保价值", "单利", "复利"],
        no_withdraw_rows,
        [
            f"总保费 US${money(child_total)}，回本约第 {payback_year} 年。",
            f"20 年约 {round(float(s20['totalSurrenderValue']) / child_total, 2)} 倍，30 年约 {round(float(s30['totalSurrenderValue']) / child_total, 2)} 倍。",
            "适合看长期累积效率。"
        ],
    )

    # Slide 13 withdraw table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    withdraw_rows = decade_rows_withdraw(sw, child_premium, child_pay_years)
    add_table_slide(
        slide,
        "第十二章 储蓄险提取数据表（每10年）",
        "口径：官方提领取用路径",
        ["年龄", "保单年度", "已交保费", "年领金额", "累计领取", "提后现价", "经济总值", "单利", "复利"],
        withdraw_rows,
        [
            f"第 {int(first_draw['policyYear'])} 年开始每年约 US${money(first_draw['annualWithdrawal'])}。",
            f"25 岁阶段累计领取约 US${money(w25['cumulativeWithdrawal'])}。",
            "适合看教育金取用效率与剩余价值。"
        ],
    )

    # Slide 14 timeline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第十三章 家庭时间轴", "把谁在什么阶段发挥作用讲清楚")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(2.35), Inches(10.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE
    add_timeline_marker(slide, Inches(0.95), "现在", ["妈妈保障生效", "孩子储蓄开始累积"])
    add_timeline_marker(slide, Inches(3.25), "孩子 6 岁", ["开始每年教育金提领", f"每年约 US${money(first_draw['annualWithdrawal'])}"])
    add_timeline_marker(slide, Inches(5.55), "孩子 18/21 岁", [f"累计约 US${money(w18['cumulativeWithdrawal'])}", f"累计约 US${money(w21['cumulativeWithdrawal'])}"])
    add_timeline_marker(slide, Inches(7.85), "孩子 25 岁", [f"累计约 US${money(w25['cumulativeWithdrawal'])}", "未来升学/创业启动金"])
    add_timeline_marker(slide, Inches(10.15), "长期", [f"30 年现价约 US${money(s30['totalSurrenderValue'])}", "长期财富继续累积"])
    add_footer(slide, "时间轴口径：现在先守住妈妈的风险，再沿孩子成长节奏释放教育金。")

    # Slide 15 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "结论", "一张保家庭不失速，一张保未来不落空")
    add_line_chart(
        slide,
        Inches(0.65), Inches(1.55), Inches(6.15), Inches(4.2),
        [1, 5, 10, 20, 30],
        [
            ("孩子储蓄险退保价值", [float(value_map[y]["totalSurrenderValue"]) for y in [1, 5, 10, 20, 30]]),
            ("妈妈重疾险赔付额", [float(next(r["ciBenefit"] for r in ci_rows if int(r["policyYear"]) == y)) for y in [1, 5, 10, 20, 30]]),
        ],
    )
    add_picture(slide, family_img, Inches(7.2), Inches(1.45), Inches(5.45), Inches(2.65))
    add_paragraphs(slide, Inches(7.15), Inches(4.35), Inches(5.1), Inches(1.85), [
        "妈妈这张单保家庭不失速：风险来时先由重疾险承担健康冲击。",
        "孩子这张单保未来不落空：教育金继续按阶段释放，不被家庭风险提前挪用。",
        "这就是家庭防火墙 + 教育金储备池的完整逻辑。"
    ], size=14.5, bullet=True)
    add_footer(slide, "后续继续扩展时，可在此基础上叠加 IUL 形成更完整的家庭资产三层结构。")

    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
