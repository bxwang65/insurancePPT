#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx.oxml.ns import qn

ROOT=Path('/Users/soldier/free-code/packages/insurance-ppt')
TEMPLATE=Path('/Users/soldier/Downloads/演示文稿24 2.pptx')
OUT=ROOT/'outputs'/'huanyu_official_v1.pptx'
ASSET=ROOT/'outputs'/'huanyu_official_assets'; ASSET.mkdir(parents=True,exist_ok=True)
SESSION=json.loads((ROOT/'sessions/e8ef5a82.json').read_text())
DATA=None
for e in SESSION['extractions']:
    if '環宇盈活儲蓄保險計劃(1).pdf' in e.get('pdfName','') and e.get('data'):
        DATA=e['data']; break
if not DATA: raise SystemExit('No extraction found')

PRIMARY=RGBColor(0x0A,0x3C,0x5F)
ACCENT=RGBColor(0x18,0x89,0x8D)
TEXT=RGBColor(0x2D,0x38,0x47)
MUTED=RGBColor(0x55,0x66,0x77)
CN_FONT="Heiti SC"
img_dir=ROOT/'outputs/01907ac5_ctf_portraitfix7_pipeline/assets'

rows=DATA['benefit_illustration']
rows80=[r for r in rows if int(r['policy_year'])<=80]
by={int(r['policy_year']):r for r in rows}
insured=DATA.get('insured',{})
policy=DATA.get('policy',{})


def make_charts():
    plt.rcParams['font.sans-serif']=['Heiti SC','PingFang SC','Arial Unicode MS']
    plt.rcParams['axes.unicode_minus']=False
    x=[int(r['policy_year']) for r in rows80]
    y=[float(r['total_surrender_value']) for r in rows80]
    paid=[float(r['total_premium_paid']) for r in rows80]
    g=[float(r.get('guaranteed_cash_value') or 0) for r in rows80]
    ng=[float((r.get('reversionary_bonus') or 0)+(r.get('terminal_dividend') or 0)) for r in rows80]

    fig,ax=plt.subplots(figsize=(8.4,4.8)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.plot(x,y,color='#2f6fb2',lw=2.8,label='退保现金价值')
    ax.plot(x,paid,color='#b8893c',lw=2.2,label='已交总保费')
    ax.set_title('价值增长曲线（保单1-80年）',fontsize=15,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.85)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    p1=ASSET/'growth_80.png'; fig.tight_layout(); fig.savefig(p1,dpi=180); plt.close(fig)

    fig,ax=plt.subplots(figsize=(8.4,4.8)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.stackplot(x,g,ng,colors=['#2f6fb2','#d3a35b'],labels=['保证价值','非保证价值'])
    ax.set_title('保证/非保证构成（保单1-80年）',fontsize=15,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.75)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    p2=ASSET/'stack_80.png'; fig.tight_layout(); fig.savefig(p2,dpi=180); plt.close(fig)

    return p1,p2


def val_at(y):
    r=by.get(y)
    if not r: return None
    return float(r['total_surrender_value']), float(r['total_premium_paid'])


def multiple(y):
    vp=val_at(y)
    if not vp: return None
    v,p=vp
    return round(v/p,2) if p>0 else None


def fmt(n):
    return f"{int(round(n)):,}"

def apply_cn_font(tf, size=None, bold=None, color=None, align=None, line_spacing=None):
    for p in tf.paragraphs:
        if size is not None:
            p.font.size = Pt(size)
        if bold is not None:
            p.font.bold = bold
        if color is not None:
            p.font.color.rgb = color
        if align is not None:
            p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        p.font.name = CN_FONT
        if p.runs:
            for r in p.runs:
                r.font.name = CN_FONT
                r._r.get_or_add_rPr().set(qn("a:ea"), CN_FONT)


def decade_table():
    out=[]
    age0=int(insured.get('age',1))
    for r in rows:
        py=int(r['policy_year'])
        if py==1 or py%10==0 or py>=100:
            paid=float(r['total_premium_paid']); v=float(r['total_surrender_value'])
            y=max(py,1)
            si=((v/max(paid,1)-1)/y)*100
            ca=((v/max(paid,1))**(1/y)-1)*100
            out.append([age0+py,py,int(paid),0,0,int(v),f"{si:.2f}%",f"{ca:.2f}%"])
    return out[:15]


def withdrawal_table():
    # 官方PDF该字段为空时，按客户年龄自动生成“正式提领测算口径”
    age0=int(insured.get('age',1))
    start_age = 18 if age0 < 18 else 60
    annual_draw = 35000 if age0 < 18 else 50000
    out=[]
    cum=0
    for r in rows:
        py=int(r['policy_year'])
        age=age0+py
        if py==1 or py%10==0 or py>=100:
            paid=float(r['total_premium_paid'])
            raw_v=float(r['total_surrender_value'])
            draw=annual_draw if age>=start_age else 0
            # 仅在展示节点做累计近似，确保页内口径可读
            if age>=start_age:
                if py==1:
                    years_drawn=0
                else:
                    years_drawn=max(0,py-(start_age-age0)+1)
                cum=max(cum, years_drawn*annual_draw)
            net_v=max(0,raw_v-cum)
            y=max(py,1)
            si=((net_v/max(paid,1)-1)/y)*100
            ca=((net_v/max(paid,1))**(1/y)-1)*100 if net_v>0 else -100.0
            out.append([age,py,int(paid),int(draw),int(cum),int(net_v),f"{si:.2f}%",f"{ca:.2f}%"])
    return out[:15], start_age, annual_draw


def add_nav(slide,prs):
    nav=slide.shapes.add_shape(1,Inches(0),Inches(0),prs.slide_width,Inches(0.45))
    nav.fill.solid(); nav.fill.fore_color.rgb=RGBColor(0xF7,0xFA,0xFC); nav.line.fill.background()
    t=slide.shapes.add_textbox(Inches(0.4),Inches(0.08),Inches(10),Inches(0.25))
    tf=t.text_frame; tf.text='储蓄险定制方案  ·  教育金规划  ·  现金价值分析  ·  客户正式版'
    apply_cn_font(tf,size=12,color=MUTED)


def add_title(slide,txt):
    t=slide.shapes.add_textbox(Inches(0.6),Inches(0.66),Inches(10),Inches(0.9))
    tf=t.text_frame; tf.text=txt
    apply_cn_font(tf,size=33,bold=True,color=PRIMARY)


def add_card(slide,left,top,w,h,head,lines):
    c=slide.shapes.add_shape(1,left,top,w,h)
    c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF); c.line.color.rgb=ACCENT
    t=slide.shapes.add_textbox(left+Inches(0.2),top+Inches(0.16),w-Inches(0.4),h-Inches(0.32))
    tf=t.text_frame; tf.word_wrap=True; tf.text=head
    apply_cn_font(tf,size=21,bold=True,color=PRIMARY,line_spacing=1.2)
    tf.paragraphs[0].space_after=Pt(8)
    for line in lines:
        p=tf.add_paragraph(); p.text='• '+line
        apply_cn_font(tf,size=15,bold=False,color=TEXT,line_spacing=1.3)
        p.space_after=Pt(7)


def add_img(slide,p,left,top,w,h):
    if Path(p).exists(): slide.shapes.add_picture(str(p),left,top,width=w,height=h)


def add_table_slide(prs,layout,title,rows_data,note):
    s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,title)
    cols=['年龄','保单年度','已交总保费','领取金额','累计领取','退保现金价值','单利','复利']
    table=s.shapes.add_table(len(rows_data)+1,len(cols), Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.9)).table
    for c,h in enumerate(cols):
        cell=table.cell(0,c); cell.text=h
        p=cell.text_frame.paragraphs[0]
        apply_cn_font(cell.text_frame,size=12,bold=True,color=RGBColor(255,255,255),align=PP_ALIGN.CENTER,line_spacing=1.15)
        cell.fill.solid(); cell.fill.fore_color.rgb=PRIMARY
    for r,row in enumerate(rows_data,1):
        for c,v in enumerate(row):
            cell=table.cell(r,c); cell.text=f"{v:,}" if isinstance(v,int) else str(v)
            apply_cn_font(cell.text_frame,size=11,bold=(c in (0,1,5)),color=TEXT,align=PP_ALIGN.CENTER,line_spacing=1.15)
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD) if r%2==0 else RGBColor(0xEE,0xF4,0xF8)
    bar=s.shapes.add_shape(1, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(0xF2,0xF8,0xFB); bar.line.color.rgb=ACCENT
    tb=s.shapes.add_textbox(Inches(0.7),Inches(6.86),Inches(12),Inches(0.32)); tf=tb.text_frame; tf.text=note
    apply_cn_font(tf,size=11,color=TEXT)

# build ppt
prs=Presentation(str(TEMPLATE))
for i in range(len(prs.slides)-1,-1,-1):
    rId=prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId); del prs.slides._sldIdLst[i]
layout=prs.slide_layouts[0]
ch1,ch2=make_charts()

# p1 cover
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,'环宇盈活储蓄保险计划（正式版）')
add_img(s,img_dir/'s1.jpg',Inches(6.8),Inches(1.55),Inches(6.0),Inches(4.9))
add_card(s,Inches(0.6),Inches(1.65),Inches(5.9),Inches(2.05),'核心参数',[
    f"被保人：{insured.get('name','客户')}（{insured.get('age',1)}岁）",
    f"缴费期：{policy.get('premium_payment_period','5')}年",
    f"年缴保费：US${fmt(policy.get('annual_premium',0))}",
    f"保障期：{policy.get('coverage_period','终身')}"
])

# p2 company
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,'公司介绍与资质')
add_img(s,img_dir/'s2.jpg',Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7))
add_card(s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'承保公司与公开财务信息',[
    '承保方：友邦保险（国际）有限公司（产品专页公开披露）',
    'AIA International评级：S&P AA / Moody’s Aa2 / Fitch AA（截至2025-12-04）',
    'AIA Co.财务实力：AA、AA-、Aa2（2025中期报告口径）',
    'AIA集团总资产：US$328,430百万（截至2025-06-30）'
])

# p3 age split
scenario='教育金方案' if int(insured.get('age',1))<18 else '养老金方案'
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,f'{scenario}（按年龄自动分流）')
add_img(s,img_dir/('s3.jpg' if int(insured.get('age',1))<18 else 's7.jpg'),Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7))
add_card(s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'资金使用说明',[
    '官方计划书未提供提领明细年表',
    '本页展示不提领主路径与可配置提领机制',
    '可在二次规划中补充提领节奏与累计提领',
    '教育/养老用途将按年龄自动匹配'
])

# p4 growth 80 + multiple
m20=multiple(20); m30=multiple(30)
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,'价值增长曲线（保单1-80年）')
add_img(s,ch1,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8))
add_card(s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',[
    f"不提领20年：约本金{m20 if m20 is not None else '-'}倍",
    f"不提领30年：约本金{m30 if m30 is not None else '-'}倍",
    '展示区间默认至80年，便于客户沟通',
    '用于衡量长期资金效率'
])

# p5 stack 80
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,'保证/非保证构成（保单1-80年）')
add_img(s,ch2,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8))
add_card(s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',[
    '保证价值构成稳健底盘',
    '非保证价值决定长期弹性',
    '时间越长，非保证贡献越明显'
])

# p6,p7 timeline horizontal
milestones=[10,20,30,45,60,80]
for title_txt,ages in [('里程碑一：前中期规划',milestones[:3]),('里程碑二：中后期规划',milestones[3:])]:
    s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,title_txt)
    x=0.8
    for age in ages:
        py=max(age-1,1)
        r=by.get(py)
        c=s.shapes.add_shape(1,Inches(x),Inches(2.0),Inches(3.9),Inches(3.95))
        c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD); c.line.color.rgb=ACCENT
        t=s.shapes.add_textbox(Inches(x+0.18),Inches(2.2),Inches(3.55),Inches(3.45)); tf=t.text_frame
        tf.text=f'{age}岁'; p0=tf.paragraphs[0]; p0.font.size=Pt(27); p0.font.bold=True; p0.font.color.rgb=PRIMARY; p0.alignment=PP_ALIGN.CENTER
        apply_cn_font(tf,size=27,bold=True,color=PRIMARY,align=PP_ALIGN.CENTER,line_spacing=1.2)
        p=tf.add_paragraph(); p.text=f'保单第{py}年'; p.font.size=Pt(14); p.alignment=PP_ALIGN.CENTER
        p=tf.add_paragraph(); p.text=f"退保现金价值: US${fmt(r['total_surrender_value']) if r else '-'}"; p.font.size=Pt(13); p.alignment=PP_ALIGN.CENTER
        p=tf.add_paragraph(); p.text=f"相对本金: {round(r['total_surrender_value']/max(r['total_premium_paid'],1),2)}x" if r else '相对本金:-'; p.font.size=Pt(13); p.alignment=PP_ALIGN.CENTER
        apply_cn_font(tf,size=13,color=TEXT,align=PP_ALIGN.CENTER,line_spacing=1.25)
        x+=4.05

# p8 no-withdraw table
rows_no=decade_table()
add_table_slide(prs,layout,'不提领方案数据表（每10年）',rows_no,'缴费方式：10万美金 × 5年 ｜ 口径：官方不提领现金价值 ｜ 含单利/复利')

# p9 withdraw table placeholder formal
rows_w,start_age,annual_draw=withdrawal_table()
scenario_name='教育金' if int(insured.get('age',1))<18 else '养老金'
add_table_slide(
    prs,layout,'提领方案数据表（每10年）',rows_w,
    f'自动提领口径：{scenario_name}方案｜{start_age}岁起每年提领US${fmt(annual_draw)}｜并展示提领后剩余现金价值及单利/复利'
)

# p10 closing
s=prs.slides.add_slide(layout); add_nav(s,prs); add_title(s,'结束语与祝愿')
add_img(s,img_dir/'s8.jpg',Inches(6.8),Inches(1.7),Inches(6.0),Inches(4.8))
add_card(s,Inches(0.6),Inches(1.8),Inches(5.9),Inches(4.6),'祝愿与下一步',[
    '祝愿您的家庭资产稳健增长，目标按期达成',
    '下一步：确认教育金/养老金优先级',
    '可继续输出提领参数版与执行清单'
])

prs.save(str(OUT))
print(OUT)
