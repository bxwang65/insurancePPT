#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

ROOT=Path('/Users/soldier/free-code/packages/insurance-ppt')
TEMPLATE=Path('/Users/soldier/Downloads/演示文稿24 2.pptx')
OUT=ROOT/'outputs'/'huanyu_template24_trial_v1.pptx'
SESSION=json.loads((ROOT/'sessions/e8ef5a82.json').read_text())
DATA=None
for e in SESSION['extractions']:
    if '環宇盈活儲蓄保險計劃(1).pdf' in e.get('pdfName','') and e.get('data'):
        DATA=e['data']; break
if not DATA: raise SystemExit('no data')

PRIMARY=RGBColor(0x0A,0x3C,0x5F); ACCENT=RGBColor(0x18,0x89,0x8D); TEXT=RGBColor(0x2D,0x38,0x47)
img_dir=ROOT/'outputs/01907ac5_ctf_portraitfix7_pipeline/assets'
asset_dir=ROOT/'outputs/huanyu_trial_assets'; asset_dir.mkdir(parents=True,exist_ok=True)

def make_charts_80y():
    rows=[r for r in DATA['benefit_illustration'] if int(r['policy_year'])<=80]
    x=[int(r['policy_year']) for r in rows]
    y=[float(r['total_surrender_value']) for r in rows]
    p=[float(r['total_premium_paid']) for r in rows]
    plt.rcParams['font.sans-serif']=['PingFang SC','Microsoft YaHei','SimHei','Arial Unicode MS']
    plt.rcParams['axes.unicode_minus']=False
    fig,ax=plt.subplots(figsize=(8.2,4.6)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.plot(x,y,color='#2f6fb2',lw=2.6,label='退保价值')
    ax.plot(x,p,color='#b8893c',lw=2.2,label='已缴保费')
    ax.set_title('价值增长曲线（保单1-80年）',fontsize=14,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.8)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    c1=asset_dir/'growth_80y.png'; fig.tight_layout(); fig.savefig(c1,dpi=180); plt.close(fig)
    g=[float(r.get('guaranteed_cash_value') or 0) for r in rows]
    ng=[float((r.get('reversionary_bonus') or 0)+(r.get('terminal_dividend') or 0)) for r in rows]
    fig,ax=plt.subplots(figsize=(8.2,4.6)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.stackplot(x,g,ng,colors=['#2f6fb2','#d3a35b'],labels=['保证价值','非保证价值'])
    ax.set_title('保证/非保证构成（保单1-80年）',fontsize=14,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.6)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    c2=asset_dir/'stack_80y.png'; fig.tight_layout(); fig.savefig(c2,dpi=180); plt.close(fig)
    return c1,c2

def add_nav(sl):
    nav=sl.shapes.add_shape(1,Inches(0),Inches(0),prs.slide_width,Inches(0.45)); nav.fill.solid(); nav.fill.fore_color.rgb=RGBColor(0xF7,0xFA,0xFC); nav.line.fill.background()
    t=sl.shapes.add_textbox(Inches(0.4),Inches(0.08),Inches(8),Inches(0.25)); t.text_frame.text='家庭理财  ·  保险理财  ·  养老保险  ·  理财案例'; t.text_frame.paragraphs[0].font.size=Pt(12)

def title(sl,s):
    tb=sl.shapes.add_textbox(Inches(0.6),Inches(0.7),Inches(9),Inches(0.8)); tb.text_frame.text=s; p=tb.text_frame.paragraphs[0]; p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=PRIMARY

def card(sl,l,t,w,h,tt,lines):
    c=sl.shapes.add_shape(1,l,t,w,h); c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF); c.line.color.rgb=ACCENT
    tb=sl.shapes.add_textbox(l+Inches(0.2),t+Inches(0.18),w-Inches(0.4),h-Inches(0.35)); tf=tb.text_frame; tf.text=tt; tf.paragraphs[0].font.size=Pt(21); tf.paragraphs[0].font.bold=True
    for ln in lines:
        p=tf.add_paragraph(); p.text='• '+ln; p.font.size=Pt(15); p.line_spacing=1.3

def img(sl,p,l,t,w,h):
    if Path(p).exists(): sl.shapes.add_picture(str(p),l,t,width=w,height=h)

def decade(withdraw=False):
    rows=DATA['benefit_illustration']; age=DATA['insured']['age']; out=[]
    for r in rows:
        py=int(r['policy_year'])
        if py==1 or py%10==0 or py>=100:
            paid=float(r['total_premium_paid']); val=float(r['total_surrender_value']); y=max(py,1)
            simp=((val/max(paid,1)-1)/y)*100; cagr=((val/max(paid,1))**(1/y)-1)*100
            out.append([age+py,py,int(paid),0,0,int(val),f'{simp:.2f}%',f'{cagr:.2f}%'])
    return out[:15]

def table(sl,title_text,rows):
    title(sl,title_text)
    cols=['年龄','保单年度','已交总保费','领取金额','累计领取','退保现金价值','单利','复利']
    t=sl.shapes.add_table(len(rows)+1,len(cols),Inches(0.5),Inches(1.7),Inches(12.3),Inches(4.9)).table
    for c,h in enumerate(cols):
        cell=t.cell(0,c); cell.text=h; p=cell.text_frame.paragraphs[0]; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=RGBColor(255,255,255); cell.fill.solid(); cell.fill.fore_color.rgb=PRIMARY
    for i,row in enumerate(rows,1):
        for c,v in enumerate(row):
            cell=t.cell(i,c); cell.text=f'{v:,}' if isinstance(v,int) else str(v); p=cell.text_frame.paragraphs[0]; p.font.size=Pt(11); p.alignment=PP_ALIGN.CENTER

prs=Presentation(str(TEMPLATE))
for i in range(len(prs.slides)-1,-1,-1):
    rId=prs.slides._sldIdLst[i].rId; prs.part.drop_rel(rId); del prs.slides._sldIdLst[i]
layout=prs.slide_layouts[0]
c1,c2=make_charts_80y()
#1
s=prs.slides.add_slide(layout); add_nav(s); title(s,'环宇盈活储蓄险方案（试跑）'); img(s,img_dir/'s1.jpg',Inches(6.8),Inches(1.55),Inches(6.0),Inches(4.9))
card(s,Inches(0.6),Inches(1.65),Inches(5.9),Inches(2.0),'核心参数',[f"缴费期：{DATA['policy']['premium_payment_period']}年",f"年缴保费：US${int(DATA['policy']['annual_premium']):,}","观察区间：保单1-80年"])
#2
s=prs.slides.add_slide(layout); add_nav(s); title(s,'公司介绍与资质'); img(s,img_dir/'s2.jpg',Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7)); card(s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'机构信息',['可接入你指定公司资料','可加入评级与偿付率','当前为模板试跑'])
#3
s=prs.slides.add_slide(layout); add_nav(s); title(s,'教育金阶段（示例）'); img(s,img_dir/'s3.jpg',Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7)); card(s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'说明',['该产品当前提领明细未在缓存中','先展示无提领价值路径','后续可补提领版计划书做细化'])
#4
s=prs.slides.add_slide(layout); add_nav(s); title(s,'价值增长曲线（保单1-80年）'); img(s,c1,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8)); card(s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',[f"20年约本金{(next(r for r in DATA['benefit_illustration'] if r['policy_year']==20)['total_surrender_value']/500000):.2f}倍",f"30年约本金{(next(r for r in DATA['benefit_illustration'] if r['policy_year']==30)['total_surrender_value']/500000):.2f}倍",'观察长期增值效率'])
#5
s=prs.slides.add_slide(layout); add_nav(s); title(s,'保证/非保证构成（保单1-80年）'); img(s,c2,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8)); card(s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',['保证价值为底盘','非保证贡献长期弹性'])
#6#7 milestones synthetic
for nm in ['里程碑一：教育与家庭现金流','里程碑二：中后期与养老金']:
    s=prs.slides.add_slide(layout); add_nav(s); title(s,nm)
    x=0.7
    for a in [10,20,30,60]:
        py=a-1; row=next((r for r in DATA['benefit_illustration'] if int(r['policy_year'])==py),None)
        if not row: continue
        c=s.shapes.add_shape(1,Inches(x),Inches(2.0),Inches(2.45),Inches(3.9)); c.fill.solid(); c.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD); c.line.color.rgb=ACCENT
        tb=s.shapes.add_textbox(Inches(x+0.15),Inches(2.15),Inches(2.15),Inches(3.5)); tf=tb.text_frame; tf.text=f'{a}岁'; tf.paragraphs[0].font.size=Pt(26); tf.paragraphs[0].font.bold=True
        for ln in [f'保单第{py}年',f"退保值: US${int(row['total_surrender_value']):,}"]:
            p=tf.add_paragraph(); p.text=ln; p.font.size=Pt(13)
        x+=2.58
#8#9 tables
s=prs.slides.add_slide(layout); add_nav(s); table(s,'提领方案数据表（每10年）',decade(True))
s=prs.slides.add_slide(layout); add_nav(s); table(s,'不提领方案数据表（每10年）',decade(False))
#10
s=prs.slides.add_slide(layout); add_nav(s); title(s,'执行建议与下一步'); img(s,img_dir/'s8.jpg',Inches(6.8),Inches(1.7),Inches(6.0),Inches(4.8)); card(s,Inches(0.6),Inches(1.8),Inches(5.9),Inches(4.6),'执行框架',['确认目标用途','确定复盘周期','按年度跟踪资产表现'])

prs.save(str(OUT)); print(OUT)
