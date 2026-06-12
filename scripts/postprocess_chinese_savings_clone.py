#!/usr/bin/env python3
from __future__ import annotations

import argparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu


BLACK = RGBColor(0x00, 0x00, 0x00)
TOP_BOUND = Emu(2.0 * 914400)


def recolor_top_title(slide) -> bool:
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(par.text.strip() for par in shape.text_frame.paragraphs).strip()
        if not text or len(text) > 60:
            continue
        top = getattr(shape, "top", Emu(999999999))
        if top > TOP_BOUND:
            continue
        candidates.append((top, shape))
    if not candidates:
        return False
    _, shape = sorted(candidates, key=lambda item: item[0])[0]
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font is not None:
            paragraph.font.color.rgb = BLACK
        for run in paragraph.runs:
            run.font.color.rgb = BLACK
    return True


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--slides", default="5,6")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    slide_numbers = [int(part.strip()) for part in args.slides.split(",") if part.strip()]
    for slide_number in slide_numbers:
      if 1 <= slide_number <= len(prs.slides):
        recolor_top_title(prs.slides[slide_number - 1])
    prs.save(args.pptx)


if __name__ == "__main__":
    main()
