#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


FONT_CN = "Microsoft YaHei"
BG = RGBColor(248, 245, 239)
PANEL = RGBColor(255, 255, 255)
PRIMARY = RGBColor(38, 38, 38)
ACCENT = RGBColor(206, 138, 44)
ACCENT_LIGHT = RGBColor(240, 230, 208)
MUTED = RGBColor(96, 96, 96)
LINE = RGBColor(223, 214, 195)
GOOD = RGBColor(24, 108, 79)


def money(value: float | int) -> str:
    return f"{round(float(value)):,}"


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=PRIMARY, align=PP_ALIGN.LEFT, font=FONT_CN,
                line_spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, left, top, width, height, lines, size=16, color=PRIMARY,
                   bullet=False, line_spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if bullet:
            p.text = f"• {line}"
        else:
            p.text = line
        if p.runs:
            run = p.runs[0]
            run.font.name = FONT_CN
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, value, subtitle="", accent=False):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_LIGHT if accent else PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.16), top + Inches(0.14), width - Inches(0.32), Inches(0.28),
                title, size=13, bold=True, color=MUTED, line_spacing=1.2)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.42), width - Inches(0.32), Inches(0.42),
                value, size=21, bold=True, color=ACCENT if accent else PRIMARY, line_spacing=1.2)
    if subtitle:
        add_textbox(slide, left + Inches(0.16), top + Inches(0.88), width - Inches(0.32), Inches(0.28),
                    subtitle, size=11, color=MUTED, line_spacing=1.2)


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(5.8), Inches(0.5), title, size=24, bold=True, line_spacing=1.2)
    if subtitle:
      add_textbox(slide, Inches(0.6), Inches(0.82), Inches(5.8), Inches(0.3), subtitle, size=11, color=MUTED, line_spacing=1.2)


def add_fact_card(slide, left, top, width, height, label, value):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.14), top + Inches(0.14), width - Inches(0.28), Inches(0.22),
                label, size=12, bold=True, color=MUTED, line_spacing=1.2)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.42), width - Inches(0.28), height - Inches(0.56),
                value, size=14, color=PRIMARY, line_spacing=1.5)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.65), Inches(7.0), Inches(12.0), Inches(0.24), text, size=10, color=MUTED, line_spacing=1.2)


def build_chart(slide, rows):
    chart_data = CategoryChartData()
    categories = [str(r["policyYear"]) for r in rows]
    chart_data.categories = categories
    chart_data.add_series("退保发还金额", [float(r["totalSurrenderValue"]) for r in rows])
    chart_data.add_series("严重疾病赔偿", [float(r["ciBenefit"]) for r in rows])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.7), Inches(1.6), Inches(6.2), Inches(3.6), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    for idx, series in enumerate(chart.series):
        series.format.line.width = Pt(2.2)
        series.format.line.color.rgb = ACCENT if idx == 0 else GOOD


def slide_cover(prs, ci, company_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, Inches(0.75), Inches(0.7), Inches(6.8), Inches(0.6),
                f"{ci['insured']['name']} 家庭重疾保障方案", size=26, bold=True, line_spacing=1.2)
    add_textbox(slide, Inches(0.75), Inches(1.45), Inches(7.8), Inches(0.5),
                f"{company_name} · {ci['productName']}", size=18, color=MUTED, line_spacing=1.2)
    add_textbox(slide, Inches(0.75), Inches(2.25), Inches(5.5), Inches(0.42),
                "第一层：健康风险防火墙", size=20, bold=True, color=ACCENT, line_spacing=1.2)
    lines = [
        f"基础保额 US${money(ci['policy']['baseSumInsured'])}",
        f"前 {ci['policy']['upgradeBenefitYears']} 年升级保障 US${money(ci['policy']['upgradeBenefitAmount'])}",
        f"年缴保费 US${money(ci['policy']['annualPremium'])}，缴费 {ci['policy']['payYears']} 年",
        f"重疾 {ci['coverageSummary']['majorCiCount']} 项，早期危疾 {ci['coverageSummary']['earlyCiCount']} 项",
    ]
    add_paragraphs(slide, Inches(0.78), Inches(2.95), Inches(5.6), Inches(2.2), lines, size=16, bullet=True)
    add_card(slide, Inches(7.55), Inches(1.5), Inches(2.2), Inches(1.25), "总保费", f"US${money(ci['policy']['totalPremium'])}", "10 年供", accent=True)
    add_card(slide, Inches(10.0), Inches(1.5), Inches(2.15), Inches(1.25), "首十年总保障", f"US${money(ci['policy']['sumInsured'] + ci['policy']['upgradeBenefitAmount'])}", "基础 + 升级")
    add_card(slide, Inches(7.55), Inches(3.0), Inches(2.2), Inches(1.25), "现金价值", "保单可退保", "兼具一定储蓄属性")
    add_card(slide, Inches(10.0), Inches(3.0), Inches(2.15), Inches(1.25), "组合定位", "家庭防火墙", "可叠加储蓄险")
    add_footer(slide, "核心口径：年缴保费、缴费年期、基础保额、升级保障、重疾责任、多重赔付、现金价值。")


def slide_company(prs, company):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "公司介绍", "以已确认的公司事实手册作为正式版展示口径")
    add_textbox(slide, Inches(0.75), Inches(1.35), Inches(5.0), Inches(1.6),
                company["companyIntro"], size=17, line_spacing=1.25)
    facts = company.get("companyFacts", [])[:4]
    positions = [
        (Inches(6.2), Inches(1.35)),
        (Inches(9.2), Inches(1.35)),
        (Inches(6.2), Inches(3.25)),
        (Inches(9.2), Inches(3.25)),
    ]
    for (left, top), fact in zip(positions, facts):
        add_fact_card(slide, left, top, Inches(2.55), Inches(1.55), fact["label"], fact["value"])
    add_paragraphs(slide, Inches(0.8), Inches(4.25), Inches(11.4), Inches(1.6), [
        "重疾险沟通重点：公司是否长期经营、偿付与评级是否稳健、理赔与服务体系是否成熟。",
        "客户听得懂的表达应聚焦：背景、实力、业务定位、评级，不展示资料来源文件名。"
    ], size=13, color=MUTED)


def slide_core(prs, ci):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "保单核心参数", "先把保费、保额、责任边界讲清楚")
    add_card(slide, Inches(0.8), Inches(1.45), Inches(3.35), Inches(1.2), "基础保额", f"US${money(ci['policy']['baseSumInsured'])}")
    add_card(slide, Inches(4.35), Inches(1.45), Inches(3.35), Inches(1.2), "升级保障", f"US${money(ci['policy']['upgradeBenefitAmount'])}", f"前 {ci['policy']['upgradeBenefitYears']} 年")
    add_card(slide, Inches(7.9), Inches(1.45), Inches(4.0), Inches(1.2), "首十年总保障", f"US${money(ci['policy']['sumInsured'] + ci['policy']['upgradeBenefitAmount'])}")
    add_card(slide, Inches(0.8), Inches(2.85), Inches(2.9), Inches(1.1), "年缴保费", f"US${money(ci['policy']['annualPremium'])}", "含征费约 US$7,015")
    add_card(slide, Inches(3.95), Inches(2.85), Inches(2.9), Inches(1.1), "缴费年期", f"{ci['policy']['payYears']} 年", "总保费 US$70,080")
    add_paragraphs(slide, Inches(0.82), Inches(4.15), Inches(5.8), Inches(2.05), [
        f"重大疾病：{ci['coverageSummary']['majorCiCount']} 项",
        f"早期危疾：{ci['coverageSummary']['earlyCiCount']} 项",
        "ICU 深切治疗保障：2 个层级",
        "配偶身故豁免保费 + 免付保费附加契约",
    ], size=16, bullet=True)
    add_paragraphs(slide, Inches(6.95), Inches(4.1), Inches(5.0), Inches(2.15), [
        "这张单的主定位不是财富增值，而是家庭重大健康风险的现金缓冲器。",
        "前 10 年升级保障把核心保额从 10 万提升至 13.5 万，是销售沟通里的第一个重点。",
        "后续若叠加储蓄险，重疾险负责防火墙，储蓄险负责现金流。"
    ], size=15, color=PRIMARY)


def slide_rules(prs, ci):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "保障责任结构", "按客户能理解的顺序讲：保什么、赔几次、附加什么")
    add_fact_card(slide, Inches(0.8), Inches(1.4), Inches(3.8), Inches(1.65), "危疾与早期危疾",
                  f"{ci['coverageItems'][0]['name']}；{ci['coverageItems'][1]['name']}")
    add_fact_card(slide, Inches(4.9), Inches(1.4), Inches(3.8), Inches(1.65), "升级保障",
                  f"前 {ci['policy']['upgradeBenefitYears']} 年额外 US${money(ci['policy']['upgradeBenefitAmount'])}")
    add_fact_card(slide, Inches(9.0), Inches(1.4), Inches(3.0), Inches(1.65), "豁免责任",
                  "免付保费附加契约 + 配偶身故豁免")
    icu_lines = [f"{item['level']}：{item.get('payoutPercentage','')}，等待 {item.get('waitingPeriodHours','')} 小时" for item in ci["icuBenefitRules"]]
    multi_lines = [f"{item['condition']}：最多 {item['claimCount']} 次，等待期 {item.get('waitingPeriod','')}" for item in ci["multiClaimRules"]]
    add_paragraphs(slide, Inches(0.85), Inches(3.45), Inches(5.5), Inches(2.6), icu_lines, size=15, bullet=True)
    add_paragraphs(slide, Inches(6.35), Inches(3.45), Inches(6.0), Inches(2.6), multi_lines, size=15, bullet=True)


def slide_cash_value(prs, ci):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "现金价值与保障路径", "这是保障型重疾险，现金价值是辅助，不是主卖点")
    rows = [r for r in ci["benefitRows"] if r["policyYear"] in (1, 5, 10, 20, 30, 65)]
    build_chart(slide, rows)
    add_paragraphs(slide, Inches(7.15), Inches(1.65), Inches(5.0), Inches(3.5), [
        f"第 10 年退保发还金额约 US${money(19976)}",
        f"第 20 年约 US${money(48002)}",
        f"第 30 年约 US${money(146065)}",
        "解读方式：重疾险的现金价值只做辅助决策，不应替代保障责任本身。",
        "客户沟通应先讲保什么、赔多少、赔几次，再讲退保价值。"
    ], size=15, bullet=True)


def slide_firewall(prs, ci):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "家庭防火墙建议", "重疾险与储蓄险组合时，各自解决不同问题")
    add_fact_card(slide, Inches(0.8), Inches(1.55), Inches(3.8), Inches(1.55), "第一层：重疾险", "解决大病、收入中断、治疗支出")
    add_fact_card(slide, Inches(4.9), Inches(1.55), Inches(3.8), Inches(1.55), "第二层：储蓄险", "解决教育金、养老金、长期现金流")
    add_fact_card(slide, Inches(9.0), Inches(1.55), Inches(3.2), Inches(1.55), "组合输出", "家庭防火墙方案")
    add_paragraphs(slide, Inches(0.85), Inches(3.55), Inches(11.4), Inches(2.5), [
        "爱伴航适合作为家庭健康风险底仓：先把重大疾病、早期危疾、ICU 和豁免责任锁住。",
        "若客户同时追求教育金或养老金，再叠加储蓄险，由储蓄险承担未来现金流目标。",
        "这就是后期组合方案里的主线：重疾险保风险，储蓄险保未来。"
    ], size=16, bullet=True)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--normalized", required=True)
    parser.add_argument("--company-context", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    ci = json.loads(Path(args.normalized).read_text(encoding="utf-8"))
    company = json.loads(Path(args.company_context).read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs, ci, company["companyName"])
    slide_company(prs, company)
    slide_core(prs, ci)
    slide_rules(prs, ci)
    slide_cash_value(prs, ci)
    slide_firewall(prs, ci)

    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
