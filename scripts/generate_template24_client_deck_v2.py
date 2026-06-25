#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json
import re
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

ROOT=Path('/Users/soldier/free-code/packages/insurance-ppt')
TEMPLATE=Path('/Users/soldier/Downloads/演示文稿24 2.pptx')
OUT=ROOT/'outputs'/'template24_client_ready_v4.pptx'
BASE=json.loads((ROOT/'outputs/1429c316_custom/base_savings.json').read_text())
ANA=json.loads((ROOT/'outputs/1429c316_custom/withdrawal_analysis.json').read_text())

img_dir=ROOT/'outputs/01907ac5_ctf_portraitfix7_pipeline/assets'
chart_dir=ROOT/'outputs/1429c316_custom'
v2_chart_dir=ROOT/'outputs'/'template24_v2_assets'
v2_chart_dir.mkdir(parents=True,exist_ok=True)

PRIMARY=RGBColor(0x0A,0x3C,0x5F)
ACCENT=RGBColor(0x18,0x89,0x8D)
TEXT=RGBColor(0x2D,0x38,0x47)


def make_chart_80y():
    base={int(r['policy_year']):r for r in BASE['benefit_illustration']}
    wr=[r for r in ANA['withdrawal_rows'] if int(r['policy_year'])<=80]
    x=[int(r['policy_year']) for r in wr]
    y_after=[float(r['surrender_value_after']) for r in wr]
    y_base=[float(base.get(py,{}).get('total_surrender_value',0)) for py in x]

    plt.rcParams['font.sans-serif']=['PingFang SC','Microsoft YaHei','SimHei','Arial Unicode MS']
    plt.rcParams['axes.unicode_minus']=False

    # growth to 80
    fig,ax=plt.subplots(figsize=(8.2,4.6)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.plot(x,y_base,color='#2f6fb2',lw=2.6,label='不提领退保价值')
    ax.plot(x,y_after,color='#b8893c',lw=2.6,label='提领后退保价值')
    ax.set_title('提领前后价值曲线（保单1-80年）',fontsize=14,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.8)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    p1=v2_chart_dir/'growth_80y.png'; fig.tight_layout(); fig.savefig(p1,dpi=180); plt.close(fig)

    # stack to 80
    b80=[r for r in BASE['benefit_illustration'] if int(r['policy_year'])<=80]
    x2=[int(r['policy_year']) for r in b80]
    g=[float(r.get('guaranteed_cash_value') or 0) for r in b80]
    ng=[float((r.get('reversionary_bonus') or 0)+(r.get('terminal_dividend') or 0)) for r in b80]
    fig,ax=plt.subplots(figsize=(8.2,4.6)); fig.patch.set_facecolor('#fff'); ax.set_facecolor('#f8fbff')
    ax.stackplot(x2,g,ng,colors=['#2f6fb2','#d3a35b'],labels=['保证价值','非保证价值'])
    ax.set_title('保证/非保证构成（保单1-80年）',fontsize=14,fontweight='bold',color='#17324d')
    ax.tick_params(colors='#35506a'); ax.grid(color='#d7e3ef',alpha=.6)
    for s in ax.spines.values(): s.set_color('#d2dfec')
    ax.legend(facecolor='#fff',edgecolor='#d2dfec')
    p2=v2_chart_dir/'stack_80y.png'; fig.tight_layout(); fig.savefig(p2,dpi=180); plt.close(fig)
    return p1,p2


def decade_rows(withdraw=False):
    insured_age=BASE['insured']['age']
    base=BASE['benefit_illustration']
    by={int(r['policy_year']):r for r in base}
    wr=ANA['withdrawal_rows'] if withdraw else []
    src=wr if withdraw else [{'age':insured_age+int(r['policy_year']),'policy_year':r['policy_year'],'total_premium_paid':r['total_premium_paid'],'annual_withdrawal':0,'cumulative_withdrawal':0,'surrender_value_after':r['total_surrender_value']} for r in base]
    out=[]
    for r in src:
        py=int(r['policy_year'])
        if py==1 or py%10==0 or py>=120:
            b=by.get(py,{})
            paid=float(b.get('total_premium_paid',r.get('total_premium_paid',0)) or 0)
            val=float(r.get('surrender_value_after',0) or 0)
            y=max(py,1)
            simp=((val/max(paid,1)-1)/y)*100
            cagr=((val/max(paid,1))**(1/y)-1)*100
            out.append([int(r.get('age',insured_age+py)),py,int(paid),int(r.get('annual_withdrawal',0)),int(r.get('cumulative_withdrawal',0)),int(val),f"{simp:.2f}%",f"{cagr:.2f}%"])
    return out[:15]


def milestones():
    key=[6,10,20,30,45,60,80]
    by={int(r['age']):r for r in ANA['withdrawal_rows']}
    rows=[]
    for a in key:
        r=by.get(a)
        if not r: continue
        rows.append((a,int(r['policy_year']),int(r['annual_withdrawal']),int(r['cumulative_withdrawal']),int(r['surrender_value_after'])))
    return rows

def first_withdraw_year():
    for r in ANA['withdrawal_rows']:
        if float(r.get('annual_withdrawal',0) or 0) > 0:
            return int(r['policy_year']), int(r.get('age',0))
    return None, None

def no_withdraw_multiple(policy_year:int):
    row = next((r for r in BASE['benefit_illustration'] if int(r['policy_year'])==policy_year), None)
    if not row: return None
    paid = float(row.get('total_premium_paid') or 0)
    val = float(row.get('total_surrender_value') or 0)
    if paid <= 0: return None
    return round(val/paid,2)

prs=Presentation(str(TEMPLATE))
for i in range(len(prs.slides)-1,-1,-1):
    rId=prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId); del prs.slides._sldIdLst[i]
layout=prs.slide_layouts[0]

def add_top_nav(slide):
    nav=slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.45))
    nav.fill.solid(); nav.fill.fore_color.rgb=RGBColor(0xF7,0xFA,0xFC); nav.line.fill.background()
    t=slide.shapes.add_textbox(Inches(0.4),Inches(0.08),Inches(8),Inches(0.25)); tf=t.text_frame; tf.text='家庭理财  ·  保险理财  ·  养老保险  ·  理财案例'
    tf.paragraphs[0].font.size=Pt(12); tf.paragraphs[0].font.color.rgb=RGBColor(0x55,0x66,0x77)

def title(slide,s):
    t=slide.shapes.add_textbox(Inches(0.6),Inches(0.7),Inches(9),Inches(0.8)); tf=t.text_frame; tf.text=s
    p=tf.paragraphs[0]; p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=PRIMARY

def add_text_card(slide,left,top,w,h,title_txt,bul):
    card=slide.shapes.add_shape(1,left,top,w,h)
    card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xFF,0xFF,0xFF); card.line.color.rgb=ACCENT
    tb=slide.shapes.add_textbox(left+Inches(0.2),top+Inches(0.18),w-Inches(0.4),h-Inches(0.35))
    tf=tb.text_frame; tf.word_wrap=True; tf.text=title_txt
    p0=tf.paragraphs[0]; p0.font.size=Pt(21); p0.font.bold=True; p0.font.color.rgb=PRIMARY; p0.line_spacing=1.2; p0.space_after=Pt(8)
    for it in bul:
        p=tf.add_paragraph(); p.text='• '+it; p.font.size=Pt(15); p.font.color.rgb=TEXT; p.space_after=Pt(7); p.line_spacing=1.3
        if re.search(r'(US\\$\\d[\\d,]*|\\d+岁|\\d+年|\\d+%)', it):
            p.font.bold=True

def add_image(slide,path,left,top,w,h):
    if Path(path).exists(): slide.shapes.add_picture(str(path),left,top,width=w,height=h)

def add_table_slide(title_text, rows):
    s=prs.slides.add_slide(layout); add_top_nav(s); title(s,title_text)
    cols=['年龄','保单年度','已交总保费','领取金额','累计领取','退保现金价值','单利','复利']
    table=s.shapes.add_table(len(rows)+1,len(cols), Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.9)).table
    for c,h in enumerate(cols):
        cell=table.cell(0,c); cell.text=h; p=cell.text_frame.paragraphs[0]; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb=PRIMARY
        p.line_spacing=1.1
    for r,row in enumerate(rows,1):
        for c,v in enumerate(row):
            cell=table.cell(r,c); cell.text=f"{v:,}" if isinstance(v,int) else str(v)
            p=cell.text_frame.paragraphs[0]; p.font.size=Pt(11); p.font.color.rgb=TEXT; p.alignment=PP_ALIGN.CENTER; p.line_spacing=1.15
            if c in (0,1,5): p.font.bold=True
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD) if r%2==0 else RGBColor(0xEE,0xF4,0xF8)

# slide1 cover
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'储蓄险家庭资产配置方案')
add_image(s,img_dir/'s1.jpg',Inches(6.8),Inches(1.55),Inches(6.0),Inches(4.9))
add_text_card(s,Inches(0.6),Inches(1.65),Inches(5.9),Inches(2.0),'核心参数',['缴费期：5年','年缴保费：US$100,000','第7年价值：US$514,498','第20年价值：US$1,366,345'])

# slide2 company - fixed overlap
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'公司介绍与资质')
add_image(s,img_dir/'s2.jpg',Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7))
add_text_card(s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'周大福人寿（CTF Life）',['Fitch 财务实力评级：A-','Moody\'s 财务实力评级：A3','香港RBC偿付能力充足率：282%','定位：长期保障+家庭财富传承'])

# slide3 education - fixed overlap
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'教育金阶段（18-21岁）')
add_image(s,img_dir/'s3.jpg',Inches(0.6),Inches(1.75),Inches(5.7),Inches(4.7))
fw_py, fw_age = first_withdraw_year()
age18 = next((r for r in ANA['withdrawal_rows'] if int(r.get('age',0))==18), None)
age21 = next((r for r in ANA['withdrawal_rows'] if int(r.get('age',0))==21), None)
add_text_card(
    s,Inches(6.55),Inches(1.75),Inches(6.2),Inches(4.7),'资金使用场景',
    [
        f'开始提领：保单第{fw_py or "-"}年（约{fw_age or "-"}岁）',
        f"18岁累计可提领：US${int(age18.get('cumulative_withdrawal',0)):,}" if age18 else '18岁累计可提领：-',
        f"21岁累计可提领：US${int(age21.get('cumulative_withdrawal',0)):,}" if age21 else '21岁累计可提领：-',
        '适配：学费、住宿、研究支出'
    ]
)

# 4,5 charts only to year 80
c1,c2=make_chart_80y()
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'价值增长曲线（保单1-80年）')
mul20 = no_withdraw_multiple(20)
mul30 = no_withdraw_multiple(30)
add_image(s,c1,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8))
add_text_card(
    s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',
    [
        '先看回本区间，再看关键节点',
        f'不提领20年：约本金{mul20}倍' if mul20 else '不提领20年：-',
        f'不提领30年：约本金{mul30}倍' if mul30 else '不提领30年：-',
        '提领后仍保留长期资产'
    ]
)

s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'保证/非保证构成（保单1-80年）')
add_image(s,c2,Inches(0.6),Inches(1.75),Inches(7.2),Inches(4.8)); add_text_card(s,Inches(8.05),Inches(1.85),Inches(4.8),Inches(4.6),'图表解读',['保证部分是底盘','非保证部分提供弹性','时间越长弹性贡献越大'])

# 6,7 horizontal milestones
ms=milestones()
for idx,(tname,focus) in enumerate([('里程碑一：教育与家庭现金流',ms[:4]),('里程碑二：中后期与养老金',ms[4:])]):
    s=prs.slides.add_slide(layout); add_top_nav(s); title(s,tname)
    x=0.7
    for a,py,ann,cum,sv in focus:
        card=s.shapes.add_shape(1,Inches(x),Inches(2.0),Inches(2.45),Inches(3.9))
        card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD); card.line.color.rgb=ACCENT
        tb=s.shapes.add_textbox(Inches(x+0.15),Inches(2.15),Inches(2.15),Inches(3.5))
        tf=tb.text_frame; tf.text=f'{a}岁'
        p=tf.paragraphs[0]; p.font.size=Pt(26); p.font.bold=True; p.font.color.rgb=PRIMARY; p.alignment=PP_ALIGN.CENTER; p.space_after=Pt(6)
        for line in [f'保单第{py}年',f'年提领: US${ann:,}',f'累计提领: US${cum:,}',f'退保值: US${sv:,}']:
            q=tf.add_paragraph(); q.text=line; q.font.size=Pt(13); q.font.color.rgb=TEXT; q.alignment=PP_ALIGN.CENTER; q.line_spacing=1.25; q.space_after=Pt(4)
            if 'US$' in line: q.font.bold=True
        x+=2.58
    # timeline line
    l=s.shapes.add_shape(1,Inches(0.75),Inches(6.2),Inches(11.8),Inches(0.08)); l.fill.solid(); l.fill.fore_color.rgb=ACCENT; l.line.fill.background()

# 8,9 tables keep mandatory
add_table_slide('提领方案数据表（每10年）',decade_rows(True))
add_table_slide('不提领方案数据表（每10年）',decade_rows(False))

# 10 close
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'执行建议与下一步')
add_image(s,img_dir/'s8.jpg',Inches(6.8),Inches(1.7),Inches(6.0),Inches(4.8))
add_text_card(s,Inches(0.6),Inches(1.8),Inches(5.9),Inches(4.6),'执行框架',['阶段1：教育金优先','阶段2：现金流补充','阶段3：养老金落地','每3年复盘提领与现价'])

prs.save(str(OUT))
print(OUT)
