#!/usr/bin/env python3
"""
极速 PPTX 渲染器 (python-pptx 原生, 无模板克隆)
目标: < 1s 出片, 76KB 输出
"""
import argparse
import json
import os
import sys
from typing import Dict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 焦糖棕 / 深海蓝 双主题
THEMES = {
    "caramel": {
        "primary": "5C3F32", "secondary": "8D624F", "accent": "D4B878",
        "bg": "EFE4DF", "card": "FEFCF7", "text": "332825", "text_light": "FBF9F6",
    },
    "deepblue": {
        "primary": "0A1628", "secondary": "1E3A5F", "accent": "C8963E",
        "bg": "F5EFE3", "card": "FFFFFF", "text": "0A1628", "text_light": "FFFFFF",
    },
    "chinese": {
        "primary": "5C3F32", "secondary": "8D624F", "accent": "C8963E",
        "bg": "F5EFE3", "card": "FEFCF7", "text": "332825", "text_light": "FBF9F6",
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _c(hex_str: str) -> RGBColor:
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x, y, w, h, fill=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = _c(fill)
    else:
        shp.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color="000000", bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(36000); tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000); tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    lines = text.split('\n') if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = _c(color)
        run.font.bold = bold
    return tb


def fill_bg(slide, color):
    bg = add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=color)
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def fmt_num(v):
    if v is None or v == 0:
        return "0"
    if abs(v) >= 1000:
        return f"{int(v):,}"
    return f"{v:,.2f}"


def slide_cover(prs, deck, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["primary"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    customer = deck["customer"]["name"]
    products = deck["products"]
    p = products[0]
    product_name = p["productName"].replace("\n", "").replace(" ", "")
    add_text(s, Inches(0.5), Inches(2), Inches(12), Inches(1),
             f"{customer} 专属方案", size=36, color=theme["text_light"], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3), Inches(12), Inches(0.6),
             product_name, size=20, color=theme["accent"], align=PP_ALIGN.CENTER)
    policy = p["policy"]
    add_text(s, Inches(0.5), Inches(4), Inches(12), Inches(0.5),
             f"{policy['currency']} {fmt_num(policy['annualPremium'])} × {policy['payYears']}年 = 总保费 {fmt_num(policy['contractualTotalPremium'])}",
             size=14, color=theme["text_light"], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5), Inches(12), Inches(0.5),
             f"保障年期: {policy['coveragePeriod']}",
             size=12, color=theme["accent"], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12), Inches(0.3),
             deck["company"]["displayName"], size=11, color=theme["text_light"], align=PP_ALIGN.CENTER)


def slide_company(prs, deck, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["bg"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             "公司介绍", size=24, color=theme["primary"], bold=True)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(0.6), Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
             deck["company"]["displayName"], size=18, color=theme["secondary"], bold=True)
    if deck["company"].get("rating"):
        add_text(s, Inches(0.5), Inches(2.1), Inches(12), Inches(0.4),
                 f"评级: {deck['company']['rating']}", size=12, color=theme["primary"])
    y = 2.7
    for ev in deck["company"].get("evidence", [])[:4]:
        add_text(s, Inches(0.5), Inches(y), Inches(12), Inches(0.5),
                 f"• {ev['text'][:120]}{'...' if len(ev['text']) > 120 else ''}", size=11, color=theme["primary"])
        y += 0.7
    # 销售洞察
    p = deck["products"][0]
    if p.get("salesInsights"):
        si = p["salesInsights"]
        if si.get("targetCustomer"):
            add_text(s, Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
                     f"目标客户: {si['targetCustomer']}", size=10, color=theme["secondary"])


def slide_growth_chart(prs, deck, theme):
    """价值增长曲线 (bar)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["bg"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             "价值增长曲线", size=24, color=theme["primary"], bold=True)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(0.6), Inches(0.04), fill=theme["accent"])
    p = deck["products"][0]
    rows = p["benefitRows"]
    # 5/10/15/20/25/30 关键年
    milestones = [r for r in rows if r["policyYear"] in [5, 10, 15, 20, 25, 30]]
    if not milestones:
        milestones = rows[::max(1, len(rows)//6)][:6]
    total_premium = p["policy"]["contractualTotalPremium"]
    # 简易柱状图
    chart_x = Inches(0.8); chart_y = Inches(2.2)
    bar_w = Inches(1.6); bar_gap = Inches(0.5)
    max_v = max((m["totalSurrenderValue"] for m in milestones), default=1) or 1
    chart_h_in = 4.0
    add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.4),
             f"退保总值（保证+非保证）/ 总保费 {fmt_num(total_premium)}", size=11, color=theme["secondary"])
    for i, m in enumerate(milestones):
        x = chart_x + i * (bar_w + bar_gap)
        h = Emu(int((m["totalSurrenderValue"] / max_v) * Inches(chart_h_in)))
        bar = add_rect(s, x, chart_y + Inches(chart_h_in) - h, bar_w, h, fill=theme["secondary"])
        # 数值标签
        add_text(s, x, chart_y + Inches(chart_h_in) - h - Inches(0.4), bar_w, Inches(0.3),
                 fmt_num(m["totalSurrenderValue"]), size=10, color=theme["primary"], bold=True, align=PP_ALIGN.CENTER)
        # 年份标签
        mult = m["totalSurrenderValue"] / total_premium if total_premium else 0
        add_text(s, x, chart_y + Inches(chart_h_in) + Inches(0.05), bar_w, Inches(0.3),
                 f"Y{m['policyYear']}", size=11, color=theme["primary"], bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, chart_y + Inches(chart_h_in) + Inches(0.35), bar_w, Inches(0.3),
                 f"{mult:.2f}x", size=10, color=theme["accent"], align=PP_ALIGN.CENTER)


def slide_withdrawal_table(prs, deck, theme):
    """提领方案表 (10/20/30)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["bg"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             "提领方案数据", size=24, color=theme["primary"], bold=True)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(0.6), Inches(0.04), fill=theme["accent"])
    p = deck["products"][0]
    rows = p.get("withdrawalRows", [])
    if not rows:
        add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(0.5),
                 "本计划未提供提领场景", size=14, color=theme["secondary"])
        return
    # 表头
    headers = ["年度", "年提取", "累计提取", "剩余价值", "总收益", "倍数"]
    col_w = [Inches(1.5), Inches(2.2), Inches(2.4), Inches(2.4), Inches(2.4), Inches(1.4)]
    table_x = Inches(0.5); table_y = Inches(2.0)
    add_rect(s, table_x, table_y, sum(col_w, Emu(0)), Inches(0.5), fill=theme["primary"])
    x = table_x
    for i, h in enumerate(headers):
        add_text(s, x, table_y, col_w[i], Inches(0.5),
                 h, size=12, color=theme["text_light"], bold=True, align=PP_ALIGN.CENTER)
        x += col_w[i]
    # 数据行 (5/10/15/20/25/30)
    total_premium = p["policy"]["contractualTotalPremium"]
    table_rows = [r for r in rows if r["policyYear"] in [5, 10, 15, 20, 25, 30]]
    if not table_rows:
        table_rows = rows[::max(1, len(rows)//6)][:6]
    for ridx, r in enumerate(table_rows):
        ry = table_y + Inches(0.5 + ridx * 0.55)
        bg = theme["card"] if ridx % 2 == 0 else theme["bg"]
        add_rect(s, table_x, ry, sum(col_w, Emu(0)), Inches(0.55), fill=bg)
        x = table_x
        total_received = r["cumulativeWithdrawal"] + r["surrenderValueAfter"]
        mult = total_received / total_premium if total_premium else 0
        values = [
            f"Y{r['policyYear']}",
            fmt_num(r["annualWithdrawal"]),
            fmt_num(r["cumulativeWithdrawal"]),
            fmt_num(r["surrenderValueAfter"]),
            fmt_num(total_received),
            f"{mult:.2f}x",
        ]
        for i, v in enumerate(values):
            add_text(s, x, ry, col_w[i], Inches(0.55),
                     v, size=11, color=theme["primary"], align=PP_ALIGN.CENTER, bold=(i == 0 or i == 5))
            x += col_w[i]


def slide_highlights(prs, deck, theme):
    """关键数字 KPI 卡"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["bg"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             "关键数字速览", size=24, color=theme["primary"], bold=True)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(0.6), Inches(0.04), fill=theme["accent"])
    p = deck["products"][0]
    total_premium = p["policy"]["contractualTotalPremium"]
    rows = p["benefitRows"]
    # 自动算 breakeven / 2x / 3x / 5x
    breakeven = None
    mult2y = None
    mult3y = None
    mult5y = None
    for r in rows:
        if breakeven is None and r["totalSurrenderValue"] >= total_premium and total_premium > 0:
            breakeven = r["policyYear"]
        if mult2y is None and r["totalSurrenderValue"] >= 2 * total_premium and total_premium > 0:
            mult2y = r["policyYear"]
        if mult3y is None and r["totalSurrenderValue"] >= 3 * total_premium and total_premium > 0:
            mult3y = r["policyYear"]
        if mult5y is None and r["totalSurrenderValue"] >= 5 * total_premium and total_premium > 0:
            mult5y = r["policyYear"]
    cards = [
        ("回本年", f"{breakeven} 年" if breakeven else "—", "本息首次超过总保费"),
        ("2 倍", f"{mult2y} 年" if mult2y else "—", "本金翻倍"),
        ("3 倍", f"{mult3y} 年" if mult3y else "—", "本金三倍"),
        ("5 倍", f"{mult5y} 年" if mult5y else "—", "本金五倍"),
    ]
    cw = Inches(2.8); ch = Inches(2.2); gx = Inches(0.3)
    sx = Inches(0.5); sy = Inches(2.2)
    for i, (label, value, sub) in enumerate(cards):
        x = sx + i * (cw + gx)
        add_rect(s, x, sy, cw, ch, fill=theme["card"])
        add_rect(s, x, sy, cw, Inches(0.5), fill=theme["primary"])
        add_text(s, x, sy, cw, Inches(0.5), label, size=14, color=theme["text_light"], bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, sy + Inches(0.7), cw, Inches(0.9), value, size=28, color=theme["accent"], bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, sy + Inches(1.6), cw, Inches(0.5), sub, size=10, color=theme["secondary"], align=PP_ALIGN.CENTER)
    # 提领快览
    if p.get("withdrawalRows"):
        wd = [r for r in p["withdrawalRows"] if r["annualWithdrawal"] > 0]
        if wd:
            add_text(s, Inches(0.5), Inches(5), Inches(12), Inches(0.5),
                     f"提领方案: Y{wd[0]['policyYear']} 起每年提取 {fmt_num(wd[0]['annualWithdrawal'])}，至 Y{wd[-1]['policyYear']} 累计 {fmt_num(wd[-1]['cumulativeWithdrawal'])}",
                     size=12, color=theme["primary"], bold=True)


def slide_closing(prs, deck, theme):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, theme["primary"])
    add_rect(s, Emu(0), Inches(0.5), SLIDE_W, Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(2.5), Inches(12), Inches(0.8),
             "愿这份规划", size=36, color=theme["text_light"], bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.4), Inches(12), Inches(0.8),
             "陪伴您的家庭稳健成长", size=36, color=theme["accent"], bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.17), Inches(4.5), Inches(1.0), Inches(0.04), fill=theme["accent"])
    add_text(s, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
             f"—— 致 {deck['customer']['name']} 与家人  ——",
             size=14, color=theme["text_light"], align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(7), Inches(12), Inches(0.3),
             f"本演示基于{deck['company']['displayName']}官方计划书数据生成 · 实际保单条款以保单文件为准",
             size=9, color=theme["accent"], align=PP_ALIGN.CENTER)


def render_pptx(deck: Dict, output_path: str, theme_name: str = "deepblue"):
    theme = THEMES.get(theme_name, THEMES["deepblue"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_cover(prs, deck, theme)
    slide_company(prs, deck, theme)
    slide_growth_chart(prs, deck, theme)
    slide_withdrawal_table(prs, deck, theme)
    slide_highlights(prs, deck, theme)
    slide_closing(prs, deck, theme)
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--deck-json", required=True, help="DeckContract JSON 路径")
    ap.add_argument("--output", required=True, help="输出 PPTX 路径")
    ap.add_argument("--theme", default="deepblue", choices=list(THEMES.keys()))
    args = ap.parse_args()
    with open(args.deck_json, "r", encoding="utf-8") as f:
        deck = json.load(f)
    out = render_pptx(deck, args.output, args.theme)
    size = os.path.getsize(out)
    print(json.dumps({"ok": True, "path": out, "size": size, "slides": 6}, ensure_ascii=False))


if __name__ == "__main__":
    main()
