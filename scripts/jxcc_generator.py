#!/usr/bin/env python3
"""
匠心传承储蓄寿险计划2（尊尚版）— 专用15页PPT生成器
数据来源: uploads/f1e275a3_匠心傳承儲蓄計劃2尊尚版.pdf
- PyMuPDF直接提取，第2页(benefit) + 第42页(withdrawal)
- Y7完整数据，年提取35,000美元（终身）
"""
import sys, os, re, json, argparse
import urllib.request
import urllib.error

# ─── 数据路径配置 ───────────────────────────────────────────────────────────
def _get_default_data_path():
    """获取默认数据路径（兼容旧逻辑）"""
    return os.path.join(os.path.dirname(__file__), "jxcc_data.json")

def load_data(path=None):
    """加载数据文件，支持自定义路径"""
    data_path = path or _get_default_data_path()
    with open(data_path) as f:
        return json.load(f)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch
import io, base64

# ─── 公司介绍（用户提供，自动搜索）────────────────────────────────────────────
# 用法：companies 传入公司名称列表，函数自动生成中英文搜索词并抓取 Wikipedia 描述
# 例如：companies=["周大福人寿", "Chow Tai Fook Life"]
WIKI_SEARCH_TERMS = {
    "周大福人寿":      ["Chow Tai Fook Life Insurance", "周大福人寿保险"],
    "富卫":           ["FWD Group Holdings", "富卫保险"],
    "AIA":            ["AIA Group"],
    "友邦":           ["AIA Group", "友邦保险"],
    "宏利":           ["Manulife", "Manulife Hong Kong"],
    "Manulife":       ["Manulife"],
    "保诚":           ["Prudential plc"],
    "Prudential":     ["Prudential plc"],
    "永明":           ["Sun Life Financial"],
    "Sunlife":        ["Sun Life Financial"],
    "安盛":           ["AXA"],
    "AXA":            ["AXA"],
    "太平":           ["Taiping Life Insurance"],
    "中银":           ["Bank of China Life Insurance"],
}

def fetch_company_info(company_slug: str, default_name: str = "") -> dict:
    """
    从 Wikipedia API 抓取公司介绍。
    company_slug: 中文或英文公司名，会自动匹配 WIKI_SEARCH_TERMS
    返回 {"name": str, "intro": str, "source": str}
    若抓取失败返回默认介绍。
    """
    # 解析搜索词
    search_variants = WIKI_SEARCH_TERMS.get(company_slug, [company_slug])

    # 尝试多个变体，取第一个成功结果
    for term in search_variants:
        url = f"https://en.wikipedia.org/w/api.php?action=query&prop=extracts&explaintext&titles={urllib.request.quote(term)}&format=json&utf8=1&exintro=1"
        try:
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 (HermesAgent/1.0)"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = json.loads(resp.read())
            pages = data.get("query", {}).get("pages", {})
            for page in pages.values():
                extract = page.get("extract", "")
                if extract:
                    # 取前 300 字符，截断到完整句
                    intro = extract[:300].strip()
                    first_period = intro.rfind("。")
                    if first_period > 50:
                        intro = intro[:first_period + 1]
                    return {
                        "name": page.get("title", default_name or term),
                        "intro": intro,
                        "source": "Wikipedia"
                    }
        except Exception:
            continue

    # 兜底默认介绍
    defaults = {
        "周大福人寿": "周大福人寿保险是香港知名保险公司，隶属于周大福集团，专注于人寿保险及财富管理服务，为客户提供全面的保障与储蓄方案。",
        "富卫":      "富卫（FWD）是香港知名保险公司，专注于人寿保险及财富管理，业务遍及亚洲多个市场，致力于为客户创造简单的保险体验。",
    }
    return {
        "name": default_name or company_slug,
        "intro": defaults.get(company_slug, f"{company_slug}，香港知名保险公司，专注于人寿保险及财富管理业务。"),
        "source": ""
    }

_company_info_cache: dict = {}

def get_company_info(company_slug: str = "周大福人寿", default_name: str = "周大福人寿保险") -> dict:
    """全局缓存：公司信息只抓取一次。"""
    if company_slug not in _company_info_cache:
        _company_info_cache[company_slug] = fetch_company_info(company_slug, default_name)
    return _company_info_cache[company_slug]

# ─── 配色 — 基于模板24设计语言 ────────────────────────────────────────────────
C = {
    # 主题色（模板24深海蓝+青绿）
    'primary':      RGBColor(0x0A, 0x3C, 0x5F),   # 深蓝 #0A3C5F
    'accent_teal':  RGBColor(0x18, 0x89, 0x8D),   # 青绿 #18898D
    'accent_gold':  RGBColor(0xC9, 0xA0, 0x27),   # 金色 #C9A027
    'mid_blue':     RGBColor(0x38, 0x5A, 0x64),   # 中蓝 #385A64
    'dark_gray':    RGBColor(0x2D, 0x38, 0x47),   # 深灰 #2D3847
    'bg_dark':      RGBColor(0x0A, 0x3C, 0x5F),   # 背景 #0A3C5F
    'bg_card':      RGBColor(0x0D, 0x2B, 0x44),   # 卡片深色背景
    'text_white':   RGBColor(0xFF, 0xFF, 0xFF),
    'text_gray':    RGBColor(0x88, 0x99, 0xA6),
    'positive':     RGBColor(0x2E, 0xD4, 0x8E),   # 增长绿
    'warning':      RGBColor(0xFF, 0xB8, 0x00),   # 警示橙
    'negative':     RGBColor(0xFF, 0x4D, 0x4D),   # 下降红
    'divider':      RGBColor(0x1A, 0x3A, 0x55),
}

# ─── 辅助函数 ─────────────────────────────────────────────────────────────────
def rgb_hex(d):
    return '#{:02X}{:02X}{:02X}'.format(d[0], d[1], d[2])

def In(s): return Inches(s)
def PtS(s): return Pt(s)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, radius=0.08):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # 全用 type 1 纯矩形，避免 LibreOffice 将 type 9 渲染为椭圆
    shape = slide.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=12, color=None, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Heiti TC"):
    txBox = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = PtS(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox

def add_accent_bar(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def kpi_card(slide, l, t, w, h, label, value, subtitle="", icon=""):
    """增强版KPI卡片 — 带顶部accent线、icon、分层清晰"""
    add_rect(slide, l, t, w, h, fill_color=C['bg_card'], radius=0.1)
    # 顶部细accent线（品牌色）
    add_accent_bar(slide, l+0.15, t+0.06, w-0.3, 0.035, C['accent_teal'])
    # icon（如果提供）
    if icon:
        add_text(slide, icon, l+0.15, t+0.12, 0.5, 0.45, font_size=22, color=C['accent_gold'])
    # label（右上对齐，配合icon偏移）
    label_x = l+0.15 if not icon else l+0.6
    label_w = w-0.3 if not icon else w-0.75
    add_text(slide, label, label_x, t+0.12, label_w, 0.32,
             font_size=10, color=C['text_gray'])
    # value（大号白字）
    add_text(slide, value, l+0.15, t+0.42, w-0.3, 0.5,
             font_size=24, color=C['text_white'], bold=True)
    # subtitle（底部青绿色）
    if subtitle:
        add_text(slide, subtitle, l+0.15, t+0.9, w-0.3, 0.32,
                 font_size=9, color=C['accent_teal'])

def make_fig(w, h, draw_fn):
    import matplotlib
    matplotlib.rcParams['font.family'] = 'Heiti TC'
    matplotlib.rcParams['axes.unicode_minus'] = False
    fig, ax = plt.subplots(figsize=(w, h))
    plt.tight_layout(pad=0.3)
    draw_fn(ax)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf

def fmt(v):
    if v >= 1_000_000: return f"${v/1_000_000:.2f}M"
    if v >= 1_000:     return f"${v/1_000:.0f}K"
    return f"${v:,.0f}"

def fmt_usd(v):
    if v >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if v >= 1_000:     return f"${v/1_000:.0f}K"
    return f"${v:,.0f}"

# ─── 数据加载 ─────────────────────────────────────────────────────────────────
def load_data():
    with open('/Users/soldier/free-code/packages/insurance-ppt/scripts/jxcc_data.json') as f:
        return json.load(f)

def get_benefit(d):
    return d['extractions'][0]['data']['benefit_illustration']

def get_withdrawal(d):
    return d['extractions'][0]['data']['withdrawal_illustration']

def get_policy(d):
    return d['extractions'][0]['data']['policy']

# ─── 图表样式 ─────────────────────────────────────────────────────────────────
def dark_axis(ax, bg='#0A3C5F'):
    spine_c = '#8899A6'
    ax.set_facecolor(bg)
    ax.spines['bottom'].set_color(spine_c)
    ax.spines['left'].set_color(spine_c)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors=spine_c, labelsize=8)
    ax.grid(axis='y', alpha=0.2, color='#FFFFFF', linewidth=0.5)

def light_axis(ax, bg='#E7E6E6'):
    spine_c = '#666666'
    ax.set_facecolor(bg)
    ax.spines['bottom'].set_color(spine_c)
    ax.spines['left'].set_color(spine_c)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors=spine_c, labelsize=8)
    ax.grid(axis='y', alpha=0.3, color='#CCCCCC', linewidth=0.5)

# ═══════════════════════════════════════════════════════════════════════════════
# 第 1 页：封面 — 大改版：左侧大胆深蓝块 + 右侧整块品牌色 + 粗体标题
# ═══════════════════════════════════════════════════════════════════════════════
def page_cover(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])

    # ── 1. 顶部粗金条（0.12寸，比原来的0.06粗一倍）───────────────────────────
    add_accent_bar(s, 0, 0, 13.33, 0.12, C['accent_gold'])

    # ── 2. 左侧深蓝整块（大面积深色块，0.4~4.6寸，视觉锚点）──────────────────
    add_rect(s, 0, 0.12, 4.8, 7.38, fill_color=RGBColor(0x06, 0x20, 0x35))

    # ── 3. 左侧金色垂直条（醒目装饰，贴在深蓝块右侧）─────────────────────────
    add_accent_bar(s, 4.8, 0.12, 0.14, 7.38, C['accent_gold'])

    # ── 4. 主标题（白字，在深蓝块内）─────────────────────────────────────────
    add_text(s, "匠心传承", 0.4, 0.35, 4.2, 1.2,
             font_size=52, color=C['text_white'], bold=True)
    add_text(s, "储蓄寿险计划 2（尊尚版）", 0.4, 1.55, 4.2, 0.7,
             font_size=22, color=C['accent_gold'], bold=True)
    add_text(s, "周大福人寿保险", 0.4, 2.25, 4.0, 0.45,
             font_size=14, color=C['accent_teal'])

    # ── 5. 三个大KPI（垂直堆叠在左侧深蓝块内，有图标有层次）─────────────────
    pol = get_policy(d)
    benefit = get_benefit(d)
    y20 = next((r for r in benefit if r['policy_year']==20), None)
    total_prem = pol.get('total_premium_with_levy') or pol.get('total_premium_paid') or 0
    if total_prem == 0:
        total_prem = (benefit[-1].get('total_premium_paid', 0) or 0) if benefit else 500000

    # KPI 1: 年缴保费
    add_rect(s, 0.3, 2.9, 4.2, 1.2, fill_color=C['bg_card'], radius=0.08)
    add_accent_bar(s, 0.3, 2.9, 0.08, 1.2, C['accent_gold'])
    add_text(s, "💰", 0.5, 2.95, 0.6, 0.5, font_size=24, color=C['accent_gold'])
    add_text(s, "年缴保费", 1.1, 2.98, 3.2, 0.3, font_size=10, color=C['text_gray'])
    annual_prem = pol.get('annual_premium') or 0
    add_text(s, f"${annual_prem:,.0f}", 1.1, 3.25, 3.2, 0.5, font_size=28, color=C['text_white'], bold=True)
    total_display = total_prem
    add_text(s, f"{pol.get('premium_payment_period', '5年')}共 ${total_display:,.0f} USD", 1.1, 3.72, 3.2, 0.3, font_size=10, color=C['accent_teal'])

    # KPI 2: 每年提取
    add_rect(s, 0.3, 4.25, 4.2, 1.2, fill_color=C['bg_card'], radius=0.08)
    add_accent_bar(s, 0.3, 4.25, 0.08, 1.2, C['accent_teal'])
    add_text(s, "📅", 0.5, 4.3, 0.6, 0.5, font_size=24, color=C['accent_gold'])
    add_text(s, "每年提取", 1.1, 4.33, 3.2, 0.3, font_size=10, color=C['text_gray'])
    add_text(s, "$35,000", 1.1, 4.6, 3.2, 0.5, font_size=28, color=C['text_white'], bold=True)
    add_text(s, "第7年起 · 终身现金流", 1.1, 5.07, 3.2, 0.3, font_size=10, color=C['accent_teal'])

    # KPI 3: 第20年退保
    if y20 and total_prem > 0:
        y20_sv = y20.get('total_surrender_value') or 0
        mult = y20_sv / total_prem if total_prem > 0 else 0
        add_rect(s, 0.3, 5.6, 4.2, 1.2, fill_color=C['bg_card'], radius=0.08)
        add_accent_bar(s, 0.3, 5.6, 0.08, 1.2, C['accent_gold'])
        add_text(s, "📈", 0.5, 5.65, 0.6, 0.5, font_size=24, color=C['accent_gold'])
        add_text(s, "第20年退保", 1.1, 5.68, 3.2, 0.3, font_size=10, color=C['text_gray'])
        add_text(s, fmt_usd(y20_sv), 1.1, 5.95, 3.2, 0.5,
                 font_size=28, color=C['text_white'], bold=True)
        add_text(s, f"{mult:.1f}x 已缴保费 · 终身保障", 1.1, 6.42, 3.2, 0.3,
                 font_size=10, color=C['accent_teal'])

    # ── 6. 公司介绍（右上区域，深蓝背景卡）──────────────────────────────────
    company = get_company_info("周大福人寿", "周大福人寿保险")
    add_rect(s, 5.2, 0.12, 7.9, 1.6, fill_color=C['bg_card'], radius=0.1)
    add_accent_bar(s, 5.2, 0.12, 0.1, 1.6, C['accent_teal'])
    add_text(s, f"关于 {company['name']}", 5.45, 0.2, 7.4, 0.38,
             font_size=12, color=C['accent_teal'], bold=True)
    add_text(s, company['intro'], 5.45, 0.58, 7.4, 0.65,
             font_size=10, color=C['text_gray'])
    if company.get('source'):
        add_text(s, f"来源: {company['source']}", 11.0, 0.2, 2.0, 0.3,
                 font_size=8, color=C['text_gray'])

    # ── 7. 计划亮点（右上卡片，大字+装饰线+图标）─────────────────────────────
    add_rect(s, 5.2, 1.85, 7.9, 3.9, fill_color=C['bg_card'], radius=0.1)
    add_text(s, "计划亮点", 5.45, 1.95, 7.5, 0.5,
             font_size=16, color=C['accent_gold'], bold=True)
    add_accent_bar(s, 5.45, 2.45, 7.4, 0.035, C['accent_gold'])

    highlights = [
        ("✅", "5年缴付，7年起每年提取 $35,000"),
        ("✅", "终身现金流，资金灵活支配"),
        ("✅", "复归红利 + 终期分红双引擎驱动"),
        ("✅", "第20年退保达保费2.7倍"),
        ("✅", "免税传承，指定受益人"),
    ]
    for i, (ico, txt) in enumerate(highlights):
        add_text(s, ico, 5.45, 2.6+i*0.6, 0.45, 0.5, font_size=14, color=C['accent_teal'])
        add_text(s, txt, 5.9, 2.6+i*0.6, 6.9, 0.5, font_size=12, color=C['text_white'])

    # ── 8. 右下角装饰：美元符号图案（用文字模拟）─────────────────────────────
    add_text(s, "$", 11.0, 5.9, 2.0, 1.5, font_size=96,
             color=RGBColor(0x0D, 0x2B, 0x44), bold=True)

    add_rect(s, 0, 7.05, 13.33, 0.45, fill_color=C['dark_gray'])
    add_text(s, f"受保人: VIP先生  |  货币: 美元  |  状态: 计划书  |  日期: 2026-05-25",
             0.5, 7.1, 12.5, 0.38, font_size=10, color=C['text_gray'])


# ═══════════════════════════════════════════════════════════════════════════════
# 第 2 页：计划概览 — 加左侧装饰条 + 参数卡加粗
# ═══════════════════════════════════════════════════════════════════════════════
def page_overview(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])

    # 顶部粗金条
    add_accent_bar(s, 0, 0, 13.33, 0.1, C['accent_gold'])
    # 左侧粗装饰条
    add_accent_bar(s, 0, 0.1, 0.12, 7.4, C['accent_teal'])

    add_text(s, "计划概览", 0.5, 0.2, 12, 0.65,
             font_size=32, color=C['text_white'], bold=True)
    add_text(s, "匠心传承储蓄寿险计划2（尊尚版）核心参数", 0.5, 0.85, 12, 0.4,
             font_size=14, color=C['text_gray'])

    pol = get_policy(d)
    params = [
        ("产品名称", "匠心传承储蓄寿险计划2（尊尚版）"),
        ("保险公司", "周大福人寿保险"),
        ("受保人", "VIP先生"),
        ("投保单位", "699,301"),
        ("缴费模式", "年缴"),
        ("年缴保费", f"${pol['annual_premium']:,} USD"),
        ("缴费年期", pol['premium_payment_period']),
        ("保障年期", pol['coverage_period']),
        ("总缴保费", f"${pol['total_premium_with_levy']:,} USD"),
    ]

    # 左列
    for i, (k, v) in enumerate(params[:5]):
        y = 1.4 + i * 0.92
        add_rect(s, 0.5, y, 5.8, 0.8, fill_color=C['bg_card'], radius=0.08)
        add_accent_bar(s, 0.5, y, 0.08, 0.8, C['accent_gold'])
        add_text(s, k, 0.72, y+0.08, 2.0, 0.35, font_size=11, color=C['text_gray'])
        add_text(s, str(v), 2.72, y+0.08, 3.4, 0.55,
                 font_size=14, color=C['text_white'], bold=True)

    # 右列
    for i, (k, v) in enumerate(params[5:]):
        y = 1.4 + i * 0.92
        add_rect(s, 6.6, y, 5.8, 0.8, fill_color=C['bg_card'], radius=0.08)
        add_accent_bar(s, 6.6, y, 0.08, 0.8, C['accent_teal'])
        add_text(s, k, 6.82, y+0.08, 2.0, 0.35, font_size=11, color=C['text_gray'])
        add_text(s, str(v), 8.82, y+0.08, 3.4, 0.55,
                 font_size=14, color=C['text_white'], bold=True)

    # 底部说明
    add_rect(s, 0.5, 6.3, 12.3, 0.65, fill_color=C['divider'], radius=0.08)
    add_text(s, "💡 此计划书数据来源于PDF计划书文件，数字仅为说明用途，实际以保险公司官方为准。",
             0.65, 6.38, 12.0, 0.5, font_size=10, color=C['text_gray'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 3 页：增长故事 — 面积图 + 左侧装饰条
# ═══════════════════════════════════════════════════════════════════════════════
def page_growth_story(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])

    add_accent_bar(s, 0, 0, 13.33, 0.1, C['accent_gold'])
    add_accent_bar(s, 0, 0.1, 0.12, 7.4, C['accent_teal'])

    add_text(s, "💰 账户价值增长 — 复利的魔力", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "不退保、持续累积 — 复归红利与终期分红双引擎驱动", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']
    key_years = [1, 3, 5, 7, 10, 15, 20, 25, 30]
    bmap = {r['policy_year']: r for r in benefit}

    sv_vals = [bmap[y]['total_surrender_value'] for y in key_years]
    gcv_vals = [bmap[y]['guaranteed_cash_value'] for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, gcv_vals, sv_vals,
                        where=[a > b for a, b in zip(sv_vals, gcv_vals)],
                        color='#C9A027', alpha=0.25, label='终期分红（非保证）', zorder=1)
        ax.plot(x, sv_vals, color='#C9A027', linewidth=2.5, marker='o', markersize=7,
                label='退保总额', zorder=3)
        ax.plot(x, gcv_vals, color='#1E6FB7', linewidth=2.0, marker='s', markersize=5,
                label='保证现金价值', zorder=3)
        # 回本线
        for i, (yv, sv) in enumerate(zip(key_years, sv_vals)):
            if sv >= total_prem and i > 0:
                ax.axvline(x=i, color='#00D4AA', linestyle='--', linewidth=1.5, alpha=0.7)
                ax.annotate(f'回本 Y{yv}', xy=(i, sv), xytext=(i+0.5, sv*1.05),
                            fontsize=9, color='#00D4AA', fontweight='bold')
                break
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.35), In(9.5), In(4.2))

    # 右侧KPI
    y7 = bmap[7]
    y20 = bmap[20]
    y30 = bmap[30]
    kpi_card(s, 10.7, 1.35, 2.8, 1.15, "Y7 退保", fmt_usd(y7['total_surrender_value']), "已回本", "📍")
    kpi_card(s, 10.7, 2.6, 2.8, 1.15, "Y20 倍数", f"{y20['total_surrender_value']/total_prem:.1f}x", "退保/已缴保费", "📈")
    kpi_card(s, 10.7, 3.85, 2.8, 1.15, "Y30 倍数", f"{y30['total_surrender_value']/total_prem:.1f}x", "退保/已缴保费", "🚀")
    kpi_card(s, 10.7, 5.1, 2.8, 1.15, "保证 vs 非保证",
             f"{y20['guaranteed_cash_value']/y20['total_surrender_value']*100:.0f}%",
             "Y20保证占比", "🛡️")

    # 底部叙事
    add_rect(s, 0.4, 5.8, 9.4, 0.95, fill_color=C['bg_card'], radius=0.1)
    add_accent_bar(s, 0.4, 5.8, 0.06, 0.95, C['accent_gold'])
    narrative = (
        "匠心传承的核心魅力在于复利的时间价值。"
        f"5年缴付$500,000，第7年账户价值已超过已缴保费。"
        f"第20年退保总额{fmt_usd(y20['total_surrender_value'])}，是已缴保费的{y20['total_surrender_value']/total_prem:.1f}倍。"
        f"第30年更达到{fmt_usd(y30['total_surrender_value'])}，{y30['total_surrender_value']/total_prem:.1f}倍保费。"
    )
    add_text(s, narrative, 0.6, 5.88, 9.1, 0.8,
             font_size=11, color=C['text_white'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 4 页：不提取每5年明细表 — 加左侧装饰条 + 表头加粗 + 交替行色
# ═══════════════════════════════════════════════════════════════════════════════
def page_no_withdrawal_table(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])

    add_accent_bar(s, 0, 0, 13.33, 0.1, C['accent_gold'])
    add_accent_bar(s, 0, 0.1, 0.12, 7.4, C['accent_teal'])

    add_text(s, "📊 不提取每5年明细 — 账户价值累积一览", 0.5, 0.18, 12, 0.6,
             font_size=26, color=C['text_white'], bold=True)
    add_text(s, "持续累积，复利驱动 — 退保总额何时回本、何时翻倍？", 0.5, 0.78, 12, 0.38,
             font_size=12, color=C['text_gray'])

    benefit = get_benefit(d)
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']
    bmap = {r['policy_year']: r for r in benefit}

    key_years = [1, 5, 10, 15, 20, 25, 30, 35, 40]
    n = len(key_years)

    # 表格：5列（左侧区域 0.4~9.9，留0.4边距，右侧KPI占10.2~13.0）
    cols = 5
    rows = n
    tbl = s.shapes.add_table(rows + 1, cols, In(0.4), In(1.3), In(9.4), In(0.52 * (rows + 1))).table

    headers = ['保单年龄', '已交保费', '退保现金价值', '倍数', '备注']
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x18, 0x89, 0x8D)  # 青绿表头
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C['text_white']
            p.alignment = PP_ALIGN.CENTER

    for i, y in enumerate(key_years):
        r = bmap[y]
        sv = r['total_surrender_value']
        mult = sv / total_prem
        paid = r['total_premium_paid']

        row_data = [f'Y{y}', f'${paid:,}', f'${sv:,}', f'{mult:.2f}x', '']
        is_break_even = (mult >= 1.0 and mult < 1.1)
        is_double = (mult >= 2.0 and mult < 2.2)

        if is_break_even: row_data[4] = '✅ 回本'
        if is_double: row_data[4] = '🚀 翻倍'

        # 交替行背景色（斑马纹）
        row_bg = RGBColor(0x0D, 0x2B, 0x44) if i % 2 == 0 else RGBColor(0x08, 0x24, 0x38)

        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                if is_double:
                    p.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
                    p.font.bold = True
                elif is_break_even:
                    p.font.color.rgb = RGBColor(0x2E, 0xD4, 0x8E)
                    p.font.bold = True
                else:
                    p.font.color.rgb = C['text_white']

    # 右侧说明
    y10 = bmap[10]
    y20 = bmap[20]
    kpi_card(s, 10.7, 1.3, 2.8, 1.15, "Y10 回本", f'${y10["total_surrender_value"]/1e6:.2f}M', "1.28x保费", "📍")
    kpi_card(s, 10.7, 2.55, 2.8, 1.15, "Y20 翻倍", f'${y20["total_surrender_value"]/1e6:.2f}M', "2.73x保费", "🚀")
    kpi_card(s, 10.7, 3.8, 2.8, 1.15, "Y30 价值", f'${bmap[30]["total_surrender_value"]/1e6:.2f}M', "5.57x保费", "📈")
    kpi_card(s, 10.7, 5.05, 2.8, 1.15, "Y40 价值", f'${bmap[40]["total_surrender_value"]/1e6:.2f}M', "10.92x保费", "🌱")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 5 页：每年提取现金流
# ═══════════════════════════════════════════════════════════════════════════════
def page_cashflow_timeline(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_teal'])

    add_text(s, "💵 终身现金流时间轴 — 第7年起每年提取", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "从第7个保单年度开始，每年提取 $35,000 美元，活到老领到老", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    withdrawal = get_withdrawal(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    # 展示前15年：Y1-Y15
    key_years = list(range(1, 16))
    bmap_benefit = {r['policy_year']: r for r in benefit}

    sv_no_withdraw = [bmap_benefit.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.bar(x, sv_no_withdraw, color='#1E6FB7', alpha=0.5, label='不退保退保总额', zorder=2)
        # 提领线：提领后
        withdrawal_vals = []
        for y in key_years:
            if y == 1:
                withdrawal_vals.append(sv_no_withdraw[0])
            elif y >= 2:
                # 从withdrawal数据取
                wr = next((r for r in withdrawal if r['policy_year']==y), None)
                if wr:
                    withdrawal_vals.append(wr['surrender_value_after'])
                else:
                    withdrawal_vals.append(sv_no_withdraw[y-1])
        ax.plot(x, withdrawal_vals, color='#00D4AA', linewidth=2.5, marker='s', markersize=5,
                label='每年提取后退保总额', zorder=3)
        # 提领注释
        ax.annotate('第7年起', xy=(6, withdrawal_vals[6]), xytext=(7.5, withdrawal_vals[6]*1.15),
                    fontsize=9, color='#C9A027', fontweight='bold',
                    arrowprops=dict(arrowstyle='->', color='#C9A027'))
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.35), In(9.5), In(4.2))

    # KPI
    y7_b = bmap.get(7, {})
    y15_b = bmap.get(15, {})
    kpi_card(s, 10.7, 1.35, 2.8, 1.2, "Y7退保价值", fmt_usd(y7_b.get('total_surrender_value',0)), "开始提取", "📍")
    kpi_card(s, 10.7, 2.65, 2.8, 1.2, "每年提取", "$35,000", "终身，活到老领到老", "💵")
    kpi_card(s, 10.7, 3.95, 2.8, 1.2, "15年累计提取", f"${35000*9:,}", "Y7-Y15共9年", "📊")
    kpi_card(s, 10.7, 5.25, 2.8, 1.2, "Y15退保价值", fmt_usd(y15_b.get('total_surrender_value',0)), "仍可退保", "📈")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 6 页：提领 vs 不提领对比
# ═══════════════════════════════════════════════════════════════════════════════
def page_withdrawal_comparison(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['warning'])

    add_text(s, "⚖️ 提领 vs 不提领 — 账户价值双线对比", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "同样本金，不同策略：持续累积 vs 终身现金流", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    withdrawal = get_withdrawal(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    # Y1 to Y15
    key_years = list(range(1, 16))
    no_wd = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    with_wd = []
    for y in key_years:
        if y == 1:
            with_wd.append(no_wd[0])
        else:
            wr = next((r for r in withdrawal if r['policy_year']==y), None)
            if wr:
                with_wd.append(wr['surrender_value_after'])
            else:
                with_wd.append(no_wd[y-1])

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, with_wd, no_wd,
                        where=[a > b for a, b in zip(no_wd, with_wd)],
                        color='#C9A027', alpha=0.2, label='提领差额（机会成本）', zorder=1)
        ax.plot(x, no_wd, color='#C9A027', linewidth=2.5, marker='o', markersize=6,
                label='不提领退保总额', zorder=3)
        ax.plot(x, with_wd, color='#00D4AA', linewidth=2.5, marker='s', markersize=6,
                label='每年提取后退保总额', zorder=3)
        # 提领起始年
        ax.axvline(x=6, color='#FFB800', linestyle='--', linewidth=1.5, alpha=0.7)
        ax.text(6.1, ax.get_ylim()[1]*0.8, '提领起始\nY7', fontsize=8, color='#FFB800')
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.35), In(9.5), In(4.2))

    # KPI
    y15_no = no_wd[14]  # Y15
    y15_wd = with_wd[14]
    kpi_card(s, 10.7, 1.35, 2.8, 1.2, "Y15 不提领", fmt_usd(y15_no), f"{y15_no/total_prem:.1f}x保费", "📈")
    kpi_card(s, 10.7, 2.65, 2.8, 1.2, "Y15 提领后", fmt_usd(y15_wd), "已领$35K×9年", "💵")
    kpi_card(s, 10.7, 3.95, 2.8, 1.2, "9年累计提取", f"${35000*9:,}", "Y7至Y15共9年", "📊")
    diff = y15_no - y15_wd
    kpi_card(s, 10.7, 5.25, 2.8, 1.2, "机会成本", fmt_usd(diff), "vs不提领", "⚖️")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 7 页：回本年份分析
# ═══════════════════════════════════════════════════════════════════════════════
def page_breakeven(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['positive'])

    add_text(s, "📍 回本分析 — 保证回本年份", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "退保总额何时超越已缴保费 $500,000？", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    # 找回本年份
    breakeven_year = None
    for r in sorted(benefit, key=lambda x: x['policy_year']):
        if r['total_surrender_value'] >= total_prem:
            breakeven_year = r['policy_year']
            break

    # 保证现金价值回本
    gcv_breakeven = None
    for r in sorted(benefit, key=lambda x: x['policy_year']):
        if r['guaranteed_cash_value'] >= total_prem:
            gcv_breakeven = r['policy_year']
            break

    # 展示Y1-Y15
    key_years = list(range(1, 16))
    prem_line = [total_prem] * len(key_years)
    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    gcv_vals = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, prem_line, sv_vals,
                        where=[sv >= total_prem for sv in sv_vals],
                        color='#2ED48E', alpha=0.15, label='已回本区域', zorder=1)
        ax.fill_between(x, prem_line, gcv_vals,
                        where=[gcv >= total_prem for gcv in gcv_vals],
                        color='#1E6FB7', alpha=0.1, label='保证回本区域', zorder=1)
        ax.plot(x, sv_vals, color='#C9A027', linewidth=2.5, marker='o', markersize=6,
                label='退保总额', zorder=3)
        ax.plot(x, gcv_vals, color='#1E6FB7', linewidth=2.0, marker='s', markersize=5,
                label='保证现金价值', zorder=3)
        ax.axhline(y=total_prem, color='#FF4D4D', linestyle='--', linewidth=1.5,
                   label=f'已缴保费 ${total_prem/1e6:.1f}M', alpha=0.8)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=9)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.5, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.5))

    # KPI
    kpi_card(s, 10.7, 1.3, 2.8, 1.2, "总价值回本", f"Y{breakeven_year}" if breakeven_year else "—",
             "退保总额≥已缴保费", "📍")
    kpi_card(s, 10.7, 2.6, 2.8, 1.2, "保证回本", f"Y{gcv_breakeven}" if gcv_breakeven else "—",
             "保证现金价值≥已缴", "🛡️")
    y15 = bmap.get(15, {})
    kpi_card(s, 10.7, 3.9, 2.8, 1.2, "Y15退保", fmt_usd(y15.get('total_surrender_value',0)),
             f"{y15.get('total_surrender_value',0)/total_prem:.1f}x保费", "📈")
    y7 = bmap.get(7, {})
    kpi_card(s, 10.7, 5.2, 2.8, 1.2, "Y7退保", fmt_usd(y7.get('total_surrender_value',0)),
             f"{y7.get('total_surrender_value',0)/total_prem:.1f}x保费", "✅")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 8 页：提取每5年明细表
# ═══════════════════════════════════════════════════════════════════════════════
def page_withdrawal_table(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_teal'])

    add_text(s, "💵 提取每5年明细 — 终身现金流方案", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "第7年起每年提取 $35,000 — 退保总额何时回本、何时翻倍？", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    withdrawal = get_withdrawal(d)
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']
    bmap = {r['policy_year']: r for r in benefit}
    wmap = {r['policy_year']: r for r in withdrawal}

    # 每5年关键节点（只展示有提领数据的）
    key_years = [1, 5, 10, 15, 20, 25, 30]
    n = len(key_years)

    cols = 6
    rows = n
    tbl = s.shapes.add_table(rows + 1, cols, In(0.4), In(1.3), In(9.4), In(0.52 * (rows + 1))).table

    headers = ['保单年龄', '已交保费', '提取金额', '累计提取', '退保现金价值', '倍数']
    col_widths = [1.5, 2.0, 2.0, 2.0, 2.5, 2.5]

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3A)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C['accent_gold']
            p.alignment = PP_ALIGN.CENTER

    for i, y in enumerate(key_years):
        b = bmap.get(y, {})
        w = wmap.get(y, {})

        paid = b.get('total_premium_paid', total_prem)
        annual_wd = w.get('annual_withdrawal', 0) if w else 0
        cum_wd = w.get('cumulative_withdrawals', 0) if w else 0
        sv_after = w.get('surrender_value_after', 0) if w else b.get('total_surrender_value', 0)
        mult = sv_after / total_prem if sv_after else 0

        is_break_even = (mult >= 1.0 and mult < 1.1)
        is_double = (mult >= 2.0 and mult < 2.2)

        row_data = [
            f'Y{y}',
            f'${paid:,}',
            f'${annual_wd:,}' if annual_wd else '—',
            f'${cum_wd:,}' if cum_wd else '—',
            f'${sv_after:,}' if sv_after else '—',
            f'{mult:.2f}x',
        ]

        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C['bg_card']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
                if is_double:
                    p.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
                    p.font.bold = True
                elif is_break_even:
                    p.font.color.rgb = RGBColor(0x2E, 0xD4, 0x8E)
                    p.font.bold = True
                else:
                    p.font.color.rgb = C['text_white']

    # 右侧说明
    y20 = wmap.get(20, {})
    kpi_card(s, 10.7, 1.3, 2.8, 1.15, "每年提取", "$35,000", "第7年起终身", "💵")
    kpi_card(s, 10.7, 2.55, 2.8, 1.15, "Y20 退保", f"${y20.get('surrender_value_after',0):,}", "1.06x保费", "📍")
    kpi_card(s, 10.7, 3.8, 2.8, 1.15, "Y20 累计领", f"${y20.get('cumulative_withdrawals',0):,}", "约100万", "📊")
    kpi_card(s, 10.7, 5.05, 2.8, 1.15, "Y30 退保", f"${wmap.get(30,{}).get('surrender_value_after','—'):,}" if wmap.get(30,{}) else "—", "1.40x保费", "📈")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 9 页：年度收益明细
# ═══════════════════════════════════════════════════════════════════════════════
def page_yearly_detail(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['mid_blue'])

    add_text(s, "📊 年度收益明细 — 前20年核心数据", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "每年退保总额、保证现金价值、复归红利、终期分红一览", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    rows = [r for r in benefit if r['policy_year'] <= 20]
    n = len(rows)

    # 表格: 6列，深色背景+白字（左侧区域）
    cols = 6
    table = s.shapes.add_table(n+1, cols, In(0.4), In(1.3), In(9.4), In(0.4*(n+1))).table

    # 设置表格默认样式：深色背景，白色字体
    table.first_row = True

    headers = ['年度', '已缴保费', '保证现金价值', '复归红利', '终期分红', '退保总额']
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3A)  # 深蓝表头背景
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = C['accent_gold']  # 金色表头字

    for i, r in enumerate(rows):
        py = r['policy_year']
        row_data = [
            f'Y{py}',
            fmt_usd(r['total_premium_paid']),
            fmt_usd(r['guaranteed_cash_value']),
            fmt_usd(r['reversionary_bonus']),
            fmt_usd(r['terminal_dividend']),
            fmt_usd(r['total_surrender_value']),
        ]
        for j, val in enumerate(row_data):
            cell = table.cell(i+1, j)
            cell.text = val
            # 给每个格子填充深色背景
            cell.fill.solid()
            cell.fill.fore_color.rgb = C['bg_card']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = C['text_white'] if j != 5 else RGBColor(0xFF, 0xD7, 0x00)  # 白色，末列金色
                if j == 5: p.font.bold = True

# ═══════════════════════════════════════════════════════════════════════════════
# 第 10 页：退保 vs 身故对比
# ═══════════════════════════════════════════════════════════════════════════════
def page_surrender_vs_death(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_teal'])

    add_text(s, "🏛️ 退保 vs 身故 — 两种选择，一个方案", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "无论选择退保变现还是代际传承，匠心传承都是您的好选择", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    key_years = [10, 20, 30, 40]
    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    # 身故赔偿 = guaranteed + terminal (最保守估计)
    death_vals = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 + bmap.get(y, {}).get('terminal_dividend', 0) or 0 for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        width = 0.35
        ax.bar(x - width/2, sv_vals, width, color='#C9A027', alpha=0.8, label='退保总额', zorder=2)
        ax.bar(x + width/2, death_vals, width, color='#1E6FB7', alpha=0.8, label='身故赔偿估算', zorder=2)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=11)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=10, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.35), In(9.5), In(4.2))

    # 右侧说明
    advantages = [
        ("💰 退保变现", "第20年可退保 $1.37M\n享受资产增值收益\n资金灵活支配"),
        ("🛡️ 身故传承", "指定受益人\n免遗产税\n财富无损代际传递"),
        ("🔄 灵活切换", "可以先提取后传承\n也可以先传承后提取\n按人生阶段自由调整"),
    ]
    for i, (title, desc) in enumerate(advantages):
        add_rect(s, 10.2, 1.35 + i*1.7, 2.8, 1.55, fill_color=C['bg_card'], radius=0.1)
        add_text(s, title, 10.3, 1.4+i*1.7, 2.6, 0.35, font_size=12, color=C['accent_gold'], bold=True)
        add_text(s, desc, 10.3, 1.75+i*1.7, 2.6, 1.0, font_size=9, color=C['text_white'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 11 页：复归红利详解
# ═══════════════════════════════════════════════════════════════════════════════
def page_reversionary_bonus(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['positive'])

    add_text(s, "📈 复归红利详解 — 稳健增值的基石", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "复归红利每年公布，以面值复归累积，是保证部分的重要组成部分", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    key_years = list(range(1, 21))
    rev = [bmap.get(y, {}).get('reversionary_bonus', 0) or 0 for y in key_years]
    gcv = [bmap.get(y, {}).get('guaranteed_cash_value', 0) or 0 for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, 0, rev, color='#2ED48E', alpha=0.5, label='复归红利', zorder=2)
        ax.plot(x, rev, color='#2ED48E', linewidth=2, marker='o', markersize=4, zorder=3)
        ax.fill_between(x, rev, [r+g for r,g in zip(rev,gcv)], color='#1E6FB7', alpha=0.3,
                        label='保证现金价值（超出部分）', zorder=2)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=8)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e3:.0f}K' if v < 1e6 else f'${v/1e6:.1f}M'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.5, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.5))

    # 右侧
    y20_rev = bmap.get(20, {}).get('reversionary_bonus', 0)
    y20_total = bmap.get(20, {}).get('total_surrender_value', 0)
    kpi_card(s, 10.7, 1.3, 2.8, 1.2, "Y20 复归红利", fmt_usd(y20_rev), "占退保总额32%", "📊")
    kpi_card(s, 10.7, 2.6, 2.8, 1.2, "Y20 保证GCV", fmt_usd(bmap.get(20,{}).get('guaranteed_cash_value',0)), "占退保总额52%", "🛡️")
    add_rect(s, 10.2, 4.0, 2.8, 1.8, fill_color=C['bg_card'], radius=0.1)
    add_text(s, "💡 复归红利特点", 10.3, 4.1, 2.6, 0.35, font_size=11, color=C['positive'], bold=True)
    add_text(s, "• 以面值复归，累积生息\n• 每年公布，稳健增值\n• 退保时可兑现\n• 是保证现金价值基础", 10.3, 4.5, 2.6, 1.2, font_size=9, color=C['text_white'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 12 页：终期分红详解
# ═══════════════════════════════════════════════════════════════════════════════
def page_terminal_dividend(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_gold'])

    add_text(s, "💰 终期分红详解 — 非保证收益的惊喜", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "终期分红在退保或保单终止时支付，是账户价值增长的主要驱动力", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    key_years = list(range(1, 21))
    term = [bmap.get(y, {}).get('terminal_dividend', 0) or 0 for y in key_years]
    rev = [bmap.get(y, {}).get('reversionary_bonus', 0) or 0 for y in key_years]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, 0, term, color='#C9A027', alpha=0.5, label='终期分红（非保证）', zorder=2)
        ax.plot(x, term, color='#C9A027', linewidth=2, marker='o', markersize=4, zorder=3)
        ax.fill_between(x, term, [t+r for t,r in zip(term,rev)], color='#2ED48E', alpha=0.25,
                        label='复归红利（参考）', zorder=2)
        ax.set_xticks(x)
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=8)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M' if v >= 1e6 else f'${v/1e3:.0f}K'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.5, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.3), In(9.5), In(4.5))

    # 右侧
    y20_term = bmap.get(20, {}).get('terminal_dividend', 0)
    y20_total = bmap.get(20, {}).get('total_surrender_value', 0)
    kpi_card(s, 10.7, 1.3, 2.8, 1.2, "Y20 终期分红", fmt_usd(y20_term), f"占{y20_term/y20_total*100:.0f}%退保额", "💰")
    add_rect(s, 10.2, 2.6, 2.8, 3.1, fill_color=C['bg_card'], radius=0.1)
    add_text(s, "⚠️ 终期分红说明", 10.3, 2.7, 2.6, 0.35, font_size=11, color=C['warning'], bold=True)
    add_text(s, "• 非保证，实际金额可能更高或更低\n• 取决于保险公司投资表现\n• 反映保险公司长期盈利分享\n• 悲观/乐观情景有差异\n• 退保时最终公布", 10.3, 3.1, 2.6, 2.4, font_size=9, color=C['text_white'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 13 页：长期增长预测（Y30-Y40）
# ═══════════════════════════════════════════════════════════════════════════════
def page_long_term(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_teal'])

    add_text(s, "🌱 长期增长预测 — Y30至Y40财富轨迹", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "匠心传承越长期越值钱 — 终身复利，世代传承", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    key_years = list(range(20, 41))
    sv_vals = [bmap.get(y, {}).get('total_surrender_value', 0) or 0 for y in key_years]
    mults = [v / total_prem for v in sv_vals]

    def draw(ax):
        x = np.arange(len(key_years))
        ax.fill_between(x, 0, sv_vals, color='#1E6FB7', alpha=0.3, zorder=1)
        ax.plot(x, sv_vals, color='#1E6FB7', linewidth=2.5, marker='o', markersize=5, zorder=3)
        for i, (yv, sv) in enumerate(zip(key_years, sv_vals)):
            if yv in [25, 30, 35, 40]:
                ax.annotate(f'Y{yv}\n{sv/1e6:.1f}M', xy=(i, sv),
                           xytext=(i, sv*1.08), fontsize=8, color='#C9A027',
                           ha='center', fontweight='bold')
        ax.set_xticks(np.arange(len(key_years)))
        ax.set_xticklabels([f'Y{y}' for y in key_years], fontsize=8)
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(
            lambda v, _: f'${v/1e6:.1f}M'))
        ax.legend(fontsize=9, loc='upper left', framealpha=0.8)
        dark_axis(ax)

    buf = make_fig(9.5, 4.2, draw)
    s.shapes.add_picture(buf, In(0.4), In(1.35), In(9.5), In(4.2))

    # 右侧
    y30 = bmap.get(30, {})
    y40 = bmap.get(40, {})
    kpi_card(s, 10.7, 1.35, 2.8, 1.2, "Y30 退保", fmt_usd(y30.get('total_surrender_value',0)), f"{y30.get('total_surrender_value',0)/total_prem:.1f}x保费", "📈")
    kpi_card(s, 10.7, 2.65, 2.8, 1.2, "Y40 退保", fmt_usd(y40.get('total_surrender_value',0)), f"{y40.get('total_surrender_value',0)/total_prem:.1f}x保费", "🚀")
    kpi_card(s, 10.7, 3.95, 2.8, 1.2, "Y30-Y40增长", f"{y40.get('total_surrender_value',0)/y30.get('total_surrender_value',1)-1:.0%}", "10年增长率", "📊")

# ═══════════════════════════════════════════════════════════════════════════════
# 第 14 页：计划总结
# ═══════════════════════════════════════════════════════════════════════════════
def page_summary(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_gold'])

    add_text(s, "📋 计划书总结 — 匠心传承核心要点", 0.5, 0.2, 12, 0.6,
             font_size=28, color=C['text_white'], bold=True)
    add_text(s, "一张保单，三代受益 — 5年播种，终身收获", 0.5, 0.82, 12, 0.38,
             font_size=13, color=C['text_gray'])

    benefit = get_benefit(d)
    bmap = {r['policy_year']: r for r in benefit}
    pol = get_policy(d)
    total_prem = pol['total_premium_with_levy']

    # 左侧总结卡片
    summary_items = [
        ("🏦 产品", "匠心传承储蓄寿险计划2（尊尚版）", C['accent_teal']),
        ("💵 年缴", f"${pol['annual_premium']:,} × 5年 = ${total_prem:,}", C['warning']),
        ("📅 提取", "第7年起每年 $35,000 终身", C['positive']),
        ("📈 回本", "Y7 退保 = $514,498（1.0x保费）", C['accent_gold']),
        ("🚀 退保", f"Y20 = {bmap[20]['total_surrender_value']/total_prem:.1f}x 保费", C['accent_gold']),
        ("🌱 退保", f"Y30 = {bmap[30]['total_surrender_value']/total_prem:.1f}x 保费", C['positive']),
        ("🛡️ 保障", "终身壂障 + 代际传承", C['accent_teal']),
        ("✅ 公司", "周大福人寿保险", C['mid_blue']),
    ]

    for i, (label, val, color) in enumerate(summary_items):
        add_rect(s, 0.5, 1.35 + i * 0.68, 7.8, 0.6, fill_color=C['bg_card'], radius=0.08)
        add_accent_bar(s, 0.5, 1.35 + i*0.68, 0.06, 0.6, color)
        add_text(s, label, 0.65, 1.4+i*0.68, 1.1, 0.5, font_size=11, color=color, bold=True)
        add_text(s, val, 1.75, 1.4+i*0.68, 6.4, 0.5, font_size=11, color=C['text_white'])

    # 右侧核心数据表
    add_rect(s, 8.6, 1.35, 4.4, 5.55, fill_color=C['bg_card'], radius=0.12)
    add_text(s, "📊 核心倍数一览", 8.75, 1.45, 4.1, 0.45,
             font_size=14, color=C['accent_gold'], bold=True)
    headers = ['年度', '倍数', '退保价值']
    for j, h in enumerate(headers):
        add_text(s, h, 8.75 + j*1.4, 1.95, 1.4, 0.35, font_size=10, color=C['text_gray'], bold=True)
    milestones = [5, 7, 10, 15, 20, 25, 30, 40]
    for i, y in enumerate(milestones):
        r = bmap.get(y, {})
        mult = r.get('total_surrender_value', 0) / total_prem
        sv = r.get('total_surrender_value', 0)
        bg_c = C['divider'] if i % 2 == 0 else None
        if bg_c:
            add_rect(s, 8.6, 2.35+i*0.48, 4.4, 0.45, fill_color=bg_c, radius=0.05)
        add_text(s, f'Y{y}', 8.75, 2.35+i*0.48, 1.4, 0.4, font_size=11, color=C['text_white'])
        add_text(s, f'{mult:.1f}x', 10.15, 2.35+i*0.48, 1.4, 0.4, font_size=11, color=C['accent_gold'], bold=True)
        add_text(s, fmt_usd(sv), 11.55, 2.35+i*0.48, 1.4, 0.4, font_size=10, color=C['text_white'])

# ═══════════════════════════════════════════════════════════════════════════════
# 第 15 页：结束页
# ═══════════════════════════════════════════════════════════════════════════════
def page_closing(prs, d):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, C['bg_dark'])
    add_accent_bar(s, 0, 0, 13.33, 0.06, C['accent_gold'])

    # 大标题
    add_text(s, "匠心传承", 1.5, 1.0, 10, 1.0,
             font_size=56, color=C['text_white'], bold=True)
    add_text(s, "储蓄寿险计划 2（尊尚版）", 1.5, 2.1, 10, 0.65,
             font_size=28, color=C['accent_gold'], bold=True)
    add_accent_bar(s, 1.5, 2.9, 5.0, 0.04, C['accent_gold'])

    # 副标题
    add_text(s, "5年播种 · 7年开始 · 终身现金流 · 代际传承", 1.5, 3.1, 10, 0.5,
             font_size=18, color=C['text_gray'])

    # 三个关键词
    kpi_card(s, 1.5, 4.0, 3.2, 1.5, "5年缴付", "$500,000", "稳健起步", "")
    kpi_card(s, 5.0, 4.0, 3.2, 1.5, "7年起领", "$35,000/年", "终身现金流", "")
    kpi_card(s, 8.5, 4.0, 3.2, 1.5, "Y20退保", "$1.37M", "2.7x保费", "")

    # 底部
    add_text(s, "此计划书数据来源于PDF文件，仅供说明用途。实际数字以保险公司官方计划书为准。",
             1.5, 6.0, 10.5, 0.45, font_size=10, color=C['text_gray'])
    add_text(s, "受保人: VIP先生  |  产品: 匠心传承储蓄寿险计划2（尊尚版）  |  周大福人寿保险  |  2026-05-25",
             1.5, 6.5, 10.5, 0.45, font_size=10, color=C['text_gray'])

# ═══════════════════════════════════════════════════════════════════════════════
# 主程序
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    parser = argparse.ArgumentParser(description="匠心传承储蓄寿险计划2 PPT生成器")
    parser.add_argument("--data", "-d", type=str, default=None,
                        help="数据JSON文件路径（默认: scripts/jxcc_data.json）")
    parser.add_argument("--output", "-o", type=str, default=None,
                        help="输出PPTX文件路径（默认: ~/Desktop/insurance_plan_jxcc.pptx）")
    args = parser.parse_args()

    d = load_data(args.data)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    page_cover(prs, d)                   # 1
    page_overview(prs, d)                # 2
    page_growth_story(prs, d)             # 3
    page_no_withdrawal_table(prs, d)      # 4 ← 不提取每5年明细表
    page_cashflow_timeline(prs, d)        # 5
    page_withdrawal_comparison(prs, d)    # 6
    page_breakeven(prs, d)                # 7
    page_withdrawal_table(prs, d)         # 8 ← 提取每5年明细表
    page_yearly_detail(prs, d)            # 9
    page_surrender_vs_death(prs, d)       # 10
    page_long_term(prs, d)                # 11
    page_summary(prs, d)                  # 12
    page_closing(prs, d)                  # 13

    out = args.output or os.path.expanduser("~/Desktop/insurance_plan_jxcc.pptx")
    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides)} slides)")
    return out

if __name__ == '__main__':
    main()
