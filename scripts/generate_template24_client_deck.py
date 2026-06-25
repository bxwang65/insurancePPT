#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json

ROOT=Path('/Users/soldier/free-code/packages/insurance-ppt')
TEMPLATE=Path('/Users/soldier/Downloads/演示文稿24 2.pptx')
OUT=ROOT/'outputs'/'template24_client_ready_v1.pptx'
BASE=json.loads((ROOT/'outputs/1429c316_custom/base_savings.json').read_text())
ANA=json.loads((ROOT/'outputs/1429c316_custom/withdrawal_analysis.json').read_text())

img_dir=ROOT/'outputs/01907ac5_ctf_portraitfix7_pipeline/assets'
chart_dir=ROOT/'outputs/1429c316_custom'

prs=Presentation(str(TEMPLATE))
# remove template slides keep theme
for i in range(len(prs.slides)-1,-1,-1):
    rId=prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[i]

layout=prs.slide_layouts[0]

PRIMARY=RGBColor(0x0A,0x3C,0x5F)
ACCENT=RGBColor(0x18,0x89,0x8D)
TEXT=RGBColor(0x2D,0x38,0x47)


def add_top_nav(slide):
    nav=slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.45))
    nav.fill.solid(); nav.fill.fore_color.rgb=RGBColor(0xF7,0xFA,0xFC)
    nav.line.fill.background()
    t=slide.shapes.add_textbox(Inches(0.4),Inches(0.08),Inches(8),Inches(0.25))
    tf=t.text_frame; tf.text='家庭理财  ·  保险理财  ·  养老保险  ·  理财案例'
    p=tf.paragraphs[0]; p.font.size=Pt(12); p.font.color.rgb=RGBColor(0x55,0x66,0x77)

def title(slide,s):
    t=slide.shapes.add_textbox(Inches(0.6),Inches(0.7),Inches(8.5),Inches(0.8))
    tf=t.text_frame; tf.text=s
    p=tf.paragraphs[0]; p.font.size=Pt(36); p.font.bold=True; p.font.color.rgb=PRIMARY


def subtitle(slide,s):
    t=slide.shapes.add_textbox(Inches(0.65),Inches(1.45),Inches(10),Inches(0.5))
    tf=t.text_frame; tf.text=s
    p=tf.paragraphs[0]; p.font.size=Pt(16); p.font.color.rgb=TEXT


def add_image(slide,path,left,top,width,height):
    if Path(path).exists():
        slide.shapes.add_picture(str(path),left,top,width=width,height=height)


def add_bullets(slide,items,left,top,width,height):
    box=slide.shapes.add_textbox(left,top,width,height)
    tf=box.text_frame; tf.clear()
    for i,it in enumerate(items):
        p=tf.add_paragraph() if i>0 else tf.paragraphs[0]
        p.text='• '+it
        p.font.size=Pt(20)
        p.font.color.rgb=TEXT
        p.space_after=Pt(10)


def add_kpis(slide,kpis):
    x=0.65
    for label,val in kpis:
        card=slide.shapes.add_shape(1, Inches(x), Inches(2.1), Inches(2.6), Inches(1.4))
        card.fill.solid(); card.fill.fore_color.rgb=RGBColor(0xF2,0xF8,0xFB)
        card.line.color.rgb=ACCENT
        tb=slide.shapes.add_textbox(Inches(x+0.15),Inches(2.25),Inches(2.3),Inches(1.1))
        tf=tb.text_frame; tf.text=label
        tf.paragraphs[0].font.size=Pt(13); tf.paragraphs[0].font.color.rgb=RGBColor(0x60,0x70,0x80)
        p=tf.add_paragraph(); p.text=val; p.font.size=Pt(24); p.font.bold=True; p.font.color.rgb=PRIMARY
        x+=2.9


def add_chart_slide(title_text, chart_file, bullets):
    s=prs.slides.add_slide(layout); add_top_nav(s); title(s,title_text)
    add_image(s,chart_dir/chart_file, Inches(0.6), Inches(1.7), Inches(7.2), Inches(4.8))
    panel=s.shapes.add_shape(1, Inches(8.05), Inches(1.75), Inches(4.9), Inches(4.75))
    panel.fill.solid(); panel.fill.fore_color.rgb=RGBColor(0xF7,0xFA,0xFC); panel.line.color.rgb=ACCENT
    add_bullets(s,bullets, Inches(8.3), Inches(2.0), Inches(4.4), Inches(4.2))


def decade_rows(withdraw=False):
    insured_age=BASE['insured']['age']
    base=BASE['benefit_illustration']
    by={int(r['policy_year']):r for r in base}
    wr=ANA['withdrawal_rows'] if withdraw else []
    src=wr if withdraw else [{'age':insured_age+int(r['policy_year']),'policy_year':r['policy_year'],'total_premium_paid':r['total_premium_paid'],'annual_withdrawal':0,'cumulative_withdrawal':0,'surrender_value_after':r['total_surrender_value']} for r in base]
    out=[]
    for r in src:
        py=int(r['policy_year'])
        if py==1 or py%10==0 or py>=120:
            b=by.get(py,{})
            paid=float(b.get('total_premium_paid',r.get('total_premium_paid',0)) or 0)
            val=float(r.get('surrender_value_after',0) or 0)
            y=max(py,1)
            simp=((val/max(paid,1)-1)/y)*100
            cagr=((val/max(paid,1))**(1/y)-1)*100
            out.append([int(r.get('age',insured_age+py)),py,int(paid),int(r.get('annual_withdrawal',0)),int(r.get('cumulative_withdrawal',0)),int(val),f"{simp:.2f}%",f"{cagr:.2f}%"])
    return out[:15]


def add_table_slide(title_text, rows):
    s=prs.slides.add_slide(layout); add_top_nav(s); title(s,title_text)
    cols=['年龄','保单年度','已交总保费','领取金额','累计领取','退保现金价值','单利','复利']
    table=s.shapes.add_table(len(rows)+1,len(cols), Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.9)).table
    for c,h in enumerate(cols):
        cell=table.cell(0,c); cell.text=h
        p=cell.text_frame.paragraphs[0]; p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=RGBColor(0xFF,0xFF,0xFF)
        cell.fill.solid(); cell.fill.fore_color.rgb=PRIMARY
    for r,row in enumerate(rows,1):
        for c,v in enumerate(row):
            cell=table.cell(r,c); cell.text=f"{v:,}" if isinstance(v,int) else str(v)
            p=cell.text_frame.paragraphs[0]; p.font.size=Pt(11); p.font.color.rgb=TEXT; p.alignment=PP_ALIGN.CENTER
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xF8,0xFB,0xFD) if r%2==0 else RGBColor(0xEE,0xF4,0xF8)
    bar=s.shapes.add_shape(1, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb=RGBColor(0xF2,0xF8,0xFB); bar.line.color.rgb=ACCENT
    tb=s.shapes.add_textbox(Inches(0.7),Inches(6.87),Inches(12),Inches(0.3))
    tf=tb.text_frame; tf.text='缴费方式：10万美金 × 5年  ｜  表格口径：每10年抽样对比  ｜  用于销售讲解“流动性 vs 长期价值”'
    tf.paragraphs[0].font.size=Pt(12); tf.paragraphs[0].font.color.rgb=TEXT

# slide 1 cover
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'储蓄险家庭资产配置方案'); subtitle(s,'匠心传承（尊尚版）| 客户可展示版（视觉重制）')
add_image(s,img_dir/'s1.jpg',Inches(6.7),Inches(1.6),Inches(6.1),Inches(4.9))
add_kpis(s,[('缴费期','5年'),('年缴保费','US$ 100,000'),('第7年价值','US$ 514,498'),('第20年价值','US$ 1,366,345')])

# slide2 company
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'公司介绍与资质');
add_image(s,img_dir/'s2.jpg',Inches(0.6),Inches(1.7),Inches(5.8),Inches(5.0))
add_bullets(s,['周大福人寿（CTF Life）','Fitch 财务实力评级：A-','Moody\'s 财务实力评级：A3','香港RBC偿付能力充足率：282%'],Inches(6.7),Inches(2.0),Inches(5.8),Inches(4.5))

# slide3 education
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'教育金阶段（18-21岁）')
add_image(s,img_dir/'s3.jpg',Inches(0.6),Inches(1.7),Inches(5.8),Inches(5.0))
add_bullets(s,['18岁累计可提领：US$ 420,004','21岁累计可提领：US$ 525,006','适配用途：学费、住宿、研究支出'],Inches(6.7),Inches(2.0),Inches(5.8),Inches(4.5))

add_chart_slide('价值增长曲线（提领前后）','withdraw_vs_base.png',['先看两条曲线是否同向上行','再看20年、30年关键节点','结论：提领后仍保留长期价值'])
add_chart_slide('保证/非保证价值构成','guarantee_stack.png',['保证部分构成底盘','非保证部分提升长期弹性','适合长期持有与传承'])
add_chart_slide('年度提领节奏','annual_withdrawal.png',['提领节奏平稳，可预算化','适合作为家庭现金流模块','可与教育/养老阶段支出匹配'])
add_chart_slide('累计提领现金流','cumulative_withdrawal.png',['累计提领持续增长','45岁、60岁为关键里程碑','用于衡量可兑现资金能力'])

add_table_slide('提领方案数据表（每10年）',decade_rows(True))
add_table_slide('不提领方案数据表（每10年）',decade_rows(False))

# closing
s=prs.slides.add_slide(layout); add_top_nav(s); title(s,'执行建议与下一步')
add_image(s,img_dir/'s8.jpg',Inches(6.7),Inches(1.7),Inches(6.1),Inches(4.9))
add_bullets(s,['先确认家庭目标：教育金 / 养老金 / 现金流','确定提领口径与复盘周期（每3年）','按年度跟踪：累计提领、剩余现金价值、配置比例'],Inches(0.7),Inches(2.0),Inches(5.8),Inches(4.5))

OUT.parent.mkdir(parents=True,exist_ok=True)
prs.save(str(OUT))
print(OUT)
